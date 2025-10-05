"""CLI 生成結果を自動検証するスクリプト。"""

from __future__ import annotations

import argparse
import json
import subprocess
from pathlib import Path

from pptx import Presentation

from pptx_generator.models import JobSpec


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="CLI 出力の検証を実行する")
    parser.add_argument(
        "--spec",
        type=Path,
        default=Path("samples/sample_spec.json"),
        help="入力 JSON 仕様のパス",
    )
    parser.add_argument(
        "--workdir",
        type=Path,
        default=Path(".pptxgen"),
        help="CLI 実行時の作業ディレクトリ",
    )
    parser.add_argument(
        "--output",
        default="proposal.pptx",
        help="出力 PPTX ファイル名",
    )
    parser.add_argument(
        "--skip-run",
        action="store_true",
        help="CLI の実行をスキップして検証のみ行う",
    )
    return parser.parse_args()


def run_cli(spec_path: Path, workdir: Path, output_filename: str) -> None:
    workdir.mkdir(parents=True, exist_ok=True)
    cmd = [
        "uv",
        "run",
        "pptx-generator",
        "run",
        str(spec_path),
        "--workdir",
        str(workdir),
        "--output",
        output_filename,
    ]
    subprocess.run(cmd, check=True)


def validate_analysis(analysis_path: Path, spec: JobSpec) -> list[str]:
    errors: list[str] = []
    if not analysis_path.exists():
        errors.append(f"analysis.json が見つかりません: {analysis_path}")
        return errors

    with analysis_path.open(encoding="utf-8") as fp:
        data = json.load(fp)

    if data.get("slides") != len(spec.slides):
        errors.append("analysis.json の slides 件数が仕様と一致しません")
    meta = data.get("meta", {})
    if meta.get("title") != spec.meta.title:
        errors.append("analysis.json のタイトルが仕様と一致しません")
    if not isinstance(data.get("issues"), list):
        errors.append("analysis.json の issues がリストではありません")
    if not isinstance(data.get("fixes"), list):
        errors.append("analysis.json の fixes がリストではありません")
    return errors


def validate_pptx(pptx_path: Path, spec: JobSpec) -> list[str]:
    errors: list[str] = []
    if not pptx_path.exists():
        errors.append(f"PPTX が見つかりません: {pptx_path}")
        return errors

    presentation = Presentation(pptx_path)
    if len(presentation.slides) != len(spec.slides):
        errors.append("PPTX のスライド枚数が仕様と一致しません")

    for slide_spec, slide in zip(spec.slides, presentation.slides, strict=False):
        expected_title = slide_spec.title
        if expected_title is None:
            continue
        actual_title = slide.shapes.title.text if slide.shapes.title else None
        if actual_title != expected_title:
            errors.append(
                f"スライド '{slide_spec.id}' のタイトルが一致しません: {actual_title} != {expected_title}"
            )
    return errors


def main() -> None:
    args = parse_args()
    spec = JobSpec.parse_file(args.spec)

    if not args.skip_run:
        run_cli(args.spec, args.workdir, args.output)

    outputs_dir = args.workdir / "outputs"
    pptx_path = outputs_dir / args.output
    analysis_path = outputs_dir / "analysis.json"

    errors = []
    errors.extend(validate_analysis(analysis_path, spec))
    errors.extend(validate_pptx(pptx_path, spec))

    if errors:
        for error in errors:
            print(f"[NG] {error}")
        raise SystemExit(1)

    print("[OK] CLI 出力の検証に成功しました")


if __name__ == "__main__":
    main()
