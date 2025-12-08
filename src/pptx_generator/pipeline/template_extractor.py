"""テンプレートファイルから図形・プレースホルダー情報を抽出する。"""

from __future__ import annotations

import json
import logging
import re
import unicodedata
from collections import defaultdict
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Literal, Optional

from pptx import Presentation
from pptx.dml.color import ColorFormat
from pptx.enum.text import PP_ALIGN
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import PlaceholderPicture, SlidePlaceholder
from zipfile import BadZipFile

from ..branding_extractor import (BrandingExtractionError,
                                  extract_branding_config)
from ..models import (FontSpec, JobSpecScaffold, JobSpecScaffoldBounds,
                      JobSpecScaffoldMeta, JobSpecScaffoldPlaceholder,
                      JobSpecScaffoldSlide, LayoutInfo, ShapeInfo,
                      TemplateBlueprint, TemplateBlueprintSlide,
                      TemplateBlueprintSlot, TemplateSpec, TextCapacity,
                      TextFramePadding, TextboxParagraph)
from ..utils.layout_metadata import (derive_usage_tags,
                                     generate_layout_description,
                                     normalise_placeholder_type,
                                     summarize_placeholders)
from ..utils.text_capacity import estimate_text_capacity
from .base import PipelineContext

logger = logging.getLogger(__name__)

# SlideBullet拡張仕様で使用される可能性のあるアンカー名パターン
SLIDE_BULLET_ANCHORS = {"bullets", "bullet_list", "content", "body"}
JOBSPEC_SCHEMA_VERSION = "0.1"
MAX_SAMPLE_TEXT_LENGTH = 200
EMU_PER_INCH = 914400

# PowerPoint 側で自動描画されるプレースホルダー種別
AUTO_DRAW_PLACEHOLDER_TYPES = {
    "SLIDE_NUMBER",
    "DATE",
    "DATETIME",
    "FOOTER",
    "HEADER",
}


@dataclass
class TemplateExtractorOptions:
    """TemplateExtractor の設定オプション。"""

    template_path: Path
    output_path: Optional[Path] = None
    layout_filter: Optional[str] = None
    anchor_filter: Optional[str] = None
    format: str = "json"  # json または yaml
    layout_mode: str = "dynamic"
    static_source: Literal["slide", "template"] = "template"


class TemplateExtractorStep:
    """テンプレートファイルから図形情報を抽出するステップ。"""
    
    name = "TemplateExtractor"
    
    def __init__(self, options: TemplateExtractorOptions) -> None:
        self.options = options
        self._slide_width_emu: int | None = None
        self._slide_height_emu: int | None = None
        self._heading_font_default = FontSpec(
            name="Meiryo UI", size_pt=32.0, color_hex="#1A1A1A"
        )
        self._body_font_default = FontSpec(
            name="Meiryo UI", size_pt=18.0, color_hex="#333333"
        )
        self._template_source: Literal["slide", "template"] = "template"
    
    def run(self, context: PipelineContext) -> None:
        """テンプレート抽出を実行する。"""
        logger.info("テンプレート抽出を開始: %s", self.options.template_path)
        
        try:
            template_spec = self.extract_template_spec()
            output_path = self._determine_output_path(context)
            self._save_template_spec(template_spec, output_path)
            jobspec_scaffold = self.build_jobspec_scaffold(template_spec, output_path)
            jobspec_path = self._determine_jobspec_path(output_path)
            self._save_jobspec_scaffold(jobspec_scaffold, jobspec_path)

            context.add_artifact("template_spec", template_spec)
            context.add_artifact("template_spec_path", output_path)
            context.add_artifact("jobspec_scaffold", jobspec_scaffold)
            context.add_artifact("jobspec_path", jobspec_path)

            logger.info("テンプレート抽出完了: %s (jobspec=%s)", output_path, jobspec_path)
            
        except Exception as exc:
            logger.error("テンプレート抽出に失敗: %s", exc)
            raise
    
    def extract_template_spec(self) -> TemplateSpec:
        """テンプレートファイルから仕様を抽出する。"""
        if not self.options.template_path.exists():
            raise FileNotFoundError(f"テンプレートファイルが見つかりません: {self.options.template_path}")

        self._load_font_defaults()

        try:
            presentation = Presentation(self.options.template_path)
        except Exception as exc:
            raise RuntimeError(f"テンプレートファイルの読み込みに失敗しました: {exc}") from exc

        try:
            slide_width = int(presentation.slide_width)
            slide_height = int(presentation.slide_height)
        except Exception as exc:  # noqa: BLE001
            logger.error("スライドサイズの取得に失敗しました: %s", exc)
            raise RuntimeError("スライドサイズの取得に失敗しました") from exc

        if slide_width <= 0 or slide_height <= 0:
            logger.error(
                "スライドサイズが不正です (width=%s, height=%s)", slide_width, slide_height
            )
            raise RuntimeError(
                "スライドサイズが不正です。テンプレートのページ設定を確認してください。"
            )

        self._slide_width_emu = slide_width
        self._slide_height_emu = slide_height

        layout_mode = (self.options.layout_mode or "dynamic").lower()
        if layout_mode not in {"dynamic", "static"}:
            layout_mode = "dynamic"
        template_source: Literal["slide", "template"] = "template"
        if layout_mode == "static":
            candidate_source = (self.options.static_source or "slide").lower()
            template_source = "slide" if candidate_source == "slide" else "template"
        self._template_source = template_source
        layouts: list[LayoutInfo] = []
        warnings: list[str] = []
        errors: list[str] = []

        if template_source == "slide":
            containers = enumerate(presentation.slides, start=1)
        else:
            containers = enumerate(presentation.slide_layouts, start=1)

        for index, container in containers:
            try:
                layout_info = self._extract_layout_info(
                    container,
                    index=index,
                    source_mode=template_source,
                )

                if self.options.layout_filter and not self._matches_filter(
                    layout_info.name, self.options.layout_filter
                ):
                    continue

                layouts.append(layout_info)

            except RuntimeError as exc:
                # アンカー名重複など致命的なエラーは即座に伝播
                if "アンカー名が重複しています" in str(exc):
                    raise
                container_name = getattr(container, "name", None) or f"index={index}"
                error_msg = f"レイアウト '{container_name}' の抽出に失敗: {exc}"
                logger.warning(error_msg)
                errors.append(error_msg)
            except Exception as exc:
                container_name = getattr(container, "name", None) or f"index={index}"
                error_msg = f"レイアウト '{container_name}' の抽出に失敗: {exc}"
                logger.warning(error_msg)
                errors.append(error_msg)

        blueprint = None

        if layout_mode == "static":
            blueprint = self._build_blueprint(layouts)

        return TemplateSpec(
            template_path=str(self.options.template_path),
            extracted_at=datetime.now(timezone.utc).isoformat(),
            template_source=template_source,
            layouts=layouts,
            warnings=warnings,
            errors=errors,
            layout_mode=layout_mode,  # type: ignore[arg-type]
            blueprint=blueprint,
        )

    def _extract_layout_info(
        self,
        container,
        *,
        index: int,
        source_mode: Literal["slide", "template"],
    ) -> LayoutInfo:
        """単一レイアウトから図形情報を抽出する。"""
        prototype_index: int | None = None
        if source_mode == "slide":
            base_name = getattr(container, "name", None)
            if not base_name:
                slide_layout = getattr(container, "slide_layout", None)
                base_name = getattr(slide_layout, "name", None) if slide_layout is not None else None
            slugified = self._slugify_layout_name(base_name)
            if slugified:
                layout_name = f"{slugified}-{index:02d}"
            else:
                layout_name = f"slide-{index:02d}"
            identifier = None
            try:
                slide_identifier = getattr(container, "slide_id", None)
            except Exception:  # noqa: BLE001
                slide_identifier = None
            if slide_identifier is not None:
                identifier = str(slide_identifier)
            prototype_index = index
            shapes_iterable = getattr(container, "shapes", ())
        else:
            layout_name = getattr(container, "name", None)
            identifier = None
            try:
                layout_identifier = getattr(container, "slide_layout_id", None)
            except Exception:  # noqa: BLE001
                layout_identifier = None
            if layout_identifier is not None:
                identifier = str(layout_identifier)
            shapes_iterable = getattr(container, "shapes", ())
            if not layout_name:
                layout_name = f"layout-{index:02d}"
        anchors = []

        for shape in shapes_iterable:
            try:
                shape_info = self._extract_shape_info(shape)
                
                # アンカーフィルタがある場合はチェック
                if self.options.anchor_filter and not self._matches_filter(
                    shape_info.name, self.options.anchor_filter
                ):
                    continue
                
                anchors.append(shape_info)
                
            except Exception as exc:
                error_msg = f"図形 '{shape.name}' の抽出エラー: {exc}"
                logger.warning(error_msg)
                
                # エラー付きの図形情報を作成
                error_shape = ShapeInfo(
                    name=getattr(shape, 'name', '不明な図形'),
                    shape_type="unknown",
                    left_in=0.0,
                    top_in=0.0,
                    width_in=0.0,
                    height_in=0.0,
                    error=error_msg,
                )
                anchors.append(error_shape)

        # アンカー名の重複チェック
        self._check_duplicate_anchors(anchors, layout_name, index, source_mode)

        placeholder_records = [
            self._build_placeholder_record(shape_info)
            for shape_info in anchors
            if self._should_include_for_summary(shape_info)
        ]
        placeholder_summary = summarize_placeholders(placeholder_records)
        heuristic_result = derive_usage_tags(layout_name or "", placeholder_records)
        heuristic_payload = {
            "tags": sorted(heuristic_result.tags),
            "reasons": heuristic_result.reasons,
            "has_title_placeholder": heuristic_result.has_title_placeholder,
            "has_body_placeholder": heuristic_result.has_body_placeholder,
            "title_from_name": heuristic_result.title_from_name,
        }
        if (
            not heuristic_payload["tags"]
            and not heuristic_payload["reasons"]
            and not heuristic_payload["has_title_placeholder"]
            and not heuristic_payload["has_body_placeholder"]
        ):
            heuristic_payload = None

        placeholder_summary_payload = placeholder_summary or None
        layout_description: dict[str, Any] | None = None
        try:
            layout_description = generate_layout_description(
                layout_name or "",
                placeholder_records,
                (self._slide_width_emu or 0, self._slide_height_emu or 0),
            )
        except Exception:  # noqa: BLE001
            layout_description = None

        return LayoutInfo(
            name=layout_name,
            identifier=identifier,
            anchors=anchors,
            prototype_index=prototype_index,
            placeholder_summary=placeholder_summary_payload,
            heuristic=heuristic_payload,
            layout_description=layout_description,
        )
    
    def _extract_shape_info(self, shape: BaseShape) -> ShapeInfo:
        """単一図形から情報を抽出する。"""
        # 基本属性の抽出
        name = getattr(shape, 'name', '')
        if not name:
            name = f"unnamed_shape_{id(shape)}"
        
        # 位置・サイズ情報（EMU単位からインチに変換）
        left_in = shape.left / 914400.0 if hasattr(shape, 'left') else 0.0
        top_in = shape.top / 914400.0 if hasattr(shape, 'top') else 0.0
        width_in = shape.width / 914400.0 if hasattr(shape, 'width') else 0.0
        height_in = shape.height / 914400.0 if hasattr(shape, 'height') else 0.0
        
        # 図形種別の判定
        shape_type = shape.__class__.__name__
        
        # テキスト内容の抽出
        text = None
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is not None:
            frame_text = getattr(text_frame, "text", None)
            if isinstance(frame_text, str):
                text = frame_text
        if text is None:
            raw_text = getattr(shape, "text", None)
            if isinstance(raw_text, str):
                text = raw_text
        
        # プレースホルダー情報の抽出
        placeholder_format = None
        try:
            placeholder_format = shape.placeholder_format  # type: ignore[attr-defined]
        except (AttributeError, ValueError):
            # python-pptx は非プレースホルダー図形にアクセスすると ValueError を送出する
            placeholder_format = None
        is_placeholder = bool(
            isinstance(shape, (SlidePlaceholder, PlaceholderPicture))
            or getattr(shape, "is_placeholder", False)
            or placeholder_format is not None
        )
        placeholder_type = None
        if placeholder_format is not None:
            placeholder_kind = getattr(placeholder_format, "type", None)
            if hasattr(placeholder_kind, "name"):
                placeholder_type = str(getattr(placeholder_kind, "name"))
            elif placeholder_kind is not None:
                placeholder_type = str(placeholder_kind)
        
        # テキスト属性（フォント/段落/余白/容量）
        font_spec: FontSpec | None = None
        paragraph_spec: TextboxParagraph | None = None
        frame_padding: TextFramePadding | None = None
        text_capacity: TextCapacity | None = None
        if self._is_text_shape(placeholder_type, shape):
            (
                font_spec,
                paragraph_spec,
                frame_padding,
                text_capacity,
            ) = self._extract_text_attributes(
                text_frame,
                placeholder_type,
                width_in,
                height_in,
            )

        # SlideBullet拡張仕様との競合チェック
        conflict = None
        if name.lower() in SLIDE_BULLET_ANCHORS:
            conflict = f"SlideBullet拡張仕様で使用される可能性のあるアンカー名: {name}"
        
        # 必須フィールドの欠落チェック
        missing_fields = []
        if not name or name.startswith("unnamed_"):
            missing_fields.append("name")
        if width_in <= 0:
            missing_fields.append("width")
        if height_in <= 0:
            missing_fields.append("height")
        
        return ShapeInfo(
            name=name,
            shape_type=shape_type,
            left_in=left_in,
            top_in=top_in,
            width_in=width_in,
            height_in=height_in,
            text=text,
            placeholder_type=placeholder_type,
            is_placeholder=is_placeholder,
            conflict=conflict,
            missing_fields=missing_fields,
            font=font_spec,
            paragraph=paragraph_spec,
            text_frame_padding=frame_padding,
            text_capacity=text_capacity,
        )

    def _check_duplicate_anchors(
        self,
        anchors: list[ShapeInfo],
        layout_name: str | None,
        index: int,
        source_mode: Literal["slide", "template"],
    ) -> None:
        """同一スライド内でアンカー名の重複をチェックする。重複が見つかった場合は RuntimeError を投げる。"""
        anchor_names: dict[str, list[int]] = {}
        
        for idx, shape_info in enumerate(anchors):
            name = shape_info.name
            # unnamed で始まる名前や空文字は無視（自動生成名）
            if not name or name.startswith("unnamed_"):
                continue
            
            if name not in anchor_names:
                anchor_names[name] = []
            anchor_names[name].append(idx)
        
        # 重複しているアンカー名を検出
        duplicates = {name: indices for name, indices in anchor_names.items() if len(indices) > 1}
        
        if duplicates:
            source_label = "スライド" if source_mode == "slide" else "レイアウト"
            layout_display = layout_name or f"{source_label}-{index:02d}"
            
            error_lines = [
                f"同一{source_label}内でアンカー名が重複しています:",
                f"  - {source_label}: {layout_display} (index={index})",
            ]
            
            for dup_name, indices in duplicates.items():
                error_lines.append(f"  - 重複アンカー: '{dup_name}' (出現回数: {len(indices)})")
            
            error_lines.extend([
                "",
                "修正方法:",
                "  PowerPoint で該当図形を選択し、図形名を一意にリネームしてください。",
                "  図形名は「ホーム」→「選択」→「オブジェクトの選択と表示」で確認・変更できます。",
            ])
            
            error_msg = "\n".join(error_lines)
            logger.error(error_msg)
            raise RuntimeError(error_msg)

    @staticmethod
    def _should_include_for_summary(shape: ShapeInfo) -> bool:
        if shape.is_placeholder:
            return True
        if shape.placeholder_type:
            return True
        if shape.name and shape.name.lower() not in {"rectangle", "textbox"}:
            return True
        return False

    @staticmethod
    def _shape_bbox_emu(shape: ShapeInfo) -> dict[str, int]:
        return {
            "x": int(round(shape.left_in * EMU_PER_INCH)),
            "y": int(round(shape.top_in * EMU_PER_INCH)),
            "width": int(round(shape.width_in * EMU_PER_INCH)),
            "height": int(round(shape.height_in * EMU_PER_INCH)),
        }

    @staticmethod
    def _build_summary_flags(shape: ShapeInfo, normalised_type: str) -> list[str]:
        flags: list[str] = []
        if normalised_type == "unknown":
            flags.append("unknown_type")
        if shape.conflict:
            flags.append("anchor_conflict")
        if shape.missing_fields:
            flags.append("missing_fields")
        return flags

    def _build_placeholder_record(self, shape: ShapeInfo) -> dict[str, Any]:
        normalised_type = normalise_placeholder_type(shape.placeholder_type, shape.name)
        record: dict[str, Any] = {
            "name": shape.name,
            "type": normalised_type,
            "bbox": self._shape_bbox_emu(shape),
            "shape_type": str(shape.shape_type or "").casefold() or None,
            "flags": self._build_summary_flags(shape, normalised_type),
        }
        if record["shape_type"] is None:
            record.pop("shape_type")
        return record

    def _is_text_shape(self, placeholder_type: str | None, shape: BaseShape) -> bool:
        placeholder = (placeholder_type or "").upper()
        if placeholder in {"TITLE", "CENTER_TITLE", "SUBTITLE", "BODY", "CONTENT", "TEXT"}:
            return True
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            return False
        text_value = getattr(text_frame, "text", None) or getattr(shape, "text", None) or ""
        return bool(text_value.strip())

    def _extract_text_attributes(
        self,
        text_frame,
        placeholder_type: str | None,
        width_in: float,
        height_in: float,
    ) -> tuple[FontSpec | None, TextboxParagraph | None, TextFramePadding | None, TextCapacity | None]:
        if text_frame is None:
            font_spec = self._resolve_font_spec(placeholder_type, None)
            text_capacity = estimate_text_capacity(
                width_in=width_in,
                height_in=height_in,
                font=font_spec,
                paragraph=None,
                padding=None,
            )
            return font_spec, None, None, text_capacity

        paragraph_obj = text_frame.paragraphs[0] if text_frame.paragraphs else None
        paragraph_spec = self._convert_paragraph(paragraph_obj)
        font_spec = self._resolve_font_spec(placeholder_type, paragraph_obj)
        padding = self._convert_text_frame_padding(text_frame)
        text_capacity = estimate_text_capacity(
            width_in=width_in,
            height_in=height_in,
            font=font_spec,
            paragraph=paragraph_spec,
            padding=padding,
        )
        return font_spec, paragraph_spec, padding, text_capacity

    def _resolve_font_spec(self, placeholder_type: str | None, paragraph) -> FontSpec:
        placeholder = (placeholder_type or "").upper()
        base = (
            self._heading_font_default
            if placeholder in {"TITLE", "CENTER_TITLE", "SUBTITLE"}
            else self._body_font_default
        )
        overrides = self._font_overrides_from_paragraph(paragraph)
        if not overrides:
            return base
        return base.model_copy(update=overrides)

    def _font_overrides_from_paragraph(self, paragraph) -> dict[str, Any]:
        if paragraph is None:
            return {}

        attributes = self._collect_font_attributes(getattr(paragraph, "font", None))
        if not self._font_attributes_complete(attributes):
            attributes = self._fill_font_attributes_from_runs(paragraph, attributes)

        overrides: dict[str, Any] = {}
        if attributes["name"]:
            overrides["name"] = attributes["name"]
        if attributes["size"] is not None:
            overrides["size_pt"] = attributes["size"]
        if attributes["color"]:
            overrides["color_hex"] = attributes["color"]
        if attributes["bold"] is not None:
            overrides["bold"] = attributes["bold"]
        if attributes["italic"] is not None:
            overrides["italic"] = attributes["italic"]
        return overrides

    @staticmethod
    def _collect_font_attributes(font) -> dict[str, Any]:
        attributes = {
            "name": None,
            "size": None,
            "color": None,
            "bold": None,
            "italic": None,
        }
        if font is None:
            return attributes
        name = getattr(font, "name", None)
        attributes["name"] = name.strip() if isinstance(name, str) else name
        attributes["size"] = _length_to_pt(getattr(font, "size", None))
        attributes["color"] = _color_to_hex(getattr(font, "color", None))
        if font.bold is not None:
            attributes["bold"] = bool(font.bold)
        if font.italic is not None:
            attributes["italic"] = bool(font.italic)
        return attributes

    @staticmethod
    def _font_attributes_complete(attributes: dict[str, Any]) -> bool:
        return all(
            [
                bool(attributes["name"]),
                attributes["size"] is not None,
                bool(attributes["color"]),
                attributes["bold"] is not None,
                attributes["italic"] is not None,
            ]
        )

    def _fill_font_attributes_from_runs(
        self, paragraph, base_attributes: dict[str, Any]
    ) -> dict[str, Any]:
        attributes = dict(base_attributes)
        for run in getattr(paragraph, "runs", []):
            run_font = getattr(run, "font", None)
            if run_font is None:
                continue
            run_attributes = self._collect_font_attributes(run_font)
            for key in ("name", "size", "color", "bold", "italic"):
                if attributes[key] is None and run_attributes[key] is not None:
                    attributes[key] = run_attributes[key]
            if self._font_attributes_complete(attributes):
                break
        return attributes

    def _convert_paragraph(self, paragraph) -> TextboxParagraph | None:
        if paragraph is None:
            return None
        fmt = getattr(paragraph, "paragraph_format", None)
        return TextboxParagraph(
            level=max(paragraph.level if paragraph.level is not None else 0, 0),
            line_spacing_pt=self._line_spacing_to_pt(paragraph),
            space_before_pt=_length_to_pt(getattr(fmt, "space_before", None)) if fmt else None,
            space_after_pt=_length_to_pt(getattr(fmt, "space_after", None)) if fmt else None,
            align=self._alignment_to_str(paragraph.alignment),
            left_indent_in=_length_to_inches(getattr(fmt, "left_margin", None)) if fmt else None,
            right_indent_in=_length_to_inches(getattr(fmt, "right_margin", None)) if fmt else None,
            first_line_indent_in=_length_to_inches(getattr(fmt, "first_line_indent", None)) if fmt else None,
        )

    @staticmethod
    def _alignment_to_str(value) -> str | None:
        if value is None:
            return None
        try:
            align = PP_ALIGN(value)
        except ValueError:
            return None
        return {
            PP_ALIGN.LEFT: "left",
            PP_ALIGN.CENTER: "center",
            PP_ALIGN.RIGHT: "right",
            PP_ALIGN.JUSTIFY: "justify",
            PP_ALIGN.DISTRIBUTE: "distributed",
        }.get(align, None)

    def _convert_text_frame_padding(self, text_frame) -> TextFramePadding | None:
        if text_frame is None:
            return None
        padding = TextFramePadding(
            left_in=_length_to_inches(getattr(text_frame, "margin_left", None)),
            right_in=_length_to_inches(getattr(text_frame, "margin_right", None)),
            top_in=_length_to_inches(getattr(text_frame, "margin_top", None)),
            bottom_in=_length_to_inches(getattr(text_frame, "margin_bottom", None)),
        )
        if all(value in (None, 0.0) for value in padding.model_dump().values()):
            return None
        return padding

    def _line_spacing_to_pt(self, paragraph) -> float | None:
        if paragraph is None:
            return None
        value = getattr(paragraph, "line_spacing", None)
        if value is None:
            return None
        if isinstance(value, (int, float)):
            font = getattr(paragraph, "font", None)
            base_size = _length_to_pt(getattr(font, "size", None)) if font else None
            if base_size is None:
                base_size = self._body_font_default.size_pt
            return float(value) * float(base_size)
        return _length_to_pt(value)

    @staticmethod
    def _matches_filter(value: str, keyword: str) -> bool:
        """フィルタ語との前方一致を確認する。"""
        if not value or not keyword:
            return True
        return value.casefold().startswith(keyword.casefold())
    
    def _determine_output_path(self, context: PipelineContext) -> Path:
        """出力パスを決定する。"""
        if self.options.output_path:
            return self.options.output_path

        context.workdir.mkdir(parents=True, exist_ok=True)
        if self.options.format == "yaml":
            return context.workdir / "template_spec.yaml"
        return context.workdir / "template_spec.json"

    @staticmethod
    def _determine_jobspec_path(spec_output_path: Path) -> Path:
        """jobspec.json の出力先パスを決定する。"""
        return spec_output_path.with_name("jobspec.json")

    def _save_template_spec(self, template_spec: TemplateSpec, output_path: Path) -> None:
        """テンプレート仕様をファイルに保存する。"""
        if self.options.format == "yaml":
            import yaml
            data = template_spec.model_dump(mode="json", exclude_none=True)
            content = yaml.dump(data, allow_unicode=True, default_flow_style=False, indent=2)
        else:
            content = json.dumps(
                template_spec.model_dump(mode="json", exclude_none=True),
                indent=2,
                ensure_ascii=False,
            )

        output_path.write_text(content, encoding="utf-8")
        logger.info("テンプレート仕様を保存: %s", output_path)

    def build_jobspec_scaffold(
        self, template_spec: TemplateSpec, template_spec_path: Path | None = None
    ) -> JobSpecScaffold:
        """テンプレート情報からジョブスペック雛形を生成する。"""
        template_path = self.options.template_path
        template_id = self._derive_template_id(template_path)
        meta = JobSpecScaffoldMeta(
            schema_version=JOBSPEC_SCHEMA_VERSION,
            template_path=str(template_path),
            template_id=template_id,
            generated_at=datetime.now(timezone.utc).isoformat(),
            layout_count=len(template_spec.layouts),
            template_spec_path=str(template_spec_path) if template_spec_path else None,
            template_source=template_spec.template_source,
        )

        counters: defaultdict[str, int] = defaultdict(int)
        slides: list[JobSpecScaffoldSlide] = []

        for layout in template_spec.layouts:
            counters[layout.name] += 1
            sequence = counters[layout.name]
            slide_id = self._resolve_slide_id(layout, sequence)

            placeholders: list[JobSpecScaffoldPlaceholder] = []
            for index, anchor in enumerate(layout.anchors, start=1):
                anchor_name = anchor.name or f"shape_{index:02d}"
                placeholder_type = (anchor.placeholder_type or "").upper()
                is_auto_draw = placeholder_type in AUTO_DRAW_PLACEHOLDER_TYPES

                bounds = JobSpecScaffoldBounds(
                    left_in=anchor.left_in,
                    top_in=anchor.top_in,
                    width_in=anchor.width_in,
                    height_in=anchor.height_in,
                )
                placeholder = JobSpecScaffoldPlaceholder(
                    anchor=anchor_name,
                    kind=self._infer_placeholder_kind(anchor),
                    placeholder_type=anchor.placeholder_type,
                    shape_type=anchor.shape_type,
                    is_placeholder=anchor.is_placeholder,
                    bounds=bounds,
                    sample_text=self._sanitize_sample_text(anchor.text),
                    notes=self._collect_placeholder_notes(anchor),
                    auto_draw=is_auto_draw,
                    font=anchor.font,
                    paragraph=anchor.paragraph,
                    text_frame_padding=anchor.text_frame_padding,
                    text_capacity=anchor.text_capacity,
                )
                placeholders.append(placeholder)

            slides.append(
                JobSpecScaffoldSlide(
                    id=slide_id,
                    layout=layout.name,
                    sequence=sequence,
                    placeholders=placeholders,
                )
            )

        return JobSpecScaffold(meta=meta, slides=slides)

    def _build_blueprint(self, layouts: list[LayoutInfo]) -> TemplateBlueprint:
        slides: list[TemplateBlueprintSlide] = []

        for index, layout in enumerate(layouts, start=1):
            slide_id = self._resolve_slide_id(layout, index)
            slot_sequence = 1
            slots: list[TemplateBlueprintSlot] = []
            for anchor in layout.anchors:
                placeholder_type = (anchor.placeholder_type or "").upper()
                if placeholder_type in AUTO_DRAW_PLACEHOLDER_TYPES:
                    logger.debug(
                        "Blueprint から自動描画プレースホルダーを除外: slide=%s anchor=%s type=%s",
                        slide_id,
                        anchor.name,
                        placeholder_type,
                    )
                    continue
                content_type = self._infer_placeholder_kind(anchor)
                slot_id = f"{slide_id}.slot{slot_sequence:02d}"
                slot_sequence += 1
                required = self._is_required_slot(anchor)
                default_text: list[str] | None = None
                default_payload: dict[str, Any] | None = None
                if content_type == "text":
                    source_text = anchor.text or ""
                    lines = [line.strip() for line in source_text.splitlines() if line.strip()]
                    if lines:
                        default_text = lines
                slots.append(
                    TemplateBlueprintSlot(
                        slot_id=slot_id,
                        anchor=anchor.name,
                        content_type=content_type,
                        required=required,
                        intent_tags=self._derive_slot_intent_tags(anchor, layout.name),
                        default_text=default_text,
                        default_payload=default_payload,
                    )
                )

            slides.append(
                TemplateBlueprintSlide(
                    slide_id=slide_id,
                    layout=layout.name,
                    prototype_index=layout.prototype_index,
                    required=True,
                    intent_tags=self._derive_layout_intent_tags(layout.name),
                    slots=slots,
                )
            )

        return TemplateBlueprint(slides=slides)

    def _save_jobspec_scaffold(self, jobspec: JobSpecScaffold, output_path: Path) -> None:
        """ジョブスペック雛形をファイルに保存する。"""
        output_path.write_text(
            json.dumps(jobspec.model_dump(), indent=2, ensure_ascii=False),
            encoding="utf-8",
        )
        logger.info("ジョブスペック雛形を保存: %s", output_path)

    def _resolve_slide_id(self, layout: LayoutInfo, sequence: int) -> str:
        base = None
        if layout.identifier:
            base = f"id_{layout.identifier}"
        if not base:
            base = self._slugify_layout_name(layout.name)
        if not base:
            base = "slide"
        suffix = f"{sequence:02d}"
        return f"{base}-{suffix}"

    def _infer_placeholder_kind(
        self, shape: ShapeInfo
    ) -> Literal["text", "image", "table", "chart", "shape", "other"]:
        placeholder_type = (shape.placeholder_type or "").upper()
        if placeholder_type in {"TITLE", "CENTER_TITLE", "SUBTITLE", "BODY", "CONTENT", "TEXT"}:
            return "text"
        if placeholder_type in {"PICTURE", "CLIP_ART", "BITMAP", "OBJECT"}:
            return "image"
        if placeholder_type in {"TABLE"}:
            return "table"
        if placeholder_type in {"CHART"}:
            return "chart"

        shape_type = (shape.shape_type or "").lower()
        if "chart" in shape_type or "graph" in shape_type:
            return "chart"
        if "table" in shape_type:
            return "table"
        if "picture" in shape_type or "image" in shape_type or "bitmap" in shape_type:
            return "image"
        if shape.text:
            return "text"
        return "other"

    def _is_required_slot(self, shape: ShapeInfo) -> bool:
        placeholder_type = (shape.placeholder_type or "").upper()
        if placeholder_type in {"TITLE", "CENTER_TITLE", "BODY"}:
            return True
        if placeholder_type in {"SUBTITLE", "CONTENT"}:
            return False
        shape_type = (shape.shape_type or "").lower()
        if "picture" in shape_type or "image" in shape_type:
            return False
        if "chart" in shape_type or "table" in shape_type:
            return False
        return True

    @staticmethod
    def _derive_slot_intent_tags(shape: ShapeInfo, layout_name: str | None) -> list[str]:
        del layout_name  # 現状は形状からの推測に限定
        placeholder_type = (shape.placeholder_type or "").upper()
        if placeholder_type in {"TITLE", "CENTER_TITLE"}:
            return ["headline"]
        if placeholder_type in {"SUBTITLE"}:
            return ["subheadline"]
        if placeholder_type in {"BODY", "CONTENT", "TEXT"}:
            return ["body"]
        return []

    @staticmethod
    def _derive_layout_intent_tags(layout_name: str | None) -> list[str]:
        name = (layout_name or "").lower()
        if "title" in name or "cover" in name:
            return ["opening"]
        if "closing" in name:
            return ["closing"]
        if "agenda" in name:
            return ["agenda"]
        if "summary" in name:
            return ["summary"]
        return []

    def _sanitize_sample_text(self, text: str | None) -> str | None:
        if text is None:
            return None
        cleaned_lines = [line.strip() for line in text.splitlines() if line.strip()]
        if not cleaned_lines:
            cleaned = text.strip()
            if not cleaned:
                return None
        else:
            cleaned = "\n".join(cleaned_lines)
        if len(cleaned) > MAX_SAMPLE_TEXT_LENGTH:
            return cleaned[:MAX_SAMPLE_TEXT_LENGTH].rstrip() + "..."
        return cleaned

    def _collect_placeholder_notes(self, shape: ShapeInfo) -> list[str]:
        notes: list[str] = []
        if shape.conflict:
            notes.append(shape.conflict)
        if shape.missing_fields:
            notes.append("missing_fields: " + ", ".join(shape.missing_fields))
        if shape.error:
            notes.append(shape.error)
        if shape.width_in <= 0 or shape.height_in <= 0:
            notes.append("size_not_positive")
        return notes

    @staticmethod
    def _slugify_layout_name(name: str | None) -> str:
        normalized = unicodedata.normalize("NFKC", (name or "").strip())
        normalized = normalized.replace(" ", "_")
        normalized = re.sub(r"[\s/\\]+", "_", normalized)
        normalized = re.sub(r"[^0-9A-Za-z_\-一-龯ぁ-んァ-ンー]+", "", normalized)
        return normalized.lower()

    @staticmethod
    def _derive_template_id(path: Path) -> str:
        stem = unicodedata.normalize("NFKC", path.stem)
        stem = re.sub(r"[^0-9A-Za-z_\-一-龯ぁ-んァ-ンー]+", "", stem)
        return stem or "template"

    def _load_font_defaults(self) -> None:
        try:
            branding = extract_branding_config(self.options.template_path)
        except (BrandingExtractionError, BadZipFile, FileNotFoundError) as exc:
            logger.debug("branding 設定の抽出に失敗したため既定フォントを使用: %s", exc)
            return

        heading_payload = branding.fonts.get("heading") if branding.fonts else None
        body_payload = branding.fonts.get("body") if branding.fonts else None
        if heading_payload:
            self._heading_font_default = self._font_spec_from_payload(
                heading_payload, self._heading_font_default
            )
        if body_payload:
            self._body_font_default = self._font_spec_from_payload(
                body_payload, self._body_font_default
            )

    @staticmethod
    def _font_spec_from_payload(payload: dict[str, Any], fallback: FontSpec) -> FontSpec:
        name = payload.get("name") or fallback.name
        size_pt = payload.get("size_pt") or fallback.size_pt
        color_hex = _normalize_hex(payload.get("color_hex")) or fallback.color_hex
        bold = payload.get("bold")
        italic = payload.get("italic")
        return fallback.model_copy(
            update={
                "name": str(name),
                "size_pt": float(size_pt),
                "color_hex": color_hex,
                "bold": bool(bold) if bold is not None else fallback.bold,
                "italic": bool(italic) if italic is not None else fallback.italic,
            }
        )


class TemplateExtractor:
    """スタンドアロンでテンプレート抽出を行うクラス。"""
    
    def __init__(self, options: TemplateExtractorOptions) -> None:
        self.options = options
        self.step = TemplateExtractorStep(options)
    
    def extract(self) -> TemplateSpec:
        """テンプレート抽出を実行してTemplateSpecを返す。"""
        return self.step.extract_template_spec()

    def build_jobspec_scaffold(
        self, template_spec: TemplateSpec, template_spec_path: Path | None = None
    ) -> JobSpecScaffold:
        """テンプレート仕様からジョブスペック雛形を構築する。"""
        return self.step.build_jobspec_scaffold(template_spec, template_spec_path)

    def save_jobspec_scaffold(self, jobspec: JobSpecScaffold, output_path: Path) -> None:
        """ジョブスペック雛形を保存する。"""
        self.step._save_jobspec_scaffold(jobspec, output_path)

    def extract_and_save(self, output_path: Optional[Path] = None) -> Path:
        """テンプレート抽出を実行してファイルに保存する。"""
        template_spec = self.extract()

        resolved_output = Path(output_path) if output_path is not None else None
        if resolved_output is None:
            if self.options.format == "yaml":
                resolved_output = Path("template_spec.yaml")
            else:
                resolved_output = Path("template_spec.json")

        jobspec_scaffold = self.build_jobspec_scaffold(template_spec, resolved_output)

        self.step._save_template_spec(template_spec, resolved_output)
        jobspec_path = self.step._determine_jobspec_path(resolved_output)
        self.step._save_jobspec_scaffold(jobspec_scaffold, jobspec_path)
        return resolved_output


def _length_to_pt(value) -> float | None:
    if value is None:
        return None
    try:
        return float(value.pt)
    except AttributeError:
        try:
            return float(value)
        except (TypeError, ValueError):
            return None


def _length_to_inches(value) -> float | None:
    if value is None:
        return None
    try:
        return float(value.inches)
    except AttributeError:
        try:
            return float(value) / EMU_PER_INCH
        except (TypeError, ValueError, ZeroDivisionError):
            return None


def _color_to_hex(color: ColorFormat | None) -> str | None:
    if color is None:
        return None
    rgb = getattr(color, "rgb", None)
    if rgb is None:
        return None
    try:
        components = tuple(rgb)
    except TypeError:
        return None
    return "#" + "".join(f"{component:02X}" for component in components)


def _normalize_hex(value) -> str | None:
    if value is None:
        return None
    text = str(value).strip()
    if not text:
        return None
    return text if text.startswith("#") else f"#{text}"
