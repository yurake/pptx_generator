"""Open XML Polisher を呼び出すステップ。"""

from __future__ import annotations

import json
import logging
import os
import shutil
import subprocess
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.util import Inches, Pt

from .base import PipelineContext

logger = logging.getLogger(__name__)


class PolisherError(RuntimeError):
    """Polisher 実行失敗時に送出される例外。"""


@dataclass(slots=True)
class PolisherOptions:
    """Polisher 呼び出しに関する設定。"""

    enabled: bool = False
    executable: Path | None = None
    rules_path: Path | None = None
    timeout_sec: int = 90
    arguments: tuple[str, ...] = ()
    working_dir: Path | None = None


class PolisherStep:
    """Open XML SDK ベースの仕上げ処理を呼び出すステップ。"""

    name = "polisher"

    def __init__(self, options: PolisherOptions | None = None) -> None:
        self.options = options or PolisherOptions()

    def run(self, context: PipelineContext) -> None:
        pptx_reference = context.require_artifact("pptx_path")
        pptx_path = Path(str(pptx_reference))
        if not pptx_path.exists():  # pragma: no cover - 異常系
            msg = f"PPTX ファイルが存在しません: {pptx_path}"
            raise PolisherError(msg)

        # AI生成フッタ付与処理（Polisher有効/無効に関わらず実行）
        footer_result = self._add_ai_footer_if_needed(pptx_path, context)

        if not self.options.enabled:
            logger.debug("Polisher は無効化されています")
            context.add_artifact(
                "polisher_metadata",
                {
                    "status": "disabled",
                    "enabled": False,
                    "ai_footer": footer_result,
                },
            )
            return

        command = self._build_command(pptx_path)
        cwd = str(self.options.working_dir) if self.options.working_dir else None

        logger.info("Polisher を実行します: %s", " ".join(command))
        start = time.perf_counter()
        try:
            completed = subprocess.run(  # noqa: S603, S607
                command,
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=self.options.timeout_sec,
                cwd=cwd,
            )
        except subprocess.TimeoutExpired as exc:  # pragma: no cover - 異常系
            raise PolisherError("Polisher の実行がタイムアウトしました") from exc
        except subprocess.CalledProcessError as exc:  # pragma: no cover - 異常系
            stdout = exc.stdout.decode(
                "utf-8", errors="replace") if exc.stdout else ""
            stderr = exc.stderr.decode(
                "utf-8", errors="replace") if exc.stderr else ""
            msg = (
                f"Polisher の実行に失敗しました (exit={exc.returncode}).\n"
                f"stdout:\n{stdout}\n"
                f"stderr:\n{stderr}"
            )
            raise PolisherError(msg) from exc
        finally:
            elapsed = time.perf_counter() - start

        metadata: dict[str, Any] = {
            "status": "success",
            "enabled": True,
            "command": command,
            "elapsed_sec": elapsed,
            "returncode": completed.returncode,
            "ai_footer": footer_result,
        }

        stdout_text = completed.stdout.decode(
            "utf-8", errors="replace") if completed.stdout else ""
        stderr_text = completed.stderr.decode(
            "utf-8", errors="replace") if completed.stderr else ""

        if stdout_text:
            metadata["stdout"] = stdout_text
            summary = self._extract_summary(stdout_text)
            if summary is not None:
                metadata["summary"] = summary
        if stderr_text:
            metadata["stderr"] = stderr_text
        if self.options.rules_path:
            metadata["rules_path"] = str(self.options.rules_path)

        context.add_artifact("polisher_metadata", metadata)

    def _build_command(self, pptx_path: Path) -> list[str]:
        executable = self._resolve_executable()
        args = self._prepare_arguments(pptx_path)
        return [str(token) for token in executable] + list(args)

    def _resolve_executable(self) -> list[str]:
        candidate = self.options.executable
        if candidate is None:
            env_value = os.environ.get(
                "POLISHER_EXECUTABLE") or os.environ.get("POLISHER_PATH")
            if env_value:
                candidate = Path(env_value)
        if candidate is None:
            msg = "Polisher の実行ファイルが指定されていません (--polisher-path または設定ファイルを確認してください)"
            raise PolisherError(msg)

        if isinstance(candidate, Path):
            path_candidate = candidate
        else:
            path_candidate = Path(str(candidate))

        if path_candidate.exists():
            if path_candidate.suffix.lower() == ".dll":
                return ["dotnet", str(path_candidate)]
            return [str(path_candidate)]

        resolved = shutil.which(str(path_candidate))
        if resolved:
            return [resolved]

        msg = f"Polisher の実行ファイルが見つかりません: {path_candidate}"
        raise PolisherError(msg)

    def _prepare_arguments(self, pptx_path: Path) -> Iterable[str]:
        rules_path = self.options.rules_path
        args = list(self.options.arguments)

        def _contains_placeholder(values: Iterable[str], placeholder: str) -> bool:
            return any(placeholder in value for value in values)

        if not _contains_placeholder(args, "{pptx}"):
            args.extend(["--input", "{pptx}"])
        if rules_path and not _contains_placeholder(args, "{rules}"):
            args.extend(["--rules", "{rules}"])

        template = {"pptx": str(pptx_path)}
        if rules_path:
            template["rules"] = str(rules_path)

        formatted: list[str] = []
        for item in args:
            try:
                formatted_item = item.format(**template)
            except KeyError as exc:
                msg = f"Polisher 引数テンプレートのプレースホルダー解決に失敗しました: {item}"
                raise PolisherError(msg) from exc
            if formatted_item:
                formatted.append(formatted_item)
        return formatted

    @staticmethod
    def _extract_summary(stdout_text: str) -> dict[str, Any] | None:
        try:
            data = json.loads(stdout_text)
        except json.JSONDecodeError:
            return None

        if isinstance(data, dict):
            return data
        return None

    def _add_ai_footer_if_needed(
        self, pptx_path: Path, context: PipelineContext
    ) -> dict[str, Any]:
        """AI生成コンテンツにフッタを付与する。

        Args:
            pptx_path: PPTX ファイルのパス
            context: パイプラインコンテキスト

        Returns:
            フッタ付与結果の辞書
        """
        result: dict[str, Any] = {
            "enabled": False,
            "slides_modified": 0,
            "error": None,
        }

        # 設定ファイル確認
        if not self.options.rules_path or not Path(self.options.rules_path).exists():
            logger.debug("ルール設定ファイルが存在しないため、AIフッタ付与をスキップします")
            return result

        try:
            with Path(self.options.rules_path).open(encoding="utf-8") as f:
                rules = json.load(f)

            footer_config = rules.get("ai_footer", {})
            if not footer_config.get("enabled", False):
                logger.debug("AIフッタ付与が無効化されています")
                return result

            result["enabled"] = True

            # AI生成メタデータ確認
            workdir = context.workdir
            meta_path = workdir / "ai_generation_meta.json"
            if not meta_path.exists():
                logger.debug("AI生成メタデータが存在しないため、フッタ付与をスキップします")
                return result

            with meta_path.open(encoding="utf-8") as f:
                ai_meta = json.load(f)

            # AI生成スライドIDを抽出
            ai_slide_ids = set()
            for card in ai_meta.get("cards", []):
                card_id = card.get("card_id")
                if card_id:
                    ai_slide_ids.add(card_id)

            if not ai_slide_ids:
                logger.debug("AI生成スライドが見つかりませんでした")
                return result

            # PPTXファイルを開いてフッタを追加
            prs = Presentation(str(pptx_path))

            footer_text = footer_config.get("text", "※本ページはAI生成コンテンツを含みます")
            font_size_pt = footer_config.get("font_size_pt", 8.0)
            color = footer_config.get("color", "#666666")
            margin_in = footer_config.get("margin_in", 0.25)

            slides_modified = 0
            for slide in prs.slides:
                # スライドIDをノート等から取得できないため、
                # 仮実装として全スライドにフッタを付与
                # (実際の実装では、slide.name や context から slide_id をマッピング)
                self._add_footer_to_slide(
                    slide, footer_text, font_size_pt, color, margin_in, prs
                )
                slides_modified += 1

            # 変更を保存
            prs.save(str(pptx_path))
            result["slides_modified"] = slides_modified
            logger.info(f"AIフッタを {slides_modified} スライドに追加しました")

        except Exception as e:  # noqa: BLE001
            logger.warning(f"AIフッタ付与中にエラーが発生しました: {e}")
            result["error"] = str(e)

        return result

    @staticmethod
    def _add_footer_to_slide(
        slide: Any,
        text: str,
        font_size_pt: float,
        color: str,
        margin_in: float,
        prs: Any,
    ) -> None:
        """スライドにフッタテキストボックスを追加する。

        Args:
            slide: スライドオブジェクト
            text: フッタテキスト
            font_size_pt: フォントサイズ（ポイント）
            color: テキスト色（16進数）
            margin_in: 余白（インチ）
            prs: プレゼンテーションオブジェクト
        """
        # スライドサイズを取得
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # テキストボックスのサイズと位置を計算
        textbox_width = Inches(3.0)
        textbox_height = Inches(0.3)
        left = slide_width - textbox_width - Inches(margin_in)
        top = slide_height - textbox_height - Inches(margin_in)

        # テキストボックスを追加
        textbox = slide.shapes.add_textbox(
            left, top, textbox_width, textbox_height)
        text_frame = textbox.text_frame
        text_frame.text = text

        # フォント設定
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(font_size_pt)

        # 色設定（16進数をRGBに変換）
        if color.startswith("#"):
            color_hex = color[1:]
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            from pptx.dml.color import RGBColor
            paragraph.font.color.rgb = RGBColor(r, g, b)
