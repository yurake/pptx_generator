"""PPTX の自動診断ステップ。"""

from __future__ import annotations

import json
import logging
from collections import Counter
from dataclasses import dataclass, field
from itertools import count
from pathlib import Path
from typing import Any, Iterator

from pptx import Presentation
from pptx.dml.color import ColorFormat
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from ..models import (
    MappingLog,
    MappingLogAnalyzerIssue,
    MappingLogAnalyzerSummary,
    Slide,
    SlideBullet,
    SlideImage,
    SlideTextbox,
)
from .base import PipelineContext

logger = logging.getLogger(__name__)

EMU_PER_INCH = 914400


@dataclass(slots=True)
class AnalyzerOptions:
    output_filename: str = "analysis.json"
    min_font_size: float = 18.0
    min_contrast_ratio: float = 4.5
    large_text_min_contrast: float = 3.0
    large_text_threshold_pt: float = 18.0
    max_bullet_level: int = 3
    margin_in: float = 0.5
    slide_width_in: float = 10.0
    slide_height_in: float = 7.5
    default_font_size: float = 18.0
    default_font_color: str = "#333333"
    preferred_text_color: str | None = None
    background_color: str = "#FFFFFF"
    grid_size_in: float = 0.125
    grid_tolerance_in: float = 0.02
    snapshot_output_filename: str | None = None


@dataclass(slots=True)
class ParagraphSnapshot:
    shape_id: int
    shape_name: str | None
    shape_type: int
    paragraph_index: int
    text: str
    level: int
    font_size_pt: float | None
    color_hex: str | None


@dataclass(slots=True)
class ShapeSnapshot:
    shape_id: int
    name: str | None
    shape_type: int
    left_in: float
    top_in: float
    width_in: float
    height_in: float
    paragraphs: list[ParagraphSnapshot] = field(default_factory=list)
    is_placeholder: bool = False
    placeholder_type: int | None = None


@dataclass(slots=True)
class SlideSnapshot:
    index: int
    shapes: list[ShapeSnapshot]
    body_placeholder_id: int | None = None

    @classmethod
    def from_slide(cls, slide, index: int) -> "SlideSnapshot":
        shapes: list[ShapeSnapshot] = []
        body_placeholder_id: int | None = None

        for shape in slide.shapes:
            shape_id = getattr(shape, "shape_id", id(shape))
            left_in = _emu_to_inches(int(getattr(shape, "left", 0)))
            top_in = _emu_to_inches(int(getattr(shape, "top", 0)))
            width_in = _emu_to_inches(int(getattr(shape, "width", 0)))
            height_in = _emu_to_inches(int(getattr(shape, "height", 0)))
            shape_name = getattr(shape, "name", None)
            shape_type = int(getattr(shape, "shape_type", MSO_SHAPE_TYPE.AUTO_SHAPE))
            is_placeholder = bool(getattr(shape, "is_placeholder", False))
            placeholder_type = None
            if is_placeholder:
                try:
                    placeholder_type = int(shape.placeholder_format.type)
                except Exception:  # noqa: BLE001
                    placeholder_type = None

            paragraphs: list[ParagraphSnapshot] = []
            if getattr(shape, "has_text_frame", False):
                text_frame = shape.text_frame
                for idx, paragraph in enumerate(text_frame.paragraphs):
                    text = paragraph.text or ""
                    font_size_pt, color_hex = _extract_font_info(paragraph)
                    level = paragraph.level if paragraph.level is not None else 0
                    paragraphs.append(
                        ParagraphSnapshot(
                            shape_id=shape_id,
                            shape_name=shape_name,
                            shape_type=shape_type,
                            paragraph_index=idx,
                            text=text,
                            level=level,
                            font_size_pt=font_size_pt,
                            color_hex=color_hex,
                        )
                    )

            snapshot = ShapeSnapshot(
                shape_id=shape_id,
                name=shape_name,
                shape_type=shape_type,
                left_in=left_in,
                top_in=top_in,
                width_in=width_in,
                height_in=height_in,
                paragraphs=paragraphs,
                is_placeholder=is_placeholder,
                placeholder_type=placeholder_type,
            )
            shapes.append(snapshot)

            if is_placeholder and placeholder_type in {
                int(PP_PLACEHOLDER.BODY),
                int(PP_PLACEHOLDER.VERTICAL_BODY),
                int(PP_PLACEHOLDER.OBJECT),
            }:
                body_placeholder_id = shape_id

        return cls(index=index, shapes=shapes, body_placeholder_id=body_placeholder_id)

    def shape_by_id(self, shape_id: int) -> ShapeSnapshot | None:
        for shape in self.shapes:
            if shape.shape_id == shape_id:
                return shape
        return None

    def find_shape_by_name(
        self, name: str, *, shape_type: int | None = None
    ) -> ShapeSnapshot | None:
        if not name:
            return None
        for shape in self.shapes:
            if shape.name != name:
                continue
            if shape_type is not None and shape.shape_type != shape_type:
                continue
            return shape
        return None

    def body_paragraphs(self) -> list[ParagraphSnapshot]:
        if self.body_placeholder_id is None:
            return []
        shape = self.shape_by_id(self.body_placeholder_id)
        return list(shape.paragraphs) if shape is not None else []


class BulletParagraphResolver:
    """スライド内の段落を箇条書きグループに対応づける。"""

    def __init__(self, snapshot: SlideSnapshot):
        self._anchor_iters: dict[str, Iterator[ParagraphSnapshot]] = {}
        for shape in snapshot.shapes:
            if not shape.name or not shape.paragraphs:
                continue
            self._anchor_iters.setdefault(shape.name, iter(shape.paragraphs))
        self._fallback_iter: Iterator[ParagraphSnapshot] = iter(snapshot.body_paragraphs())

    def resolve(self, anchor: str | None) -> ParagraphSnapshot | None:
        if anchor:
            iterator = self._anchor_iters.get(anchor)
            if iterator is None:
                return None
            return next(iterator, None)
        return next(self._fallback_iter, None)


class SimpleAnalyzerStep:
    """PPTX の実データを解析して analysis.json を生成する。"""

    name = "analyzer"

    def __init__(
        self,
        options: AnalyzerOptions | None = None,
        *,
        artifact_key: str = "analysis_path",
        register_default_artifact: bool = True,
        allow_missing_artifact: bool = False,
    ) -> None:
        self.options = options or AnalyzerOptions()
        self._sequence = count(1)
        self._artifact_key = artifact_key
        self._register_default_artifact = register_default_artifact
        self._allow_missing_artifact = allow_missing_artifact

    def run(self, context: PipelineContext) -> None:
        pptx_reference = context.artifacts.get("pptx_path")
        if pptx_reference is None:
            if self._allow_missing_artifact:
                logger.info(
                    "Analyzer (%s) をスキップします: pptx_path artifact が見つかりません",
                    self.options.output_filename,
                )
                return
            msg = "解析対象の PPTX が存在しません。renderer の実行順序を確認してください。"
            raise RuntimeError(msg)

        pptx_path = Path(str(pptx_reference))
        if not pptx_path.exists():  # pragma: no cover - 異常系
            raise FileNotFoundError(f"PPTX ファイルが存在しません: {pptx_path}")

        presentation = Presentation(pptx_path)
        issues: list[dict[str, Any]] = []
        fixes: list[dict[str, Any]] = []
        snapshot_slides: list[dict[str, Any]] = []

        spec_slides = context.spec.slides
        if len(presentation.slides) < len(spec_slides):
            logger.warning(
                "PPTX のスライド数が不足しています: spec=%s, pptx=%s",
                len(spec_slides),
                len(presentation.slides),
            )

        presentation_width_in = _emu_to_inches(int(getattr(presentation, "slide_width", 0)))
        presentation_height_in = _emu_to_inches(int(getattr(presentation, "slide_height", 0)))
        if presentation_width_in <= 0:
            presentation_width_in = self.options.slide_width_in
        if presentation_height_in <= 0:
            presentation_height_in = self.options.slide_height_in

        for index, slide_spec in enumerate(spec_slides):
            if index >= len(presentation.slides):
                break
            slide = presentation.slides[index]
            slide_width_in = _emu_to_inches(int(getattr(slide, "slide_width", 0)))
            slide_height_in = _emu_to_inches(int(getattr(slide, "slide_height", 0)))
            if slide_width_in <= 0:
                slide_width_in = presentation_width_in
            if slide_height_in <= 0:
                slide_height_in = presentation_height_in
            snapshot = SlideSnapshot.from_slide(slide, index)
            slide_issues, slide_fixes = self._analyze_slide(
                slide_spec, snapshot, slide_width_in, slide_height_in
            )
            issues.extend(slide_issues)
            fixes.extend(slide_fixes)
            if self.options.snapshot_output_filename:
                snapshot_slides.append(
                    self._export_snapshot_slide(slide_spec, snapshot)
                )

        analysis = {
            "slides": len(spec_slides),
            "meta": context.spec.meta.model_dump(),
            "issues": issues,
            "fixes": fixes,
        }
        output_path = self._save(analysis, context.workdir)
        context.add_artifact(self._artifact_key, output_path)
        if self._register_default_artifact and self._artifact_key != "analysis_path":
            context.add_artifact("analysis_path", output_path)
        elif self._register_default_artifact and self._artifact_key == "analysis_path":
            # 既定キーと同じ場合は追加登録不要
            pass
        self._sync_mapping_log(context, analysis)
        logger.info("%s を出力しました: %s", self.options.output_filename, output_path)
        if self.options.snapshot_output_filename:
            snapshot_path = self._save_snapshot(
                snapshot_slides, context.workdir
            )
            context.add_artifact("analyzer_snapshot_path", snapshot_path)
            logger.info(
                "構造スナップショットを出力しました: %s", snapshot_path
            )

    def _analyze_slide(
        self,
        slide_spec: Slide,
        snapshot: SlideSnapshot,
        slide_width_in: float,
        slide_height_in: float,
    ) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
        issues: list[dict[str, Any]] = []
        fixes: list[dict[str, Any]] = []

        resolver = BulletParagraphResolver(snapshot)
        applied_level: int | None = None
        previous_level: int | None = None

        for group in slide_spec.bullets:
            for bullet in group.items:
                paragraph = resolver.resolve(group.anchor)
                actual_level = paragraph.level if paragraph else bullet.level
                target = {
                    "slide_id": slide_spec.id,
                    "element_id": bullet.id,
                    "element_type": "bullet",
                }

                depth = self._check_bullet_depth(slide_spec, bullet, actual_level, target)
                if depth:
                    issue, fix = depth
                    issues.append(issue)
                    if fix:
                        fixes.append(fix)

                font = self._check_font_size(slide_spec, bullet, paragraph, target)
                if font:
                    issue, fix = font
                    issues.append(issue)
                    if fix:
                        fixes.append(fix)

                contrast = self._check_contrast(slide_spec, bullet, paragraph, target)
                if contrast:
                    issue, fix = contrast
                    issues.append(issue)
                    if fix:
                        fixes.append(fix)

                allowed_level = (
                    0 if applied_level is None else min(applied_level + 1, self.options.max_bullet_level)
                )
                if actual_level > allowed_level:
                    issue_id = self._next_issue_id("layout_consistency", slide_spec.id, bullet.id)
                    fix = {
                        "id": f"fix-{issue_id}",
                        "issue_id": issue_id,
                        "type": "bullet_reindent",
                        "target": target,
                        "payload": {"level": allowed_level},
                    }
                    issue = self._make_issue(
                        issue_id=issue_id,
                        issue_type="layout_consistency",
                        severity="warning",
                        message=(
                            f"スライド '{slide_spec.id}' の箇条書き '{bullet.id}' のレベル {actual_level} が"
                            f" 許容ステップ {allowed_level} を超えています"
                        ),
                        target=target,
                        metrics={
                            "level": actual_level,
                            "allowed_level": allowed_level,
                            "previous_level": previous_level,
                        },
                        fix=fix,
                    )
                    issues.append(issue)
                    fixes.append(fix)
                    applied_level = allowed_level
                else:
                    applied_level = actual_level
                previous_level = actual_level

                if paragraph is None:
                    logger.debug(
                        "箇条書き '%s' に対応する PPTX 段落が見つかりませんでした (slide=%s, anchor=%s)",
                        bullet.id,
                        slide_spec.id,
                        group.anchor,
                    )

        for image_spec in slide_spec.images:
            shape = self._locate_image_shape(snapshot, image_spec)
            if shape is None:
                logger.debug("画像 '%s' の図形が見つかりません", image_spec.id)
                continue
            result = self._check_margins(
                slide_spec,
                image_spec,
                shape,
                slide_width_in=slide_width_in,
                slide_height_in=slide_height_in,
            )
            if result:
                issue, fix = result
                issues.append(issue)
                if fix:
                    fixes.append(fix)
            grid = self._check_grid_alignment(slide_spec, image_spec.id, "image", shape)
            if grid:
                issue, fix = grid
                issues.append(issue)
                if fix:
                    fixes.append(fix)

        for textbox in slide_spec.textboxes:
            shape = self._locate_textbox_shape(snapshot, textbox)
            if shape is None:
                logger.debug("テキストボックス '%s' の図形が見つかりません", textbox.id)
                continue
            grid = self._check_grid_alignment(slide_spec, textbox.id, "textbox", shape)
            if grid:
                issue, fix = grid
                issues.append(issue)
                if fix:
                    fixes.append(fix)

        return issues, fixes

    def _export_snapshot_slide(
        self, slide_spec: Slide, snapshot: SlideSnapshot
    ) -> dict[str, Any]:
        named_shapes: list[dict[str, Any]] = []
        placeholders: list[dict[str, Any]] = []

        for shape in snapshot.shapes:
            base_record = {
                "shape_id": shape.shape_id,
                "name": shape.name or "",
                "shape_type": self._shape_type_name(shape.shape_type),
                "left_in": shape.left_in,
                "top_in": shape.top_in,
                "width_in": shape.width_in,
                "height_in": shape.height_in,
            }
            if shape.is_placeholder or shape.placeholder_type is not None:
                placeholders.append(
                    {
                        **base_record,
                        "is_placeholder": shape.is_placeholder,
                        "placeholder_type": self._placeholder_type_name(
                            shape.placeholder_type
                        ),
                    }
                )
            if shape.name:
                named_shapes.append(base_record)

        spec_anchors = sorted(
            {
                *(group.anchor for group in slide_spec.bullets if group.anchor),
                *(image.anchor for image in slide_spec.images if image.anchor),
                *(table.anchor for table in slide_spec.tables if table.anchor),
                *(chart.anchor for chart in slide_spec.charts if chart.anchor),
                *(textbox.anchor for textbox in slide_spec.textboxes if textbox.anchor),
            }
        )

        return {
            "index": snapshot.index,
            "slide_id": slide_spec.id,
            "layout": slide_spec.layout,
            "placeholders": placeholders,
            "named_shapes": named_shapes,
            "spec_anchors": spec_anchors,
        }

    def _save_snapshot(
        self, slides: list[dict[str, Any]], workdir: Path
    ) -> Path:
        if not self.options.snapshot_output_filename:
            raise ValueError("snapshot_output_filename が設定されていません")
        payload = {
            "schema_version": "1.0.0",
            "slides": slides,
        }
        path = workdir / self.options.snapshot_output_filename
        path.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        return path

    @staticmethod
    def _shape_type_name(shape_type: int | None) -> str:
        if shape_type is None:
            return "unknown"
        try:
            return MSO_SHAPE_TYPE(shape_type).name
        except ValueError:  # pragma: no cover - 予期しない値
            return str(shape_type)

    @staticmethod
    def _placeholder_type_name(
        placeholder_type: int | None,
    ) -> str | None:
        if placeholder_type is None:
            return None
        try:
            return PP_PLACEHOLDER(placeholder_type).name
        except ValueError:  # pragma: no cover - 予期しない値
            return str(placeholder_type)

    def _locate_image_shape(
        self, snapshot: SlideSnapshot, image_spec: SlideImage
    ) -> ShapeSnapshot | None:
        if image_spec.anchor:
            shape = snapshot.find_shape_by_name(
                image_spec.anchor, shape_type=int(MSO_SHAPE_TYPE.PICTURE)
            )
            if shape:
                return shape
        shape = snapshot.find_shape_by_name(image_spec.id, shape_type=int(MSO_SHAPE_TYPE.PICTURE))
        if shape:
            return shape
        for candidate in snapshot.shapes:
            if candidate.shape_type == int(MSO_SHAPE_TYPE.PICTURE):
                return candidate
        return None

    def _locate_textbox_shape(
        self, snapshot: SlideSnapshot, textbox: SlideTextbox
    ) -> ShapeSnapshot | None:
        candidates = []
        if textbox.anchor:
            shape = snapshot.find_shape_by_name(textbox.anchor)
            if shape:
                candidates.append(shape)
        if not candidates:
            shape = snapshot.find_shape_by_name(textbox.id)
            if shape:
                candidates.append(shape)
        for shape in candidates:
            if shape.shape_type in {int(MSO_SHAPE_TYPE.TEXT_BOX), int(MSO_SHAPE_TYPE.PLACEHOLDER)}:
                return shape
        for shape in snapshot.shapes:
            if shape.shape_type in {int(MSO_SHAPE_TYPE.TEXT_BOX), int(MSO_SHAPE_TYPE.PLACEHOLDER)}:
                return shape
        return None

    def _check_bullet_depth(
        self,
        slide: Slide,
        bullet: SlideBullet,
        actual_level: int,
        target: dict[str, Any],
    ) -> tuple[dict[str, Any], dict[str, Any]] | None:
        if actual_level <= self.options.max_bullet_level:
            return None

        issue_id = self._next_issue_id("bullet_depth", slide.id, bullet.id)
        fix = {
            "id": f"fix-{issue_id}",
            "issue_id": issue_id,
            "type": "bullet_cap",
            "target": target,
            "payload": {"level": self.options.max_bullet_level},
        }
        issue = self._make_issue(
            issue_id=issue_id,
            issue_type="bullet_depth",
            severity="warning",
            message=(
                f"スライド '{slide.id}' の箇条書き '{bullet.id}' のレベルが"
                f" 上限 {self.options.max_bullet_level} を超えています"
            ),
            target=target,
            metrics={
                "level": actual_level,
                "max_level": self.options.max_bullet_level,
            },
            fix=fix,
        )
        return issue, fix

    def _check_font_size(
        self,
        slide: Slide,
        bullet: SlideBullet,
        paragraph: ParagraphSnapshot | None,
        target: dict[str, Any],
    ) -> tuple[dict[str, Any], dict[str, Any]] | None:
        size = paragraph.font_size_pt if paragraph and paragraph.font_size_pt is not None else None
        if size is None:
            size = self.options.default_font_size

        if size >= self.options.min_font_size:
            return None

        issue_id = self._next_issue_id("font_min", slide.id, bullet.id)
        fix = {
            "id": f"fix-{issue_id}",
            "issue_id": issue_id,
            "type": "font_raise",
            "target": target,
            "payload": {"size_pt": self.options.min_font_size},
        }
        issue = self._make_issue(
            issue_id=issue_id,
            issue_type="font_min",
            severity="warning",
            message=(
                f"スライド '{slide.id}' の箇条書き '{bullet.id}' のフォントサイズ {size:.1f}pt が"
                f" 下限 {self.options.min_font_size:.1f}pt を下回っています"
            ),
            target=target,
            metrics={
                "size_pt": size,
                "min_size_pt": self.options.min_font_size,
                "shape_name": paragraph.shape_name if paragraph else None,
            },
            fix=fix,
        )
        return issue, fix

    def _check_contrast(
        self,
        slide: Slide,
        bullet: SlideBullet,
        paragraph: ParagraphSnapshot | None,
        target: dict[str, Any],
    ) -> tuple[dict[str, Any], dict[str, Any]] | None:
        color_hex = (
            paragraph.color_hex if paragraph and paragraph.color_hex else self.options.default_font_color
        )
        font_size = (
            paragraph.font_size_pt if paragraph and paragraph.font_size_pt else self.options.default_font_size
        )
        try:
            ratio = _contrast_ratio(color_hex, self.options.background_color)
        except ValueError:
            logger.debug("無効なカラーコードのためコントラスト判定をスキップ: %s", color_hex)
            return None

        required_ratio = self.options.min_contrast_ratio
        if font_size >= self.options.large_text_threshold_pt:
            required_ratio = min(required_ratio, self.options.large_text_min_contrast)

        if ratio >= required_ratio:
            return None

        issue_id = self._next_issue_id("contrast_low", slide.id, bullet.id)
        suggested_color = self.options.preferred_text_color or self.options.default_font_color
        fix = {
            "id": f"fix-{issue_id}",
            "issue_id": issue_id,
            "type": "color_adjust",
            "target": target,
            "payload": {"color_hex": suggested_color},
        }
        issue = self._make_issue(
            issue_id=issue_id,
            issue_type="contrast_low",
            severity="warning",
            message=(
                f"スライド '{slide.id}' の箇条書き '{bullet.id}' の文字色と背景色のコントラスト比"
                f" {ratio:.2f} が基準 {required_ratio:.2f} を下回っています"
            ),
            target=target,
            metrics={
                "color_hex": _normalize_hex(color_hex),
                "background_hex": _normalize_hex(self.options.background_color),
                "contrast_ratio": ratio,
                "required_ratio": required_ratio,
                "font_size_pt": font_size,
                "shape_name": paragraph.shape_name if paragraph else None,
            },
            fix=fix,
        )
        return issue, fix

    def _check_margins(
        self,
        slide: Slide,
        image: SlideImage,
        shape: ShapeSnapshot,
        *,
        slide_width_in: float,
        slide_height_in: float,
    ) -> tuple[dict[str, Any], dict[str, Any]] | None:
        left = shape.left_in
        top = shape.top_in
        width = shape.width_in
        height = shape.height_in
        margin = self.options.margin_in
        base_width = slide_width_in if slide_width_in > 0 else self.options.slide_width_in
        base_height = slide_height_in if slide_height_in > 0 else self.options.slide_height_in

        if base_width is None or base_height is None:
            return None

        violations: list[str] = []
        if left < margin:
            violations.append("left")
        if top < margin:
            violations.append("top")
        right = left + width
        bottom = top + height
        if right > base_width - margin:
            violations.append("right")
        if bottom > base_height - margin:
            violations.append("bottom")

        if not violations:
            return None

        issue_id = self._next_issue_id("margin", slide.id, image.id)
        target = {
            "slide_id": slide.id,
            "element_id": image.id,
            "element_type": "image",
        }

        target_left = min(
            max(left, margin),
            max(base_width - margin - width, margin),
        )
        target_top = min(
            max(top, margin),
            max(base_height - margin - height, margin),
        )

        fix_payload: dict[str, float] = {}
        if abs(target_left - left) > 1e-3:
            fix_payload["left_in"] = round(target_left, 3)
        if abs(target_top - top) > 1e-3:
            fix_payload["top_in"] = round(target_top, 3)

        fix = None
        if fix_payload:
            fix = {
                "id": f"fix-{issue_id}",
                "issue_id": issue_id,
                "type": "move",
                "target": target,
                "payload": fix_payload,
            }

        issue = self._make_issue(
            issue_id=issue_id,
            issue_type="margin",
            severity="warning",
            message=(
                f"スライド '{slide.id}' の画像 '{image.id}' が余白基準 {margin:.1f}in を外れています"
            ),
            target=target,
            metrics={
                "left_in": left,
                "top_in": top,
                "width_in": width,
                "height_in": height,
                "margin_in": margin,
                "slide_width_in": base_width,
                "slide_height_in": base_height,
                "violations": violations,
                "shape_name": shape.name,
            },
            fix=fix,
        )
        return issue, fix

    def _check_grid_alignment(
        self,
        slide: Slide,
        element_id: str,
        element_type: str,
        shape: ShapeSnapshot,
    ) -> tuple[dict[str, Any], dict[str, Any]] | None:
        grid = self.options.grid_size_in
        tolerance = self.options.grid_tolerance_in

        deviations = {
            "left": _grid_deviation(shape.left_in, grid),
            "top": _grid_deviation(shape.top_in, grid),
        }
        out_of_grid = {axis: dev for axis, dev in deviations.items() if dev > tolerance}
        if not out_of_grid:
            return None

        target = {
            "slide_id": slide.id,
            "element_id": element_id,
            "element_type": element_type,
        }
        issue_id = self._next_issue_id("grid_misaligned", slide.id, element_id)
        fix_payload: dict[str, float] = {}
        snapped_left = _snap_to_grid(shape.left_in, grid)
        snapped_top = _snap_to_grid(shape.top_in, grid)
        if "left" in out_of_grid:
            fix_payload["left_in"] = round(snapped_left, 3)
        if "top" in out_of_grid:
            fix_payload["top_in"] = round(snapped_top, 3)

        fix = None
        if fix_payload:
            fix = {
                "id": f"fix-{issue_id}",
                "issue_id": issue_id,
                "type": "move",
                "target": target,
                "payload": fix_payload,
            }

        issue = self._make_issue(
            issue_id=issue_id,
            issue_type="grid_misaligned",
            severity="warning",
            message=(
                f"スライド '{slide.id}' の要素 '{element_id}' がグリッド {grid:.3f}in に揃っていません"
            ),
            target=target,
            metrics={
                "left_in": shape.left_in,
                "top_in": shape.top_in,
                "width_in": shape.width_in,
                "height_in": shape.height_in,
                "grid_in": grid,
                "tolerance_in": tolerance,
                "deviations_in": out_of_grid,
                "shape_name": shape.name,
            },
            fix=fix,
        )
        return issue, fix

    def _next_issue_id(self, issue_type: str, slide_id: str, element_id: str | None) -> str:
        sequence = next(self._sequence)
        parts: list[str] = [issue_type, slide_id]
        if element_id:
            parts.append(element_id)
        parts.append(str(sequence))
        return "-".join(parts)

    def _make_issue(
        self,
        *,
        issue_id: str,
        issue_type: str,
        severity: str,
        message: str,
        target: dict[str, Any],
        metrics: dict[str, Any],
        fix: dict[str, Any] | None,
    ) -> dict[str, Any]:
        issue = {
            "id": issue_id,
            "type": issue_type,
            "severity": severity,
            "message": message,
            "target": target,
            "metrics": metrics,
        }
        if fix:
            issue["fix"] = fix
        return issue

    def _save(self, payload: dict[str, Any], workdir: Path) -> Path:
        workdir.mkdir(parents=True, exist_ok=True)
        output_path = workdir / self.options.output_filename
        output_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
        return output_path

    def _sync_mapping_log(self, context: PipelineContext, analysis: dict[str, Any]) -> None:
        mapping_log_ref = context.artifacts.get("mapping_log_path")
        if mapping_log_ref is None:
            return
        try:
            mapping_log_path = Path(str(mapping_log_ref))
        except Exception:  # noqa: BLE001
            return
        if not mapping_log_path.exists():
            return
        try:
            mapping_log = MappingLog.model_validate_json(
                mapping_log_path.read_text(encoding="utf-8")
            )
        except Exception as exc:  # noqa: BLE001
            logger.warning(
                "mapping_log.json の読み込みに失敗したため Analyzer 連携をスキップします: %s",
                exc,
            )
            return

        issues_payload = analysis.get("issues")
        if not isinstance(issues_payload, list):
            return

        overall_by_type: Counter[str] = Counter()
        overall_by_severity: Counter[str] = Counter()
        for slide in mapping_log.slides:
            slide_issues = [
                issue
                for issue in issues_payload
                if isinstance(issue, dict)
                and self._target_slide_id(issue) == slide.ref_id
            ]
            summary = self._build_analyzer_summary(slide_issues)
            slide.analyzer = summary
            overall_by_type.update(summary.issue_counts_by_type)
            overall_by_severity.update(summary.issue_counts_by_severity)

        mapping_log.meta.analyzer_issue_count = sum(overall_by_type.values())
        mapping_log.meta.analyzer_issue_counts_by_type = dict(overall_by_type)
        mapping_log.meta.analyzer_issue_counts_by_severity = dict(overall_by_severity)

        mapping_log_path.write_text(
            json.dumps(
                mapping_log.model_dump(mode="json"),
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )
        context.add_artifact("mapping_log", mapping_log)

    @staticmethod
    def _target_slide_id(issue: dict[str, Any]) -> str | None:
        target = issue.get("target")
        if not isinstance(target, dict):
            return None
        slide_id = target.get("slide_id")
        if isinstance(slide_id, str) and slide_id:
            return slide_id
        return None

    def _build_analyzer_summary(
        self, issues_payload: list[dict[str, Any]]
    ) -> MappingLogAnalyzerSummary:
        type_counter: Counter[str] = Counter()
        severity_counter: Counter[str] = Counter()
        summary_issues: list[MappingLogAnalyzerIssue] = []

        for issue in issues_payload:
            if not isinstance(issue, dict):
                continue

            issue_id = str(issue.get("id") or "")
            issue_type = str(issue.get("type") or "")
            severity = str(issue.get("severity") or "")
            message = str(issue.get("message") or "")

            target = issue.get("target")
            if not isinstance(target, dict):
                target = {}
            metrics = issue.get("metrics")
            if not isinstance(metrics, dict):
                metrics = {}

            fix = issue.get("fix")
            fix_type: str | None = None
            fix_payload: dict[str, Any] | None = None
            if isinstance(fix, dict):
                fix_type_value = fix.get("type")
                if isinstance(fix_type_value, str):
                    fix_type = fix_type_value
                payload_value = fix.get("payload")
                if isinstance(payload_value, dict):
                    fix_payload = payload_value

            summary_issues.append(
                MappingLogAnalyzerIssue(
                    issue_id=issue_id,
                    issue_type=issue_type,
                    severity=severity,
                    message=message,
                    target=target,
                    metrics=metrics,
                    fix_type=fix_type,
                    fix_payload=fix_payload,
                )
            )
            type_counter[issue_type] += 1
            severity_counter[severity] += 1

        return MappingLogAnalyzerSummary(
            issue_count=sum(type_counter.values()),
            issue_counts_by_type=dict(type_counter),
            issue_counts_by_severity=dict(severity_counter),
            issues=summary_issues,
        )

def _emu_to_inches(value: int) -> float:
    return value / EMU_PER_INCH


def _length_to_pt(length) -> float | None:
    if length is None:
        return None
    try:
        return float(length.pt)
    except AttributeError:
        return None


def _color_to_hex(color: ColorFormat | None) -> str | None:
    if color is None:
        return None
    try:
        rgb = color.rgb
    except AttributeError:
        return None
    if rgb is None:
        return None
    components = tuple(rgb)
    return "#" + "".join(f"{component:02X}" for component in components)


def _extract_font_info(paragraph) -> tuple[float | None, str | None]:
    font = paragraph.font
    size = _length_to_pt(getattr(font, "size", None))
    color = _color_to_hex(getattr(font, "color", None))

    if size is None or color is None:
        for run in paragraph.runs:
            run_font = run.font
            if size is None:
                size = _length_to_pt(getattr(run_font, "size", None))
            if color is None:
                color = _color_to_hex(getattr(run_font, "color", None))
            if size is not None and color is not None:
                break

    return size, color


def _normalize_hex(value: str) -> str:
    return value if value.startswith("#") else f"#{value}"


def _hex_to_rgb(value: str) -> tuple[float, float, float]:
    hex_value = _normalize_hex(value).lstrip("#")
    if len(hex_value) != 6:
        raise ValueError("hex color must be 6 characters")
    r = int(hex_value[0:2], 16)
    g = int(hex_value[2:4], 16)
    b = int(hex_value[4:6], 16)
    return r / 255.0, g / 255.0, b / 255.0


def _relative_luminance(rgb: tuple[float, float, float]) -> float:
    def linearize(channel: float) -> float:
        return channel / 12.92 if channel <= 0.03928 else ((channel + 0.055) / 1.055) ** 2.4

    r, g, b = (linearize(c) for c in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrast_ratio(foreground_hex: str, background_hex: str) -> float:
    fg_lum = _relative_luminance(_hex_to_rgb(foreground_hex))
    bg_lum = _relative_luminance(_hex_to_rgb(background_hex))
    lighter = max(fg_lum, bg_lum)
    darker = min(fg_lum, bg_lum)
    return (lighter + 0.05) / (darker + 0.05)


def _grid_deviation(value: float, grid_size: float) -> float:
    remainder = value % grid_size
    return min(remainder, grid_size - remainder)


def _snap_to_grid(value: float, grid_size: float) -> float:
    cells = round(value / grid_size)
    return cells * grid_size
