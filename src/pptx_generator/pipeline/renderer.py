"""PPTX 生成を担うステップ。"""

from __future__ import annotations

import logging
import os
import tempfile
import time
from dataclasses import dataclass
from pathlib import Path
from urllib.parse import urlparse
from urllib.request import urlopen

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..models import (
    ChartSeries,
    FontSpec,
    JobSpec,
    Slide,
    SlideBullet,
    SlideBulletGroup,
    SlideTextbox,
    TextboxParagraph,
)
from ..settings import BrandingConfig, BrandingFont, BoxSpec, ParagraphStyle
from .base import PipelineContext

logger = logging.getLogger(__name__)


@dataclass(slots=True)
class LayoutBox:
    left_in: float
    top_in: float
    width_in: float
    height_in: float

    def to_emu(self) -> tuple[int, int, int, int]:
        return (
            int(Inches(self.left_in)),
            int(Inches(self.top_in)),
            int(Inches(self.width_in)),
            int(Inches(self.height_in)),
        )


@dataclass(slots=True)
class AnchorResolution:
    shape: object | None
    left: int
    top: int
    width: int
    height: int
    is_placeholder: bool = False

    def as_box(self) -> tuple[int, int, int, int]:
        return self.left, self.top, self.width, self.height


@dataclass(slots=True)
class RenderingOptions:
    template_path: Path | None = None
    output_filename: str = "proposal.pptx"
    branding: BrandingConfig | None = None
    branding: BrandingConfig | None = None


class SimpleRendererStep:
    """最小機能の PPTX レンダラー。"""

    name = "renderer"

    def __init__(self, options: RenderingOptions | None = None) -> None:
        self.options = options or RenderingOptions()
        if self.options.branding is None:
            self.options.branding = BrandingConfig.default()
        self._branding: BrandingConfig = self.options.branding
        self._temp_files: list[Path] = []

    def run(self, context: PipelineContext) -> None:
        presentation = self._load_template()
        start = time.perf_counter()
        try:
            self._render_slides(presentation, context.spec)
            output_path = self._save(presentation, context.workdir)
            context.add_artifact("pptx_path", output_path)
            logger.info("PPTX を出力しました: %s", output_path)
        finally:
            self._cleanup_temp_files()

        elapsed_ms = int((time.perf_counter() - start) * 1000)
        context.add_artifact("renderer_stats", {"rendering_time_ms": elapsed_ms})

    def _load_template(self) -> Presentation:
        if self.options.template_path and self.options.template_path.exists():
            logger.debug("テンプレートを使用: %s", self.options.template_path)
            return Presentation(self.options.template_path)
        logger.debug("既定テンプレートを利用")
        return Presentation()

    def _render_slides(self, presentation: Presentation, spec: JobSpec) -> None:
        for slide_spec in spec.slides:
            layout = self._resolve_layout(presentation, slide_spec)
            slide = presentation.slides.add_slide(layout)
            self._apply_title(slide, slide_spec)
            self._apply_subtitle(slide, slide_spec)
            self._apply_bullets(slide, slide_spec)
            self._apply_textboxes(slide, slide_spec)
            self._apply_tables(slide, slide_spec)
            self._apply_images(slide, slide_spec)
            self._apply_charts(slide, slide_spec)
            self._apply_notes(slide, slide_spec)

    def _resolve_layout(self, presentation: Presentation, slide_spec: Slide):
        for layout in presentation.slide_layouts:
            if layout.name == slide_spec.layout:
                return layout
        logger.debug("レイアウト '%s' が見つからないため既定を使用", slide_spec.layout)
        try:
            return presentation.slide_layouts[1]
        except IndexError:
            if len(presentation.slide_layouts) == 0:
                raise RuntimeError("テンプレートに利用可能なレイアウトが存在しません")
            logger.warning(
                "テンプレートにレイアウト index=1 が存在しないため、index=0 を使用します"
            )
            return presentation.slide_layouts[0]

    def _apply_title(self, slide, slide_spec: Slide) -> None:
        if slide_spec.title is None:
            return
        title_shape = slide.shapes.title
        if title_shape is not None:
            text_frame = title_shape.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.text = slide_spec.title
            self._apply_brand_font(paragraph, self._branding.heading_font)
            return
        textbox = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.5), Inches(8.0), Inches(1.0)
        )
        text_frame = textbox.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = slide_spec.title
        self._apply_brand_font(paragraph, self._branding.heading_font)

    def _apply_bullets(self, slide, slide_spec: Slide) -> None:
        groups = slide_spec.bullets
        if not groups:
            return

        fallback_items: list[SlideBullet] = []
        used_anchors: set[str] = set()
        default_paragraph_style = self._branding.components.textbox.paragraph

        for group in groups:
            anchor_name = group.anchor
            if anchor_name:
                if anchor_name in used_anchors:
                    raise ValueError(
                        f"箇条書きのアンカー '{anchor_name}' が複数のグループで指定されています。"
                        "図形名はグループごとに一意にしてください。"
                    )
                used_anchors.add(anchor_name)
                paragraph_style = self._branding.resolve_layout_paragraph(
                    layout=slide_spec.layout,
                    placement_key=anchor_name,
                    default=default_paragraph_style,
                )
                self._render_bullet_group_to_anchor(
                    slide,
                    anchor_name,
                    group.items,
                    paragraph_style=paragraph_style,
                )
            else:
                fallback_items.extend(group.items)

        if fallback_items:
            target_shape = self._find_body_placeholder(slide)
            self._write_bullets_to_text_frame(
                target_shape.text_frame,
                fallback_items,
                paragraph_style=default_paragraph_style,
            )

    def _render_bullet_group_to_anchor(
        self,
        slide,
        anchor_name: str,
        bullets: list[SlideBullet],
        *,
        paragraph_style: ParagraphStyle,
    ) -> None:
        fallback_box = self._box_spec_to_layout_box(
            self._branding.components.textbox.fallback_box
        )
        text_frame = self._obtain_text_frame(
            slide=slide,
            anchor_name=anchor_name,
            fallback_box=fallback_box,
            strict_anchor=True,
        )

        if text_frame is None:
            raise ValueError(
                f"Shape with name '{anchor_name}' not found in slide. "
                "テンプレートの図形名を確認してください。"
            )
        self._write_bullets_to_text_frame(
            text_frame,
            bullets,
            paragraph_style=paragraph_style,
        )

    def _apply_subtitle(self, slide, slide_spec: Slide) -> None:
        if not slide_spec.subtitle:
            return

        subtitle_placeholder = None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                subtitle_placeholder = placeholder
                break

        if subtitle_placeholder is not None:
            text_frame = subtitle_placeholder.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.text = slide_spec.subtitle
            self._apply_brand_font(paragraph, self._branding.body_font)
            return

        textbox = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.5), Inches(8.0), Inches(1.0)
        )
        text_frame = textbox.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = slide_spec.subtitle
        self._apply_brand_font(paragraph, self._branding.body_font)

    def _apply_notes(self, slide, slide_spec: Slide) -> None:
        if not slide_spec.notes:
            return
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame

        text_frame.clear()
        lines = slide_spec.notes.splitlines() or [slide_spec.notes]
        for index, line in enumerate(lines):
            paragraph = (
                text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            )
            paragraph.text = line
            self._apply_brand_font(paragraph, self._branding.body_font)
            paragraph.level = 0

    def _apply_textboxes(self, slide, slide_spec: Slide) -> None:
        if not slide_spec.textboxes:
            return

        for textbox_spec in slide_spec.textboxes:
            fallback_box = self._resolve_textbox_fallback(slide_spec, textbox_spec)
            text_frame = self._obtain_text_frame(
                slide=slide,
                anchor_name=textbox_spec.anchor,
                fallback_box=fallback_box,
                strict_anchor=False,
            )
            if text_frame is None:
                msg = (
                    f"Shape with name '{textbox_spec.anchor}' not found in slide."
                    " テンプレートの図形名を確認してください。"
                )
                raise ValueError(msg)
            shape = getattr(text_frame, "_parent", None)
            if shape is not None:
                target_name = textbox_spec.anchor or textbox_spec.id
                if target_name:
                    try:
                        shape.name = target_name
                    except ValueError:
                        logger.debug("テキストボックス名 '%s' の設定に失敗", target_name, exc_info=True)
            self._write_textbox_content(slide_spec, textbox_spec, text_frame)

    def _write_bullets_to_text_frame(
        self,
        text_frame,
        bullets: list[SlideBullet],
        *,
        paragraph_style: ParagraphStyle,
    ) -> None:
        text_frame.clear()
        text_frame.word_wrap = True
        for index, bullet in enumerate(bullets):
            paragraph = (
                text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            )
            paragraph.text = bullet.text
            paragraph.level = bullet.level
            self._apply_font(
                paragraph,
                bullet.font,
                fallback=self._branding.body_font,
            )
            self._apply_paragraph_style(
                paragraph,
                None,
                fallback=paragraph_style,
                preserve_level=True,
            )

    def _write_textbox_content(
        self, slide_spec: Slide, textbox_spec: SlideTextbox, text_frame
    ) -> None:
        text_frame.clear()
        text_frame.word_wrap = True
        default_font = self._branding.resolve_layout_font(
            layout=slide_spec.layout,
            placement_key=textbox_spec.id,
            default=self._branding.components.textbox.font,
        )
        default_paragraph = self._branding.resolve_layout_paragraph(
            layout=slide_spec.layout,
            placement_key=textbox_spec.id,
            default=self._branding.components.textbox.paragraph,
        )

        lines = textbox_spec.text.splitlines() or [""]
        for index, line in enumerate(lines):
            paragraph = (
                text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            )
            paragraph.text = line
            self._apply_font(
                paragraph,
                textbox_spec.font,
                fallback=default_font,
            )
            self._apply_paragraph_style(
                paragraph,
                textbox_spec.paragraph,
                fallback=default_paragraph,
                preserve_level=False,
            )

    def _resolve_textbox_fallback(
        self, slide_spec: Slide, textbox_spec: SlideTextbox
    ) -> LayoutBox:
        position = textbox_spec.position
        if position is not None:
            return LayoutBox(
                position.left_in,
                position.top_in,
                position.width_in,
                position.height_in,
            )
        box_spec = self._branding.resolve_fallback_box(
            "textbox", layout=slide_spec.layout, placement_key=textbox_spec.id
        )
        return self._box_spec_to_layout_box(box_spec)

    def _apply_tables(self, slide, slide_spec: Slide) -> None:
        if not slide_spec.tables:
            return

        for table_spec in slide_spec.tables:
            header = list(table_spec.columns)
            rows = [list(row) for row in table_spec.rows]
            column_count = len(header) or (len(rows[0]) if rows else 0)
            if column_count == 0:
                logger.debug("テーブル '%s' に列情報がないためスキップ", table_spec.id)
                continue

            row_count = len(rows) + (1 if header else 0)
            if row_count == 0:
                logger.debug("テーブル '%s' にデータがないためスキップ", table_spec.id)
                continue

            fallback_box = self._box_spec_to_layout_box(
                self._branding.resolve_fallback_box(
                    "table", layout=slide_spec.layout, placement_key=table_spec.id
                )
            )
            resolution = self._resolve_anchor(slide, table_spec.anchor, fallback_box)
            anchor_shape = resolution.shape
            if resolution.is_placeholder:
                self._prepare_placeholder(anchor_shape)
            left, top, width, height = resolution.as_box()
            table_shape = slide.shapes.add_table(
                row_count, column_count, left, top, width, height
            )
            table = table_shape.table

            total_width = table_shape.width
            for idx in range(column_count):
                table.columns[idx].width = total_width // column_count

            start_row = 0
            if header:
                self._fill_table_row(
                    table.rows[0],
                    header,
                    is_header=True,
                    table_spec=table_spec,
                )
                start_row = 1

            for offset, row_values in enumerate(rows):
                target_row = table.rows[start_row + offset]
                padded = row_values + [""] * (column_count - len(row_values))
                self._fill_table_row(
                    target_row,
                    padded,
                    is_header=False,
                    table_spec=table_spec,
                    zebra_index=offset,
                )

            if anchor_shape is not None:
                self._remove_shape(anchor_shape)

    def _apply_images(self, slide, slide_spec: Slide) -> None:
        if not slide_spec.images:
            return

        for image_spec in slide_spec.images:
            fallback_box = self._box_spec_to_layout_box(
                self._branding.resolve_fallback_box(
                    "image", layout=slide_spec.layout, placement_key=image_spec.id
                )
            )
            resolution = self._resolve_anchor(slide, image_spec.anchor, fallback_box)
            anchor_shape = resolution.shape
            if resolution.is_placeholder:
                self._prepare_placeholder(anchor_shape)
            image_path = self._resolve_image_source(image_spec.source, image_spec.id)
            left = self._override_emu(resolution.left, image_spec.left_in)
            top = self._override_emu(resolution.top, image_spec.top_in)
            target_width = self._override_emu(resolution.width, image_spec.width_in)
            target_height = self._override_emu(resolution.height, image_spec.height_in)
            picture = slide.shapes.add_picture(str(image_path), left, top)
            self._resize_picture(
                picture, target_width, target_height, image_spec.sizing
            )
            target_name = image_spec.anchor or image_spec.id
            if target_name:
                try:
                    picture.name = target_name
                except ValueError:
                    logger.debug("画像図形名 '%s' の設定に失敗", target_name, exc_info=True)

            if anchor_shape is not None:
                self._remove_shape(anchor_shape)

    def _apply_charts(self, slide, slide_spec: Slide) -> None:
        if not slide_spec.charts:
            return

        for chart_spec in slide_spec.charts:
            if not chart_spec.series:
                logger.debug("チャート '%s' に系列がないためスキップ", chart_spec.id)
                continue

            data = CategoryChartData()
            categories = chart_spec.categories or [
                str(index + 1) for index in range(len(chart_spec.series[0].values))
            ]
            data.categories = categories
            for series in chart_spec.series:
                data.add_series(series.name, series.values)

            chart_type = self._resolve_chart_type(chart_spec.type)
            fallback_box = self._box_spec_to_layout_box(
                self._branding.resolve_fallback_box(
                    "chart", layout=slide_spec.layout, placement_key=chart_spec.id
                )
            )
            resolution = self._resolve_anchor(slide, chart_spec.anchor, fallback_box)
            anchor_shape = resolution.shape
            if resolution.is_placeholder:
                self._prepare_placeholder(anchor_shape)
            left, top, width, height = resolution.as_box()
            chart_shape = slide.shapes.add_chart(
                chart_type, left, top, width, height, data
            )
            chart = chart_shape.chart

            self._apply_chart_series_colors(chart.series, chart_spec.series)
            self._style_chart(chart, chart_spec.options)

            if anchor_shape is not None:
                self._remove_shape(anchor_shape)

    def _find_body_placeholder(self, slide):
        for shape in slide.placeholders:
            if shape.placeholder_format.type in {
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.VERTICAL_BODY,
                PP_PLACEHOLDER.OBJECT,
            }:
                return shape
        logger.debug("本文用プレースホルダがないためテキストボックスを追加")
        return slide.shapes.add_textbox(
            Inches(1.0), Inches(1.5), Inches(8.0), Inches(4.5)
        )

    def _apply_font(
        self,
        paragraph,
        font_spec: FontSpec | BrandingFont | None,
        *,
        fallback: BrandingFont,
    ) -> None:
        if isinstance(font_spec, BrandingFont):
            self._apply_brand_font(paragraph, font_spec)
            return
        if isinstance(font_spec, FontSpec):
            font = paragraph.font
            font.name = font_spec.name
            font.size = Pt(font_spec.size_pt)
            font.bold = font_spec.bold
            font.italic = font_spec.italic
            font.color.rgb = RGBColor.from_string(font_spec.color_hex.lstrip("#"))
            return
        self._apply_brand_font(paragraph, fallback)

    def _apply_paragraph_style(
        self,
        paragraph,
        paragraph_spec: TextboxParagraph | None,
        *,
        fallback: ParagraphStyle,
        preserve_level: bool,
    ) -> None:
        level = (
            paragraph_spec.level
            if paragraph_spec and paragraph_spec.level is not None
            else fallback.level
        )
        if not preserve_level or paragraph.level is None:
            paragraph.level = level if level is not None else (paragraph.level or 0)

        align = (
            paragraph_spec.align
            if paragraph_spec and paragraph_spec.align
            else fallback.align
        )
        if align:
            paragraph.alignment = self._resolve_alignment(align)

        line_spacing = (
            paragraph_spec.line_spacing_pt
            if paragraph_spec and paragraph_spec.line_spacing_pt is not None
            else fallback.line_spacing_pt
        )
        if line_spacing is not None:
            paragraph.line_spacing = Pt(line_spacing)

        space_before = (
            paragraph_spec.space_before_pt
            if paragraph_spec and paragraph_spec.space_before_pt is not None
            else fallback.space_before_pt
        )
        if space_before is not None:
            paragraph.space_before = Pt(space_before)

        space_after = (
            paragraph_spec.space_after_pt
            if paragraph_spec and paragraph_spec.space_after_pt is not None
            else fallback.space_after_pt
        )
        if space_after is not None:
            paragraph.space_after = Pt(space_after)

        paragraph_properties = paragraph._p.get_or_add_pPr()
        left_indent = (
            paragraph_spec.left_indent_in
            if paragraph_spec and paragraph_spec.left_indent_in is not None
            else fallback.left_indent_in
        )
        if left_indent is not None:
            paragraph_properties.set("marL", str(int(Inches(left_indent))))

        right_indent = (
            paragraph_spec.right_indent_in
            if paragraph_spec and paragraph_spec.right_indent_in is not None
            else fallback.right_indent_in
        )
        if right_indent is not None:
            paragraph_properties.set("marR", str(int(Inches(right_indent))))

        first_line_indent = (
            paragraph_spec.first_line_indent_in
            if paragraph_spec and paragraph_spec.first_line_indent_in is not None
            else fallback.first_line_indent_in
        )
        if first_line_indent is not None:
            paragraph_properties.set("indent", str(int(Inches(first_line_indent))))

    def _resolve_alignment(self, align: str) -> PP_ALIGN:
        mapping = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
            "distributed": PP_ALIGN.DISTRIBUTE,
        }
        return mapping.get(align.lower(), PP_ALIGN.LEFT)

    def _box_spec_to_layout_box(self, spec: BoxSpec) -> LayoutBox:
        return LayoutBox(spec.left_in, spec.top_in, spec.width_in, spec.height_in)

    def _save(self, presentation: Presentation, workdir: Path) -> Path:
        workdir.mkdir(parents=True, exist_ok=True)
        output_path = workdir / self.options.output_filename
        presentation.save(output_path)
        return output_path

    def _apply_brand_font(self, paragraph, branding_font: BrandingFont) -> None:
        font = paragraph.font
        font.name = branding_font.name
        font.size = Pt(branding_font.size_pt)
        font.color.rgb = RGBColor.from_string(branding_font.color_hex.lstrip("#"))
        font.bold = branding_font.bold
        font.italic = branding_font.italic

    def _find_shape_by_name(self, slide, name: str):
        for shape in slide.shapes:
            if shape.name == name:
                return shape
        return None

    def _find_placeholder_by_name(self, slide, name: str):
        layout = getattr(slide, "slide_layout", None)
        if layout is None:
            return None
        target_idx: int | None = None
        for layout_shape in layout.shapes:
            if (
                getattr(layout_shape, "is_placeholder", False)
                and layout_shape.name == name
            ):
                target_idx = layout_shape.placeholder_format.idx
                break
        if target_idx is None:
            return None
        for shape in slide.shapes:
            if (
                getattr(shape, "is_placeholder", False)
                and shape.placeholder_format.idx == target_idx
            ):
                return shape
        return None

    def _resolve_anchor(
        self, slide, anchor: str | None, fallback_box: LayoutBox
    ) -> AnchorResolution:
        if anchor:
            shape = self._find_shape_by_name(slide, anchor)
            if shape is not None:
                return AnchorResolution(
                    shape,
                    int(shape.left),
                    int(shape.top),
                    int(shape.width),
                    int(shape.height),
                    getattr(shape, "is_placeholder", False),
                )
            placeholder = self._find_placeholder_by_name(slide, anchor)
            if placeholder is not None:
                return AnchorResolution(
                    placeholder,
                    int(placeholder.left),
                    int(placeholder.top),
                    int(placeholder.width),
                    int(placeholder.height),
                    True,
                )
        left, top, width, height = fallback_box.to_emu()
        return AnchorResolution(None, left, top, width, height)

    def _obtain_text_frame(
        self,
        *,
        slide,
        anchor_name: str | None,
        fallback_box: LayoutBox,
        strict_anchor: bool,
    ):
        resolution = self._resolve_anchor(slide, anchor_name, fallback_box)
        shape = resolution.shape

        if strict_anchor and anchor_name and shape is None:
            return None

        if shape is not None and getattr(shape, "has_text_frame", False):
            if resolution.is_placeholder:
                self._prepare_placeholder(shape)
            else:
                shape.text_frame.clear()

        left, top, width, height = resolution.as_box()
        text_shape = slide.shapes.add_textbox(left, top, width, height)
        if shape is not None:
            self._remove_shape(shape)
        if anchor_name:
            try:
                text_shape.name = anchor_name
            except ValueError:
                logger.debug("アンカー名 '%s' の再設定に失敗", anchor_name, exc_info=True)
        return text_shape.text_frame

    def _prepare_placeholder(self, placeholder) -> None:
        if placeholder is None:
            return
        try:
            if getattr(placeholder, "has_text_frame", False):
                placeholder.text_frame.clear()
        except Exception:  # noqa: BLE001
            logger.debug("プレースホルダーの初期化に失敗", exc_info=True)

    def _override_emu(self, default: int, value_in: float | None) -> int:
        if value_in is None:
            return default
        return int(Inches(value_in))

    def _fill_table_row(
        self,
        row,
        values: list[object],
        *,
        is_header: bool,
        table_spec: SlideTable,
        zebra_index: int | None = None,
    ) -> None:
        table_style = self._branding.components.table
        spec_style = table_spec.style

        for idx, value in enumerate(values):
            cell = row.cells[idx]
            text_frame = cell.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.text = str(value)
            if is_header:
                self._apply_font(
                    paragraph,
                    None,
                    fallback=table_style.header.font,
                )
                header_color = (
                    spec_style.header_fill if spec_style and spec_style.header_fill else table_style.header.fill_color
                )
                fill_color = header_color
            else:
                self._apply_font(
                    paragraph,
                    None,
                    fallback=table_style.body.font,
                )
                zebra_enabled = bool(spec_style and spec_style.zebra)
                use_zebra = (
                    zebra_enabled
                    and table_style.body.zebra_fill_color
                    and zebra_index is not None
                    and zebra_index % 2 == 1
                )
                fill_color = (
                    table_style.body.zebra_fill_color
                    if use_zebra
                    else table_style.body.fill_color
                )

            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip("#"))

    def _resolve_image_source(self, source: str, image_id: str) -> Path:
        parsed = urlparse(str(source))
        if parsed.scheme in {"http", "https"}:
            return self._download_remote_image(str(source), image_id)

        path = Path(source).expanduser()
        if not path.is_absolute():
            path = Path.cwd() / path
        if not path.exists():
            msg = f"画像ファイルが見つかりません: {source}"
            raise FileNotFoundError(msg)
        return path

    def _download_remote_image(self, url: str, image_id: str) -> Path:
        tmp = tempfile.NamedTemporaryFile(
            delete=False, suffix=Path(url).suffix or ".img"
        )
        try:
            with urlopen(url) as response:
                tmp.write(response.read())
        finally:
            tmp.close()
        path = Path(tmp.name)
        self._temp_files.append(path)
        logger.info("画像をダウンロードしました: id=%s, path=%s", image_id, path)
        return path

    def _resize_picture(self, picture, width: int, height: int, sizing: str) -> None:
        if width is None and height is None:
            return

        target_width = width or picture.width
        target_height = height or picture.height

        if sizing == "stretch":
            picture.width = target_width
            picture.height = target_height
            return

        original_width = picture.width
        original_height = picture.height
        if original_width == 0 or original_height == 0:
            picture.width = target_width
            picture.height = target_height
            return

        width_ratio = target_width / original_width
        height_ratio = target_height / original_height

        if sizing == "fill":
            scale = max(width_ratio, height_ratio)
        else:  # fit
            scale = min(width_ratio, height_ratio)

        picture.width = int(original_width * scale)
        picture.height = int(original_height * scale)

        if width is not None:
            picture.left += (target_width - picture.width) // 2
        if height is not None:
            picture.top += (target_height - picture.height) // 2

    def _apply_chart_series_colors(
        self, chart_series, series_specs: list[ChartSeries]
    ) -> None:
        palette = self._branding.components.chart.palette
        if not palette:
            palette = (
                self._branding.accent_color,
                self._branding.primary_color,
                self._branding.secondary_color,
            )
        for index, (series, spec) in enumerate(
            zip(chart_series, series_specs, strict=False)
        ):
            fill = series.format.fill
            fill.solid()
            color = spec.color_hex or palette[index % len(palette)]
            fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))

    def _style_chart(self, chart, options) -> None:
        component = self._branding.components.chart
        labels_enabled = component.data_labels.enabled
        labels_format = component.data_labels.format

        if options is not None:
            labels_enabled = options.data_labels
            labels_format = options.y_axis_format or labels_format

        for plot in chart.plots:
            plot.has_data_labels = bool(labels_enabled)
            if plot.has_data_labels:
                data_labels = plot.data_labels
                data_labels.show_value = True
                if labels_format:
                    data_labels.number_format = labels_format

        if labels_format and hasattr(chart, "value_axis"):
            chart.value_axis.tick_labels.number_format = labels_format

        if getattr(chart, "has_legend", False):
            if chart.has_legend:
                chart.legend.include_in_layout = False

        axis_font = component.axis.font
        if hasattr(chart, "category_axis"):
            font = chart.category_axis.tick_labels.font
            font.name = axis_font.name
            font.size = Pt(axis_font.size_pt)
            font.color.rgb = RGBColor.from_string(axis_font.color_hex.lstrip("#"))
            font.bold = axis_font.bold
            font.italic = axis_font.italic
        if hasattr(chart, "value_axis"):
            font = chart.value_axis.tick_labels.font
            font.name = axis_font.name
            font.size = Pt(axis_font.size_pt)
            font.color.rgb = RGBColor.from_string(axis_font.color_hex.lstrip("#"))
            font.bold = axis_font.bold
            font.italic = axis_font.italic

    def _resolve_chart_type(self, chart_type: str) -> XL_CHART_TYPE:
        mapping = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }
        return mapping.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)

    def _remove_shape(self, shape) -> None:
        if shape is None:
            return
        try:
            shape.element.getparent().remove(shape.element)
        except Exception:  # noqa: BLE001
            logger.debug(
                "shape の削除に失敗しました: %s",
                shape.name if hasattr(shape, "name") else shape,
            )

    def _cleanup_temp_files(self) -> None:
        for path in self._temp_files:
            try:
                os.unlink(path)
            except OSError:
                logger.debug("一時ファイル削除に失敗: %s", path)
        self._temp_files.clear()
