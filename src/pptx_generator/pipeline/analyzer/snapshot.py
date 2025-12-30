from __future__ import annotations

from dataclasses import dataclass, field
from typing import Iterator

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from . import utils


@dataclass(slots=True)
class ParagraphSnapshot:
    shape_id: int
    shape_name: str | None
    shape_type: int
    paragraph_index: int
    text: str
    level: int
    font_size_pt: float | None
    color_hex: str | None
    font_name: str | None = None
    bold: bool | None = None
    italic: bool | None = None
    alignment: str | None = None
    line_spacing_pt: float | None = None
    space_before_pt: float | None = None
    space_after_pt: float | None = None
    left_indent_in: float | None = None
    right_indent_in: float | None = None
    first_line_indent_in: float | None = None


@dataclass(slots=True)
class ShapeSnapshot:
    shape_id: int
    name: str | None
    shape_type: int
    left_in: float
    top_in: float
    width_in: float
    height_in: float
    paragraphs: list[ParagraphSnapshot] = field(default_factory=list)
    parent_shape_id: int | None = None
    is_placeholder: bool = False
    placeholder_type: int | None = None
    placeholder_index: int | None = None
    z_order: int | None = None
    rotation_deg: float | None = None
    text_frame_padding: dict[str, float | None] | None = None
    text_frame_word_wrap: bool | None = None
    text_frame_vertical_anchor: str | None = None
    text_frame_auto_size: str | None = None
    table_cell: dict[str, int] | None = None


@dataclass(slots=True)
class SlideSnapshot:
    index: int
    shapes: list[ShapeSnapshot]
    body_placeholder_id: int | None = None

    @classmethod
    def from_slide(cls, slide, index: int) -> "SlideSnapshot":
        shapes: list[ShapeSnapshot] = []
        body_placeholder_id: int | None = None

        for snapshot in _collect_shape_snapshots(slide.shapes):
            shapes.append(snapshot)
            if _is_body_placeholder(snapshot) and body_placeholder_id is None:
                body_placeholder_id = snapshot.shape_id

        return cls(index=index, shapes=shapes, body_placeholder_id=body_placeholder_id)

    def shape_by_id(self, shape_id: int) -> ShapeSnapshot | None:
        for shape in self.shapes:
            if shape.shape_id == shape_id:
                return shape
        return None

    def find_shape_by_name(self, name: str, *, shape_type: int | None = None) -> ShapeSnapshot | None:
        if not name:
            return None
        for shape in self.shapes:
            if shape.name != name:
                continue
            if shape_type is not None and shape.shape_type != shape_type:
                continue
            return shape
        return None

    def body_paragraphs(self) -> list[ParagraphSnapshot]:
        if self.body_placeholder_id is None:
            return []
        shape = self.shape_by_id(self.body_placeholder_id)
        return list(shape.paragraphs) if shape is not None else []


class BulletParagraphResolver:
    """スライド内の段落を箇条書きグループに対応づける。"""

    def __init__(self, snapshot: SlideSnapshot):
        self._anchor_iters: dict[str, Iterator[ParagraphSnapshot]] = {}
        for shape in snapshot.shapes:
            if not shape.name or not shape.paragraphs:
                continue
            self._anchor_iters.setdefault(shape.name, iter(shape.paragraphs))
        self._fallback_iter: Iterator[ParagraphSnapshot] = iter(snapshot.body_paragraphs())

    def resolve(self, anchor: str | None) -> ParagraphSnapshot | None:
        if anchor:
            iterator = self._anchor_iters.get(anchor)
            if iterator is None:
                return None
            return next(iterator, None)
        return next(self._fallback_iter, None)


def _collect_shape_snapshots(
    shapes, *, parent_shape_id: int | None = None, left_offset_in: float = 0.0, top_offset_in: float = 0.0
) -> list[ShapeSnapshot]:
    snapshots: list[ShapeSnapshot] = []
    for shape in shapes:
        snapshot = _build_shape_snapshot(
            shape,
            parent_shape_id=parent_shape_id,
            left_offset_in=left_offset_in,
            top_offset_in=top_offset_in,
        )
        snapshots.append(snapshot)

        if getattr(shape, "has_table", False):
            snapshots.extend(
                _build_table_cell_snapshots(
                    shape,
                    parent_shape_id=snapshot.shape_id,
                    left_offset_in=left_offset_in,
                    top_offset_in=top_offset_in,
                )
            )

        is_group = int(getattr(shape, "shape_type", MSO_SHAPE_TYPE.AUTO_SHAPE)) == int(MSO_SHAPE_TYPE.GROUP)
        if is_group:
            group_left, group_top, _, _ = _extract_geometry(shape, left_offset_in, top_offset_in)
            snapshots.extend(
                _collect_shape_snapshots(
                    getattr(shape, "shapes", []),
                    parent_shape_id=snapshot.shape_id,
                    left_offset_in=group_left,
                    top_offset_in=group_top,
                )
            )
    return snapshots


def _build_shape_snapshot(
    shape,
    *,
    parent_shape_id: int | None = None,
    left_offset_in: float = 0.0,
    top_offset_in: float = 0.0,
    table_cell: dict[str, int] | None = None,
    override_shape_id: int | None = None,
) -> ShapeSnapshot:
    shape_id = override_shape_id if override_shape_id is not None else _shape_id(shape)
    left_in, top_in, width_in, height_in = _extract_geometry(shape, left_offset_in, top_offset_in)
    shape_name = getattr(shape, "name", None)
    shape_type = int(getattr(shape, "shape_type", MSO_SHAPE_TYPE.AUTO_SHAPE))
    is_placeholder = bool(getattr(shape, "is_placeholder", False))
    placeholder_type, placeholder_index = _extract_placeholder(shape, is_placeholder)
    z_order = _extract_z_order(shape)
    rotation = _extract_rotation(shape)
    text_frame = _get_text_frame(shape)
    paragraphs = _build_paragraph_snapshots(shape_id, shape_name, shape_type, text_frame)
    (
        text_frame_padding,
        text_frame_word_wrap,
        text_frame_vertical_anchor,
        text_frame_auto_size,
    ) = _extract_text_frame_metadata(text_frame)

    return ShapeSnapshot(
        shape_id=shape_id,
        name=shape_name,
        shape_type=shape_type,
        left_in=left_in,
        top_in=top_in,
        width_in=width_in,
        height_in=height_in,
        paragraphs=paragraphs,
        parent_shape_id=parent_shape_id,
        is_placeholder=is_placeholder,
        placeholder_type=placeholder_type,
        placeholder_index=placeholder_index,
        z_order=z_order,
        rotation_deg=rotation,
        text_frame_padding=text_frame_padding,
        text_frame_word_wrap=text_frame_word_wrap,
        text_frame_vertical_anchor=text_frame_vertical_anchor,
        text_frame_auto_size=text_frame_auto_size,
        table_cell=table_cell,
    )


def _shape_id(shape) -> int:
    return getattr(shape, "shape_id", id(shape))


def _extract_geometry(shape, left_offset_in: float = 0.0, top_offset_in: float = 0.0) -> tuple[float, float, float, float]:
    left = left_offset_in + utils.emu_to_inches(int(getattr(shape, "left", 0)))
    top = top_offset_in + utils.emu_to_inches(int(getattr(shape, "top", 0)))
    width = utils.emu_to_inches(int(getattr(shape, "width", 0)))
    height = utils.emu_to_inches(int(getattr(shape, "height", 0)))
    return left, top, width, height


def _extract_placeholder(shape, is_placeholder: bool) -> tuple[int | None, int | None]:
    if not is_placeholder:
        return None, None
    try:
        placeholder_type = int(shape.placeholder_format.type)
        placeholder_index = int(getattr(shape.placeholder_format, "idx", 0))
        return placeholder_type, placeholder_index
    except Exception:  # noqa: BLE001
        return None, None


def _extract_z_order(shape) -> int | None:
    z_order = getattr(shape, "z_order_position", None)
    if z_order is None:
        return None
    try:
        return int(z_order)
    except (TypeError, ValueError):
        return None


def _extract_rotation(shape) -> float | None:
    rotation = getattr(shape, "rotation", None)
    if rotation is None:
        return None
    try:
        return float(rotation)
    except (TypeError, ValueError):
        return None


def _get_text_frame(shape):
    if getattr(shape, "has_text_frame", False):
        return getattr(shape, "text_frame", None)
    if hasattr(shape, "text_frame"):
        return getattr(shape, "text_frame", None)
    return None


def _build_paragraph_snapshots(
    shape_id: int,
    shape_name: str | None,
    shape_type: int,
    text_frame,
) -> list[ParagraphSnapshot]:
    if text_frame is None:
        return []
    paragraphs: list[ParagraphSnapshot] = []
    for idx, paragraph in enumerate(text_frame.paragraphs):
        text = paragraph.text or ""
        (
            font_size_pt,
            color_hex,
            font_name,
            bold,
            italic,
        ) = utils.extract_font_info(paragraph)
        paragraph_style = utils.extract_paragraph_style(paragraph)
        level = paragraph.level if paragraph.level is not None else 0
        paragraphs.append(
            ParagraphSnapshot(
                shape_id=shape_id,
                shape_name=shape_name,
                shape_type=shape_type,
                paragraph_index=idx,
                text=text,
                level=level,
                font_size_pt=font_size_pt,
                color_hex=color_hex,
                font_name=font_name,
                bold=bold,
                italic=italic,
                alignment=paragraph_style.get("alignment"),
                line_spacing_pt=paragraph_style.get("line_spacing_pt"),
                space_before_pt=paragraph_style.get("space_before_pt"),
                space_after_pt=paragraph_style.get("space_after_pt"),
                left_indent_in=paragraph_style.get("left_indent_in"),
                right_indent_in=paragraph_style.get("right_indent_in"),
                first_line_indent_in=paragraph_style.get("first_line_indent_in"),
            )
        )
    return paragraphs


def _extract_text_frame_metadata(
    text_frame,
) -> tuple[
    dict[str, float | None] | None,
    bool | None,
    str | None,
    str | None,
]:
    padding = utils.extract_text_frame_padding(text_frame)
    if text_frame is None:
        return padding, None, None, None

    word_wrap_value = getattr(text_frame, "word_wrap", None)
    word_wrap = bool(word_wrap_value) if word_wrap_value is not None else None
    vertical_anchor = utils.enum_name(getattr(text_frame, "vertical_anchor", None))
    auto_size = utils.enum_name(getattr(text_frame, "auto_size", None))
    return padding, word_wrap, vertical_anchor, auto_size


def _is_body_placeholder(shape: ShapeSnapshot) -> bool:
    if not shape.is_placeholder or shape.placeholder_type is None:
        return False
    body_placeholders = {
        int(PP_PLACEHOLDER.BODY),
        int(PP_PLACEHOLDER.VERTICAL_BODY),
        int(PP_PLACEHOLDER.OBJECT),
    }
    return shape.placeholder_type in body_placeholders


__all__ = [
    "ParagraphSnapshot",
    "ShapeSnapshot",
    "SlideSnapshot",
    "BulletParagraphResolver",
]


def _build_table_cell_snapshots(table_shape, *, parent_shape_id: int, left_offset_in: float, top_offset_in: float) -> list[ShapeSnapshot]:
    snapshots: list[ShapeSnapshot] = []
    table = getattr(table_shape, "table", None)
    if table is None:
        return snapshots

    for row_idx, row in enumerate(getattr(table, "rows", [])):
        for col_idx, cell in enumerate(getattr(row, "cells", [])):
            shape_id = _table_cell_shape_id(parent_shape_id, row_idx, col_idx)
            text_frame = getattr(cell, "text_frame", None)
            paragraphs = _build_paragraph_snapshots(
                shape_id,
                getattr(table_shape, "name", None),
                int(getattr(table_shape, "shape_type", MSO_SHAPE_TYPE.TABLE)),
                text_frame,
            )
            padding, word_wrap, vertical_anchor, auto_size = _extract_text_frame_metadata(text_frame)
            left_in, top_in, width_in, height_in = _extract_geometry(cell, left_offset_in, top_offset_in)

            snapshots.append(
                ShapeSnapshot(
                    shape_id=shape_id,
                    name=f"{getattr(table_shape, 'name', '')}:r{row_idx}c{col_idx}",
                    shape_type=int(MSO_SHAPE_TYPE.TABLE),
                    left_in=left_in,
                    top_in=top_in,
                    width_in=width_in,
                    height_in=height_in,
                    paragraphs=paragraphs,
                    parent_shape_id=parent_shape_id,
                    is_placeholder=False,
                    placeholder_type=None,
                    placeholder_index=None,
                    z_order=_extract_z_order(cell),
                    rotation_deg=_extract_rotation(cell),
                    text_frame_padding=padding,
                    text_frame_word_wrap=word_wrap,
                    text_frame_vertical_anchor=vertical_anchor,
                    text_frame_auto_size=auto_size,
                    table_cell={"row": row_idx, "col": col_idx},
                )
            )
    return snapshots


def _table_cell_shape_id(parent_shape_id: int, row_idx: int, col_idx: int) -> int:
    return parent_shape_id * 10000 + row_idx * 100 + col_idx
