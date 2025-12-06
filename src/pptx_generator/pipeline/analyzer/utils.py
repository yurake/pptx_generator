from __future__ import annotations

import math
from typing import Tuple

from pptx.dml.color import ColorFormat

EMU_PER_INCH = 914400


def emu_to_inches(value: int) -> float:
    return value / EMU_PER_INCH


def length_to_inches(length) -> float | None:
    if length is None:
        return None
    try:
        return float(length.inches)
    except AttributeError:
        try:
            return float(length) / EMU_PER_INCH
        except (TypeError, ValueError, ZeroDivisionError):
            return None


def length_to_pt(length) -> float | None:
    if length is None:
        return None
    try:
        return float(length.pt)
    except AttributeError:
        try:
            return float(length)
        except (TypeError, ValueError):
            return None


def color_to_hex(color: ColorFormat | None) -> str | None:
    if color is None:
        return None
    try:
        rgb = color.rgb
    except AttributeError:
        return None
    if rgb is None:
        return None
    components = tuple(rgb)
    return "#" + "".join(f"{component:02X}" for component in components)


def extract_font_info(paragraph) -> tuple[float | None, str | None, str | None, bool | None, bool | None]:
    font = paragraph.font
    size = length_to_pt(getattr(font, "size", None))
    color = color_to_hex(getattr(font, "color", None))
    name = getattr(font, "name", None)
    bold = getattr(font, "bold", None)
    italic = getattr(font, "italic", None)

    if size is None or color is None or name is None or bold is None or italic is None:
        for run in paragraph.runs:
            run_font = run.font
            if size is None:
                size = length_to_pt(getattr(run_font, "size", None))
            if color is None:
                color = color_to_hex(getattr(run_font, "color", None))
            if name is None:
                name = getattr(run_font, "name", None)
            if bold is None:
                bold = getattr(run_font, "bold", None)
            if italic is None:
                italic = getattr(run_font, "italic", None)
            if all(value is not None for value in (size, color, name, bold, italic)):
                break

    return size, color, name, bold, italic


def enum_name(value) -> str | None:
    if value is None:
        return None
    name = getattr(value, "name", None)
    if isinstance(name, str) and name:
        return name.lower()
    try:
        text = str(value)
    except Exception:  # noqa: BLE001
        return None
    text = text.strip()
    return text.lower() if text else None


def extract_paragraph_style(paragraph) -> dict[str, float | str | None]:
    paragraph_format = getattr(paragraph, "paragraph_format", None)
    line_spacing_pt = None
    space_before_pt = None
    space_after_pt = None
    left_indent_in = None
    right_indent_in = None
    first_line_indent_in = None
    if paragraph_format is not None:
        line_spacing_pt = length_to_pt(getattr(paragraph_format, "line_spacing", None))
        space_before_pt = length_to_pt(getattr(paragraph_format, "space_before", None))
        space_after_pt = length_to_pt(getattr(paragraph_format, "space_after", None))
        left_indent_in = length_to_inches(getattr(paragraph_format, "left_indent", None))
        right_indent_in = length_to_inches(getattr(paragraph_format, "right_indent", None))
        first_line_indent_in = length_to_inches(getattr(paragraph_format, "first_line_indent", None))

    return {
        "alignment": enum_name(getattr(paragraph, "alignment", None)),
        "line_spacing_pt": line_spacing_pt,
        "space_before_pt": space_before_pt,
        "space_after_pt": space_after_pt,
        "left_indent_in": left_indent_in,
        "right_indent_in": right_indent_in,
        "first_line_indent_in": first_line_indent_in,
    }


def extract_text_frame_padding(text_frame) -> dict[str, float | None] | None:
    if text_frame is None:
        return None
    padding = {
        "left_in": length_to_inches(getattr(text_frame, "margin_left", None)),
        "right_in": length_to_inches(getattr(text_frame, "margin_right", None)),
        "top_in": length_to_inches(getattr(text_frame, "margin_top", None)),
        "bottom_in": length_to_inches(getattr(text_frame, "margin_bottom", None)),
    }
    if all(value is None for value in padding.values()):
        return None
    return padding


def normalize_hex(value: str) -> str:
    return value if value.startswith("#") else f"#{value}"


def hex_to_rgb(value: str) -> tuple[float, float, float]:
    hex_value = normalize_hex(value).lstrip("#")
    if len(hex_value) != 6:
        raise ValueError("hex color must be 6 characters")
    r = int(hex_value[0:2], 16)
    g = int(hex_value[2:4], 16)
    b = int(hex_value[4:6], 16)
    return r / 255.0, g / 255.0, b / 255.0


def relative_luminance(rgb: Tuple[float, float, float]) -> float:
    def linearize(channel: float) -> float:
        return channel / 12.92 if channel <= 0.03928 else ((channel + 0.055) / 1.055) ** 2.4

    r, g, b = (linearize(c) for c in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def contrast_ratio(foreground_hex: str, background_hex: str) -> float:
    fg_lum = relative_luminance(hex_to_rgb(foreground_hex))
    bg_lum = relative_luminance(hex_to_rgb(background_hex))
    lighter = max(fg_lum, bg_lum)
    darker = min(fg_lum, bg_lum)
    return (lighter + 0.05) / (darker + 0.05)


def grid_deviation(value: float, grid_size: float) -> float:
    remainder = math.fmod(value, grid_size)
    if remainder < 0:
        remainder += grid_size
    return min(remainder, grid_size - remainder)


def snap_to_grid(value: float, grid_size: float) -> float:
    cells = round(value / grid_size)
    return cells * grid_size
