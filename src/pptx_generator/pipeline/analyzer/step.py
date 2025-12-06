from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Tuple

from pptx import Presentation

from ...models import Slide
from ..base import PipelineContext
from .bullet_analysis import (
    _check_bullet_depth as _check_bullet_depth_impl,
    _check_contrast as _check_contrast_impl,
    _check_font_size as _check_font_size_impl,
    _evaluate_bullet as _evaluate_bullet_impl,
    analyze_bullet_groups,
)
from .image_analysis import analyze_images
from .issues import IssueTracker
from .layout_checks import (
    check_grid_alignment as check_grid_alignment_impl,
    check_margins as check_margins_impl,
    locate_image_shape as locate_image_shape_impl,
    locate_textbox_shape as locate_textbox_shape_impl,
)
from .mapping_sync import sync_mapping_log
from .options import AnalyzerOptions
from .snapshot import SlideSnapshot
from .snapshot_export import (
    _placeholder_type_name as placeholder_type_name_impl,
    _shape_type_name as shape_type_name_impl,
    export_snapshot_slide,
    save_analysis,
    save_snapshot,
)
from .textbox_analysis import analyze_textboxes
from .utils import emu_to_inches

logger = logging.getLogger(__name__)


class SimpleAnalyzerStep:
    """PPTX の実データを解析して analysis.json を生成する。"""

    name = "analyzer"

    def __init__(
        self,
        options: AnalyzerOptions | None = None,
        *,
        artifact_key: str = "analysis_path",
        register_default_artifact: bool = True,
        allow_missing_artifact: bool = False,
    ) -> None:
        self.options = options or AnalyzerOptions()
        self._issue_tracker = IssueTracker()
        self._artifact_key = artifact_key
        self._register_default_artifact = register_default_artifact
        self._allow_missing_artifact = allow_missing_artifact

    def run(self, context: PipelineContext) -> None:
        pptx_reference = context.artifacts.get("pptx_path")
        if pptx_reference is None:
            if self._allow_missing_artifact:
                logger.info(
                    "Analyzer (%s) をスキップします: pptx_path artifact が見つかりません",
                    self.options.output_filename,
                )
                return
            msg = "解析対象の PPTX が存在しません。renderer の実行順序を確認してください。"
            raise RuntimeError(msg)

        pptx_path = Path(str(pptx_reference))
        if not pptx_path.exists():  # pragma: no cover - 異常系
            raise FileNotFoundError(f"PPTX ファイルが存在しません: {pptx_path}")

        presentation = Presentation(pptx_path)
        issues: list[dict[str, Any]] = []
        fixes: list[dict[str, Any]] = []
        snapshot_slides: list[dict[str, Any]] = []

        spec_slides = context.spec.slides
        if len(presentation.slides) < len(spec_slides):
            logger.warning(
                "PPTX のスライド数が不足しています: spec=%s, pptx=%s",
                len(spec_slides),
                len(presentation.slides),
            )

        presentation_width_in = emu_to_inches(int(getattr(presentation, "slide_width", 0)))
        presentation_height_in = emu_to_inches(int(getattr(presentation, "slide_height", 0)))
        if presentation_width_in <= 0:
            presentation_width_in = self.options.slide_width_in
        if presentation_height_in <= 0:
            presentation_height_in = self.options.slide_height_in

        for index, slide_spec in enumerate(spec_slides):
            if index >= len(presentation.slides):
                break
            slide = presentation.slides[index]
            slide_width_in = emu_to_inches(int(getattr(slide, "slide_width", 0)))
            slide_height_in = emu_to_inches(int(getattr(slide, "slide_height", 0)))
            if slide_width_in <= 0:
                slide_width_in = presentation_width_in
            if slide_height_in <= 0:
                slide_height_in = presentation_height_in
            snapshot = SlideSnapshot.from_slide(slide, index)
            slide_issues, slide_fixes = self._analyze_slide(
                slide_spec, snapshot, slide_width_in, slide_height_in
            )
            issues.extend(slide_issues)
            fixes.extend(slide_fixes)
            if self.options.snapshot_output_filename:
                snapshot_slides.append(export_snapshot_slide(slide_spec, snapshot))

        analysis = {
            "slides": len(spec_slides),
            "meta": context.spec.meta.model_dump(),
            "issues": issues,
            "fixes": fixes,
        }
        output_path = save_analysis(analysis, context.workdir, self.options.output_filename)
        context.add_artifact(self._artifact_key, output_path)
        if self._register_default_artifact and self._artifact_key != "analysis_path":
            context.add_artifact("analysis_path", output_path)
        elif self._register_default_artifact and self._artifact_key == "analysis_path":
            pass
        sync_mapping_log(context, analysis)
        logger.info("%s を出力しました: %s", self.options.output_filename, output_path)
        if self.options.snapshot_output_filename:
            snapshot_path = save_snapshot(
                snapshot_slides, context.workdir, self.options.snapshot_output_filename
            )
            context.add_artifact("analyzer_snapshot_path", snapshot_path)
            logger.info("構造スナップショットを出力しました: %s", snapshot_path)

    def _analyze_slide(
        self,
        slide_spec: Slide,
        snapshot: SlideSnapshot,
        slide_width_in: float,
        slide_height_in: float,
    ) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
        issues: list[dict[str, Any]] = []
        fixes: list[dict[str, Any]] = []

        bullet_issues, bullet_fixes = analyze_bullet_groups(
            self.options, self._issue_tracker, slide_spec, snapshot
        )
        issues.extend(bullet_issues)
        fixes.extend(bullet_fixes)

        image_issues, image_fixes = analyze_images(
            self.options,
            self._issue_tracker,
            slide_spec,
            snapshot,
            slide_width_in=slide_width_in,
            slide_height_in=slide_height_in,
        )
        issues.extend(image_issues)
        fixes.extend(image_fixes)

        textbox_issues, textbox_fixes = analyze_textboxes(
            self.options, self._issue_tracker, slide_spec, snapshot
        )
        issues.extend(textbox_issues)
        fixes.extend(textbox_fixes)

        return issues, fixes

    # --- Compatibility wrappers for legacy tests ---

    def _locate_image_shape(self, snapshot: SlideSnapshot, image_spec):
        return locate_image_shape_impl(snapshot, image_spec)

    def _locate_textbox_shape(self, snapshot: SlideSnapshot, textbox):
        return locate_textbox_shape_impl(snapshot, textbox)

    def _analyze_images(
        self,
        slide_spec: Slide,
        snapshot: SlideSnapshot,
        *,
        slide_width_in: float,
        slide_height_in: float,
    ) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
        return analyze_images(
            self.options,
            self._issue_tracker,
            slide_spec,
            snapshot,
            slide_width_in=slide_width_in,
            slide_height_in=slide_height_in,
        )

    def _analyze_textboxes(
        self,
        slide_spec: Slide,
        snapshot: SlideSnapshot,
    ) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
        return analyze_textboxes(self.options, self._issue_tracker, slide_spec, snapshot)

    def _evaluate_bullet(
        self,
        slide_spec: Slide,
        bullet,
        anchor: str | None,
        resolver,
        applied_level: int | None,
        previous_level: int | None,
    ):
        return _evaluate_bullet_impl(
            self.options,
            self._issue_tracker,
            slide_spec,
            bullet,
            anchor,
            resolver,
            applied_level,
            previous_level,
        )

    def _check_bullet_depth(self, slide, bullet, actual_level, target):
        return _check_bullet_depth_impl(self.options, self._issue_tracker, slide, bullet, actual_level, target)

    def _check_font_size(self, slide, bullet, paragraph, target):
        return _check_font_size_impl(self.options, self._issue_tracker, slide, bullet, paragraph, target)

    def _check_contrast(self, slide, bullet, paragraph, target):
        return _check_contrast_impl(self.options, self._issue_tracker, slide, bullet, paragraph, target)

    def _check_margins(
        self,
        slide,
        image,
        shape,
        *,
        slide_width_in: float,
        slide_height_in: float,
    ):
        return check_margins_impl(
            self.options,
            self._issue_tracker,
            slide,
            image,
            shape,
            slide_width_in=slide_width_in,
            slide_height_in=slide_height_in,
        )

    def _check_grid_alignment(self, slide, element_id, element_type, shape):
        return check_grid_alignment_impl(
            self.options,
            self._issue_tracker,
            slide,
            element_id,
            element_type,
            shape,
        )

    def _save_snapshot(self, slides: list[dict[str, Any]], workdir: Path) -> Path:
        if not self.options.snapshot_output_filename:
            raise ValueError("snapshot_output_filename が設定されていません")
        return save_snapshot(slides, workdir, self.options.snapshot_output_filename)

    def _sync_mapping_log(self, context: PipelineContext, analysis: dict[str, Any]) -> None:
        sync_mapping_log(context, analysis)

    def _extend_results(self, issues, fixes, outcome) -> None:
        self._issue_tracker.extend_results(issues, fixes, outcome)

    def _shape_type_name(self, shape_type: int | None) -> str:
        return shape_type_name_impl(shape_type)

    def _placeholder_type_name(self, placeholder_type: int | None) -> str | None:
        return placeholder_type_name_impl(placeholder_type)
