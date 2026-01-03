from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from .snapshot import SlideSnapshot


def _write_json_payload(payload: dict[str, Any], workdir: Path, filename: str) -> Path:
    workdir.mkdir(parents=True, exist_ok=True)
    path = workdir / filename
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def export_snapshot_slide(slide_spec, snapshot: SlideSnapshot) -> dict[str, Any]:
    named_shapes: list[dict[str, Any]] = []
    placeholders: list[dict[str, Any]] = []

    for shape in snapshot.shapes:
        paragraphs = [
            {
                "index": paragraph.paragraph_index,
                "text": paragraph.text,
                "level": paragraph.level,
                "font_size_pt": paragraph.font_size_pt,
                "color_hex": paragraph.color_hex,
                "font_name": paragraph.font_name,
                "bold": paragraph.bold,
                "italic": paragraph.italic,
                "alignment": paragraph.alignment,
                "line_spacing_pt": paragraph.line_spacing_pt,
                "space_before_pt": paragraph.space_before_pt,
                "space_after_pt": paragraph.space_after_pt,
                "left_indent_in": paragraph.left_indent_in,
                "right_indent_in": paragraph.right_indent_in,
                "first_line_indent_in": paragraph.first_line_indent_in,
            }
            for paragraph in shape.paragraphs
        ]
        base_record = {
            "shape_id": shape.shape_id,
            "name": shape.name or "",
            "shape_type": _shape_type_name(shape.shape_type),
            "left_in": shape.left_in,
            "top_in": shape.top_in,
            "width_in": shape.width_in,
            "height_in": shape.height_in,
            "paragraphs": paragraphs,
            "is_placeholder": shape.is_placeholder,
            "placeholder_type": _placeholder_type_name(shape.placeholder_type),
            "placeholder_index": shape.placeholder_index,
            "z_order": shape.z_order,
            "rotation_deg": shape.rotation_deg,
            "text_frame_padding": shape.text_frame_padding,
            "text_frame_word_wrap": shape.text_frame_word_wrap,
            "text_frame_vertical_anchor": shape.text_frame_vertical_anchor,
            "text_frame_auto_size": shape.text_frame_auto_size,
        }
        if shape.is_placeholder or shape.placeholder_type is not None:
            placeholders.append(dict(base_record))
        if shape.name:
            named_shapes.append(dict(base_record))

    spec_anchors = sorted(
        {
            *(group.anchor for group in slide_spec.bullets if group.anchor),
            *(image.anchor for image in slide_spec.images if image.anchor),
            *(table.anchor for table in slide_spec.tables if table.anchor),
            *(chart.anchor for chart in slide_spec.charts if chart.anchor),
            *(textbox.anchor for textbox in slide_spec.textboxes if textbox.anchor),
        }
    )

    return {
        "index": snapshot.index,
        "slide_id": slide_spec.id,
        "layout": slide_spec.layout,
        "placeholders": placeholders,
        "named_shapes": named_shapes,
        "spec_anchors": spec_anchors,
    }


def save_snapshot(slides: list[dict[str, Any]], workdir: Path, filename: str) -> Path:
    payload = {
        "schema_version": "1.0.0",
        "slides": slides,
    }
    return _write_json_payload(payload, workdir, filename)


def save_analysis(payload: dict[str, Any], workdir: Path, filename: str) -> Path:
    return _write_json_payload(payload, workdir, filename)


def _shape_type_name(shape_type: int | None) -> str:
    if shape_type is None:
        return "unknown"
    try:
        return MSO_SHAPE_TYPE(shape_type).name
    except ValueError:  # pragma: no cover - 予期しない値
        return str(shape_type)


def _placeholder_type_name(placeholder_type: int | None) -> str | None:
    if placeholder_type is None:
        return None
    try:
        return PP_PLACEHOLDER(placeholder_type).name
    except ValueError:  # pragma: no cover - 予期しない値
        return str(placeholder_type)
