from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_SHAPE_TYPE

from ...models import Slide, SlideImage, SlideTextbox
from .issues import IssueTracker
from .options import AnalyzerOptions
from .snapshot import ShapeSnapshot, SlideSnapshot
from .utils import grid_deviation, snap_to_grid


def locate_image_shape(snapshot: SlideSnapshot, image_spec: SlideImage) -> ShapeSnapshot | None:
    if image_spec.anchor:
        shape = snapshot.find_shape_by_name(image_spec.anchor, shape_type=int(MSO_SHAPE_TYPE.PICTURE))
        if shape:
            return shape
    shape = snapshot.find_shape_by_name(image_spec.id, shape_type=int(MSO_SHAPE_TYPE.PICTURE))
    if shape:
        return shape
    for candidate in snapshot.shapes:
        if candidate.shape_type == int(MSO_SHAPE_TYPE.PICTURE):
            return candidate
    return None


def locate_textbox_shape(snapshot: SlideSnapshot, textbox: SlideTextbox) -> ShapeSnapshot | None:
    candidates: list[ShapeSnapshot] = []
    if textbox.anchor:
        shape = snapshot.find_shape_by_name(textbox.anchor)
        if shape:
            candidates.append(shape)
    if not candidates:
        shape = snapshot.find_shape_by_name(textbox.id)
        if shape:
            candidates.append(shape)
    for shape in candidates:
        if shape.shape_type in {int(MSO_SHAPE_TYPE.TEXT_BOX), int(MSO_SHAPE_TYPE.PLACEHOLDER)}:
            return shape
    for shape in snapshot.shapes:
        if shape.shape_type in {int(MSO_SHAPE_TYPE.TEXT_BOX), int(MSO_SHAPE_TYPE.PLACEHOLDER)}:
            return shape
    return None


def check_margins(
    options: AnalyzerOptions,
    issue_tracker: IssueTracker,
    slide: Slide,
    image: SlideImage,
    shape: ShapeSnapshot,
    *,
    slide_width_in: float,
    slide_height_in: float,
) -> tuple[dict[str, Any], dict[str, Any]] | None:
    left = shape.left_in
    top = shape.top_in
    width = shape.width_in
    height = shape.height_in
    margin = options.margin_in
    base_width = slide_width_in if slide_width_in > 0 else options.slide_width_in
    base_height = slide_height_in if slide_height_in > 0 else options.slide_height_in

    if base_width is None or base_height is None:
        return None

    violations: list[str] = []
    if left < margin:
        violations.append("left")
    if top < margin:
        violations.append("top")
    right = left + width
    bottom = top + height
    if right > base_width - margin:
        violations.append("right")
    if bottom > base_height - margin:
        violations.append("bottom")

    if not violations:
        return None

    issue_id = issue_tracker.next_issue_id("margin", slide.id, image.id)
    target = {
        "slide_id": slide.id,
        "element_id": image.id,
        "element_type": "image",
    }

    target_left = min(max(left, margin), max(base_width - margin - width, margin))
    target_top = min(max(top, margin), max(base_height - margin - height, margin))

    fix_payload: dict[str, float] = {}
    if abs(target_left - left) > 1e-3:
        fix_payload["left_in"] = round(target_left, 3)
    if abs(target_top - top) > 1e-3:
        fix_payload["top_in"] = round(target_top, 3)

    fix = None
    if fix_payload:
        fix = {
            "id": f"fix-{issue_id}",
            "issue_id": issue_id,
            "type": "move",
            "target": target,
            "payload": fix_payload,
        }

    issue = issue_tracker.make_issue(
        issue_id=issue_id,
        issue_type="margin",
        severity="warning",
        message=(
            f"スライド '{slide.id}' の画像 '{image.id}' が余白基準 {margin:.1f}in を外れています"
        ),
        target=target,
        metrics={
            "left_in": left,
            "top_in": top,
            "width_in": width,
            "height_in": height,
            "margin_in": margin,
            "slide_width_in": base_width,
            "slide_height_in": base_height,
            "violations": violations,
            "shape_name": shape.name,
        },
        fix=fix,
    )
    return issue, fix


def check_grid_alignment(
    options: AnalyzerOptions,
    issue_tracker: IssueTracker,
    slide: Slide,
    element_id: str,
    element_type: str,
    shape: ShapeSnapshot,
) -> tuple[dict[str, Any], dict[str, Any]] | None:
    grid = options.grid_size_in
    tolerance = options.grid_tolerance_in

    deviations = {
        "left": grid_deviation(shape.left_in, grid),
        "top": grid_deviation(shape.top_in, grid),
    }
    out_of_grid = {axis: dev for axis, dev in deviations.items() if dev > tolerance}
    if not out_of_grid:
        return None

    target = {
        "slide_id": slide.id,
        "element_id": element_id,
        "element_type": element_type,
    }
    issue_id = issue_tracker.next_issue_id("grid_misaligned", slide.id, element_id)
    fix_payload: dict[str, float] = {}
    snapped_left = snap_to_grid(shape.left_in, grid)
    snapped_top = snap_to_grid(shape.top_in, grid)
    if "left" in out_of_grid:
        fix_payload["left_in"] = round(snapped_left, 3)
    if "top" in out_of_grid:
        fix_payload["top_in"] = round(snapped_top, 3)

    fix = None
    if fix_payload:
        fix = {
            "id": f"fix-{issue_id}",
            "issue_id": issue_id,
            "type": "move",
            "target": target,
            "payload": fix_payload,
        }

    issue = issue_tracker.make_issue(
        issue_id=issue_id,
        issue_type="grid_misaligned",
        severity="warning",
        message=(
            f"スライド '{slide.id}' の要素 '{element_id}' がグリッド {grid:.3f}in に揃っていません"
        ),
        target=target,
        metrics={
            "left_in": shape.left_in,
            "top_in": shape.top_in,
            "width_in": shape.width_in,
            "height_in": shape.height_in,
            "grid_in": grid,
            "tolerance_in": tolerance,
            "deviations_in": out_of_grid,
            "shape_name": shape.name,
        },
        fix=fix,
    )
    return issue, fix
