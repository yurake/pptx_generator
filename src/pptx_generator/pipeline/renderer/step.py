from __future__ import annotations

import logging
import time
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches

from ...models import JobSpec, Slide, TemplateStyle
from .autodraw import AutoDrawMixin
from .bullets import BulletMixin
from .charts import ChartMixin
from .images import ImageMixin
from .layout import LayoutBox, LayoutMixin
from .notes import NotesMixin
from .options import RenderingOptions
from .styling import StylingMixin
from .tables import TableMixin
from .textboxes import TextboxMixin

logger = logging.getLogger(__name__)


class SimpleRendererStep(
    AutoDrawMixin,
    NotesMixin,
    TableMixin,
    ChartMixin,
    ImageMixin,
    TextboxMixin,
    BulletMixin,
    LayoutMixin,
    StylingMixin,
):
    """最小機能の PPTX レンダラー。"""

    name = "renderer"

    def __init__(self, options: RenderingOptions | None = None) -> None:
        self.options = options or RenderingOptions()
        if self.options.template_style is None:
            self.options.template_style = TemplateStyle.default()
        self._style = self.options.template_style
        self._temp_files: list[Path] = []

    def run(self, context) -> None:
        presentation = self._load_template()
        start = time.perf_counter()
        try:
            self._render_slides(presentation, context.spec)
            output_path = self._save(presentation, context.workdir)
            context.add_artifact("pptx_path", output_path)
            logger.info("PPTX を出力しました: %s", output_path)
        finally:
            self._cleanup_temp_files()

        elapsed_ms = int((time.perf_counter() - start) * 1000)
        context.add_artifact("renderer_stats", {"rendering_time_ms": elapsed_ms})

    def _load_template(self) -> Presentation:
        if self.options.template_path and self.options.template_path.exists():
            logger.debug("テンプレートを使用: %s", self.options.template_path)
            return Presentation(self.options.template_path)
        logger.debug("既定テンプレートを利用")
        return Presentation()

    def _render_slides(self, presentation: Presentation, spec: JobSpec) -> None:
        if self.options.template_source == "slide":
            self._render_using_prototype_slides(presentation, spec)
        else:
            self._render_using_layouts(presentation, spec)

    def _render_using_layouts(self, presentation: Presentation, spec: JobSpec) -> None:
        for page_number, slide_spec in enumerate(spec.slides, start=1):
            layout = self._resolve_layout(presentation, slide_spec)
            slide = presentation.slides.add_slide(layout)
            self._apply_title(slide, slide_spec)
            self._apply_subtitle(slide, slide_spec)
            self._apply_bullets(slide, slide_spec)
            self._apply_textboxes(slide, slide_spec)
            self._apply_tables(slide, slide_spec)
            self._apply_images(slide, slide_spec)
            self._apply_charts(slide, slide_spec)
            self._apply_notes(slide, slide_spec)
            self._apply_auto_draw(slide, slide_spec, page_number)

    def _render_using_prototype_slides(self, presentation: Presentation, spec: JobSpec) -> None:
        total_required = len(spec.slides)
        mapped_indices = self.options.prototype_mapping or []
        if mapped_indices and len(mapped_indices) < total_required:
            raise RuntimeError(
                "プロトタイプスライドの指定数が不足しています。"
            )
        if not mapped_indices:
            mapped_indices = list(range(1, total_required + 1))
        expected_sequence = list(range(1, total_required + 1))
        if mapped_indices[:total_required] != expected_sequence:
            raise RuntimeError(
                "プロトタイプスライドの順序指定には現在対応していません。テンプレートの先頭から順に実スライドを配置してください。"
            )
        if len(presentation.slides) < total_required:
            raise RuntimeError(
                "テンプレートに含まれる実スライド数が不足しています。"
            )

        for page_number, slide_spec in enumerate(spec.slides, start=1):
            slide = presentation.slides[page_number - 1]
            self._apply_title(slide, slide_spec)
            self._apply_subtitle(slide, slide_spec)
            self._apply_bullets(slide, slide_spec)
            self._apply_textboxes(slide, slide_spec)
            self._apply_tables(slide, slide_spec)
            self._apply_images(slide, slide_spec)
            self._apply_charts(slide, slide_spec)
            self._apply_notes(slide, slide_spec)
            self._apply_auto_draw(slide, slide_spec, page_number)

        self._truncate_slides(presentation, total_required)

    def _resolve_layout(self, presentation: Presentation, slide_spec: Slide):
        for layout in presentation.slide_layouts:
            if layout.name == slide_spec.layout:
                return layout
        logger.debug("レイアウト '%s' が見つからないため既定を使用", slide_spec.layout)
        try:
            return presentation.slide_layouts[1]
        except IndexError:
            if len(presentation.slide_layouts) == 0:
                raise RuntimeError("テンプレートに利用可能なレイアウトが存在しません")
            logger.warning(
                "テンプレートにレイアウト index=1 が存在しないため、index=0 を使用します"
            )
            return presentation.slide_layouts[0]

    def _apply_title(self, slide, slide_spec: Slide) -> None:
        if slide_spec.title is None:
            return
        title_shape = slide.shapes.title
        if title_shape is not None:
            text_frame = title_shape.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.text = slide_spec.title
            self._set_font(paragraph, self._style.heading_font)
            return
        textbox = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.5), Inches(8.0), Inches(1.0)
        )
        text_frame = textbox.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = slide_spec.title
        self._set_font(paragraph, self._style.heading_font)

    def _apply_subtitle(self, slide, slide_spec: Slide) -> None:
        if not slide_spec.subtitle:
            return

        subtitle_placeholder = None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                subtitle_placeholder = placeholder
                break

        if subtitle_placeholder is not None:
            text_frame = subtitle_placeholder.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.text = slide_spec.subtitle
            self._set_font(paragraph, self._style.body_font)
            return

        textbox = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.5), Inches(8.0), Inches(1.0)
        )
        text_frame = textbox.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = slide_spec.subtitle
        self._set_font(paragraph, self._style.body_font)

    def _find_body_placeholder(self, slide):
        for shape in slide.placeholders:
            if shape.placeholder_format.type in {
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.VERTICAL_BODY,
                PP_PLACEHOLDER.OBJECT,
            }:
                return shape
        logger.debug("本文用プレースホルダがないためテキストボックスを追加")
        return slide.shapes.add_textbox(
            Inches(1.0), Inches(1.5), Inches(8.0), Inches(4.5)
        )

    def _truncate_slides(self, presentation: Presentation, keep_count: int) -> None:
        total = len(presentation.slides)
        if keep_count >= total:
            return
        for index in range(total - 1, keep_count - 1, -1):
            self._delete_slide(presentation, index)

    @staticmethod
    def _delete_slide(presentation: Presentation, index: int) -> None:
        slides = presentation.slides
        slide_id = slides._sldIdLst[index]
        r_id = slide_id.rId
        presentation.part.drop_rel(r_id)
        slides._sldIdLst.remove(slide_id)

    def _save(self, presentation: Presentation, workdir: Path) -> Path:
        workdir.mkdir(parents=True, exist_ok=True)
        output_path = workdir / self.options.output_filename
        presentation.save(output_path)
        return output_path
