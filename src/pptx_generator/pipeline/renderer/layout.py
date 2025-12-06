from __future__ import annotations

import logging
from dataclasses import dataclass
from typing import Any

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches

from ...models import PipelineFallbackError

logger = logging.getLogger(__name__)


@dataclass(slots=True)
class LayoutBox:
    left_in: float
    top_in: float
    width_in: float
    height_in: float

    def to_emu(self) -> tuple[int, int, int, int]:
        return (
            int(Inches(self.left_in)),
            int(Inches(self.top_in)),
            int(Inches(self.width_in)),
            int(Inches(self.height_in)),
        )


@dataclass(slots=True)
class AnchorResolution:
    shape: Any | None
    left: int
    top: int
    width: int
    height: int
    is_placeholder: bool = False

    def as_box(self) -> tuple[int, int, int, int]:
        return self.left, self.top, self.width, self.height


class LayoutMixin:
    def _find_shape_by_name(self, slide, name: str):
        for shape in slide.shapes:
            if shape.name == name:
                return shape
        return None

    def _find_placeholder_by_name(self, slide, name: str):
        layout = getattr(slide, "slide_layout", None)
        if layout is None:
            return None
        target_idx: int | None = None
        for layout_shape in layout.shapes:
            if (
                getattr(layout_shape, "is_placeholder", False)
                and layout_shape.name == name
            ):
                target_idx = layout_shape.placeholder_format.idx
                break
        if target_idx is None:
            return None
        for shape in slide.shapes:
            if (
                getattr(shape, "is_placeholder", False)
                and shape.placeholder_format.idx == target_idx
            ):
                return shape
        return None

    def _resolve_anchor(
        self,
        slide,
        anchor: str | None,
        fallback_box: LayoutBox,
        *,
        owner_description: str | None = None,
    ) -> AnchorResolution:
        if anchor:
            shape = self._find_shape_by_name(slide, anchor)
            if shape is not None:
                return AnchorResolution(
                    shape,
                    int(shape.left),
                    int(shape.top),
                    int(shape.width),
                    int(shape.height),
                    getattr(shape, "is_placeholder", False),
                )
            placeholder = self._find_placeholder_by_name(slide, anchor)
            if placeholder is not None:
                return AnchorResolution(
                    placeholder,
                    int(placeholder.left),
                    int(placeholder.top),
                    int(placeholder.width),
                    int(placeholder.height),
                    True,
                )
            message = (
                f"{owner_description or '指定された要素'} のアンカー '{anchor}' がテンプレートで見つかりません。"
            )
            logger.error(message)
            raise PipelineFallbackError(message)
        left, top, width, height = fallback_box.to_emu()
        return AnchorResolution(None, left, top, width, height)

    def _prepare_placeholder(self, placeholder) -> None:
        if placeholder is None:
            return
        try:
            if getattr(placeholder, "has_text_frame", False):
                placeholder.text_frame.clear()
        except Exception:  # noqa: BLE001
            logger.debug("プレースホルダーの初期化に失敗", exc_info=True)

    def _remove_shape(self, shape) -> None:
        if shape is None:
            return
        try:
            shape.element.getparent().remove(shape.element)
        except Exception:  # noqa: BLE001
            logger.debug(
                "shape の削除に失敗しました: %s",
                shape.name if hasattr(shape, "name") else shape,
            )

    def _override_emu(self, default: int, value_in: float | None) -> int:
        if value_in is None:
            return default
        return int(Inches(value_in))

    def _obtain_text_frame(
        self,
        *,
        slide,
        anchor_name: str | None,
        fallback_box: LayoutBox,
        strict_anchor: bool,
    ):
        owner_description = (
            f"テキストボックス anchor '{anchor_name}'" if anchor_name else "テキストボックス要素"
        )
        resolution = self._resolve_anchor(
            slide,
            anchor_name,
            fallback_box,
            owner_description=owner_description,
        )
        shape = resolution.shape

        if strict_anchor and anchor_name and shape is None:
            return None

        if shape is not None and getattr(shape, "has_text_frame", False):
            if resolution.is_placeholder:
                self._prepare_placeholder(shape)
            else:
                shape.text_frame.clear()

        left, top, width, height = resolution.as_box()
        text_shape = slide.shapes.add_textbox(left, top, width, height)
        if shape is not None:
            self._remove_shape(shape)
        if anchor_name:
            try:
                text_shape.name = anchor_name
            except ValueError:
                logger.debug("アンカー名 '%s' の再設定に失敗", anchor_name, exc_info=True)
        return text_shape.text_frame
