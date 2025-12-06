from __future__ import annotations

import logging

from pptx.dml.color import RGBColor

from ...models import Slide, SlideTable, TemplateTableDefaults
from .layout import LayoutBox

logger = logging.getLogger(__name__)


class TableMixin:
    def _apply_tables(self, slide, slide_spec: Slide) -> None:
        if not slide_spec.tables:
            return

        table_defaults = self._style.table
        for table_spec in slide_spec.tables:
            header = list(table_spec.columns)
            rows = [list(row) for row in table_spec.rows]
            column_count = len(header) or (len(rows[0]) if rows else 0)
            if column_count == 0:
                logger.debug("テーブル '%s' に列情報がないためスキップ", table_spec.id)
                continue

            row_count = len(rows) + (1 if header else 0)
            if row_count == 0:
                logger.debug("テーブル '%s' にデータがないためスキップ", table_spec.id)
                continue

            fallback_position = table_defaults.fallback_box
            if fallback_position is not None:
                fallback_box = LayoutBox(
                    fallback_position.left_in,
                    fallback_position.top_in,
                    fallback_position.width_in,
                    fallback_position.height_in,
                )
            else:
                fallback_box = LayoutBox(1.0, 1.5, 8.5, 3.0)
            element_label = table_spec.id or (table_spec.anchor or "table")
            resolution = self._resolve_anchor(
                slide,
                table_spec.anchor,
                fallback_box,
                owner_description=f"テーブル要素 '{element_label}' (slide_id='{slide_spec.id}')",
            )
            anchor_shape = resolution.shape
            if resolution.is_placeholder:
                self._prepare_placeholder(anchor_shape)
            left, top, width, height = resolution.as_box()
            table_shape = slide.shapes.add_table(
                row_count, column_count, left, top, width, height
            )
            table = table_shape.table

            total_width = table_shape.width
            for idx in range(column_count):
                table.columns[idx].width = total_width // column_count

            start_row = 0
            if header:
                self._fill_table_row(
                    table.rows[0],
                    header,
                    is_header=True,
                    table_spec=table_spec,
                    defaults=table_defaults,
                )
                start_row = 1

            for offset, row_values in enumerate(rows):
                target_row = table.rows[start_row + offset]
                padded = row_values + [""] * (column_count - len(row_values))
                self._fill_table_row(
                    target_row,
                    padded,
                    is_header=False,
                    table_spec=table_spec,
                    defaults=table_defaults,
                    zebra_index=offset,
                )

            if anchor_shape is not None:
                self._remove_shape(anchor_shape)

    def _fill_table_row(
        self,
        row,
        values: list[object],
        *,
        is_header: bool,
        table_spec: SlideTable,
        defaults: TemplateTableDefaults,
        zebra_index: int | None = None,
    ) -> None:
        spec_style = table_spec.style
        header_font = defaults.header_font or self._style.body_font
        body_font = defaults.body_font or self._style.body_font
        header_fill_color = (
            spec_style.header_fill
            if spec_style and spec_style.header_fill
            else defaults.header_fill_color
        )
        body_fill_color = defaults.body_fill_color or "#FFFFFF"
        zebra_fill_color = defaults.zebra_fill_color

        for idx, value in enumerate(values):
            cell = row.cells[idx]
            text_frame = cell.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.text = str(value)
            if is_header:
                self._set_font(paragraph, header_font)
                fill_color = header_fill_color or body_fill_color
            else:
                self._apply_font(
                    paragraph,
                    None,
                    fallback=body_font,
                )
                zebra_enabled = bool(spec_style and spec_style.zebra)
                use_zebra = (
                    zebra_enabled
                    and zebra_fill_color
                    and zebra_index is not None
                    and zebra_index % 2 == 1
                )
                fill_color = (
                    zebra_fill_color
                    if use_zebra
                    else body_fill_color
                )

            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip("#"))
