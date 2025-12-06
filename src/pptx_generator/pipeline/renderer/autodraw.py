from __future__ import annotations

import logging

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from ...models import Slide

logger = logging.getLogger(__name__)


class AutoDrawMixin:
    def _apply_auto_draw(self, slide, slide_spec: Slide, page_number: int) -> None:
        if not slide_spec.auto_draw_boxes:
            return

        for anchor, position in slide_spec.auto_draw_boxes.items():
            left = int(Inches(position.left_in))
            top = int(Inches(position.top_in))
            width = int(Inches(position.width_in))
            height = int(Inches(position.height_in))

            try:
                existing = self._find_shape_by_name(slide, anchor)
            except Exception:  # noqa: BLE001
                existing = None
            if existing is not None:
                continue

            textbox = slide.shapes.add_textbox(left, top, width, height)
            try:
                textbox.name = anchor
            except ValueError:
                logger.debug("自動描画アンカー '%s' の命名に失敗", anchor, exc_info=True)
            text_frame = textbox.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.text = str(page_number)
            paragraph.alignment = PP_ALIGN.RIGHT
            self._set_font(paragraph, self._style.body_font)
