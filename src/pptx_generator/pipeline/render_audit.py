"""レンダリング結果の軽量監査ステップ。"""

from __future__ import annotations

import json
import logging
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from .base import PipelineContext

logger = logging.getLogger(__name__)


@dataclass(slots=True)
class RenderingAuditOptions:
    """レンダリング監査に関する設定。"""

    enabled: bool = True
    output_filename: str = "rendering_log.json"


class RenderingAuditStep:
    """レンダリング済み PPTX の整合性を検査し、監査ログを生成する。"""

    name = "rendering_audit"

    def __init__(self, options: RenderingAuditOptions | None = None) -> None:
        self.options = options or RenderingAuditOptions()

    def run(self, context: PipelineContext) -> None:
        if not self.options.enabled:
            logger.debug("Rendering audit is disabled")
            return

        try:
            pptx_reference = context.require_artifact("pptx_path")
        except KeyError:
            logger.warning("Rendering audit skipped because pptx_path artifact is missing")
            return

        pptx_path = Path(str(pptx_reference))
        if not pptx_path.exists():
            logger.warning("Rendering audit skipped because PPTX was not found: %s", pptx_path)
            return

        presentation = Presentation(pptx_path)
        spec = context.spec

        slides_payload: list[dict[str, Any]] = []
        warnings_total = 0
        empty_placeholder_total = 0

        for index, slide in enumerate(presentation.slides, start=1):
            spec_slide = spec.slides[index - 1] if index - 1 < len(spec.slides) else None
            detected, warnings, empty_count = self._inspect_slide(slide, spec_slide)
            warnings_total += len(warnings)
            empty_placeholder_total += empty_count
            slides_payload.append(
                {
                    "page_no": index,
                    "layout_id": spec_slide.layout if spec_slide else None,
                    "detected": detected,
                    "warnings": warnings,
                }
            )

        if len(spec.slides) > len(presentation.slides):
            for offset, spec_slide in enumerate(spec.slides[len(presentation.slides):], start=len(presentation.slides) + 1):
                warning = {
                    "code": "missing_slide",
                    "message": f"spec slide '{spec_slide.id}' is missing in rendered PPTX",
                }
                warnings_total += 1
                slides_payload.append(
                    {
                        "page_no": offset,
                        "layout_id": spec_slide.layout,
                        "detected": {
                            "title": False,
                            "subtitle": False,
                            "body": False,
                            "notes": False,
                        },
                        "warnings": [warning],
                    }
                )

        template_version: str | None = None
        rendering_ready = context.artifacts.get("rendering_ready")
        if rendering_ready is not None:
            template_version = getattr(getattr(rendering_ready, "meta", None), "template_version", None)

        renderer_stats = context.artifacts.get("renderer_stats") or {}

        meta: dict[str, Any] = {
            "generated_at": datetime.now(timezone.utc).isoformat(),
            "template_version": template_version,
            "rendering_time_ms": renderer_stats.get("rendering_time_ms"),
            "warnings_total": warnings_total,
            "empty_placeholders": empty_placeholder_total,
        }
        if len(presentation.slides) != len(spec.slides):
            meta["slide_count_actual"] = len(presentation.slides)
            meta["slide_count_expected"] = len(spec.slides)

        rendering_log = {
            "meta": meta,
            "slides": slides_payload,
        }

        output_dir = context.workdir
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / self.options.output_filename
        output_path.write_text(
            json.dumps(rendering_log, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

        context.add_artifact("rendering_log", rendering_log)
        context.add_artifact("rendering_log_path", str(output_path))
        context.add_artifact(
            "rendering_summary",
            {
                "warnings_total": warnings_total,
                "empty_placeholders": empty_placeholder_total,
            },
        )

        logger.info(
            "Rendering audit completed: warnings=%d empty_placeholders=%d log=%s",
            warnings_total,
            empty_placeholder_total,
            output_path,
        )

    def _inspect_slide(self, slide, spec_slide) -> tuple[dict[str, bool], list[dict[str, Any]], int]:
        detected = {
            "title": False,
            "subtitle": False,
            "body": False,
            "notes": False,
        }
        warnings: list[dict[str, Any]] = []
        empty_placeholder_count = 0
        flagged_placeholders: set[str] = set()

        title_shape = getattr(slide.shapes, "title", None)
        title_text = self._shape_text(title_shape)
        detected["title"] = bool(title_text)
        if spec_slide and spec_slide.title and not detected["title"]:
            warnings.append({"code": "missing_title", "message": "title プレースホルダーが空です"})
            flagged_placeholders.add(self._shape_name(title_shape))

        subtitle_shape = self._find_placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
        subtitle_text = self._shape_text(subtitle_shape)
        detected["subtitle"] = bool(subtitle_text)
        if spec_slide and spec_slide.subtitle and not detected["subtitle"]:
            warnings.append({"code": "missing_subtitle", "message": "subtitle プレースホルダーが空です"})
            flagged_placeholders.add(self._shape_name(subtitle_shape))

        detected_body = self._has_body_content(slide)
        detected["body"] = detected_body
        if spec_slide and self._expects_body(spec_slide) and not detected_body:
            warnings.append({"code": "missing_body", "message": "本文プレースホルダーに内容が挿入されていません"})

        detected_notes = self._has_notes(slide)
        detected["notes"] = detected_notes
        if spec_slide and spec_slide.notes and not detected_notes:
            warnings.append({"code": "missing_notes", "message": "ノートが空です"})

        for placeholder in self._iter_placeholders(slide):
            name = self._shape_name(placeholder)
            if name in flagged_placeholders:
                continue
            text = self._shape_text(placeholder)
            if text:
                continue
            empty_placeholder_count += 1
            warnings.append(
                {
                    "code": "empty_placeholder",
                    "message": f"プレースホルダー '{name or placeholder.placeholder_format.type}' が空です",
                }
            )

        return detected, warnings, empty_placeholder_count

    @staticmethod
    def _shape_text(shape) -> str:
        if shape is None or not getattr(shape, "has_text_frame", False):
            return ""
        texts: list[str] = []
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.text:
                texts.append(paragraph.text.strip())
        return "\n".join(filter(None, texts)).strip()

    @staticmethod
    def _shape_name(shape) -> str | None:
        if shape is None:
            return None
        name = getattr(shape, "name", None)
        if isinstance(name, str) and name:
            return name
        return None

    @staticmethod
    def _iter_placeholders(slide):
        for shape in slide.shapes:
            if getattr(shape, "is_placeholder", False):
                yield shape

    @staticmethod
    def _find_placeholder(slide, placeholder_type) -> Any:
        for shape in slide.shapes:
            if not getattr(shape, "is_placeholder", False):
                continue
            try:
                if int(shape.placeholder_format.type) == int(placeholder_type):
                    return shape
            except Exception:  # noqa: BLE001
                continue
        return None

    def _has_body_content(self, slide) -> bool:
        for shape in self._iter_placeholders(slide):
            try:
                placeholder_type = int(shape.placeholder_format.type)
            except Exception:  # noqa: BLE001
                placeholder_type = None
            if placeholder_type in {
                int(PP_PLACEHOLDER.BODY),
                int(PP_PLACEHOLDER.VERTICAL_BODY),
                int(PP_PLACEHOLDER.OBJECT),
            }:
                if self._shape_text(shape):
                    return True
        for shape in slide.shapes:
            if getattr(shape, "is_placeholder", False):
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            if self._shape_text(shape):
                return True
        return False

    @staticmethod
    def _expects_body(spec_slide) -> bool:
        if spec_slide is None:
            return False
        for group in spec_slide.bullets:
            if not group.items:
                continue
            if group.anchor is None:
                return True
        return False

    @staticmethod
    def _has_notes(slide) -> bool:
        try:
            frame = slide.notes_slide.notes_text_frame
        except AttributeError:
            return False
        text = frame.text or ""
        return bool(text.strip())
