"""Shape extraction related mixins."""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING, Any

from pptx.enum.text import PP_ALIGN
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import PlaceholderPicture, SlidePlaceholder

from ...utils.text_capacity import estimate_text_capacity
from .constants import EMU_PER_INCH, SLIDE_BULLET_ANCHORS
from .helpers import color_to_hex, length_to_inches, length_to_pt

if TYPE_CHECKING:
    from ...models import (
        FontSpec,
        ShapeInfo,
        TextCapacity,
        TextFramePadding,
        TextboxParagraph,
    )

logger = logging.getLogger(__name__)

__all__ = ["ShapeExtractionMixin"]


class ShapeExtractionMixin:
    """Provides utilities for extracting shapes from slides/layouts."""

    _heading_font_default: "FontSpec"
    _body_font_default: "FontSpec"

    def _extract_shape_info(self, shape: BaseShape) -> "ShapeInfo":
        from ...models import ShapeInfo  # local import to avoid cycles

        name = getattr(shape, "name", "")
        if not name:
            name = f"unnamed_shape_{id(shape)}"

        left_in = shape.left / EMU_PER_INCH if hasattr(shape, "left") else 0.0
        top_in = shape.top / EMU_PER_INCH if hasattr(shape, "top") else 0.0
        width_in = shape.width / EMU_PER_INCH if hasattr(shape, "width") else 0.0
        height_in = shape.height / EMU_PER_INCH if hasattr(shape, "height") else 0.0

        shape_type = shape.__class__.__name__

        text = None
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is not None:
            frame_text = getattr(text_frame, "text", None)
            if isinstance(frame_text, str):
                text = frame_text
        if text is None:
            raw_text = getattr(shape, "text", None)
            if isinstance(raw_text, str):
                text = raw_text

        placeholder_format = None
        try:
            placeholder_format = shape.placeholder_format  # type: ignore[attr-defined]
        except (AttributeError, ValueError):
            placeholder_format = None
        is_placeholder = bool(
            isinstance(shape, (SlidePlaceholder, PlaceholderPicture))
            or getattr(shape, "is_placeholder", False)
            or placeholder_format is not None
        )
        placeholder_type = None
        if placeholder_format is not None:
            placeholder_kind = getattr(placeholder_format, "type", None)
            if hasattr(placeholder_kind, "name"):
                placeholder_type = str(getattr(placeholder_kind, "name"))
            elif placeholder_kind is not None:
                placeholder_type = str(placeholder_kind)

        (
            font_spec,
            paragraph_spec,
            frame_padding,
            text_capacity,
        ) = (None, None, None, None)
        if self._is_text_shape(placeholder_type, shape):
            (
                font_spec,
                paragraph_spec,
                frame_padding,
                text_capacity,
            ) = self._extract_text_attributes(
                text_frame,
                placeholder_type,
                width_in,
                height_in,
            )

        conflict = None
        if name.lower() in SLIDE_BULLET_ANCHORS:
            conflict = f"SlideBullet拡張仕様で使用される可能性のあるアンカー名: {name}"

        missing_fields = []
        if not name or name.startswith("unnamed_"):
            missing_fields.append("name")
        if width_in <= 0:
            missing_fields.append("width")
        if height_in <= 0:
            missing_fields.append("height")

        return ShapeInfo(
            name=name,
            shape_type=shape_type,
            left_in=left_in,
            top_in=top_in,
            width_in=width_in,
            height_in=height_in,
            text=text,
            placeholder_type=placeholder_type,
            is_placeholder=is_placeholder,
            conflict=conflict,
            missing_fields=missing_fields,
            font=font_spec,
            paragraph=paragraph_spec,
            text_frame_padding=frame_padding,
            text_capacity=text_capacity,
        )

    def _is_text_shape(self, placeholder_type: str | None, shape: BaseShape) -> bool:
        placeholder = (placeholder_type or "").upper()
        if placeholder in {
            "TITLE",
            "CENTER_TITLE",
            "SUBTITLE",
            "BODY",
            "CONTENT",
            "TEXT",
        }:
            return True
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            return False
        text_value = (
            getattr(text_frame, "text", None) or getattr(shape, "text", None) or ""
        )
        return bool(text_value.strip())

    def _extract_text_attributes(
        self,
        text_frame,
        placeholder_type: str | None,
        width_in: float,
        height_in: float,
    ) -> tuple[
        "FontSpec | None",
        "TextboxParagraph | None",
        "TextFramePadding | None",
        "TextCapacity | None",
    ]:
        if text_frame is None:
            font_spec = self._resolve_font_spec(placeholder_type, None)
            text_capacity = estimate_text_capacity(
                width_in=width_in,
                height_in=height_in,
                font=font_spec,
                paragraph=None,
                padding=None,
            )
            return font_spec, None, None, text_capacity

        paragraph_obj = text_frame.paragraphs[0] if text_frame.paragraphs else None
        paragraph_spec = self._convert_paragraph(paragraph_obj)
        font_spec = self._resolve_font_spec(placeholder_type, paragraph_obj)
        padding = self._convert_text_frame_padding(text_frame)
        text_capacity = estimate_text_capacity(
            width_in=width_in,
            height_in=height_in,
            font=font_spec,
            paragraph=paragraph_spec,
            padding=padding,
        )
        return font_spec, paragraph_spec, padding, text_capacity

    def _resolve_font_spec(self, placeholder_type: str | None, paragraph) -> "FontSpec":
        placeholder = (placeholder_type or "").upper()
        base = (
            self._heading_font_default
            if placeholder in {"TITLE", "CENTER_TITLE", "SUBTITLE"}
            else self._body_font_default
        )
        overrides = self._font_overrides_from_paragraph(paragraph)
        if not overrides:
            return base
        return base.model_copy(update=overrides)

    def _font_overrides_from_paragraph(self, paragraph) -> dict[str, Any]:
        if paragraph is None:
            return {}

        attributes = self._collect_font_attributes(getattr(paragraph, "font", None))
        if not self._font_attributes_complete(attributes):
            attributes = self._fill_font_attributes_from_runs(paragraph, attributes)

        overrides: dict[str, Any] = {}
        if attributes["name"]:
            overrides["name"] = attributes["name"]
        if attributes["size"] is not None:
            overrides["size_pt"] = attributes["size"]
        if attributes["color"]:
            overrides["color_hex"] = attributes["color"]
        if attributes["bold"] is not None:
            overrides["bold"] = attributes["bold"]
        if attributes["italic"] is not None:
            overrides["italic"] = attributes["italic"]
        return overrides

    def _collect_font_attributes(self, font) -> dict[str, Any]:
        attributes = {
            "name": None,
            "size": None,
            "color": None,
            "bold": None,
            "italic": None,
        }
        if font is None:
            return attributes
        name = getattr(font, "name", None)
        attributes["name"] = name.strip() if isinstance(name, str) else name
        attributes["size"] = length_to_pt(getattr(font, "size", None))
        attributes["color"] = color_to_hex(getattr(font, "color", None))
        if font.bold is not None:
            attributes["bold"] = bool(font.bold)
        if font.italic is not None:
            attributes["italic"] = bool(font.italic)
        return attributes

    @staticmethod
    def _font_attributes_complete(attributes: dict[str, Any]) -> bool:
        return all(
            [
                bool(attributes["name"]),
                attributes["size"] is not None,
                bool(attributes["color"]),
                attributes["bold"] is not None,
                attributes["italic"] is not None,
            ]
        )

    def _fill_font_attributes_from_runs(
        self, paragraph, base_attributes: dict[str, Any]
    ) -> dict[str, Any]:
        attributes = dict(base_attributes)
        for run in getattr(paragraph, "runs", []):
            run_font = getattr(run, "font", None)
            if run_font is None:
                continue
            run_attributes = self._collect_font_attributes(run_font)
            for key in ("name", "size", "color", "bold", "italic"):
                if attributes[key] is None and run_attributes[key] is not None:
                    attributes[key] = run_attributes[key]
            if self._font_attributes_complete(attributes):
                break
        return attributes

    def _convert_paragraph(self, paragraph) -> "TextboxParagraph | None":
        from ...models import TextboxParagraph

        if paragraph is None:
            return None
        fmt = getattr(paragraph, "paragraph_format", None)
        return TextboxParagraph(
            level=max(paragraph.level if paragraph.level is not None else 0, 0),
            line_spacing_pt=self._line_spacing_to_pt(paragraph),
            space_before_pt=(
                length_to_pt(getattr(fmt, "space_before", None)) if fmt else None
            ),
            space_after_pt=(
                length_to_pt(getattr(fmt, "space_after", None)) if fmt else None
            ),
            align=self._alignment_to_str(paragraph.alignment),
            left_indent_in=(
                length_to_inches(getattr(fmt, "left_margin", None)) if fmt else None
            ),
            right_indent_in=(
                length_to_inches(getattr(fmt, "right_margin", None)) if fmt else None
            ),
            first_line_indent_in=(
                length_to_inches(getattr(fmt, "first_line_indent", None))
                if fmt
                else None
            ),
        )

    @staticmethod
    def _alignment_to_str(value) -> str | None:
        if value is None:
            return None
        try:
            align = PP_ALIGN(value)
        except ValueError:
            return None
        return {
            PP_ALIGN.LEFT: "left",
            PP_ALIGN.CENTER: "center",
            PP_ALIGN.RIGHT: "right",
            PP_ALIGN.JUSTIFY: "justify",
            PP_ALIGN.DISTRIBUTE: "distributed",
        }.get(align, None)

    def _convert_text_frame_padding(self, text_frame) -> "TextFramePadding | None":
        from ...models import TextFramePadding

        if text_frame is None:
            return None
        padding = TextFramePadding(
            left_in=length_to_inches(getattr(text_frame, "margin_left", None)),
            right_in=length_to_inches(getattr(text_frame, "margin_right", None)),
            top_in=length_to_inches(getattr(text_frame, "margin_top", None)),
            bottom_in=length_to_inches(getattr(text_frame, "margin_bottom", None)),
        )
        if all(value in (None, 0.0) for value in padding.model_dump().values()):
            return None
        return padding

    def _line_spacing_to_pt(self, paragraph) -> float | None:
        if paragraph is None:
            return None
        value = getattr(paragraph, "line_spacing", None)
        if value is None:
            return None
        if isinstance(value, (int, float)):
            font = getattr(paragraph, "font", None)
            base_size = length_to_pt(getattr(font, "size", None)) if font else None
            if base_size is None:
                base_size = self._body_font_default.size_pt
            return float(value) * float(base_size)
        return length_to_pt(value)

    @staticmethod
    def _should_include_for_summary(shape: "ShapeInfo") -> bool:
        if shape.is_placeholder:
            return True
        if shape.placeholder_type:
            return True
        if shape.name and shape.name.lower() not in {"rectangle", "textbox"}:
            return True
        return False

    @staticmethod
    def _shape_bbox_emu(shape: "ShapeInfo") -> dict[str, int]:
        return {
            "x": int(round(shape.left_in * EMU_PER_INCH)),
            "y": int(round(shape.top_in * EMU_PER_INCH)),
            "width": int(round(shape.width_in * EMU_PER_INCH)),
            "height": int(round(shape.height_in * EMU_PER_INCH)),
        }

    @staticmethod
    def _build_summary_flags(shape: "ShapeInfo", normalised_type: str) -> list[str]:
        flags: list[str] = []
        if normalised_type == "unknown":
            flags.append("unknown_type")
        if shape.conflict:
            flags.append("anchor_conflict")
        if shape.missing_fields:
            flags.append("missing_fields")
        return flags

    def _build_placeholder_record(self, shape: "ShapeInfo") -> dict[str, Any]:
        from ...utils.layout_metadata import normalise_placeholder_type

        normalised_type = normalise_placeholder_type(shape.placeholder_type, shape.name)
        record: dict[str, Any] = {
            "name": shape.name,
            "type": normalised_type,
            "bbox": self._shape_bbox_emu(shape),
            "shape_type": str(shape.shape_type or "").casefold() or None,
            "flags": self._build_summary_flags(shape, normalised_type),
        }
        if record["shape_type"] is None:
            record.pop("shape_type")
        return record
