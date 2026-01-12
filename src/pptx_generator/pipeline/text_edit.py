from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List

from pptx import Presentation
from pptx.presentation import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

from .analyzer.snapshot import table_cell_shape_id


def overwrite_text_frame_preserving_style(text_frame, new_text: str) -> None:
    """
    text_frame のテキストを差し替える。既存のラン／段落の書式をベースに再適用する。
    """
    if text_frame is None:
        return

    paragraphs = list(getattr(text_frame, "paragraphs", []))
    if paragraphs:
        base_paragraph = paragraphs[0]
    else:
        base_paragraph = text_frame.add_paragraph()
    base_run = base_paragraph.runs[0] if base_paragraph.runs else base_paragraph.add_run()

    run_style = _snapshot_run_style(base_run.font)
    paragraph_style = _snapshot_paragraph_style(base_paragraph)
    frame_style = _snapshot_text_frame_style(text_frame)

    lines = new_text.splitlines() or [""]

    text_frame.clear()
    _apply_text_frame_style(text_frame, frame_style)

    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        _apply_paragraph_style(paragraph, paragraph_style)
        if not paragraph.runs:
            paragraph.add_run()
        run = paragraph.runs[0]
        _apply_run_style(run.font, run_style)


def apply_shape_text_edits(
    pptx_path: Path | str,
    edits: Iterable[dict[str, object]],
    *,
    output_path: Path | str | None = None,
) -> tuple[int, list[str]]:
    """
    shape_id (任意で slide_index/name) をキーにテキスト差し替えを適用する。書式は overwrite_text_frame_preserving_style で保持する。
    戻り値: (適用件数, 未適用キーリスト)。slide_index 指定ありの場合は "<slide_index>:<shape_id>" 形式。
    """
    pptx_path = Path(pptx_path)
    presentation = Presentation(pptx_path)
    output_path = Path(output_path) if output_path is not None else pptx_path

    normalized = _normalize_edits(edits)
    applied, missing = _apply_edits_to_presentation(presentation, normalized)
    presentation.save(output_path)
    return applied, missing


def generate_edits_template(pptx_path: Path | str, output_path: Path | str | None = None) -> Path:
    """
    PPTX から shape_id ごとのテキストを抽出し、差分定義のテンプレート(JSON)を出力する。
    """
    pptx_path = Path(pptx_path)
    presentation = Presentation(pptx_path)
    output_path = Path(output_path) if output_path is not None else pptx_path.with_name(f"{pptx_path.stem}_edits.json")

    edits = snapshot_shapes_for_edit(presentation)
    payload = {"edits": edits}
    output_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return output_path


def snapshot_shapes_for_edit(pptx: Path | Presentation) -> List[dict[str, object]]:
    """LLM 入力やテンプレ生成向けに shape_id とテキストを抽出する。"""
    presentation = pptx if isinstance(pptx, PptxPresentation) else Presentation(pptx)
    edits: list[dict[str, object]] = []
    for slide_index, slide in enumerate(presentation.slides):
        for shape in _iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                edits.append(
                    {
                        "shape_id": int(shape.shape_id),
                        "slide_index": slide_index,
                        "name": getattr(shape, "name", ""),
                        "table_cell": None,
                        "edit": False,
                        "contents": shape.text,
                    }
                )
            if getattr(shape, "has_table", False):
                for row_idx, row in enumerate(getattr(shape.table, "rows", [])):
                    for col_idx, cell in enumerate(getattr(row, "cells", [])):
                        cell_id = table_cell_shape_id(int(shape.shape_id), row_idx, col_idx)
                        edits.append(
                            {
                                "shape_id": cell_id,
                                "slide_index": slide_index,
                                "name": getattr(shape, "name", ""),
                                "table_cell": {"row": row_idx, "col": col_idx},
                                "edit": False,
                                "contents": cell.text,
                            }
                        )
    return edits


def _snapshot_run_style(font) -> dict[str, object | None]:
    color = getattr(getattr(font, "color", None), "rgb", None)
    return {
        "name": getattr(font, "name", None),
        "size": getattr(font, "size", None),
        "bold": getattr(font, "bold", None),
        "italic": getattr(font, "italic", None),
        "underline": getattr(font, "underline", None),
        "color_rgb": color if isinstance(color, RGBColor) else None,
    }


def _apply_run_style(font, style: dict[str, object | None]) -> None:
    if style.get("name") is not None:
        font.name = style["name"]
    if style.get("size") is not None:
        font.size = style["size"]
    if style.get("bold") is not None:
        font.bold = style["bold"]
    if style.get("italic") is not None:
        font.italic = style["italic"]
    if style.get("underline") is not None:
        font.underline = style["underline"]
    if style.get("color_rgb") is not None:
        font.color.rgb = style["color_rgb"]


def _snapshot_paragraph_style(paragraph) -> dict[str, object | None]:
    return {
        "alignment": getattr(paragraph, "alignment", None),
        "level": getattr(paragraph, "level", None),
        "line_spacing": getattr(paragraph, "line_spacing", None),
        "space_before": getattr(paragraph, "space_before", None),
        "space_after": getattr(paragraph, "space_after", None),
        "left_indent": getattr(paragraph, "left_indent", None),
        "right_indent": getattr(paragraph, "right_indent", None),
        "first_line_indent": getattr(paragraph, "first_line_indent", None),
    }


def _apply_paragraph_style(paragraph, style: dict[str, object | None]) -> None:
    if style.get("alignment") is not None:
        paragraph.alignment = style["alignment"]
    if style.get("level") is not None:
        paragraph.level = style["level"]
    if style.get("line_spacing") is not None and hasattr(paragraph, "line_spacing"):
        paragraph.line_spacing = style["line_spacing"]
    if style.get("space_before") is not None and hasattr(paragraph, "space_before"):
        paragraph.space_before = style["space_before"]
    if style.get("space_after") is not None and hasattr(paragraph, "space_after"):
        paragraph.space_after = style["space_after"]
    if style.get("left_indent") is not None and hasattr(paragraph, "left_indent"):
        paragraph.left_indent = style["left_indent"]
    if style.get("right_indent") is not None and hasattr(paragraph, "right_indent"):
        paragraph.right_indent = style["right_indent"]
    if style.get("first_line_indent") is not None and hasattr(paragraph, "first_line_indent"):
        paragraph.first_line_indent = style["first_line_indent"]


def _snapshot_text_frame_style(text_frame) -> dict[str, object | None]:
    return {
        "word_wrap": getattr(text_frame, "word_wrap", None),
        "vertical_anchor": getattr(text_frame, "vertical_anchor", None),
        "auto_size": getattr(text_frame, "auto_size", None),
        "margin_left": getattr(text_frame, "margin_left", None),
        "margin_right": getattr(text_frame, "margin_right", None),
        "margin_top": getattr(text_frame, "margin_top", None),
        "margin_bottom": getattr(text_frame, "margin_bottom", None),
    }


def _apply_text_frame_style(text_frame, style: dict[str, object | None]) -> None:
    if style.get("word_wrap") is not None:
        text_frame.word_wrap = style["word_wrap"]
    if style.get("vertical_anchor") is not None:
        text_frame.vertical_anchor = style["vertical_anchor"]
    if style.get("auto_size") is not None:
        text_frame.auto_size = style["auto_size"]
    if style.get("margin_left") is not None:
        text_frame.margin_left = style["margin_left"]
    if style.get("margin_right") is not None:
        text_frame.margin_right = style["margin_right"]
    if style.get("margin_top") is not None:
        text_frame.margin_top = style["margin_top"]
    if style.get("margin_bottom") is not None:
        text_frame.margin_bottom = style["margin_bottom"]


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if int(getattr(shape, "shape_type", 0)) == int(MSO_SHAPE_TYPE.GROUP) and hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)


def _edit_applied(target_shape_id: int, presentation, target_slide_index: int | None, target_name: str | None) -> bool:
    for slide_index, slide in enumerate(presentation.slides):
        if target_slide_index is not None and slide_index != target_slide_index:
            continue
        for shape in _iter_shapes(slide.shapes):
            if getattr(shape, "shape_id", None) == target_shape_id:
                if target_name is None or target_name == getattr(shape, "name", None):
                    return True
            if getattr(shape, "has_table", False):
                for row_idx, row in enumerate(getattr(shape.table, "rows", [])):
                    for col_idx, _cell in enumerate(getattr(row, "cells", [])):
                        if table_cell_shape_id(int(shape.shape_id), row_idx, col_idx) == target_shape_id:
                            if target_name is None or target_name == getattr(shape, "name", None):
                                return True
    return False


@dataclass
class _NormalizedEdit:
    shape_id: int
    slide_index: int | None
    name: str | None
    contents: str


def _normalize_edits(edits: Iterable[dict[str, object]]) -> list[_NormalizedEdit]:
    normalized: list[_NormalizedEdit] = []
    for edit in edits:
        if not isinstance(edit, dict):
            continue
        if not edit.get("edit", True):
            continue
        shape_id = edit.get("shape_id")
        contents = edit.get("contents")
        if shape_id is None or contents is None:
            continue
        try:
            shape_id_int = int(shape_id)
        except (TypeError, ValueError):
            continue
        try:
            slide_idx = int(edit.get("slide_index")) if edit.get("slide_index") is not None else None
        except (TypeError, ValueError):
            slide_idx = None
        name_val = edit.get("name")
        normalized.append(
            _NormalizedEdit(
                shape_id=shape_id_int,
                slide_index=slide_idx,
                name=str(name_val) if name_val is not None else None,
                contents=str(contents),
            )
        )
    return normalized


def _apply_edits_to_presentation(
    presentation: PptxPresentation, edits: list[_NormalizedEdit]
) -> tuple[int, list[str]]:
    applied = 0
    missing: list[str] = []
    for edit in edits:
        if _apply_single_edit(presentation, edit):
            applied += 1
        else:
            missing.append(_format_missing_key(edit))
    return applied, missing


def _apply_single_edit(presentation: PptxPresentation, edit: _NormalizedEdit) -> bool:
    for slide_index, slide in enumerate(presentation.slides):
        if edit.slide_index is not None and slide_index != edit.slide_index:
            continue
        for shape_id, shape_name, text_frame in _iter_slide_text_frames(slide):
            if shape_id != edit.shape_id:
                continue
            if edit.name is not None and edit.name != shape_name:
                continue
            overwrite_text_frame_preserving_style(text_frame, edit.contents)
            return True
    return False


def _iter_slide_text_frames(slide):
    for shape in _iter_shapes(slide.shapes):
        shape_id = getattr(shape, "shape_id", None)
        shape_name = getattr(shape, "name", None)
        if getattr(shape, "has_text_frame", False):
            yield shape_id, shape_name, shape.text_frame
        if getattr(shape, "has_table", False):
            for row_idx, row in enumerate(getattr(shape.table, "rows", [])):
                for col_idx, cell in enumerate(getattr(row, "cells", [])):
                    cell_id = table_cell_shape_id(int(shape.shape_id), row_idx, col_idx)
                    yield cell_id, shape_name, cell.text_frame


def _format_missing_key(edit: _NormalizedEdit) -> str:
    return f"{edit.slide_index}:{edit.shape_id}" if edit.slide_index is not None else str(edit.shape_id)


__all__ = ["overwrite_text_frame_preserving_style", "apply_shape_text_edits", "generate_edits_template"]
