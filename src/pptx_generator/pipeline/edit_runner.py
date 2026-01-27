from __future__ import annotations

import json
import logging
import os
from pathlib import Path
from typing import Iterable, Callable

from pptx import Presentation

from pptx_generator.edit_ai import create_edit_ai_client, EditAIRequest, build_user_prompt
from pptx_generator.edit_ai.client import EditAIImage, EditAIResponseFormatError
from pptx_generator.pipeline.analyzer.snapshot import SlideSnapshot, ShapeSnapshot
from pptx_generator.pipeline.analyzer.utils import emu_to_inches
from pptx_generator.pipeline.slide_image_exporter import (
    SlideImageAsset,
    SlideImageExportOptions,
    SlideImageExporter,
)
from pptx_generator.pipeline.text_edit import apply_shape_text_edits, snapshot_shapes_for_edit


logger = logging.getLogger(__name__)


class EditRunError(RuntimeError):
    """edit 実行時の共通エラー。"""


def load_edits(edits_path: Path, *, error_cls: type[Exception] = EditRunError) -> list[dict]:
    payload = json.loads(edits_path.read_text(encoding="utf-8"))
    if isinstance(payload, dict) and "edits" in payload:
        return payload.get("edits", [])
    if isinstance(payload, list):
        return payload
    raise error_cls('edits ファイルの形式が不正です。リストまたは {"edits": [...]} を指定してください。')


def resolve_explicit_edits(
    edits_json: str | Path | None,
    edits_inline: object,
    *,
    error_cls: type[Exception] = EditRunError,
) -> list[dict] | None:
    if edits_json:
        json_path = Path(edits_json).expanduser()
        if not json_path.exists():
            raise error_cls(f"edits_json not found: {json_path}")
        edits = load_edits(json_path, error_cls=error_cls)
        if not isinstance(edits, list):
            raise error_cls("edits はリストである必要があります")
        return edits
    if edits_inline is not None:
        if not isinstance(edits_inline, list):
            raise error_cls("edits must be a list when provided inline")
        return edits_inline
    return None


def generate_edits_via_llm(
    pptx_path: Path,
    *,
    snapshot_fn: Callable[[Path], list[dict]] = snapshot_shapes_for_edit,
    client_factory: Callable[[], object] = create_edit_ai_client,
    image_exporter: SlideImageExporter | None = None,
    image_output_dir: Path | None = None,
) -> tuple[list[dict], set[str]]:
    shapes = snapshot_fn(pptx_path)
    client = client_factory()
    all_edits: list[dict[str, object]] = []
    models: set[str] = set()

    slide_images: dict[int, list[SlideImageAsset]] = {}
    slide_snapshots: dict[int, SlideSnapshot] = {}
    slide_size_in: dict[str, float] | None = None

    if image_exporter and image_output_dir and pptx_path.exists():
        export_result = image_exporter.export(pptx_path, image_output_dir)
        slide_images = export_result.images_by_slide
        if export_result.errors:
            for err in export_result.errors:
                logger.warning("edit screenshot export skipped: %s", err)
        if slide_images:
            try:
                slide_snapshots, slide_size_in = _build_slide_snapshots(pptx_path)
                metadata_path = _save_edit_image_metadata(
                    image_output_dir,
                    pptx_path=pptx_path,
                    slide_snapshots=slide_snapshots,
                    slide_size_in=slide_size_in,
                    images_by_slide=slide_images,
                )
                logger.info("edit screenshot metadata saved: %s", metadata_path)
            except Exception as exc:  # noqa: BLE001
                logger.warning("edit screenshot metadata generation failed: %s", exc)

    slides: dict[int, list[dict[str, object]]] = {}
    for shape in shapes:
        slides.setdefault(int(shape.get("slide_index", 0)), []).append(shape)

    for slide_idx, contexts in slides.items():
        slide_snapshot = slide_snapshots.get(slide_idx)
        slide_contexts = _attach_shape_geometry(contexts, slide_snapshot)
        slide_image = _select_slide_image(slide_images.get(slide_idx))
        screenshot_payload = _build_screenshot_payload(slide_image, image_output_dir)
        prompt = build_user_prompt(
            slide_title=None,
            shape_contexts=slide_contexts,
            slide_index=slide_idx,
            slide_size_in=slide_size_in,
            screenshot=screenshot_payload,
        )
        request_images = _load_edit_images([slide_image] if slide_image else [])
        request = EditAIRequest(prompt=prompt, shape_contexts=slide_contexts, images=request_images)
        response = client.rewrite(request)
        models.add(response.model)
        for edit in response.edits:
            if isinstance(edit, dict):
                edit["slide_index"] = slide_idx
            all_edits.append(edit)

    return all_edits, models


def apply_and_save_edits(
    pptx_path: Path, edits: list[dict], *, output_path: Path, models: Iterable[str] | set[str] | list[str]
):
    applied, missing = apply_shape_text_edits(pptx_path, edits, output_path=output_path)
    edits_path = _save_applied_edits(output_path, _normalize_edits_for_save(edits))
    return {
        "artifacts": {"pptx_url": str(output_path)},
        "applied": applied,
        "missing": missing,
        "models": sorted(models),
        "edits_path": str(edits_path),
    }


def run_edit_job(
    *,
    pptx_path: Path,
    edits_json: str | Path | None,
    edits_inline: object,
    output_path: Path,
    snapshot_fn: Callable[[Path], list[dict]] = snapshot_shapes_for_edit,
    client_factory: Callable[[], object] = create_edit_ai_client,
    error_cls: type[Exception] = EditRunError,
    apply_fn: Callable[[Path, list[dict]], dict] | None = None,
):
    """差分適用の共通実行フロー（CLI/API両用）。"""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    explicit = resolve_explicit_edits(edits_json, edits_inline, error_cls=error_cls)
    if explicit is not None:
        return (apply_fn or apply_and_save_edits)(pptx_path, explicit, output_path=output_path, models=[])
    try:
        image_exporter = _resolve_edit_image_exporter()
        image_output_dir = _resolve_edit_image_output_dir(output_path) if image_exporter else None
        llm_edits, models = generate_edits_via_llm(
            pptx_path,
            snapshot_fn=snapshot_fn,
            client_factory=client_factory,
            image_exporter=image_exporter,
            image_output_dir=image_output_dir,
        )
    except EditAIResponseFormatError as exc:
        raise error_cls(str(exc)) from exc
    return (apply_fn or apply_and_save_edits)(pptx_path, llm_edits, output_path=output_path, models=models)


def _save_applied_edits(output_path: Path, applied: list[dict]) -> Path:
    path = output_path.parent / "applied_edits.json"
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps({"edits": applied}, ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def _normalize_edits_for_save(edits: list[dict] | tuple | set) -> list[dict]:
    result: list[dict] = []
    for edit in edits:
        if not isinstance(edit, dict) or not edit.get("edit", True):
            continue
        contents = edit.get("contents")
        try:
            shape_id = int(edit.get("shape_id"))
        except (TypeError, ValueError):
            continue
        if contents is None:
            continue
        slide_raw = edit.get("slide_index")
        try:
            slide_idx = int(slide_raw) if slide_raw is not None else None
        except (TypeError, ValueError):
            slide_idx = None
        name_val = edit.get("name")
        result.append(
            {
                "shape_id": shape_id,
                "slide_index": slide_idx,
                "name": None if name_val is None else str(name_val),
                "contents": str(contents),
            }
        )
    return result


def _resolve_edit_image_exporter() -> SlideImageExporter | None:
    if not _env_flag("PPTX_EDIT_IMAGE_INPUT", default=False):
        return None
    formats = _env_list("PPTX_EDIT_IMAGE_FORMATS", default=["png", "jpg"])
    prefer_first = _env_flag("PPTX_EDIT_IMAGE_PREFER_FIRST", default=True)
    timeout_sec = _env_int("PPTX_EDIT_IMAGE_TIMEOUT_SEC", default=120)
    max_retries = _env_int("PPTX_EDIT_IMAGE_RETRIES", default=2)
    soffice_path = os.getenv("PPTX_EDIT_IMAGE_SOFFICE_PATH")
    options = SlideImageExportOptions(
        enabled=True,
        formats=tuple(formats),
        prefer_first_success=prefer_first,
        soffice_path=Path(soffice_path) if soffice_path else None,
        timeout_sec=timeout_sec,
        max_retries=max_retries,
    )
    return SlideImageExporter(options)


def _resolve_edit_image_output_dir(output_path: Path) -> Path:
    image_dir = output_path.parent / "images"
    image_dir.mkdir(parents=True, exist_ok=True)
    return image_dir


def _build_slide_snapshots(pptx_path: Path) -> tuple[dict[int, SlideSnapshot], dict[str, float]]:
    presentation = Presentation(pptx_path)
    slide_size_in = {
        "width": emu_to_inches(int(presentation.slide_width)),
        "height": emu_to_inches(int(presentation.slide_height)),
    }
    snapshots: dict[int, SlideSnapshot] = {}
    for slide_idx, slide in enumerate(presentation.slides):
        snapshots[slide_idx] = SlideSnapshot.from_slide(slide, slide_idx)
    return snapshots, slide_size_in


def _attach_shape_geometry(
    contexts: list[dict[str, object]],
    slide_snapshot: SlideSnapshot | None,
) -> list[dict[str, object]]:
    if slide_snapshot is None:
        return list(contexts)
    shapes_by_id = {shape.shape_id: shape for shape in slide_snapshot.shapes}
    enriched: list[dict[str, object]] = []
    for context in contexts:
        record = dict(context)
        shape_id = context.get("shape_id")
        try:
            shape_id_int = int(shape_id) if shape_id is not None else None
        except (TypeError, ValueError):
            shape_id_int = None
        if shape_id_int is not None:
            shape = shapes_by_id.get(shape_id_int)
            if shape:
                record["geometry"] = {
                    "left_in": shape.left_in,
                    "top_in": shape.top_in,
                    "width_in": shape.width_in,
                    "height_in": shape.height_in,
                    "rotation_deg": shape.rotation_deg,
                    "z_order": shape.z_order,
                    "shape_type": shape.shape_type,
                    "parent_shape_id": shape.parent_shape_id,
                    "table_cell": shape.table_cell,
                }
        enriched.append(record)
    return enriched


def _select_slide_image(images: list[SlideImageAsset] | None) -> SlideImageAsset | None:
    if not images:
        return None
    order = _env_list("PPTX_EDIT_IMAGE_FORMAT_ORDER", default=["png", "jpg", "jpeg"])
    order = [_normalize_image_format(fmt) for fmt in order if _normalize_image_format(fmt)]
    by_format = {asset.format: asset for asset in images}
    for fmt in order:
        asset = by_format.get(fmt)
        if asset:
            return asset
    return images[0]


def _build_screenshot_payload(slide_image: SlideImageAsset | None, image_output_dir: Path | None) -> dict[str, object] | None:
    if slide_image is None:
        return None
    path = slide_image.path
    path_str = str(path)
    if image_output_dir:
        try:
            path_str = str(path.relative_to(image_output_dir))
        except ValueError:
            pass
    return {
        "path": path_str,
        "format": slide_image.format,
        "media_type": slide_image.media_type,
    }


def _load_edit_images(images: Iterable[SlideImageAsset]) -> list[EditAIImage]:
    payloads: list[EditAIImage] = []
    for image in images:
        try:
            data = image.path.read_bytes()
        except OSError as exc:
            logger.warning("failed to read edit image: %s", exc)
            continue
        payloads.append(EditAIImage(media_type=image.media_type, data=data))
    return payloads


def _save_edit_image_metadata(
    image_output_dir: Path,
    *,
    pptx_path: Path,
    slide_snapshots: dict[int, SlideSnapshot],
    slide_size_in: dict[str, float] | None,
    images_by_slide: dict[int, list[SlideImageAsset]],
) -> Path:
    slides_payload: list[dict[str, object]] = []
    for slide_index, snapshot in slide_snapshots.items():
        images = images_by_slide.get(slide_index, [])
        slides_payload.append(
            {
                "slide_index": slide_index,
                "images": [_format_slide_image_record(image_output_dir, image) for image in images],
                "shapes": [_format_shape_snapshot(shape) for shape in snapshot.shapes],
            }
        )
    payload = {
        "schema_version": "1.0.0",
        "pptx_path": str(pptx_path),
        "slide_size_in": slide_size_in,
        "slides": slides_payload,
    }
    image_output_dir.mkdir(parents=True, exist_ok=True)
    metadata_path = image_output_dir / "edit_slide_images.json"
    metadata_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return metadata_path


def _format_slide_image_record(base_dir: Path, image: SlideImageAsset) -> dict[str, object]:
    path_str = str(image.path)
    try:
        path_str = str(image.path.relative_to(base_dir))
    except ValueError:
        pass
    return {
        "format": image.format,
        "media_type": image.media_type,
        "path": path_str,
    }


def _format_shape_snapshot(shape: ShapeSnapshot) -> dict[str, object]:
    return {
        "shape_id": shape.shape_id,
        "name": shape.name or "",
        "shape_type": shape.shape_type,
        "left_in": shape.left_in,
        "top_in": shape.top_in,
        "width_in": shape.width_in,
        "height_in": shape.height_in,
        "rotation_deg": shape.rotation_deg,
        "z_order": shape.z_order,
        "parent_shape_id": shape.parent_shape_id,
        "is_placeholder": shape.is_placeholder,
        "placeholder_type": shape.placeholder_type,
        "placeholder_index": shape.placeholder_index,
        "table_cell": shape.table_cell,
    }


def _normalize_image_format(fmt: str) -> str:
    trimmed = fmt.strip().lower().lstrip(".")
    return "jpg" if trimmed == "jpeg" else trimmed


def _env_flag(name: str, *, default: bool = False) -> bool:
    raw = os.getenv(name)
    if raw is None:
        return default
    return raw.strip().lower() in {"1", "true", "yes", "on"}


def _env_int(name: str, *, default: int) -> int:
    raw = os.getenv(name)
    if raw is None:
        return default
    try:
        return int(raw)
    except ValueError:
        return default


def _env_list(name: str, *, default: list[str]) -> list[str]:
    raw = os.getenv(name)
    if raw is None:
        return list(default)
    parts = [part.strip() for part in raw.split(",")]
    return [part for part in parts if part]


__all__ = [
    "EditRunError",
    "load_edits",
    "resolve_explicit_edits",
    "generate_edits_via_llm",
    "apply_and_save_edits",
]
