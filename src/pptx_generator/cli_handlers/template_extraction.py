from __future__ import annotations

import json
import logging
import os
from dataclasses import dataclass
from datetime import datetime, timezone
from importlib import resources as importlib_resources
from pathlib import Path
from typing import Any, Optional, Sequence

import click
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from pptx_generator.branding_extractor import extract_branding_config
from pptx_generator.cli_handlers.prepare import (
    PROMPT_TEMPLATE_FILENAME_PATTERN,
    PROMPT_USER_SECTION_END,
    PROMPT_USER_SECTION_START,
    SLIDE_INPUTS_FILENAME,
    build_prompt_identifier,
    slugify_prompt_layout,
)
from pptx_generator.layout_validation import (
    LayoutValidationError,
    LayoutValidationResult,
    LayoutValidationOptions,
    LayoutValidationSuite,
)
from pptx_generator.models import (
    JobSpecScaffold,
    TemplateBlueprint,
    TemplateBlueprintSlide,
    TemplateSpec,
)
from pptx_generator.pipeline import TemplateExtractor, TemplateExtractorOptions
from pptx_generator.pipeline.analyzer import SlideSnapshot

logger = logging.getLogger(__name__)

PROMPT_TEMPLATE_DIRNAME = Path("prompts")


@dataclass(slots=True)
class TemplateExtractionResult:
    template_spec: TemplateSpec
    jobspec_scaffold: JobSpecScaffold
    template_spec_path: Path
    branding_path: Path
    jobspec_path: Path
    validation_result: LayoutValidationResult | None
    output_dir: Path
    slide_snapshot_path: Path | None
    prompt_templates_dir: Path | None
    prompt_templates_created: int
    slide_inputs_path: Path | None


def run_template_extraction(
    *,
    template_path: Path,
    output_dir: Path,
    layout: str | None,
    anchor: str | None,
    output_format: str,
    template_ai_policy: Path | None,
    template_ai_policy_id: str | None,
    disable_template_ai: bool,
    layout_mode: str,
    skip_validation: bool = False,
    emit_slide_snapshot: bool = False,
) -> TemplateExtractionResult:
    fmt = output_format.lower()
    extractor_options = TemplateExtractorOptions(
        template_path=template_path,
        output_path=None,
        layout_filter=layout,
        anchor_filter=anchor,
        format=fmt,
        layout_mode=layout_mode,
    )

    output_dir.mkdir(parents=True, exist_ok=True)

    ai_policy_path = template_ai_policy
    if ai_policy_path is None and not disable_template_ai:
        env_policy = os.getenv("PPTX_TEMPLATE_AI_POLICY")
        if env_policy:
            ai_policy_path = Path(env_policy)
        else:
            ai_policy_path = discover_template_ai_policy()
    ai_policy_id = template_ai_policy_id or os.getenv("PPTX_TEMPLATE_AI_POLICY_ID")
    effective_disable = disable_template_ai or ai_policy_path is None
    if effective_disable:
        ai_policy_path = None
        if not disable_template_ai:
            logger.info("Template AI validation disabled: no policy file available")

    extractor = TemplateExtractor(extractor_options)
    template_spec = extractor.extract()
    jobspec_scaffold = extractor.build_jobspec_scaffold(template_spec)
    branding_result = extract_branding_config(template_path)

    if fmt == "yaml":
        import yaml

        spec_path = output_dir / "template_spec.yaml"
        spec_content = yaml.dump(
            template_spec.model_dump(mode="json", exclude_none=True),
            allow_unicode=True,
            default_flow_style=False,
            indent=2,
        )
    else:
        spec_path = output_dir / "template_spec.json"
        spec_content = json.dumps(
            template_spec.model_dump(mode="json", exclude_none=True),
            indent=2,
            ensure_ascii=False,
        )

    spec_path.write_text(spec_content, encoding="utf-8")
    logger.info("Saved template spec to %s", spec_path.resolve())

    branding_path = output_dir / "branding.json"
    branding_payload = branding_result.to_branding_payload()
    branding_text = json.dumps(branding_payload, ensure_ascii=False, indent=2)
    branding_path.write_text(branding_text, encoding="utf-8")
    logger.info("Saved branding payload to %s", branding_path.resolve())

    validation_result: LayoutValidationResult | None = None
    if not skip_validation:
        logger.info("Starting layout validation for %s", template_path)
        validation_options = LayoutValidationOptions(
            template_path=template_path,
            output_dir=output_dir,
            template_ai_policy_path=ai_policy_path,
            template_ai_policy_id=ai_policy_id,
            disable_template_ai=effective_disable,
        )
        validation_suite = LayoutValidationSuite(validation_options)
        validation_result = validation_suite.run()
        logger.info(
            "Layout validation finished: warnings=%d errors=%d",
            validation_result.warnings_count,
            validation_result.errors_count,
        )

        try:
            layouts_relative = str(
                validation_result.layouts_path.relative_to(output_dir)
            )
        except ValueError:
            layouts_relative = str(validation_result.layouts_path)
    else:
        validation_result = None
        layouts_relative = None

    template_spec_relative: str | None = None
    try:
        template_spec_relative = str(spec_path.relative_to(output_dir))
    except ValueError:
        template_spec_relative = str(spec_path)

    meta_update: dict[str, str | None] = {
        "template_spec_path": template_spec_relative,
    }
    if layouts_relative is not None:
        meta_update["layouts_path"] = layouts_relative

    jobspec_scaffold.meta = jobspec_scaffold.meta.model_copy(update=meta_update)

    jobspec_path = output_dir / "jobspec.json"
    extractor.save_jobspec_scaffold(jobspec_scaffold, jobspec_path)
    logger.info("Saved jobspec scaffold to %s", jobspec_path.resolve())

    slide_snapshot_path: Path | None = None
    if emit_slide_snapshot:
        slide_snapshot_path = _generate_slide_snapshot(
            template_path=template_path,
            output_dir=output_dir,
        )

    prompt_templates_dir: Path | None = None
    prompt_templates_created = 0
    if template_spec.layout_mode == "static" and template_spec.blueprint:
        try:
            prompt_templates_dir, prompt_templates_created = _ensure_prompt_templates(
                output_dir=output_dir,
                template_spec=template_spec,
            )
        except Exception as exc:  # noqa: BLE001
            logger.warning("プロンプト雛形の生成に失敗しました: %s", exc)
            prompt_templates_dir = None
            prompt_templates_created = 0

    slide_inputs_path: Path | None = None
    if template_spec.layout_mode == "static" and template_spec.blueprint:
        try:
            slide_inputs_path = _ensure_slide_inputs_manifest(
                output_dir=output_dir,
                template_spec=template_spec,
            )
        except Exception as exc:  # noqa: BLE001
            logger.warning("スライド入力マニフェストの生成に失敗しました: %s", exc)
            slide_inputs_path = None

    return TemplateExtractionResult(
        template_spec=template_spec,
        jobspec_scaffold=jobspec_scaffold,
        template_spec_path=spec_path,
        branding_path=branding_path,
        jobspec_path=jobspec_path,
        validation_result=validation_result,
        output_dir=output_dir,
        slide_snapshot_path=slide_snapshot_path,
        prompt_templates_dir=prompt_templates_dir,
        prompt_templates_created=prompt_templates_created,
        slide_inputs_path=slide_inputs_path,
    )


def echo_template_extraction_result(result: TemplateExtractionResult) -> None:
    template_spec = result.template_spec
    jobspec_scaffold = result.jobspec_scaffold
    validation_result = result.validation_result

    click.echo(f"テンプレート抽出が完了しました: {result.template_spec_path}")
    click.echo(f"ブランド設定を出力しました: {result.branding_path}")
    click.echo(f"ジョブスペック雛形を出力しました: {result.jobspec_path}")
    click.echo(f"抽出されたレイアウト数: {len(template_spec.layouts)}")
    click.echo(f"Layout Mode: {template_spec.layout_mode}")
    if template_spec.blueprint:
        click.echo(f"Blueprint Slides: {len(template_spec.blueprint.slides)}")

    total_anchors = sum(len(layout.anchors) for layout in template_spec.layouts)
    click.echo(f"抽出された図形・アンカー数: {total_anchors}")
    click.echo(f"ジョブスペックのスライド数: {len(jobspec_scaffold.slides)}")

    if validation_result is not None:
        click.echo(f"Layouts: {validation_result.layouts_path}")
        click.echo(f"Diagnostics: {validation_result.diagnostics_path}")
        if validation_result.diff_report_path is not None:
            click.echo(f"Diff: {validation_result.diff_report_path}")
        click.echo(
            "検出結果: warnings=%d, errors=%d"
            % (validation_result.warnings_count, validation_result.errors_count)
        )
    else:
        click.echo("検証をスキップしました (--force)")

    if result.slide_snapshot_path is not None:
        click.echo(f"スライドスナップショットを出力しました: {result.slide_snapshot_path}")
    if result.prompt_templates_dir is not None:
        click.echo(
            f"カスタムプロンプト雛形フォルダ: {result.prompt_templates_dir}"
        )
        if result.prompt_templates_created:
            click.echo(
                f"  -> {result.prompt_templates_created} 件のスライド雛形を生成しました。必要に応じて Markdown の user-editable 節を編集してください。"
            )
        else:
            click.echo("  -> 既存雛形を保持しました。変更が不要な場合はそのままご利用ください。")
    if result.slide_inputs_path is not None:
        click.echo(f"スライド入力マニフェストを出力しました: {result.slide_inputs_path}")

    if template_spec.warnings:
        click.echo(f"警告: {len(template_spec.warnings)} 件")
        for warning in template_spec.warnings:
            click.echo(f"  - {warning}", err=True)

    if template_spec.errors:
        click.echo(f"エラー: {len(template_spec.errors)} 件")
        for error in template_spec.errors:
            click.echo(f"  - {error}", err=True)


def discover_template_ai_policy() -> Path | None:
    env_policy = os.getenv("PPTX_TEMPLATE_AI_POLICY")
    if env_policy:
        candidate = Path(env_policy)
        if candidate.is_file():
            return candidate.resolve()

    cwd_candidate = Path.cwd() / "config" / "template_ai_policies.json"
    if cwd_candidate.exists():
        logger.info("Detected default template AI policy: %s", cwd_candidate)
        return cwd_candidate.resolve()

    try:
        resource = importlib_resources.files("pptx_generator").joinpath(
            "config/template_ai_policies.json"
        )
        if resource.is_file():
            try:
                text = resource.read_text(encoding="utf-8")
            except FileNotFoundError:
                text = None
            if text is not None:
                cache_dir = Path(".pptx/cache")
                cache_dir.mkdir(parents=True, exist_ok=True)
                cached_path = cache_dir / "template_ai_policies.json"
                cached_path.write_text(text, encoding="utf-8")
                logger.info("Using bundled template AI policy: %s", cached_path)
                return cached_path.resolve()
    except ModuleNotFoundError:
        logger.debug("Bundled template AI policy not found in package resources")
    return None


def _generate_slide_snapshot(*, template_path: Path, output_dir: Path) -> Path | None:
    logger.info("Generating slide snapshot for %s", template_path)
    try:
        presentation = Presentation(template_path)
    except Exception as exc:  # noqa: BLE001
        logger.warning("スライドスナップショットの生成に失敗しました: %s", exc)
        return None

    slides_payload: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides):
        snapshot = SlideSnapshot.from_slide(slide, index)
        slides_payload.append(_serialize_slide_snapshot(slide, snapshot))

    if not slides_payload:
        logger.info("スライドが存在しないためスナップショットを出力しません: %s", template_path)
        return None

    payload = {
        "template_path": str(template_path),
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "slides": slides_payload,
    }

    output_dir.mkdir(parents=True, exist_ok=True)
    path = output_dir / "slide_snapshot.json"
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    logger.info("Saved slide snapshot to %s", path.resolve())
    return path


def _serialize_slide_snapshot(slide, snapshot: SlideSnapshot) -> dict[str, Any]:
    layout_name = getattr(getattr(slide, "slide_layout", None), "name", None)
    slide_identifier = getattr(slide, "slide_id", None)

    shapes_payload: list[dict[str, Any]] = []
    for shape in snapshot.shapes:
        paragraphs = [
            {
                "index": paragraph.paragraph_index,
                "text": paragraph.text,
                "level": paragraph.level,
                "font_size_pt": paragraph.font_size_pt,
                "color_hex": paragraph.color_hex,
            }
            for paragraph in shape.paragraphs
        ]
        shapes_payload.append(
            {
                "shape_id": shape.shape_id,
                "name": shape.name or "",
                "shape_type": _shape_type_name(shape.shape_type),
                "left_in": shape.left_in,
                "top_in": shape.top_in,
                "width_in": shape.width_in,
                "height_in": shape.height_in,
                "is_placeholder": shape.is_placeholder,
                "placeholder_type": _placeholder_type_name(shape.placeholder_type),
                "paragraphs": paragraphs,
            }
        )

    return {
        "index": snapshot.index,
        "slide_id": slide_identifier,
        "layout": layout_name,
        "shapes": shapes_payload,
    }


def _shape_type_name(shape_type: int | None) -> str:
    if shape_type is None:
        return "unknown"
    try:
        return MSO_SHAPE_TYPE(shape_type).name
    except ValueError:
        return str(shape_type)


def _placeholder_type_name(placeholder_type: int | None) -> str | None:
    if placeholder_type is None:
        return None
    try:
        return PP_PLACEHOLDER(placeholder_type).name
    except ValueError:
        return str(placeholder_type)


def _ensure_prompt_templates(
    *,
    output_dir: Path,
    template_spec: TemplateSpec,
) -> tuple[Path, int]:
    prompts_dir = output_dir / PROMPT_TEMPLATE_DIRNAME
    prompts_dir.mkdir(parents=True, exist_ok=True)

    blueprint: TemplateBlueprint | None = template_spec.blueprint
    if blueprint is None:
        return prompts_dir, 0

    created = 0
    for index, slide in enumerate(blueprint.slides, start=1):
        slug_source = slide.layout or slide.slide_id or f"slide{index:02}"
        slug = slugify_prompt_layout(slug_source)
        filename = prompts_dir / f"{index:02}_{slug}.md"
        if filename.exists():
            logger.debug("Prompt template already exists, skipping: %s", filename)
            continue
        content = _render_prompt_template(slide=slide, index=index)
        filename.write_text(content, encoding="utf-8")
        created += 1

    return prompts_dir, created


def _render_prompt_template(*, slide: TemplateBlueprintSlide, index: int) -> str:
    layout_label = slide.layout or "Unnamed Layout"
    slide_id = slide.slide_id or f"slide-{index:02}"
    required_marker = "必須" if slide.required else "任意"
    intent_tags = ", ".join(slide.intent_tags) if slide.intent_tags else "(なし)"

    slot_lines: list[str] = []
    for slot in slide.slots:
        slot_required = "必須" if slot.required else "任意"
        slot_tags = ", ".join(slot.intent_tags) if slot.intent_tags else "(なし)"
        slot_lines.append(
            f"- `{slot.slot_id}` (anchor: {slot.anchor or '-'}, type: {slot.content_type}, {slot_required}, intent_tags: {slot_tags})"
        )
    if not slot_lines:
        slot_lines.append("- (slot 未定義)")

    fixed_section = "\n".join(
        [
            f"# Slide {index:02d}: {layout_label}",
            "",
            "## システム指定 (編集不可)",
            f"- slide_id: {slide_id}",
            f"- layout: {layout_label}",
            f"- スライド必須: {required_marker}",
            f"- intent_tags: {intent_tags}",
            "",
            "### slot 一覧",
            *slot_lines,
            "",
            "## 編集方法",
            "以下の user-editable セクションのみ編集してください。Markdown の構造を壊さないよう注意してください。",
            "",
        ]
    )

    user_section = "\n".join(
        [
            PROMPT_USER_SECTION_START,
            "- 例: このスライドでは ROI の定量値を箇条書きで入れる",
            "- 例: リスクを最低 2 点列挙する",
            PROMPT_USER_SECTION_END,
            "",
            "<!-- 編集しない場合は user-editable セクションを空のままにしてください -->",
        ]
    )

    return f"{fixed_section}\n{user_section}\n"


def _ensure_slide_inputs_manifest(*, output_dir: Path, template_spec: TemplateSpec) -> Path | None:
    blueprint = template_spec.blueprint
    if blueprint is None:
        return None

    base_dir = output_dir.parent if output_dir.parent != output_dir else output_dir
    manifest_path = base_dir / SLIDE_INPUTS_FILENAME
    if manifest_path.exists():
        logger.debug("Slide inputs manifest already exists: %s", manifest_path)
        return manifest_path

    lines = [
        "# Slide Inputs Manifest",
        "# 記法: <01_system-layout>: <data file path>",
        "# 例: 01_system-layout: samples/input/pitch.md",
        "",
    ]
    for index, slide in enumerate(blueprint.slides, start=1):
        identifier = build_prompt_identifier(index, slide)
        lines.append(f"{identifier}: <data file path>")

    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    manifest_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    logger.info("Saved slide inputs manifest to %s", manifest_path.resolve())
    return manifest_path
