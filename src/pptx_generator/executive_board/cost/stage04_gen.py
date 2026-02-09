#!/usr/bin/env python3
"""Stage4 hook: generate_ready.json の elements を PPTX へ反映する。"""
from __future__ import annotations

import json
import os
from copy import deepcopy
from pathlib import Path
import re
from typing import Any, Iterable

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from common.stage_shared import CONTEXT_PATH, is_keep_template_value, load_mapping_config, resolve_local_path
from common.rendering.title_style import apply_title_font_size, apply_title_prefix
from common.rendering.text_cleanup import apply_strip_bullet_prefix

TABLE_CELL_MAPPINGS: dict[str, dict[str, tuple[int, int]]] = {
    "表 8": {
        "初期計": (2, 4),
        "初期計利益": (2, 5),
        "期間請負": (3, 4),
        "期間請負利益": (3, 5),
        "社員": (4, 4),
        "社員利益": (4, 5),
        "内製開発": (5, 4),
        "内製開発利益": (5, 5),
        "pj管理": (6, 4),
        "pj管理利益": (6, 5),
        "pnibm": (7, 4),
        "pnibm利益": (7, 5),
        "si": (8, 4),
        "si利益": (8, 5),
        "機器": (14, 4),
        "ソフト使用料": (15, 4),
        "運用": (16, 4),
        "ランニング計": (13, 4),
    },
}

STATIC_TEXT_ANCHORS = {
    "Rectangle 23",
}


def _load_context() -> dict[str, str]:
    if CONTEXT_PATH.exists():
        return json.loads(CONTEXT_PATH.read_text(encoding="utf-8"))
    fallback = Path(".pptx/hook_context.json").resolve()
    if fallback.exists():
        return json.loads(fallback.read_text(encoding="utf-8"))
    return {}


def _resolve_template_path(path_str: str, base_dir: Path | None = None) -> Path:
    resolved = resolve_local_path(path_str, base_dir or Path(__file__).resolve().parent)
    if resolved.exists():
        try:
            Presentation(str(resolved))
            return resolved
        except Exception:
            pass
    return resolved


def _find_layout(presentation, target_name: str | None, blueprint_slots: list[dict[str, Any]] | None = None):
    if target_name:
        target_norm = target_name.strip().lower()
        for layout in presentation.slide_layouts:
            if layout.name.lower() == target_norm:
                return layout
    if blueprint_slots:
        anchors = {str(slot.get("anchor")) for slot in blueprint_slots if isinstance(slot, dict) and slot.get("anchor")}
        best_layout = None
        best_score = -1
        for layout in presentation.slide_layouts:
            layout_shape_names = {getattr(shape, "name", "") for shape in layout.shapes}
            score = len(anchors & layout_shape_names)
            if score > best_score:
                best_layout = layout
                best_score = score
        if best_layout is not None and best_score > 0:
            return best_layout
    return None


def _load_payload(generate_ready_path: Path) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    payload = json.loads(generate_ready_path.read_text(encoding="utf-8"))
    doc_meta = payload.get("meta") if isinstance(payload, dict) else {}
    slides = payload.get("slides")
    slides = slides if isinstance(slides, list) else []
    return [slide for slide in slides if isinstance(slide, dict)], doc_meta


def _pick_slide_payload(
    slides: list[dict[str, Any]], index: int
) -> tuple[dict[str, Any], dict[str, Any], dict[str, Any]]:
    if index < 0 or index >= len(slides):
        return {}, {}, {}
    slide = slides[index]
    elements = slide.get("elements") if isinstance(slide, dict) else {}
    elements = elements if isinstance(elements, dict) else {}
    meta = slide.get("meta") if isinstance(slide, dict) else {}
    meta = meta if isinstance(meta, dict) else {}
    return elements, meta, slide


def _normalize_name(value: str) -> str:
    normalized = value.lower()
    replacements = {
        "テキストプレースホルダー": "textplaceholder",
        "テキスト プレースホルダー": "textplaceholder",
        "テキストホルダー": "textplaceholder",
        "テキストボックス": "textbox",
        "テキスト ボックス": "textbox",
        "タイトル": "title",
        "表": "table",
        "図": "picture",
    }
    for src, dest in replacements.items():
        normalized = normalized.replace(src, dest)
    normalized = normalized.replace(" ", "")
    normalized = normalized.replace("_", "")
    normalized = normalized.replace("-", "")
    normalized = normalized.replace("（", "(").replace("）", ")")
    normalized = re.sub(r"[0-9\.]+", "", normalized)
    normalized = re.sub(r"[^a-z]+", "", normalized)
    return normalized


def _anchor_aliases(key: str) -> list[str]:
    aliases: list[str] = []
    if key in {"textplaceholder", "textbox"}:
        aliases.extend(["textplaceholder", "textbox", "rectangle"])
    if key == "title":
        aliases.extend(["title", "textbox"])
    if key == "table":
        aliases.extend(["table", "chart"])
    if key == "picture":
        aliases.extend(["picture", "image"])
    aliases.append(key)
    seen = set()
    ordered: list[str] = []
    for alias in aliases:
        if alias not in seen and alias:
            seen.add(alias)
            ordered.append(alias)
    return ordered


def _build_shape_groups(slide) -> dict[str, list[Any]]:
    groups: dict[str, list[Any]] = {}
    for shape in slide.shapes:
        name = getattr(shape, "name", "")
        key = _normalize_name(name)
        groups.setdefault(key, []).append(shape)
    return groups


def _pop_shape(groups: dict[str, list[Any]], anchor: str):
    exact_key = anchor.strip()
    for shapes in groups.values():
        for idx, shape in enumerate(shapes):
            if getattr(shape, "name", "").strip() == exact_key:
                return shapes.pop(idx), _normalize_name(anchor)
    key = _normalize_name(anchor)
    for alias in _anchor_aliases(key):
        shapes = groups.get(alias)
        if shapes:
            return shapes.pop(0), alias
    # fallback: substring matching
    raw_target = anchor.replace(" ", "").lower()
    for alias, shapes in groups.items():
        if raw_target and raw_target in alias and shapes:
            return shapes.pop(0), alias
    return None, None


def _set_shape_text(shape, value: Any) -> None:
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    text_frame.clear()
    values: list[str]
    if isinstance(value, list):
        values = [str(item) for item in value]
    elif value is None:
        values = [""]
    else:
        values = [str(value)]

    for idx, text in enumerate(values):
        if idx == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        paragraph.text = text
        paragraph.level = 0
        for run in paragraph.runs:
            run.font.size = Pt(12)


def _normalize_target_name(value: str) -> str:
    text = value.lower()
    replacements = {
        "（": "",
        "）": "",
        " ": "",
        "　": "",
        "・": "",
        "ー": "",
        "‐": "",
        "−": "",
        "_": "",
    }
    for src, dest in replacements.items():
        text = text.replace(src, dest)
    text = text.replace("計", "計")
    text = re.sub(r"[^0-9a-z一-龯ぁ-んァ-ン]", "", text)
    return text


def _apply_table_from_targets(shape, table_payload: dict[str, Any]) -> bool:
    if not shape.has_table or not isinstance(table_payload, dict):
        return False
    table = shape.table
    applied = False
    mapping = TABLE_CELL_MAPPINGS.get(getattr(shape, "name", ""))
    for payload in table_payload.values():
        target = payload.get("target")
        formatted = payload.get("formatted_value")
        if target is None or formatted is None:
            continue
        target_str = str(target)
        formatted_str = str(formatted)
        if mapping:
            normalized = _normalize_target_name(target_str)
            if normalized in mapping:
                row_idx, col_idx = mapping[normalized]
                if row_idx < len(table.rows) and col_idx < len(table.columns):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = formatted_str
                    applied = True
                    continue
        target_col = 5 if "利益" in target_str else 4
        target_col = min(target_col, len(table.columns) - 1)
        for row_idx in range(len(table.rows)):
            row_texts = " ".join(
                table.cell(row_idx, col).text for col in range(min(4, len(table.columns)))
            )
            if not row_texts:
                continue
            if target_str in row_texts:
                table.cell(row_idx, target_col).text = formatted_str
                applied = True
                break
            if "初期" in target_str and "（１）初期" in row_texts:
                table.cell(row_idx, target_col).text = formatted_str
                applied = True
                break
            if "ランニング" in target_str and "（２）ランニング" in row_texts:
                table.cell(row_idx, target_col).text = formatted_str
                applied = True
                break
            if "期間請負" in target_str and "期間請負" in row_texts:
                table.cell(row_idx, target_col).text = formatted_str
                applied = True
                break
            if "社員" in target_str and "社員" in row_texts and "SI" not in target_str:
                table.cell(row_idx, target_col).text = formatted_str
                applied = True
                break
            if "内製開発" in target_str and "内製開発" in row_texts:
                table.cell(row_idx, target_col).text = formatted_str
                applied = True
                break
            if "PJ管理" in target_str and "PJ管理" in row_texts:
                table.cell(row_idx, target_col).text = formatted_str
                applied = True
                break
            if "PN" in target_str and "PN" in row_texts:
                table.cell(row_idx, target_col).text = formatted_str
                applied = True
                break
            if "機器" in target_str and "機器" in row_texts:
                table.cell(row_idx, target_col).text = formatted_str
                applied = True
                break
            if "ソフト使用料" in target_str and "ソフト使用料" in row_texts:
                table.cell(row_idx, target_col).text = formatted_str
                applied = True
                break
            if "運用" in target_str and "運用" in row_texts:
                table.cell(row_idx, target_col).text = formatted_str
                applied = True
                break
    if applied:
        for row in table.rows:
            for cell in row.cells:
                if not cell.text:
                    continue
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
    return applied


def _apply_table_payload(shape, payload: dict[str, Any]) -> None:
    if not shape.has_table:
        return
    table = shape.table
    rows = payload.get("rows")
    if not isinstance(rows, list):
        return
    for row_idx in range(len(table.rows)):
        if row_idx >= len(rows):
            break
        row_values = rows[row_idx]
        if not isinstance(row_values, Iterable):
            continue
        for col_idx in range(len(table.columns)):
            if col_idx >= len(row_values):
                break
            cell = table.cell(row_idx, col_idx)
            cell.text = str(row_values[col_idx])
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)


def _apply_elements(
    slide,
    elements: dict[str, Any],
    blueprint_slots: list[dict[str, Any]],
    context_table_payload: dict[str, Any] | None,
    layout_defaults: dict[str, list[str]],
    detail_table_rows: list[list[Any]] | None,
) -> None:
    shape_groups = _build_shape_groups(slide)
    processed_keys: set[str] = set()

    def reinsert(alias: str | None, shape: Any | None) -> None:
        if alias and shape is not None:
            shape_groups.setdefault(alias, []).insert(0, shape)

    def apply_to_anchor(anchor: str, value: Any) -> None:
        shape, alias = _pop_shape(shape_groups, anchor)
        if shape is None:
            return
        if isinstance(value, dict) and ("rows" in value or value.get("headers")):
            rows = value.get("rows") if isinstance(value.get("rows"), list) else None
            if rows:
                _apply_table_payload(shape, value)
                return
            if context_table_payload and _apply_table_from_targets(shape, context_table_payload):
                return
            return
        _set_shape_text(shape, value)

    for slot in blueprint_slots:
        anchor = slot.get("anchor")
        if not isinstance(anchor, str):
            continue
        if anchor in STATIC_TEXT_ANCHORS:
            shape, alias = _pop_shape(shape_groups, anchor)
            default_value = layout_defaults.get(anchor)
            if shape is not None and default_value:
                _set_shape_text(shape, default_value)
            processed_keys.add(anchor)
            continue
        value = elements.get(anchor)
        slot_type = str(slot.get("content_type") or "").lower()
        if is_keep_template_value(value):
            processed_keys.add(anchor)
            continue
        if not slot_type:
            normalized_anchor_type = _normalize_name(anchor)
            if normalized_anchor_type in {"table", "chart"}:
                slot_type = "table"
        if value is None:
            if anchor == "表 10" and detail_table_rows:
                shape, alias = _pop_shape(shape_groups, anchor)
                if shape is None:
                    continue
                _apply_table_payload(shape, {"rows": detail_table_rows})
                processed_keys.add(anchor)
                continue
            if slot_type in {"chart", "table"} and context_table_payload:
                shape, alias = _pop_shape(shape_groups, anchor)
                if shape is None:
                    continue
                applied = _apply_table_from_targets(shape, context_table_payload)
                if applied:
                    processed_keys.add(anchor)
                else:
                    reinsert(alias, shape)
            continue
        processed_keys.add(anchor)
        apply_to_anchor(anchor, value)

    for anchor, value in elements.items():
        if anchor in processed_keys:
            continue
        apply_to_anchor(anchor, value)


def _clear_table(slide, table_name: str) -> None:
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        if getattr(shape, "name", "").strip() != table_name:
            continue
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                cell.text = ""


def _max_shape_id(slide) -> int:
    max_id = 0
    for shape in slide.shapes:
        cNvPr = getattr(shape._element, "cNvPr", None)
        if cNvPr is not None and cNvPr.get("id"):
            try:
                max_id = max(max_id, int(cNvPr.get("id")))
            except ValueError:
                continue
    return max_id


def _clear_slide_shapes(slide) -> None:
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


def _copy_shapes_from_slide(source_slide, target_slide) -> None:
    next_id = _max_shape_id(target_slide) + 1
    for shape in source_slide.shapes:
        new_element = deepcopy(shape.element)
        for cNvPr in new_element.iter(qn("p:cNvPr")):
            cNvPr.set("id", str(next_id))
            next_id += 1
        target_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")


def _find_cost_template_slide(template_pptx: Presentation):
    for slide in template_pptx.slides:
        names = {getattr(shape, "name", "") for shape in slide.shapes}
        if "表 8" in names or "表 10" in names:
            return slide
    return None


def _build_message_lines_from_summary(summary: dict[str, Any]) -> list[str]:
    if not isinstance(summary, dict):
        return []
    initial = summary.get("initial_amount", {}).get("formatted_value", "N/A")
    running_amount = summary.get("running_amount", {}).get("formatted_value", "N/A")
    running_years = summary.get("running_years", {}).get("formatted_value", "-")
    return [
        f"提示金額（初期）：{initial}億円",
        f"提示金額（ランニング）：{running_years} {running_amount}億円",
    ]


def _apply_text_to_named_shape(slide, target_name: str, lines: list[str]) -> None:
    for shape in slide.shapes:
        if getattr(shape, "name", "").strip() != target_name:
            continue
        if not shape.has_text_frame:
            return
        text_frame = shape.text_frame
        text_frame.clear()
        for idx, text in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            paragraph.text = text
            paragraph.level = 0
            for run in paragraph.runs:
                run.font.size = Pt(12)
        return


def _apply_tables_for_cost(slide, context_table_payload: dict[str, Any] | None, detail_rows: list[list[Any]] | None) -> None:
    if context_table_payload:
        for shape in slide.shapes:
            if getattr(shape, "name", "").strip() == "表 8":
                _apply_table_from_targets(shape, context_table_payload)
                break
    if detail_rows:
        for shape in slide.shapes:
            if getattr(shape, "name", "").strip() == "表 10":
                _apply_table_payload(shape, {"rows": detail_rows})
                break


def main() -> int:
    context = _load_context()
    mapping_config_env = os.environ.get("JRI_MAPPING_CONFIG")
    mapping_config_path: Path | None = None
    if isinstance(context, dict):
        mapping_config_path_str = context.get("mapping_config_path")
        if isinstance(mapping_config_path_str, str) and mapping_config_path_str.strip():
            mapping_config_path = resolve_local_path(mapping_config_path_str.strip(), Path(__file__).resolve().parent)
    if mapping_config_path is None and mapping_config_env:
        mapping_config_path = resolve_local_path(mapping_config_env.strip(), Path(__file__).resolve().parent)
    if mapping_config_path is None:
        default_mapping = Path(__file__).resolve().parent / "mapping_config.json"
        if default_mapping.exists():
            mapping_config_path = default_mapping.resolve()

    mapping_template_path: Path | None = None
    if mapping_config_path and mapping_config_path.exists():
        mapping_config = load_mapping_config(mapping_config_path)
        template_source_value = mapping_config.get("template_source_path")
        if isinstance(template_source_value, str) and template_source_value.strip():
            mapping_template_path = resolve_local_path(
                template_source_value.strip(), mapping_config_path.parent
            )

    generate_ready_path_env = os.environ.get("PPTX_GENERATE_READY_PATH") or context.get("generate_ready_path")
    if not generate_ready_path_env:
        raise RuntimeError("generate_ready path is not provided")
    generate_ready_path = Path(generate_ready_path_env).expanduser().resolve()
    if not generate_ready_path.exists():
        raise FileNotFoundError(f"generate_ready.json not found: {generate_ready_path}")

    slides_payload, doc_meta = _load_payload(generate_ready_path)

    template_candidates: list[str] = []
    if isinstance(doc_meta, dict):
        meta_template_path = doc_meta.get("template_path")
        if isinstance(meta_template_path, str) and meta_template_path.strip():
            template_candidates.append(meta_template_path.strip())
    env_template = os.environ.get("PPTX_TEMPLATE_PATH")
    if isinstance(env_template, str) and env_template.strip():
        template_candidates.append(env_template.strip())
    context_template = context.get("template_path") if isinstance(context, dict) else None
    if isinstance(context_template, str) and context_template.strip():
        template_candidates.append(context_template.strip())
    if mapping_template_path is not None:
        template_candidates.append(str(mapping_template_path))
    default_exec = Path(__file__).resolve().parents[3] / "templates/executive_board.pptx"
    template_candidates.append(str(default_exec))

    template_path: Path | None = None
    missing_templates: list[Path] = []
    for candidate in template_candidates:
        base_dir = mapping_template_path.parent if mapping_template_path and str(mapping_template_path) == candidate else Path(__file__).resolve().parent
        resolved = _resolve_template_path(candidate, base_dir)
        template_path = resolved
        if resolved.exists():
            break
        missing_templates.append(resolved)

    if template_path is None:
        raise FileNotFoundError("Template PPTX path could not be determined. Provide PPTX_TEMPLATE_PATH or update mapping_config.json.")
    if not template_path.exists():
        missing_repr = ", ".join(str(path) for path in missing_templates) if missing_templates else str(template_path)
        raise FileNotFoundError(f"template pptx not found: {missing_repr}")

    output_dir_env = os.environ.get("PPTX_OUTPUT_DIR")
    output_dir = Path(output_dir_env) if output_dir_env else Path(".pptx/gen")
    output_dir = output_dir.expanduser().resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_name = os.environ.get("PPTX_PPTX_NAME", "external_generated.pptx")
    output_path = output_dir / pptx_name

    slide_index_env = os.environ.get("PPTX_SLIDE_INDEX") or os.environ.get("PPTX_SLIDE_PAGE_NO")
    slide_index = int(slide_index_env) - 1 if slide_index_env and slide_index_env.isdigit() else 0
    elements, slide_meta, slide_payload = _pick_slide_payload(slides_payload, slide_index)
    blueprint_slots = slide_meta.get("blueprint_slots")
    if not isinstance(blueprint_slots, list):
        blueprint_slots = []
    target_layout = (
        (slide_payload.get("layout_id") if isinstance(slide_payload, dict) else None)
        or slide_meta.get("layout_name")
        or slide_meta.get("layout_id")
    )
    if output_path.exists():
        presentation = Presentation(output_path)
    else:
        presentation = Presentation(template_path)
    template_source = (
        doc_meta.get("template_source")
        if isinstance(doc_meta, dict)
        else None
    ) or slide_meta.get("template_source") or "template"
    prototype_index = slide_meta.get("prototype_index")
    layout = None
    slide = None

    layout = _find_layout(presentation, target_layout, blueprint_slots)
    if output_path.exists():
        if layout is None:
            layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
    else:
        if template_source == "slide":
            if isinstance(prototype_index, int) and 1 <= prototype_index <= len(presentation.slides):
                slide = presentation.slides[prototype_index - 1]
                layout = slide.slide_layout
            else:
                if presentation.slides:
                    slide = presentation.slides[0]
                else:
                    if layout is None:
                        layout = presentation.slide_layouts[0]
                    slide = presentation.slides.add_slide(layout)
        else:
            if presentation.slides:
                slide = presentation.slides[0]
                if layout is not None and slide.slide_layout != layout:
                    slide = presentation.slides.add_slide(layout)
            else:
                if layout is None:
                    layout = presentation.slide_layouts[0]
                slide = presentation.slides.add_slide(layout)

        # 余分なスライドは削除し、対象スライドのみ残す（関連リレーションも破棄）
        if slide is not None and len(presentation.slides) > 1:
            sld_id_lst = presentation.slides._sldIdLst  # type: ignore[attr-defined]
            slides_snapshot = list(presentation.slides)
            package = presentation.part.package
            for idx in reversed(range(len(slides_snapshot))):
                candidate = slides_snapshot[idx]
                if candidate is slide:
                    continue
                rel_id = sld_id_lst[idx].rId
                rel = presentation.part.rels.get(rel_id)
                target_part = rel.target_part if rel is not None else None
                presentation.part.drop_rel(rel_id)
                if target_part is not None:
                    try:
                        package.drop_part(target_part.partname)
                    except Exception:
                        pass
                sld_id_lst.remove(sld_id_lst[idx])
            slide = presentation.slides[0]

    layout_defaults: dict[str, list[str]] = {}
    if slide is not None:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_lines = [paragraph.text for paragraph in shape.text_frame.paragraphs]
            layout_defaults[getattr(shape, "name", "")] = text_lines
    elif layout is not None:
        for shape in layout.shapes:
            if not shape.has_text_frame:
                continue
            text_lines = [paragraph.text for paragraph in shape.text_frame.paragraphs]
            layout_defaults[getattr(shape, "name", "")] = text_lines

    if (
        template_source != "slide"
        and layout is not None
        and slide.shapes
        and len(slide.shapes) < len(layout.shapes)
    ):
        next_id = _max_shape_id(slide) + 1
        for layout_shape in layout.shapes:
            if getattr(layout_shape, "is_placeholder", False):
                continue
            new_element = deepcopy(layout_shape.element)
            for cNvPr in new_element.iter(qn("p:cNvPr")):
                cNvPr.set("id", str(next_id))
                next_id += 1
            slide.shapes._spTree.insert_element_before(new_element, "p:extLst")

    extract_summary = context.get("extract_summary") if isinstance(context, dict) else None
    table_payload = extract_summary.get("table") if isinstance(extract_summary, dict) else None
    detail_rows = context.get("table_detail_rows") if isinstance(context, dict) else None
    if detail_rows and not isinstance(detail_rows, list):
        detail_rows = None

    template_id = None
    if isinstance(doc_meta, dict):
        job_meta = doc_meta.get("job_meta")
        if isinstance(job_meta, dict):
            template_id = job_meta.get("template_id")

    if template_id == "executive_board" and len(slides_payload) >= 3:
        if output_path.exists():
            presentation = Presentation(output_path)
        else:
            presentation = Presentation(template_path)

        target_index = 2
        template_pptx = Presentation(template_path)
        source_slide = _find_cost_template_slide(template_pptx)
        if source_slide is None:
            raise RuntimeError("cost template slide not found (表 8/表 10)")

        if len(presentation.slides) > target_index:
            slide = presentation.slides[target_index]
        else:
            layout = _find_layout(presentation, getattr(source_slide.slide_layout, "name", None))
            if layout is None:
                layout = presentation.slide_layouts[0]
            slide = presentation.slides.add_slide(layout)

        _clear_slide_shapes(slide)
        _copy_shapes_from_slide(source_slide, slide)

        message_lines = _build_message_lines_from_summary(
            extract_summary.get("message_line") if isinstance(extract_summary, dict) else {}
        )
        if message_lines:
            _apply_text_to_named_shape(slide, "テキスト プレースホルダー 3", message_lines)
        _apply_tables_for_cost(slide, table_payload, detail_rows)
        apply_title_font_size(slide)
        apply_title_prefix(slide, 3)
        apply_strip_bullet_prefix(slide)

        presentation.save(output_path)
        print(f"[stage04_gen] PPTX generated -> {output_path}")
        return 0

    _apply_elements(slide, elements, blueprint_slots, table_payload, layout_defaults, detail_rows)
    if slide is not None:
        apply_title_font_size(slide)
        apply_title_prefix(slide, 3)
        apply_strip_bullet_prefix(slide)

    presentation.save(output_path)
    if output_path.exists():
        refreshed = Presentation(output_path)
        if refreshed.slides:
            _clear_table(refreshed.slides[0], "表 10")
        refreshed.save(output_path)

    print(f"[stage04_gen] PPTX generated -> {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
