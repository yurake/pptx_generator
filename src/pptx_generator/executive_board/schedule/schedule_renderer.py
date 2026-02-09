"""スケジュールガントチャート（表形式＋ホームベース）のレンダリング機能。"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt

from schedule.models import ScheduleGantt, ScheduleProject, ScheduleTask
from common.settings import BrandingConfig

logger = logging.getLogger(__name__)


def _get_fiscal_year(year: int, month: int) -> int:
    """暦年と月から年度を計算する（4月始まり）。
    
    Args:
        year: 暦年
        month: 月（1-12）
    
    Returns:
        年度（4-12月は当該年、1-3月は前年度）
    
    Examples:
        >>> _get_fiscal_year(2025, 4)  # 2025年4月
        2025
        >>> _get_fiscal_year(2025, 1)  # 2025年1月
        2024
    """
    if 1 <= month <= 3:
        return year - 1
    return year


def _get_task_color(task_name: str) -> str:
    """タスク名からタスクの背景色を判定する。
    
    Args:
        task_name: タスク名
    
    Returns:
        RGB色コード（16進数文字列）
    """
    task_name_lower = task_name.lower()
    
    # 白色：拠点要件精緻化、拠点展開、計画策定、契約締結手続
    white_keywords = ['拠点要件精緻化', '拠点展開', '計画策定', '契約締結手続']
    if any(keyword in task_name for keyword in white_keywords):
        return "#FFFFFF"
    
    # 明るく薄い緑色：移行関連
    light_green_keywords = ['移行', 'データ移行']
    if any(keyword in task_name for keyword in light_green_keywords):
        return "#C6E0B4"  # 明るく薄い緑
    
    # 明るい緑色：要件定義、設計、UAT
    green_keywords = ['要件定義', '設計', 'uat']
    if any(keyword in task_name_lower for keyword in green_keywords):
        return "#92D050"  # 明るい緑
    
    # ピンク色：開発、製造、構築、単体テスト
    pink_keywords = ['開発', '製造', '構築', '単体テスト']
    if any(keyword in task_name for keyword in pink_keywords):
        return "#FFC7CE"  # ピンク
    
    # 薄い黄色：結合テスト、性能テスト、障害テスト
    yellow_keywords = ['結合テスト', '性能テスト', '障害テスト']
    if any(keyword in task_name for keyword in yellow_keywords):
        return "#FFF2CC"  # 薄い黄色
    
    # 薄い青色：システムテスト、運用テスト
    cyan_keywords = ['システムテスト', '運用テスト']
    if any(keyword in task_name for keyword in cyan_keywords):
        return "#DEEBF7"  # 薄い青色
    
    # デフォルト：白色
    return "#FFFFFF"


@dataclass(slots=True)
class GanttRenderConfig:
    """ガントチャート描画の設定。"""

    slide_width_in: float = 10.0
    slide_height_in: float = 7.5
    title_top_in: float = 0.4
    title_height_in: float = 0.5
    title_left_in: float = 0.2
    title_width_in: float = 9.6
    table_top_in: float = 1.1
    table_left_in: float = 0.2
    table_width_in: float = 9.6
    project_col_width_in: float = 1.2
    task_col_width_in: float = 1.5


class ScheduleGanttRenderer:
    """スケジュールガントチャート（表形式＋ホームベース）のレンダラー。"""

    # プロジェクトごとに使用するカラーパレット
    COLOR_PALETTE = (
        "#4472C4",  # 青
        "#ED7D31",  # オレンジ
        "#A5A5A5",  # グレー
        "#FFC000",  # 黄
        "#5B9BD5",  # 水色
        "#70AD47",  # 緑
    )

    def __init__(
        self,
        *,
        template_path: Path | None = None,
        branding: BrandingConfig | None = None,
        config: GanttRenderConfig | None = None,
    ) -> None:
        """
        Args:
            template_path: テンプレートPPTXファイルパス
            branding: ブランド設定
            config: ガントチャート描画設定
        """
        self.template_path = template_path
        self.branding = branding or BrandingConfig.default()
        self.config = config or GanttRenderConfig()

    def render(
        self, schedule: ScheduleGantt, output_path: Path, layout_name: str = "System_layout"
    ) -> None:
        """スケジュールガントチャートをPPTXファイルとして保存する。

        Args:
            schedule: スケジュールデータ
            output_path: 出力PPTXファイルパス
            layout_name: 使用するレイアウト名（デフォルト: "System_layout"）
        """
        presentation = self._load_template()
        
        # テンプレート内の既存スライドを削除
        self._clear_existing_slides(presentation)
        
        layout = self._find_layout(presentation, layout_name)
        slide = presentation.slides.add_slide(layout)

        # 背景を白に設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Titleプレースホルダーにタイトルを設定
        self._set_title(slide, "プロジェクトのスケジュール")

        # Date_deptプレースホルダーに日付と部門を設定
        self._set_date_dept(slide, schedule)

        # メッセージラインプレースホルダーにメッセージを設定
        # JSONで指定されていればそれを使用、未指定なら自動生成
        message = schedule.meta.message_line or self._generate_message_line(schedule)
        self._set_message_line(slide, message)

        # Schedule_boxプレースホルダーを削除して表を配置
        self._remove_schedule_box_placeholder(slide)
        
        # 表形式でスケジュールを描画
        self._draw_schedule_table(slide, schedule)

        # 保存
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        logger.info("スケジュールガントチャートを保存しました: %s", output_path)

    def _load_template(self) -> Presentation:
        """テンプレートをロードする。"""
        if self.template_path and self.template_path.exists():
            logger.debug("テンプレートを使用: %s", self.template_path)
            return Presentation(self.template_path)
        logger.debug("既定テンプレートを利用")
        return Presentation()

    def _find_layout(self, presentation: Presentation, layout_name: str):
        """レイアウトを検索する。"""
        for layout in presentation.slide_layouts:
            if layout.name == layout_name:
                return layout
        logger.warning("レイアウト '%s' が見つからないため最初のレイアウトを使用", layout_name)
        if len(presentation.slide_layouts) == 0:
            msg = "テンプレートに利用可能なレイアウトが存在しません"
            raise RuntimeError(msg)
        return presentation.slide_layouts[0]

    def _clear_existing_slides(self, presentation: Presentation) -> None:
        """プレゼンテーション内の既存スライドを全て削除する。"""
        # スライドを逆順で削除（インデックスのずれを防ぐため）
        slide_count = len(presentation.slides)
        for i in range(slide_count - 1, -1, -1):
            rId = presentation.slides._sldIdLst[i].rId
            presentation.part.drop_rel(rId)
            del presentation.slides._sldIdLst[i]

    def _generate_message_line(self, schedule: ScheduleGantt) -> str:
        """スケジュールの概要を説明するメッセージラインを生成する。
        
        チェックポイント、移行リリース判定、移行時期などを含める。
        """
        start_year = schedule.meta.year
        end_year = schedule.meta.end_year or start_year
        start_month = schedule.meta.start_month
        end_month = schedule.meta.end_month
        
        # 期間の表記（年度形式）
        start_fiscal_year = _get_fiscal_year(start_year, start_month)
        end_fiscal_year = _get_fiscal_year(end_year, end_month)
        
        if start_fiscal_year == end_fiscal_year:
            period = f"{start_fiscal_year}年度{start_month}月～{end_month}月"
        else:
            period = f"{start_fiscal_year}年度{start_month}月～{end_fiscal_year}年度{end_month}月"
        
        # マイルストーン情報を収集
        milestones = schedule.meta.milestones
        checkpoint_info = []
        migration_info = []
        
        for milestone in milestones:
            milestone_date = datetime.strptime(milestone.date, "%Y-%m-%d")
            fiscal_year = _get_fiscal_year(milestone_date.year, milestone_date.month)
            date_str = f"{fiscal_year}年度{milestone_date.month}月"
            
            # 完了・承認系のマイルストーン
            if '完了' in milestone.name:
                checkpoint_info.append(f"{milestone.name}は{date_str}")
            
            # 移行・稼働系のマイルストーン
            if '稼働' in milestone.name or 'リリース' in milestone.name:
                migration_info.append((milestone.name, date_str))
        
        # 移行タスクの開始・終了情報を収集
        migration_tasks = []
        for project in schedule.projects:
            for task in project.tasks:
                if 'データ移行' in task.name:
                    start_date = datetime.strptime(task.start, "%Y-%m-%d")
                    end_date = datetime.strptime(task.end, "%Y-%m-%d")
                    start_fiscal_year = _get_fiscal_year(start_date.year, start_date.month)
                    end_fiscal_year = _get_fiscal_year(end_date.year, end_date.month)
                    start_str = f"{start_fiscal_year}年度{start_date.month}月"
                    end_str = f"{end_fiscal_year}年度{end_date.month}月"
                    if '実施' in task.name:
                        migration_tasks.append(('実施', start_str, end_str))
                    elif '準備' in task.name:
                        migration_tasks.append(('準備', start_str, end_str))
        
        # メッセージを構築
        message_parts = []
        
        # 基本情報
        message_parts.append(f"本スケジュールは{period}のプロジェクト計画を示しています。")
        
        # チェックポイント
        if checkpoint_info:
            checkpoints = '、'.join(checkpoint_info)
            message_parts.append(f"主要なチェックポイントとして、{checkpoints}を設定しています。")
        
        # 移行情報
        migration_text = []
        for task_type, start_str, end_str in migration_tasks:
            if task_type == '準備':
                migration_text.append(f"移行準備は{start_str}から開始し{end_str}に完了")
            elif task_type == '実施':
                migration_text.append(f"移行実施は{start_str}から{end_str}まで")
        
        if migration_text:
            message_parts.append('、'.join(migration_text) + "を予定しています。")
        
        # 移行リリース判定
        if migration_info:
            for name, date_str in migration_info:
                message_parts.append(f"{name}は{date_str}を予定しています。")
        
        message = '　'.join(message_parts)
        
        return message

    def _set_title(self, slide, title: str) -> None:
        """Titleプレースホルダーにタイトルを設定する（idx: 0）。"""
        for shape in slide.shapes:
            if shape.is_placeholder and hasattr(shape, 'text_frame'):
                ph_format = shape.placeholder_format
                if ph_format.idx == 0:  # Title
                    text_frame = shape.text_frame
                    text_frame.clear()
                    paragraph = text_frame.paragraphs[0]
                    paragraph.text = title
                    logger.info("Titleを設定しました")
                    return
        logger.warning("Titleプレースホルダー（idx: 0）が見つかりませんでした")

    def _set_date_dept(self, slide, schedule: ScheduleGantt) -> None:
        """Date_deptプレースホルダーに日付と部門を設定する（idx: 10）。"""
        # 現在の日付を取得
        today = datetime.now().strftime("%Y年%m月%d日")
        # 部門情報（実際には schedule.meta から取得することも可能）
        dept = "システム部"  # デフォルト値
        date_dept_text = f"{today} {dept}"
        
        for shape in slide.shapes:
            if shape.is_placeholder and hasattr(shape, 'text_frame'):
                ph_format = shape.placeholder_format
                if ph_format.idx == 10:  # Date_dept
                    text_frame = shape.text_frame
                    text_frame.clear()
                    paragraph = text_frame.paragraphs[0]
                    paragraph.text = date_dept_text
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
                    logger.info("Date_deptを設定しました")
                    return
        logger.warning("Date_deptプレースホルダー（idx: 10）が見つかりませんでした")

    def _set_message_line(self, slide, message: str) -> None:
        """Message_lineプレースホルダーにメッセージを設定する（idx: 12）。"""
        for shape in slide.shapes:
            if shape.is_placeholder and hasattr(shape, 'text_frame'):
                ph_format = shape.placeholder_format
                if ph_format.idx == 12:  # Message_line
                    text_frame = shape.text_frame
                    text_frame.clear()
                    paragraph = text_frame.paragraphs[0]
                    paragraph.text = message
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    # フォントを設定
                    font = paragraph.font
                    font.size = Pt(16)
                    font.bold = True
                    logger.info("Message_lineにメッセージを設定しました")
                    return
        
        # プレースホルダーが見つからない場合はテキストボックスを作成（フォールバック）
        logger.info("Message_lineプレースホルダーが見つからないため、テキストボックスを作成します")
        textbox = slide.shapes.add_textbox(
            Inches(self.config.title_left_in),
            Inches(self.config.title_top_in),
            Inches(self.config.title_width_in),
            Inches(self.config.title_height_in),
        )
        text_frame = textbox.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = message
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        font = paragraph.font
        font.size = Pt(16)
        font.bold = True

    def _remove_schedule_box_placeholder(self, slide) -> None:
        """Schedule_boxプレースホルダーを削除する。"""
        for shape in list(slide.shapes):
            if shape.is_placeholder and shape.name in ("Schedule_box", "Schedule_box_title"):
                sp = shape.element
                sp.getparent().remove(sp)

    def _draw_schedule_table(self, slide, schedule: ScheduleGantt) -> None:
        """表形式でスケジュールを描画する。"""
        projects = schedule.projects
        milestones = schedule.meta.milestones
        
        # 表示単位を取得
        display_unit = schedule.meta.get_display_unit()
        
        # 表示単位に応じた列数を計算
        if display_unit == "month":
            # 月単位の場合、年を跨ぐ可能性を考慮
            start_year = schedule.meta.year
            end_year = schedule.meta.end_year or start_year
            start_month = schedule.meta.start_month
            end_month = schedule.meta.end_month
            
            period_count = (end_year - start_year) * 12 + (end_month - start_month + 1)
        elif display_unit == "quarter":
            # 四半期単位
            period_count = self._calculate_quarter_count(schedule.meta)
        else:  # year
            # 年単位
            start_year = schedule.meta.year
            end_year = schedule.meta.end_year or start_year
            period_count = end_year - start_year + 1

        # プロジェクト数を計算（プロジェクトごとに1行）
        project_count = len(projects)
        
        # 四半期表示の場合はヘッダーが2行、それ以外は1行
        header_rows = 2 if display_unit == "quarter" else 1
        
        # ヘッダー行 + マイルストーン行 + プロジェクト行
        rows = header_rows + 1 + project_count
        cols = 1 + period_count  # プロジェクト列 + 期間列

        # 行の高さを計算
        # ヘッダー行とマイルストーン行は固定の小さい高さ
        header_row_height_in = 0.3  # ヘッダー行の固定高さ
        milestone_row_height_in = 0.3  # マイルストーン行の固定高さ
        
        # プロジェクト行の高さは矢羽の高さに合わせて固定
        # 矢羽の高さを0.8cm（約0.315インチ）に制限
        project_row_height_in = 0.315  # 0.8cm = 0.315インチ
        
        # 実際のテーブル高さを計算
        fixed_height_in = header_row_height_in * header_rows + milestone_row_height_in
        table_height_in = fixed_height_in + project_row_height_in * project_count
        
        logger.info(f"表の高さ: {table_height_in:.3f}インチ ({table_height_in * 2.54:.2f}cm)")
        
        # フォントサイズを設定
        header_font_size_pt = 9  # ヘッダー用フォントサイズ
        milestone_font_size_pt = 8  # マイルストーン用フォントサイズ
        
        # プロジェクト行のフォントサイズを行高に応じて調整
        if project_row_height_in >= 0.4:
            project_font_size_pt = 10
        elif project_row_height_in <= 0.2:
            project_font_size_pt = 6
        else:
            # 線形補間: 0.2インチで6pt、0.4インチで10pt
            project_font_size_pt = 6 + (project_row_height_in - 0.2) / 0.2 * 4
        project_font_size_pt = max(6, min(10, project_font_size_pt))

        # Schedule_boxプレースホルダーの位置を取得
        table_left_in = self.config.table_left_in
        table_top_in = self.config.table_top_in
        table_width_in = self.config.table_width_in
        
        # 最大幅を26cm（約10.24インチ）に制限
        max_table_width_in = 10.24  # 26cm = 10.24インチ
        table_width_in = min(table_width_in, max_table_width_in)
        
        for shape in slide.shapes:
            if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.idx == 27:  # Schedule_box
                    # Schedule_boxの位置とサイズを使用
                    table_left_in = shape.left / 914400  # EMUからインチに変換
                    table_top_in = shape.top / 914400
                    table_width_in = shape.width / 914400
                    logger.info("Schedule_boxの位置を使用: left=%.2f, top=%.2f, width=%.2f",
                              table_left_in, table_top_in, table_width_in)
                    break

        # 表を作成
        table_shape = slide.shapes.add_table(
            rows,
            cols,
            int(table_left_in * 914400),
            int(table_top_in * 914400),
            int(table_width_in * 914400),
            int(table_height_in * 914400),
        )
        table = table_shape.table

        # 列幅を設定
        table.columns[0].width = Inches(self.config.project_col_width_in)
        period_col_width_in = (
            self.config.table_width_in - self.config.project_col_width_in
        ) / period_count
        for i in range(1, cols):
            table.columns[i].width = Inches(period_col_width_in)

        # 行の高さを設定（ヘッダー、マイルストーン、プロジェクトで異なる高さ）
        # ヘッダー行
        for i in range(header_rows):
            table.rows[i].height = Inches(header_row_height_in)
        # マイルストーン行
        table.rows[header_rows].height = Inches(milestone_row_height_in)
        # プロジェクト行
        for i in range(header_rows + 1, rows):
            table.rows[i].height = Inches(project_row_height_in)

        # ヘッダー行を設定（表示単位に応じて変更）
        if display_unit == "quarter":
            # 四半期表示：2行ヘッダー
            # 1行目と2行目のプロジェクト列を統合して「スケジュール」と表記
            
            # 1行目と2行目のプロジェクト列を統合
            cell_0_0 = table.cell(0, 0)
            cell_1_0 = table.cell(1, 0)
            cell_0_0.merge(cell_1_0)
            cell_0_0.text = "スケジュール"
            
            # 期間ヘッダーを設定（1行目：年度、2行目：四半期）
            self._set_period_headers(table, schedule.meta, display_unit, period_count)
            
            # 1行目のプロジェクト列（統合されたセル）のスタイル設定
            self._style_header_cell(cell_0_0, header_font_size_pt)
            
            # 2行目のヘッダー行のスタイル設定（プロジェクト列以外）
            for col_idx in range(1, cols):
                cell = table.cell(1, col_idx)
                self._style_header_cell(cell, header_font_size_pt)
            
            # マイルストーン行（3行目）
            milestone_row_idx = 2
        else:
            # 月表示・年表示：1行ヘッダー
            table.cell(0, 0).text = "スケジュール"
            self._set_period_headers(table, schedule.meta, display_unit, period_count)
            
            # ヘッダー行のスタイル設定
            for col_idx in range(cols):
                cell = table.cell(0, col_idx)
                self._style_header_cell(cell, header_font_size_pt)
            
            # マイルストーン行（2行目）
            milestone_row_idx = 1

        # マイルストーン行を設定（文言は空欄）
        milestone_cell = table.cell(milestone_row_idx, 0)
        milestone_cell.text = ""  # マイルストーンの文言を削除
        milestone_cell.vertical_anchor = 1
        self._style_data_cell(milestone_cell, milestone_font_size_pt)
        
        # マイルストーン行の空セルをスタイル設定
        for i in range(1, cols):
            cell = table.cell(milestone_row_idx, i)
            self._style_data_cell(cell, milestone_font_size_pt)

        # 各プロジェクトを行に追加（プロジェクトごとに1行）
        current_row = milestone_row_idx + 1
        for project_idx, project in enumerate(projects):
            # プロジェクト名セル
            project_cell = table.cell(current_row, 0)
            project_cell.text = project.name
            project_cell.vertical_anchor = 1
            self._style_project_cell(project_cell, project_font_size_pt)

            # 期間列セルのスタイル設定
            for i in range(1, cols):
                self._style_data_cell(table.cell(current_row, i), project_font_size_pt)

            # プロジェクトの全タスクを矢羽で配置（重複を考慮してトラック配置）
            self._draw_project_tasks_with_tracks(
                slide, table_shape, table, current_row, project, schedule.meta,
                period_col_width_in, project_row_height_in, display_unit, project_font_size_pt
            )

            current_row += 1
        
        # マイルストーンを下向き三角形とテキストボックスで配置
        self._draw_milestones_as_triangles(
            slide, table_shape, table, schedule.meta, milestones,
            period_col_width_in, milestone_row_height_in, display_unit, milestone_font_size_pt
        )

    def _calculate_quarter_count(self, meta) -> int:
        """四半期の数を計算する。"""
        start_year = meta.year
        end_year = meta.end_year or start_year
        start_quarter = (meta.start_month - 1) // 3 + 1
        end_quarter = (meta.end_month - 1) // 3 + 1
        
        years = end_year - start_year
        quarters = years * 4 + (end_quarter - start_quarter + 1)
        return quarters

    def _set_period_headers(self, table, meta, display_unit: str, period_count: int) -> None:
        """表示単位に応じたヘッダーを設定する。"""
        start_year = meta.year
        start_month = meta.start_month
        
        if display_unit == "month":
            # 月単位のヘッダー（1行）
            current_year = start_year
            current_month = start_month
            for i in range(period_count):
                if current_year == start_year:
                    header_text = f"{current_month}月"
                else:
                    header_text = f"{current_year}年{current_month}月"
                table.cell(0, i + 1).text = header_text
                
                current_month += 1
                if current_month > 12:
                    current_month = 1
                    current_year += 1
                    
        elif display_unit == "quarter":
            # 四半期単位のヘッダー（2行：1行目=年度、2行目=四半期）
            start_quarter = (start_month - 1) // 3 + 1
            current_year = start_year
            current_quarter = start_quarter
            
            # 各四半期の年度を計算してグループ化
            fiscal_year_groups = []  # [(fiscal_year, start_col, count), ...]
            current_fiscal_year = None
            fiscal_year_start_col = None
            fiscal_year_col_count = 0
            
            for i in range(period_count):
                # 四半期の開始月を計算
                # Q1=1月, Q2=4月, Q3=7月, Q4=10月
                quarter_start_month = (current_quarter - 1) * 3 + 1
                
                # 年度を計算
                fiscal_year = _get_fiscal_year(current_year, quarter_start_month)
                
                # 年度グループの開始・終了を記録
                if fiscal_year != current_fiscal_year:
                    # 前のグループを保存
                    if current_fiscal_year is not None:
                        fiscal_year_groups.append((current_fiscal_year, fiscal_year_start_col, fiscal_year_col_count))
                    # 新しいグループ開始
                    current_fiscal_year = fiscal_year
                    fiscal_year_start_col = i + 1
                    fiscal_year_col_count = 1
                else:
                    fiscal_year_col_count += 1
                
                # 2行目：四半期ラベル（"1Q", "2Q", etc.）
                table.cell(1, i + 1).text = f"{current_quarter}Q"
                
                # 次の四半期へ
                current_quarter += 1
                if current_quarter > 4:
                    current_quarter = 1
                    current_year += 1
            
            # 最後の年度グループを追加
            if current_fiscal_year is not None:
                fiscal_year_groups.append((current_fiscal_year, fiscal_year_start_col, fiscal_year_col_count))
            
            # 1行目：年度ラベル（同じ年度の列をマージ）
            for fiscal_year, start_col, col_count in fiscal_year_groups:
                if col_count > 1:
                    # 複数列をマージ
                    cell = table.cell(0, start_col)
                    end_col = start_col + col_count - 1
                    cell.merge(table.cell(0, end_col))
                    cell.text = f"{fiscal_year}年度"
                else:
                    # 単一列
                    table.cell(0, start_col).text = f"{fiscal_year}年度"
                
                # 1行目のスタイル設定
                cell = table.cell(0, start_col)
                self._style_header_cell(cell, 10)  # フォントサイズは後で調整される
                    
        else:  # year
            # 年単位のヘッダー（1行）
            for i in range(period_count):
                year = start_year + i
                table.cell(0, i + 1).text = f"{year}年"

    def _draw_milestones_as_triangles(
        self,
        slide,
        table_shape,
        table,
        meta,
        milestones,
        period_col_width_in: float,
        row_height_in: float,
        display_unit: str,
        font_size_pt: float,
    ) -> None:
        """マイルストーンを下向き三角形とテキストボックスで表示する。"""
        start_year = meta.year
        end_year = meta.end_year or start_year
        start_month = meta.start_month
        end_month = meta.end_month
        
        # 表全体の位置
        table_left = table_shape.left
        table_top = table_shape.top
        
        # プロジェクト列の幅
        col_0_width = table.columns[0].width
        
        # マイルストーン行のインデックス（四半期表示は2行目、それ以外は1行目）
        milestone_row_idx = 2 if display_unit == "quarter" else 1
        
        for milestone in milestones:
            milestone_date = datetime.strptime(milestone.date, "%Y-%m-%d")
            milestone_year = milestone_date.year
            milestone_month = milestone_date.month
            
            # スケジュール範囲外はスキップ
            if milestone_year < start_year or milestone_year > end_year:
                continue
            if milestone_year == start_year and milestone_month < start_month:
                continue
            if milestone_year == end_year and milestone_month > end_month:
                continue
            
            col_idx = self._get_period_column_index(
                milestone_date, meta, display_unit
            )
            if col_idx is None:
                continue
            
            # マイルストーン列のセル位置を計算
            cell_left = table_left + col_0_width + Inches(period_col_width_in * (col_idx - 1))
            cell_top = table_top + Inches(row_height_in * milestone_row_idx)  # マイルストーン行
            cell_width = Inches(period_col_width_in)
            cell_height = Inches(row_height_in)
            
            # 三角形のサイズ（フォントサイズに基づいて小さく）
            triangle_size = Pt(font_size_pt * 1.2)  # フォントサイズの1.2倍
            
            # 下向き三角形を作成（セルの中央上部に配置）
            triangle_left = cell_left + (cell_width - triangle_size) / 2
            triangle_top = cell_top + Inches(0.05)
            
            triangle = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                triangle_left,
                triangle_top,
                triangle_size,
                triangle_size,
            )
            
            # 三角形を180度回転（下向きに）
            triangle.rotation = 180
            
            # 塗りつぶし（白色）
            fill = triangle.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色
            
            # 枠線（黒色）
            line = triangle.line
            line.color.rgb = RGBColor(0, 0, 0)  # 黒色
            line.width = Pt(1)
            
            # テキストボックスを三角形の右隣に近接配置
            # 間隔を0.1cm（約0.039インチ）に設定
            textbox_left = cell_left + (cell_width + triangle_size) / 2 + Inches(0.039)
            # 高さをマイルストーン行の高さに合わせる
            textbox_top = cell_top
            textbox_width = Inches(1.5)  # 固定幅（やや広く）
            textbox_height = cell_height
            
            textbox = slide.shapes.add_textbox(
                textbox_left,
                textbox_top,
                textbox_width,
                textbox_height,
            )
            
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.vertical_anchor = 1  # 中央揃え
            paragraph = text_frame.paragraphs[0]
            paragraph.text = milestone.name
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            
            # フォント設定
            font = paragraph.font
            font.size = Pt(font_size_pt)
            font.bold = True
            font.color.rgb = RGBColor(0, 0, 0)  # 黒色

    def _get_period_column_index(self, date: datetime, meta, display_unit: str) -> int | None:
        """日付から対応する列インデックスを計算する。"""
        start_year = meta.year
        start_month = meta.start_month
        
        if display_unit == "month":
            # 月単位
            months_diff = (date.year - start_year) * 12 + (date.month - start_month)
            return 1 + months_diff
            
        elif display_unit == "quarter":
            # 四半期単位
            start_quarter = (start_month - 1) // 3 + 1
            date_quarter = (date.month - 1) // 3 + 1
            quarters_diff = (date.year - start_year) * 4 + (date_quarter - start_quarter)
            return 1 + quarters_diff
            
        else:  # year
            # 年単位
            years_diff = date.year - start_year
            return 1 + years_diff

    def _draw_project_tasks_with_tracks(
        self,
        slide,
        table_shape,
        table,
        row_idx: int,
        project: ScheduleProject,
        meta,
        period_col_width_in: float,
        row_height_in: float,
        display_unit: str,
        font_size_pt: float,
    ) -> None:
        """プロジェクトの全タスクをトラック配置で描画する。
        
        期間が重複するタスクは縦にずらして配置し、行の高さを調整する。
        """
        # タスクの期間重複を検出してトラックを割り当て
        tasks_with_tracks = self._assign_tracks_to_tasks(project.tasks, meta, display_unit)
        
        # 必要なトラック数を計算
        max_track = max(track for _, track in tasks_with_tracks)
        tracks_needed = max_track + 1
        
        # トラックの高さは固定（矢羽の最大高さ0.8cm = 0.315インチ）
        track_height = 0.315  # インチ
        
        # 行の高さを調整（トラック数に応じて）
        if tracks_needed > 1:
            new_row_height = track_height * tracks_needed
            table.rows[row_idx].height = Inches(new_row_height)
        else:
            new_row_height = row_height_in
            table.rows[row_idx].height = Inches(new_row_height)
        
        for task, track in tasks_with_tracks:
            self._draw_task_arrow(
                slide, table_shape, table, row_idx, task, meta,
                period_col_width_in, track_height, track, display_unit, font_size_pt
            )
    
    def _assign_tracks_to_tasks(
        self,
        tasks: list[ScheduleTask],
        meta,
        display_unit: str,
    ) -> list[tuple[ScheduleTask, int]]:
        """タスクの期間重複を検出してトラック番号を割り当てる。
        
        Returns:
            タスクとトラック番号のペアのリスト
        """
        # タスクを開始日順にソート
        sorted_tasks = sorted(
            tasks,
            key=lambda t: datetime.strptime(t.start, "%Y-%m-%d")
        )
        
        # トラック割り当て（グリーディアルゴリズム）
        tracks: list[list[tuple[datetime, datetime]]] = []  # 各トラックの占有期間リスト
        task_tracks: list[tuple[ScheduleTask, int]] = []
        
        for task in sorted_tasks:
            task_start = datetime.strptime(task.start, "%Y-%m-%d")
            task_end = datetime.strptime(task.end, "%Y-%m-%d")
            
            # 既存のトラックで配置可能か確認
            placed = False
            for track_idx, track_periods in enumerate(tracks):
                # このトラックで重複がないか確認
                overlap = False
                for period_start, period_end in track_periods:
                    if task_start <= period_end and task_end >= period_start:
                        overlap = True
                        break
                
                if not overlap:
                    # このトラックに配置可能
                    track_periods.append((task_start, task_end))
                    task_tracks.append((task, track_idx))
                    placed = True
                    break
            
            if not placed:
                # 新しいトラックを作成
                tracks.append([(task_start, task_end)])
                task_tracks.append((task, len(tracks) - 1))
        
        return task_tracks
    
    def _draw_task_arrow(
        self,
        slide,
        table_shape,
        table,
        row_idx: int,
        task: ScheduleTask,
        meta,
        period_col_width_in: float,
        track_height_in: float,
        track_idx: int,
        display_unit: str,
        font_size_pt: float,
    ) -> None:
        """タスクを矢羽図形で描画する（トラック配置対応、日付ベースの精緻な位置計算）。"""
        start_date = datetime.strptime(task.start, "%Y-%m-%d")
        end_date = datetime.strptime(task.end, "%Y-%m-%d")
        
        start_year = meta.year
        end_year = meta.end_year or start_year
        start_month = meta.start_month
        end_month = meta.end_month

        # スケジュール範囲外のタスクはスキップ
        if end_date.year < start_year or start_date.year > end_year:
            logger.warning("タスク '%s' がスケジュール範囲外です", task.name)
            return
        if start_date.year == end_year and start_date.month > end_month:
            logger.warning("タスク '%s' がスケジュール範囲外です", task.name)
            return
        if end_date.year == start_year and end_date.month < start_month:
            logger.warning("タスク '%s' がスケジュール範囲外です", task.name)
            return

        # 表全体の位置
        table_left = table_shape.left
        table_top = table_shape.top

        # プロジェクト列の幅
        col_0_width = table.columns[0].width

        # トラック内の位置を計算
        track_top_offset = Inches(track_height_in * track_idx)
        arrow_top = table_top + Inches(sum(table.rows[i].height / 914400 for i in range(row_idx))) + track_top_offset

        # 精緻な位置計算（日付ベース）
        start_offset, end_offset = self._calculate_precise_position(
            start_date, end_date, meta, display_unit, period_col_width_in
        )
        
        # 矢羽の開始位置
        arrow_left = table_left + col_0_width + start_offset
        # 矢羽の幅
        arrow_width = end_offset - start_offset
        arrow_height = Inches(track_height_in)

        # タスク名から背景色を取得
        task_color = _get_task_color(task.name)

        # 矢羽図形（五角形）を作成
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.PENTAGON,
            arrow_left + Inches(0.02),
            arrow_top + Inches(0.05),
            arrow_width - Inches(0.04),
            arrow_height - Inches(0.1),
        )

        # 塗りつぶし（タスク名に応じた色）
        fill = arrow.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(task_color.lstrip("#"))

        # 枠線（青色）
        line = arrow.line
        line.color.rgb = RGBColor(0, 112, 192)  # 青色 #0070C0
        line.width = Pt(1.5)
        
        # タスク名をテキストボックスとして矢羽内に追加
        text_frame = arrow.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        paragraph = text_frame.paragraphs[0]
        paragraph.text = task.name
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        # フォント設定
        font = paragraph.font
        font.size = Pt(max(6, font_size_pt * 0.8))
        font.bold = False
        font.color.rgb = RGBColor(0, 0, 0)  # 黒色
    
    def _calculate_precise_position(
        self,
        start_date: datetime,
        end_date: datetime,
        meta,
        display_unit: str,
        period_col_width_in: float,
    ) -> tuple[int, int]:
        """タスクの開始・終了位置を日付ベースで精緻に計算する（EMU単位）。
        
        Returns:
            (start_offset_emu, end_offset_emu): 列の開始位置からのオフセット（EMU単位）
        """
        start_year = meta.year
        start_month = meta.start_month
        
        if display_unit == "month":
            # 月表示：月の中の日数で位置を計算
            # 開始位置：開始月の何日目か
            start_col_idx = (start_date.year - start_year) * 12 + (start_date.month - start_month)
            days_in_start_month = self._days_in_month(start_date.year, start_date.month)
            start_day_ratio = (start_date.day - 1) / days_in_start_month
            
            # 終了位置：終了月の何日目か
            end_col_idx = (end_date.year - start_year) * 12 + (end_date.month - start_month)
            days_in_end_month = self._days_in_month(end_date.year, end_date.month)
            end_day_ratio = end_date.day / days_in_end_month
            
            # EMU単位で計算（1インチ = 914400 EMU）
            col_width_emu = int(period_col_width_in * 914400)
            start_offset_emu = int(start_col_idx * col_width_emu + start_day_ratio * col_width_emu)
            end_offset_emu = int(end_col_idx * col_width_emu + end_day_ratio * col_width_emu)
            
        elif display_unit == "quarter":
            # 四半期表示：四半期内の日数で位置を計算
            start_quarter = (start_month - 1) // 3 + 1
            start_col_idx = (start_date.year - start_year) * 4 + ((start_date.month - 1) // 3 + 1 - start_quarter)
            
            # 四半期の開始日を計算
            quarter_start_month = ((start_date.month - 1) // 3) * 3 + 1
            quarter_start = datetime(start_date.year, quarter_start_month, 1)
            
            # 四半期内の経過日数
            days_from_quarter_start = (start_date - quarter_start).days
            # 四半期の総日数（約90日）
            days_in_quarter = self._days_in_quarter(start_date.year, quarter_start_month)
            start_day_ratio = days_from_quarter_start / days_in_quarter
            
            # 終了位置
            end_col_idx = (end_date.year - start_year) * 4 + ((end_date.month - 1) // 3 + 1 - start_quarter)
            end_quarter_start_month = ((end_date.month - 1) // 3) * 3 + 1
            end_quarter_start = datetime(end_date.year, end_quarter_start_month, 1)
            
            # 四半期内の経過日数（終了日まで）
            days_to_end = (end_date - end_quarter_start).days + 1  # 終了日を含む
            days_in_end_quarter = self._days_in_quarter(end_date.year, end_quarter_start_month)
            end_day_ratio = days_to_end / days_in_end_quarter
            
            # EMU単位で計算
            col_width_emu = int(period_col_width_in * 914400)
            start_offset_emu = int(start_col_idx * col_width_emu + start_day_ratio * col_width_emu)
            end_offset_emu = int(end_col_idx * col_width_emu + end_day_ratio * col_width_emu)
            
        else:  # year
            # 年表示：年内の日数で位置を計算
            start_col_idx = start_date.year - start_year
            # 年の開始からの経過日数
            year_start = datetime(start_date.year, 1, 1)
            days_from_year_start = (start_date - year_start).days
            days_in_year = 366 if self._is_leap_year(start_date.year) else 365
            start_day_ratio = days_from_year_start / days_in_year
            
            # 終了位置
            end_col_idx = end_date.year - start_year
            end_year_start = datetime(end_date.year, 1, 1)
            days_to_end = (end_date - end_year_start).days + 1
            days_in_end_year = 366 if self._is_leap_year(end_date.year) else 365
            end_day_ratio = days_to_end / days_in_end_year
            
            # EMU単位で計算
            col_width_emu = int(period_col_width_in * 914400)
            start_offset_emu = int(start_col_idx * col_width_emu + start_day_ratio * col_width_emu)
            end_offset_emu = int(end_col_idx * col_width_emu + end_day_ratio * col_width_emu)
        
        return start_offset_emu, end_offset_emu
    
    def _days_in_month(self, year: int, month: int) -> int:
        """指定月の日数を返す。"""
        if month == 12:
            next_month = datetime(year + 1, 1, 1)
        else:
            next_month = datetime(year, month + 1, 1)
        current_month = datetime(year, month, 1)
        return (next_month - current_month).days
    
    def _days_in_quarter(self, year: int, quarter_start_month: int) -> int:
        """四半期の日数を返す（3ヶ月分）。"""
        total_days = 0
        for i in range(3):
            month = quarter_start_month + i
            if month > 12:
                month -= 12
                year += 1
            total_days += self._days_in_month(year, month)
        return total_days
    
    def _is_leap_year(self, year: int) -> bool:
        """うるう年判定。"""
        return year % 4 == 0 and (year % 100 != 0 or year % 400 == 0)

    def _style_header_cell(self, cell, font_size_pt: float) -> None:
        """ヘッダーセルのスタイルを設定する（青背景・白文字）。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)  # 青色 #4472C4
        
        text_frame = cell.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        font = paragraph.font
        font.bold = True
        font.size = Pt(font_size_pt)
        font.color.rgb = RGBColor(255, 255, 255)  # 白色

    def _style_data_cell(self, cell, font_size_pt: float) -> None:
        """データセルのスタイルを設定する（白背景・黒文字）。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色
        
        text_frame = cell.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        font = paragraph.font
        font.size = Pt(font_size_pt)
        font.color.rgb = RGBColor(0, 0, 0)  # 黒色

    def _style_project_cell(self, cell, font_size_pt: float) -> None:
        """プロジェクトセルのスタイルを設定する（薄い水色背景・黒文字）。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(217, 225, 242)  # 薄い水色 #D9E1F2
        
        text_frame = cell.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        font = paragraph.font
        font.size = Pt(font_size_pt)
        font.color.rgb = RGBColor(0, 0, 0)  # 黒色

    def _apply_font(self, paragraph, branding_font) -> None:
        """ブランドフォントを適用する。"""
        font = paragraph.font
        font.name = branding_font.name
        font.size = Pt(branding_font.size_pt)
        font.color.rgb = RGBColor.from_string(branding_font.color_hex.lstrip("#"))
        font.bold = branding_font.bold
        font.italic = branding_font.italic
