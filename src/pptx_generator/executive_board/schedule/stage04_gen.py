#!/usr/bin/env python3
"""Stage4 hook: schedule_data.json を用いて PPTX を生成する。"""
from __future__ import annotations

import json
import os
import sys
from pathlib import Path
from typing import Any

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from common.stage_shared import load_context, resolve_local_path
from common.rendering.title_style import apply_title_font_size, apply_title_prefix
from common.rendering.text_cleanup import apply_strip_bullet_prefix
from common.rendering.text_style import apply_text_style_by_name
from pptx import Presentation


def _configure_sys_path() -> None:
    src_dir = Path(__file__).resolve().parents[3] / "src"
    if str(src_dir) not in sys.path:
        sys.path.insert(0, str(src_dir))


_configure_sys_path()

from schedule_renderer import ScheduleGanttRenderer  # noqa: E402
from schedule_parser import parse_schedule_markdown, save_schedule_json  # noqa: E402


DEFAULT_TEMPLATE = Path(__file__).resolve().parents[3] / "templates/schedule.pptx"


class AppendRenderer(ScheduleGanttRenderer):  # type: ignore[misc]
    """既存 PPTX をベースに新規スライドを追加するレンダラー。"""

    def _prepare_slide(self, presentation, layout_name: str):
        layout = self._find_layout(presentation, layout_name)
        return presentation.slides.add_slide(layout)

    def _clear_existing_slides(self, presentation) -> None:  # override
        # 既存スライドは残したまま末尾に追加する
        return


def _resolve_template_path(context: dict[str, Any]) -> Path:
    candidates: list[str] = []
    env_template = os.environ.get("PPTX_TEMPLATE_PATH")
    if isinstance(env_template, str) and env_template.strip():
        candidates.append(env_template.strip())
    context_template = context.get("template_path")
    if isinstance(context_template, str) and context_template.strip():
        candidates.append(context_template.strip())
    candidates.append(str(DEFAULT_TEMPLATE))

    for candidate in candidates:
        resolved = resolve_local_path(candidate, Path(__file__).resolve().parent)
        if resolved.exists():
            return resolved
    raise FileNotFoundError("Template PPTX not found. Provide PPTX_TEMPLATE_PATH or ensure default exists.")


def _load_schedule(path: Path):
    text = path.read_text(encoding="utf-8")
    try:
        data = json.loads(text)
    except json.JSONDecodeError:
        return parse_schedule_markdown(path)
    from schedule.models import ScheduleGantt

    return ScheduleGantt.model_validate(data)


def _load_schedule_from_context(context: dict[str, Any]) -> ScheduleGantt:
    schedule_md = context.get("schedule_md_path")
    env_md = os.environ.get("PPTX_SCHEDULE_MD")
    schedule_json = Path(context.get("schedule_json_path", "")).expanduser() if context.get("schedule_json_path") else None
    candidate = None
    for raw in (env_md, schedule_md):
        if isinstance(raw, str) and raw.strip():
            path = Path(raw).expanduser().resolve()
            if path.exists():
                candidate = path
                break
    if candidate is None and schedule_json and schedule_json.exists():
        return _load_schedule(schedule_json)
    if candidate is None:
        raise FileNotFoundError("schedule markdown path not found (env/context)")
    return parse_schedule_markdown(candidate)


def main(argv: list[str] | None = None) -> int:
    parser = None  # placeholder
    _ = argv, parser

    context = load_context()

    output_dir_env = os.environ.get("PPTX_OUTPUT_DIR")
    pptx_name_env = os.environ.get("PPTX_PPTX_NAME")
    output_pptx_env = os.environ.get("PPTX_OUTPUT_PPTX_PATH")

    if output_pptx_env:
        output_path = Path(output_pptx_env).expanduser().resolve()
        output_dir = output_path.parent
        pptx_name = output_path.name
    else:
        output_dir = Path(output_dir_env) if output_dir_env else Path(".pptx/gen")
        output_dir = output_dir.expanduser().resolve()
        pptx_name = pptx_name_env or "schedule_gantt.pptx"
        output_path = output_dir / pptx_name

    # 既存の PPTX（cost フックが先に生成）をテンプレートとして再利用し、2枚目以降に追記する
    if output_path.exists():
        template_path = output_path
    else:
        template_path = _resolve_template_path(context)

    # compose/gen から呼ばれる場合は generate_ready.json が渡るが、schedule_data が必要なので
    # context もしくは env に記録された schedule.md から再パースする。
    try:
        schedule = _load_schedule_from_context(context)
    except FileNotFoundError:
        # フォールバック: generate_ready を ScheduleGantt として読む（外部呼び出し手動時など）
        env_path = os.environ.get("PPTX_GENERATE_READY_PATH") or context.get("generate_ready_path")
        if not env_path:
            raise
        schedule = _load_schedule(Path(env_path).expanduser().resolve())

    output_dir.mkdir(parents=True, exist_ok=True)

    renderer = AppendRenderer(template_path=template_path)
    layout_name = getattr(schedule.meta, "layout", None) or "System_layout"
    renderer.render(schedule, output_path, layout_name=layout_name)

    prs = Presentation(output_path)
    if prs.slides:
        apply_title_font_size(prs.slides[-1])
        apply_title_prefix(prs.slides[-1], 4)
        apply_text_style_by_name(prs.slides[-1], "Text Placeholder 1", font_size=14, bold=False)
        apply_strip_bullet_prefix(prs.slides[-1])
        prs.save(output_path)

    json_path = output_dir / "schedule_data.json"
    save_schedule_json(schedule, json_path)

    print(f"[stage04_gen] PPTX -> {output_path}")
    print(f"[stage04_gen] schedule JSON -> {json_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
