#!/usr/bin/env python3
from __future__ import annotations

import json
import os
import re
from pathlib import Path
from typing import Any

from pptx.util import Pt


def load_payload(generate_ready_path: Path) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    payload = json.loads(generate_ready_path.read_text(encoding="utf-8"))
    doc_meta = payload.get("meta") if isinstance(payload, dict) else {}
    slides = payload.get("slides")
    slides = slides if isinstance(slides, list) else []
    return [slide for slide in slides if isinstance(slide, dict)], doc_meta


def pick_slide_payload(slides: list[dict[str, Any]], index: int) -> tuple[dict[str, Any], dict[str, Any]]:
    if index < 0 or index >= len(slides):
        return {}, {}
    slide = slides[index]
    elements = slide.get("elements") if isinstance(slide.get("elements"), dict) else {}
    meta = slide.get("meta") if isinstance(slide.get("meta"), dict) else {}
    return elements, meta


def resolve_source_path(source: str, generate_ready_path: Path) -> Path:
    candidate = Path(source).expanduser()
    if candidate.is_absolute():
        return candidate
    repo_root = generate_ready_path.resolve().parent.parent.parent
    return (repo_root / candidate).resolve()


def normalize_name(value: str) -> str:
    text = value.lower()
    text = text.replace(" ", "").replace("_", "").replace("-", "")
    text = text.replace("（", "(").replace("）", ")")
    text = re.sub(r"[0-9\.]+", "", text)
    text = re.sub(r"[^a-z]+", "", text)
    return text


def build_shape_groups(slide) -> dict[str, list[Any]]:
    groups: dict[str, list[Any]] = {}
    for shape in slide.shapes:
        name = getattr(shape, "name", "")
        key = normalize_name(name)
        groups.setdefault(key, []).append(shape)
    return groups


def pop_shape(groups: dict[str, list[Any]], anchor: str):
    exact = anchor.strip()
    for shapes in groups.values():
        for idx, shape in enumerate(shapes):
            if getattr(shape, "name", "").strip() == exact:
                return shapes.pop(idx)
    key = normalize_name(anchor)
    shapes = groups.get(key)
    if shapes:
        return shapes.pop(0)
    return None


def set_text(shape, value: Any) -> None:
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    text_frame.clear()

    if isinstance(value, dict) and "text" in value:
        text_frame.text = str(value.get("text", ""))
        return
    elif isinstance(value, list) and value and isinstance(value[0], dict) and "text" in value[0]:
        values = [str(item.get("text", "")) for item in value]
        levels = [int(item.get("level", 0)) for item in value]
    elif isinstance(value, list):
        values = [str(item) for item in value]
        levels = [0] * len(values)
    elif value is None:
        values = [""]
        levels = [0]
    else:
        values = [str(value)]
        levels = [0]

    for idx, text in enumerate(values):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = text
        paragraph.level = levels[idx] if idx < len(levels) else 0


def apply_table_payload(shape, payload: dict[str, Any]) -> None:
    if not shape.has_table:
        return
    table = shape.table
    headers = payload.get("headers") if isinstance(payload.get("headers"), list) else []
    rows = payload.get("rows") if isinstance(payload.get("rows"), list) else []

    for row in table.rows:
        for cell in row.cells:
            cell.text = ""

    if headers:
        for col_idx, header in enumerate(headers):
            if col_idx >= len(table.columns):
                break
            table.cell(0, col_idx).text = str(header)

    for row_idx, row_values in enumerate(rows, start=1):
        if row_idx >= len(table.rows):
            break
        if not isinstance(row_values, list):
            continue
        for col_idx, value in enumerate(row_values):
            if col_idx >= len(table.columns):
                break
            table.cell(row_idx, col_idx).text = str(value)


def apply_table_font_size(shape, font_size: int) -> None:
    if not shape.has_table:
        return
    table = shape.table
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)


def apply_table_row_height(shape, height_emu: int) -> None:
    if not shape.has_table:
        return
    table = shape.table
    for row in table.rows:
        row.height = height_emu


def apply_table_column_widths(shape, widths_emu: list[int]) -> None:
    if not shape.has_table:
        return
    table = shape.table
    if len(widths_emu) != len(table.columns):
        return
    for idx, width in enumerate(widths_emu):
        table.columns[idx].width = width


def apply_text_font_size(shape, font_size: int) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        for run in paragraph.runs:
            run.font.size = Pt(font_size)


def _create_table_from_payload(slide, shape, payload: dict[str, Any], anchor_name: str):
    headers = payload.get("headers") if isinstance(payload.get("headers"), list) else []
    rows = payload.get("rows") if isinstance(payload.get("rows"), list) else []
    col_count = len(headers)
    for row in rows:
        if isinstance(row, list):
            col_count = max(col_count, len(row))
    row_count = len(rows) + (1 if headers else 0)
    if row_count <= 0 or col_count <= 0:
        return None

    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    try:
        shape.element.getparent().remove(shape.element)
    except Exception:
        pass
    table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
    try:
        table_shape.name = anchor_name
    except Exception:
        pass
    apply_table_payload(table_shape, payload)
    return table_shape


def apply_image(slide, shape, payload: dict[str, Any], generate_ready_path: Path, anchor: str | None = None) -> None:
    source = payload.get("source")
    if not isinstance(source, str) or not source:
        return
    image_path = resolve_source_path(source, generate_ready_path)
    if not image_path.exists():
        return
    try:
        if hasattr(shape, "insert_picture"):
            shape.insert_picture(str(image_path))
            return
    except Exception:
        pass
    picture = slide.shapes.add_picture(str(image_path), shape.left, shape.top, shape.width, shape.height)
    try:
        shape.element.getparent().remove(shape.element)
    except Exception:
        pass
    if anchor:
        try:
            picture.name = anchor
        except Exception:
            pass


ANCHOR_ALIAS_BY_LAYOUT: dict[str, dict[str, str]] = {
    "project_background_layout-01": {
        "Solusion": "コンテンツ プレースホルダー 1",
        "Solusion_Message_line": "テキスト プレースホルダー 2",
        "Solusion_title": "テキスト プレースホルダー 3",
        "Problem": "表プレースホルダー 4",
        "Problem_Message_line": "テキスト プレースホルダー 5",
        "Problem_title": "テキスト プレースホルダー 6",
        "Message_line": "テキスト プレースホルダー 7",
        "Date_dept": "テキスト プレースホルダー 8",
        "title": "Title 9",
    },
    "3_system_layout-02": {
        "Items_titile": "Text Placeholder 1",
        "Image_title": "Text Placeholder 2",
        "Items": "Table Placeholder 3",
        "Message_line": "Text Placeholder 4",
        "Date_dept": "Text Placeholder 5",
        "title": "Title 6",
        "Image": "Content Placeholder 7",
    },
}

ANCHOR_GEOMETRY_BY_LAYOUT: dict[str, dict[str, tuple[int, int, int, int]]] = {
    "project_background_layout-01": {
        "Problem": (164996, 2440028, 4686134, 4192547),
    },
    "3_system_layout-02": {
        "Items": (5054770, 2288312, 4686132, 16301193),
    },
}

ANCHOR_RENAME_BY_LAYOUT: dict[str, dict[str, str]] = {
    "project_background_layout-01": {
        "Solusion": "コンテンツ プレースホルダー 1",
        "Solusion_Message_line": "テキスト プレースホルダー 2",
        "Solusion_title": "テキスト プレースホルダー 3",
        "Problem": "表プレースホルダー 4",
        "Problem_Message_line": "テキスト プレースホルダー 5",
        "Problem_title": "テキスト プレースホルダー 6",
        "Message_line": "テキスト プレースホルダー 7",
        "Date_dept": "テキスト プレースホルダー 8",
        "title": "Title 9",
    },
    "3_system_layout-02": {
        "Items_titile": "Text Placeholder 1",
        "Image_title": "Text Placeholder 2",
        "Items": "Table Placeholder 3",
        "Message_line": "Text Placeholder 4",
        "Date_dept": "Text Placeholder 5",
        "title": "Title 6",
        "Image": "Content Placeholder 7",
    },
}

ANCHOR_FALLBACKS_BY_LAYOUT: dict[str, dict[str, list[str]]] = {
    "project_background_layout-01": {
        "Solusion": ["Content Placeholder 1", "コンテンツ プレースホルダー 1"],
        "Solusion_Message_line": ["Text Placeholder 2", "テキスト プレースホルダー 2"],
        "Solusion_title": ["Text Placeholder 3", "テキスト プレースホルダー 3"],
        "Problem": ["Table Placeholder 4", "表プレースホルダー 4"],
        "Problem_Message_line": ["Text Placeholder 5", "テキスト プレースホルダー 5"],
        "Problem_title": ["Text Placeholder 6", "テキスト プレースホルダー 6"],
        "Message_line": ["Text Placeholder 7", "テキスト プレースホルダー 7"],
        "Date_dept": ["Text Placeholder 8", "テキスト プレースホルダー 8"],
        "title": ["タイトル 9", "Title 9"],
    },
    "3_system_layout-02": {
        "Items_titile": ["Items_titile"],
        "Image_title": ["Image_title"],
        "Items": ["Table 12", "Table Placeholder 3"],
        "Message_line": ["Message_line"],
        "Date_dept": ["Date_dept"],
        "title": ["タイトル 6", "Title 6"],
        "Image": ["Image", "Content Placeholder 7"],
    },
}


def _resolve_anchor_alias(anchor: str, layout_id: str | None) -> str:
    if not layout_id:
        return anchor
    aliases = ANCHOR_ALIAS_BY_LAYOUT.get(layout_id)
    if not aliases:
        return anchor
    return aliases.get(anchor, anchor)


def _ensure_textbox(slide, shape, anchor_name: str):
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    try:
        shape.element.getparent().remove(shape.element)
    except Exception:
        pass
    textbox = slide.shapes.add_textbox(left, top, width, height)
    try:
        textbox.name = anchor_name
    except Exception:
        pass
    return textbox


def apply_elements(slide, elements: dict[str, Any], generate_ready_path: Path, layout_id: str | None = None) -> None:
    shape_groups = build_shape_groups(slide)
    for anchor, value in elements.items():
        resolved_anchor = _resolve_anchor_alias(anchor, layout_id)
        shape = pop_shape(shape_groups, resolved_anchor)
        if shape is None and layout_id:
            fallback_map = ANCHOR_FALLBACKS_BY_LAYOUT.get(layout_id, {})
            for fallback in fallback_map.get(anchor, []):
                shape = pop_shape(shape_groups, fallback)
                if shape is not None:
                    break
        if shape is None:
            continue
        if layout_id:
            rename_map = ANCHOR_RENAME_BY_LAYOUT.get(layout_id)
            if rename_map and anchor in rename_map:
                try:
                    shape.name = rename_map[anchor]
                except Exception:
                    pass
        if layout_id:
            geometry_map = ANCHOR_GEOMETRY_BY_LAYOUT.get(layout_id)
            if geometry_map and anchor in geometry_map:
                left, top, width, height = geometry_map[anchor]
                shape.left = left
                shape.top = top
                shape.width = width
                shape.height = height
        if isinstance(value, dict) and ("headers" in value or "rows" in value):
            if getattr(shape, "has_table", False):
                apply_table_payload(shape, value)
            else:
                created = _create_table_from_payload(slide, shape, value, resolved_anchor)
                if created is None:
                    continue
                shape = created
            anchor_name = getattr(shape, "name", "").strip()
            if anchor_name == "表プレースホルダー 4":
                apply_table_font_size(shape, 10)
            elif anchor_name == "Table Placeholder 3":
                apply_table_font_size(shape, 9)
                apply_table_row_height(shape, 300000)
                total_width = shape.width
                widths = [
                    int(total_width * 0.08),
                    int(total_width * 0.22),
                    int(total_width * 0.70),
                ]
                apply_table_column_widths(shape, widths)
            continue
        if isinstance(value, dict) and "source" in value:
            apply_image(slide, shape, value, generate_ready_path, anchor)
            continue
        if not getattr(shape, "has_text_frame", False):
            shape = _ensure_textbox(slide, shape, anchor)
        set_text(shape, value)
        if getattr(shape, "name", "").strip() == "コンテンツ プレースホルダー 1":
            apply_text_font_size(shape, 10)
