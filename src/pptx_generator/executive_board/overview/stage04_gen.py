#!/usr/bin/env python3
from __future__ import annotations

import os
import sys
import re
from pathlib import Path
from typing import Any

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from common.stage_shared import load_context, resolve_local_path
from overview.slide_renderer import apply_elements, load_payload, pick_slide_payload
from common.rendering.title_style import apply_title_font_size, apply_title_prefix
from common.rendering.text_cleanup import apply_strip_bullet_prefix


def _resolve_output_path() -> Path:
    output_dir_env = os.environ.get("PPTX_OUTPUT_DIR")
    pptx_name_env = os.environ.get("PPTX_PPTX_NAME")
    output_pptx_env = os.environ.get("PPTX_OUTPUT_PPTX_PATH")

    if output_pptx_env:
        return Path(output_pptx_env).expanduser().resolve()

    output_dir = Path(output_dir_env) if output_dir_env else Path(".pptx/gen")
    output_dir = output_dir.expanduser().resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_name = pptx_name_env or "executive_board_generated.pptx"
    return output_dir / pptx_name


def _resolve_template_path(context: dict[str, Any], generate_ready_path: Path) -> Path:
    payload_slides, doc_meta = load_payload(generate_ready_path)
    if isinstance(doc_meta, dict):
        template_path = doc_meta.get("template_path")
        if isinstance(template_path, str) and template_path.strip():
            return Path(template_path).expanduser().resolve()
        job_meta = doc_meta.get("job_meta") if isinstance(doc_meta, dict) else None
        template_path = job_meta.get("template_path") if isinstance(job_meta, dict) else None
        if isinstance(template_path, str) and template_path.strip():
            return (generate_ready_path.parent.parent.parent / template_path).resolve()
    if isinstance(context, dict):
        template_path = context.get("template_path")
        if isinstance(template_path, str) and template_path.strip():
            return resolve_local_path(template_path.strip(), Path(__file__).resolve().parents[1])
    raise FileNotFoundError("template_path is not provided")


def _find_layout(presentation: Presentation, layout_name: str | None):
    if not layout_name:
        return None
    for layout in presentation.slide_layouts:
        if getattr(layout, "name", "").strip() == layout_name:
            return layout
    normalized = layout_name.replace(" ", "").lower()
    normalized_trim = re.sub(r"-\d+$", "", normalized)
    for layout in presentation.slide_layouts:
        layout_norm = getattr(layout, "name", "").replace(" ", "").lower()
        if normalized in layout_norm or (normalized_trim and normalized_trim in layout_norm):
            return layout
    return None


def _remove_shapes_by_name(slide, names: set[str]) -> None:
    for shape in list(slide.shapes):
        if getattr(shape, "name", "").strip() in names:
            slide.shapes._spTree.remove(shape._element)


def main() -> int:
    context = load_context()
    generate_ready_path_env = os.environ.get("PPTX_GENERATE_READY_PATH")
    if not generate_ready_path_env:
        raise RuntimeError("generate_ready path is not provided")
    generate_ready_path = Path(generate_ready_path_env).expanduser().resolve()
    if not generate_ready_path.exists():
        raise FileNotFoundError(f"generate_ready.json not found: {generate_ready_path}")

    output_path = _resolve_output_path()
    template_path = _resolve_template_path(context, generate_ready_path)
    presentation = Presentation(template_path)

    slides_payload, _ = load_payload(generate_ready_path)
    elements, _ = pick_slide_payload(slides_payload, 0)
    target_layout = slides_payload[0].get("layout_id") if slides_payload else None
    layout = _find_layout(presentation, target_layout)

    if presentation.slides:
        slide = presentation.slides[0]
        if layout is not None and slide.slide_layout != layout:
            slide = presentation.slides.add_slide(layout)
    else:
        if layout is None:
            layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)

    if len(presentation.slides) > 1:
        sld_id_lst = presentation.slides._sldIdLst  # type: ignore[attr-defined]
        slides_snapshot = list(presentation.slides)
        package = presentation.part.package
        for idx in reversed(range(len(slides_snapshot))):
            candidate = slides_snapshot[idx]
            if candidate is slide:
                continue
            rel_id = sld_id_lst[idx].rId
            rel = presentation.part.rels.get(rel_id)
            target_part = rel.target_part if rel is not None else None
            presentation.part.drop_rel(rel_id)
            if target_part is not None:
                try:
                    package.drop_part(target_part.partname)
                except Exception:
                    pass
            sld_id_lst.remove(sld_id_lst[idx])
        slide = presentation.slides[0]

    apply_elements(slide, elements, generate_ready_path, target_layout)
    apply_title_font_size(slide)
    apply_title_prefix(slide, 1)
    apply_strip_bullet_prefix(slide)
    _remove_shapes_by_name(slide, {"表 10"})
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    print(f"[stage04_gen] overview slide generated -> {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
