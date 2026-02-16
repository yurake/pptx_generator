#!/usr/bin/env python3
"""Stage4 hook: organization_data を用いて PPTX に組織図スライドを追加する。"""
from __future__ import annotations

import json
import os
import sys
from pathlib import Path

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from common.stage_shared import load_context, resolve_local_path
from common.rendering.title_style import apply_title_font_size, apply_title_prefix
from common.rendering.text_cleanup import apply_strip_bullet_prefix
from pptx import Presentation


def _configure_sys_path() -> None:
    src_dir = Path(__file__).resolve().parents[3] / "src"
    if str(src_dir) not in sys.path:
        sys.path.insert(0, str(src_dir))
    organization_dir = Path(__file__).resolve().parent
    if str(organization_dir) not in sys.path:
        sys.path.insert(0, str(organization_dir))


_configure_sys_path()

from organization_parser import parse_organization_markdown  # noqa: E402
from organization_renderer import OrganizationChartRenderer  # noqa: E402
from organization.models import OrganizationChart  # noqa: E402


DEFAULT_MD = Path(__file__).resolve().parents[1] / "input/organization.md"
DEFAULT_TEMPLATE = Path(__file__).resolve().parents[3] / "templates/organization.pptx"


class AppendOrganizationRenderer(OrganizationChartRenderer):  # type: ignore[misc]
    """既存の PPTX にスライドを追加するレンダラー。"""

    def _clear_existing_slides(self, presentation) -> None:  # override
        return  # 既存スライドは維持する


def _load_chart(context: dict[str, str]) -> OrganizationChart:
    # 優先順位: JSON -> Markdown
    json_path_str = context.get("organization_json_path")
    if json_path_str:
        json_path = Path(json_path_str).expanduser().resolve()
        if json_path.exists():
            try:
                data = json.loads(json_path.read_text(encoding="utf-8"))
                return OrganizationChart.model_validate(data)
            except json.JSONDecodeError:
                pass

    md_env = os.environ.get("PPTX_ORGANIZATION_MD")
    candidates = [
        Path(md_env).expanduser().resolve() if md_env else None,
        Path(context.get("organization_md_path", "")).expanduser().resolve() if context.get("organization_md_path") else None,
        DEFAULT_MD,
    ]
    for candidate in candidates:
        if candidate and candidate.exists():
            return parse_organization_markdown(candidate)

    raise FileNotFoundError("organization markdown/json not found. Set PPTX_ORGANIZATION_MD or run stage02_prepare.")


def _resolve_template_path(context: dict[str, str], output_path: Path) -> Path:
    if output_path.exists():
        return output_path

    candidates: list[str] = []
    env_template = os.environ.get("PPTX_TEMPLATE_PATH")
    if isinstance(env_template, str) and env_template.strip():
        candidates.append(env_template.strip())
    if context.get("template_path"):
        candidates.append(context["template_path"])
    candidates.append(str(DEFAULT_TEMPLATE))

    for candidate in candidates:
        resolved = resolve_local_path(candidate, Path(__file__).resolve().parent)
        if resolved.exists():
            return resolved
    raise FileNotFoundError("Template PPTX not found for organization slide.")


def main(argv: list[str] | None = None) -> int:
    _ = argv  # unused

    context = load_context()

    output_dir_env = os.environ.get("PPTX_OUTPUT_DIR")
    pptx_name_env = os.environ.get("PPTX_PPTX_NAME")
    output_pptx_env = os.environ.get("PPTX_OUTPUT_PPTX_PATH")

    if output_pptx_env:
        output_path = Path(output_pptx_env).expanduser().resolve()
    else:
        output_dir = Path(output_dir_env) if output_dir_env else Path(".pptx/gen")
        output_dir = output_dir.expanduser().resolve()
        output_dir.mkdir(parents=True, exist_ok=True)
        pptx_name = pptx_name_env or "executive_board_generated.pptx"
        output_path = output_dir / pptx_name

    template_path = _resolve_template_path(context, output_path)
    chart = _load_chart(context)

    layout_name = os.environ.get("PPTX_ORGANIZATION_LAYOUT") or context.get("organization_layout") or "System_layout"

    renderer = AppendOrganizationRenderer(template_path=template_path)
    renderer.render(chart, output_path, layout_name=layout_name)

    prs = Presentation(output_path)
    if prs.slides:
        apply_title_font_size(prs.slides[-1])
        apply_title_prefix(prs.slides[-1], 5)
        apply_strip_bullet_prefix(prs.slides[-1])
        prs.save(output_path)

    print(f"[stage04_gen] organization slide appended -> {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
