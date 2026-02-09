"""組織図のレンダリング機能。"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt

from organization.models import OrganizationCategory, OrganizationChart
from common.settings import BrandingConfig

logger = logging.getLogger(__name__)


# 色定義
COLORS = {
    "light_green": RGBColor(198, 224, 180),  # #C6E0B4
    "green": RGBColor(146, 208, 80),  # #92D050
    "light_blue": RGBColor(222, 235, 247),  # #DEEBF7
    "blue": RGBColor(68, 114, 196),  # #4472C4
    "white": RGBColor(255, 255, 255),  # #FFFFFF
    "gray": RGBColor(128, 128, 128),  # #808080
    "black": RGBColor(0, 0, 0),  # #000000
}

# テンプレートのプレースホルダー位置とサイズ（インチ単位）
TEMPLATE_PLACEHOLDERS = [
    {"left": 0.18, "top": 0.65, "width": 3.23, "height": 6.48},  # 左列（SMBC）
    {"left": 3.80, "top": 0.65, "width": 3.23, "height": 6.48},  # 中央列（JRI）
    {"left": 7.42, "top": 0.65, "width": 3.23, "height": 6.48},  # 右列（開発ベンダー）
]


@dataclass(slots=True)
class OrganizationRenderConfig:
    """組織図描画の設定。"""

    slide_width_in: float = 10.0
    slide_height_in: float = 7.5
    title_top_in: float = 0.4
    title_height_in: float = 0.5
    title_left_in: float = 0.2
    title_width_in: float = 9.6
    content_top_in: float = 1.2
    content_left_in: float = 0.3
    content_width_in: float = 9.4
    content_height_in: float = 6.0
    category_spacing_in: float = 0.2
    group_spacing_in: float = 0.15
    box_title_height_in: float = 0.3937  # タイトル部の高さ 1.0cm
    box_member_height_per_line_in: float = 0.15  # メンバー1行あたりの高さ（約0.38cm）
    box_base_height_in: float = 0.4528  # 基本高さ 1.15cm
    box_extra_height_in: float = 0.1181  # 追加高さ 0.3cm
    font_size_title_pt: float = 11
    font_size_member_pt: float = 9


@dataclass(slots=True)
class BoxLayout:
    """箱の配置情報。"""

    left_in: float
    top_in: float
    width_in: float
    height_in: float
    title: str
    members: list[str]
    title_color: RGBColor
    member_color: RGBColor
    border_color: RGBColor


@dataclass(slots=True)
class CategoryLayout:
    """カテゴリーの配置情報。"""

    left_in: float
    top_in: float
    width_in: float
    height_in: float
    name: str
    boxes: list[BoxLayout]
    background_color: RGBColor
    border_color: RGBColor


class OrganizationChartRenderer:
    """組織図のレンダラー。"""

    def __init__(
        self,
        *,
        template_path: Path | None = None,
        branding: BrandingConfig | None = None,
        config: OrganizationRenderConfig | None = None,
    ) -> None:
        """
        Args:
            template_path: テンプレートPPTXファイルパス
            branding: ブランド設定
            config: 組織図描画設定
        """
        self.template_path = template_path
        self.branding = branding or BrandingConfig.default()
        self.config = config or OrganizationRenderConfig()

    def render(
        self, chart: OrganizationChart, output_path: Path, layout_name: str = "白紙"
    ) -> None:
        """組織図をPPTXファイルとして保存する。

        Args:
            chart: 組織図データ
            output_path: 出力PPTXファイルパス
            layout_name: 使用するレイアウト名（デフォルト: "白紙"）
        """
        presentation = self._load_template()

        # テンプレート内の既存スライドを削除
        self._clear_existing_slides(presentation)

        layout = self._find_layout(presentation, layout_name)
        slide = presentation.slides.add_slide(layout)

        # 背景を白に設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS["white"]

        # タイトルをプレースホルダーに設定
        self._set_title_to_placeholder(slide, "開発体制")

        # 日付と部署をプレースホルダーに設定
        self._set_date_dept_to_placeholder(slide)

        # カテゴリーレイアウトを計算
        category_layouts = self._calculate_category_layouts(chart.categories)

        # カテゴリーを描画
        for category_layout in category_layouts:
            self._draw_category(slide, category_layout)

        # 接続線を描画
        self._draw_connections(slide, category_layouts)

        # 空のプレースホルダーを削除
        self._remove_empty_placeholders(slide)

        # 保存
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        logger.info("組織図を保存しました: %s", output_path)

    def _load_template(self) -> Presentation:
        """テンプレートをロードする。"""
        if self.template_path and self.template_path.exists():
            logger.debug("テンプレートを使用: %s", self.template_path)
            return Presentation(self.template_path)
        logger.debug("既定テンプレートを利用")
        return Presentation()

    def _find_layout(self, presentation: Presentation, layout_name: str):
        """レイアウトを検索する。"""
        for layout in presentation.slide_layouts:
            if layout.name == layout_name:
                return layout
        logger.warning("レイアウト '%s' が見つからないため最初のレイアウトを使用", layout_name)
        if len(presentation.slide_layouts) == 0:
            msg = "テンプレートに利用可能なレイアウトが存在しません"
            raise RuntimeError(msg)
        return presentation.slide_layouts[0]

    def _clear_existing_slides(self, presentation: Presentation) -> None:
        """プレゼンテーション内の既存スライドを全て削除する。"""
        slide_count = len(presentation.slides)
        for i in range(slide_count - 1, -1, -1):
            rId = presentation.slides._sldIdLst[i].rId
            presentation.part.drop_rel(rId)
            del presentation.slides._sldIdLst[i]

    def _set_title_to_placeholder(self, slide, title: str) -> None:
        """スライドのタイトルプレースホルダーにタイトルを設定する。"""
        # タイトルプレースホルダーを検索
        for shape in slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # タイトルプレースホルダー（type 1）を探す
                if phf.type == 1:  # PP_PLACEHOLDER.TITLE
                    text_frame = shape.text_frame
                    text_frame.clear()
                    paragraph = text_frame.paragraphs[0]
                    paragraph.text = title
                    return
        
        # プレースホルダーが見つからない場合はログ出力
        logger.warning("タイトルプレースホルダーが見つかりませんでした")

    def _set_date_dept_to_placeholder(self, slide) -> None:
        """スライドのDate_dept図形に日付と部署を設定する。"""
        # 今日の日付を取得（日本語形式）
        today = datetime.now()
        date_str = today.strftime("%Y年%m月%d日")
        dept_str = "システム情報部"
        content = f"{date_str}\n{dept_str}"
        
        # Date_dept図形を検索（優先順位順）
        target_shape = None
        search_log: list[str] = []
        
        # スライドのすべての図形を検索
        logger.debug(f"スライド内の図形総数: {len(slide.shapes)}")
        for idx, shape in enumerate(slide.shapes):
            shape_name = getattr(shape, 'name', None)
            is_placeholder = shape.is_placeholder
            shape_type = type(shape).__name__
            
            # 詳細情報をログに記録
            if is_placeholder:
                phf = shape.placeholder_format
                placeholder_idx = getattr(phf, 'idx', 'N/A')
                search_log.append(
                    f"[{idx}] name={shape_name}, type={shape_type}, "
                    f"is_placeholder=True, ph_type={phf.type}, ph_idx={placeholder_idx}"
                )
            else:
                search_log.append(
                    f"[{idx}] name={shape_name}, type={shape_type}, is_placeholder=False"
                )
            
            # 第1優先：図形名が完全一致で「Date_dept」
            if shape_name == "Date_dept":
                target_shape = shape
                logger.debug(f"Date_dept図形を検出（完全一致）: {shape_name}, is_placeholder={is_placeholder}")
                break
            
            # 第2優先：図形名に「Date_dept」を含む（大文字小文字無視）
            if shape_name and "date_dept" in shape_name.lower():
                target_shape = shape
                logger.debug(f"Date_dept図形を検出（部分一致）: {shape_name}, is_placeholder={is_placeholder}")
                break
            
            # 第3優先（プレースホルダーの場合のみ）：Date/Footerタイプまたは名前に'date'を含む
            if target_shape is None and is_placeholder:
                phf = shape.placeholder_format
                is_date_type = phf.type in (15, 16)  # FOOTER or DATE
                is_date_name = shape_name and 'date' in shape_name.lower()
                if is_date_type or is_date_name:
                    target_shape = shape
                    logger.debug(f"Date_dept図形候補を検出（タイプ/名前マッチ）: {shape_name}, type={phf.type}")
        
        # 検索ログを出力
        logger.debug(f"検索した図形一覧:\n" + "\n".join(search_log))
        
        if target_shape:
            if hasattr(target_shape, 'text_frame'):
                text_frame = target_shape.text_frame
                text_frame.clear()
                paragraph = text_frame.paragraphs[0]
                paragraph.text = content
                shape_name = getattr(target_shape, 'name', 'unknown')
                logger.info(f"Date_dept図形に日付を設定しました: {shape_name}, 内容={content}")
            else:
                logger.warning(f"Date_dept図形が見つかりましたがtext_frameがありません: {getattr(target_shape, 'name', 'unknown')}")
        else:
            # プレースホルダーが見つからない場合は、右上にテキストボックスを作成
            logger.info("Date_dept図形が見つからないため、右上にテキストボックスを作成します")
            
            # 右上の位置（template_spec.jsonの座標を参考）
            left = Inches(8.75)
            top = Inches(0.38)
            width = Inches(1.90)
            height = Inches(0.21)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            paragraph = text_frame.paragraphs[0]
            paragraph.text = content
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
            
            # フォント設定
            font = paragraph.font
            font.size = Pt(9)
            font.name = "メイリオ"
            
            logger.info(f"右上にDate_deptテキストボックスを作成しました: 内容={content}")

    def _calculate_category_layouts(
        self, categories: list[OrganizationCategory]
    ) -> list[CategoryLayout]:
        """カテゴリーのレイアウトを計算する。
        
        テンプレートのプレースホルダーと同じ位置・サイズで配置する。
        """
        category_count = len(categories)
        if category_count == 0:
            return []

        # テンプレートのプレースホルダー数と照合
        if category_count > len(TEMPLATE_PLACEHOLDERS):
            logger.warning(
                "カテゴリー数(%d)がテンプレートのプレースホルダー数(%d)を超えています",
                category_count,
                len(TEMPLATE_PLACEHOLDERS),
            )

        layouts: list[CategoryLayout] = []

        for idx, category in enumerate(categories):
            # テンプレートのプレースホルダー位置を使用
            if idx < len(TEMPLATE_PLACEHOLDERS):
                placeholder = TEMPLATE_PLACEHOLDERS[idx]
                left = placeholder["left"]
                top = placeholder["top"]
                width = placeholder["width"]
                height = placeholder["height"]
            else:
                # プレースホルダーが不足している場合はデフォルト位置
                left = self.config.content_left_in + idx * 3.5
                top = self.config.content_top_in
                width = 3.0
                height = 6.0
                logger.warning("プレースホルダー不足: カテゴリー %s", category.name)

            # グループの箱を計算
            boxes = self._calculate_box_layouts(category, left, top, width)

            # 背景色を取得
            bg_color = COLORS.get(category.color, COLORS["light_blue"])

            layouts.append(
                CategoryLayout(
                    left_in=left,
                    top_in=top,
                    width_in=width,
                    height_in=height,
                    name=category.name,
                    boxes=boxes,
                    background_color=bg_color,
                    border_color=COLORS["gray"],
                )
            )

        return layouts

    def _calculate_box_layouts(
        self, category: OrganizationCategory, left: float, top: float, width: float
    ) -> list[BoxLayout]:
        """カテゴリー内の箱のレイアウトを計算する。"""
        boxes: list[BoxLayout] = []
        current_top = top + 0.4  # 上部マージン（カテゴリータイトル用のスペースを確保）

        title_color = COLORS.get(category.box_title_color, COLORS["blue"])
        member_color = COLORS["white"]
        border_color = COLORS["black"]
        
        # 箱の幅を6cm（約2.36インチ）に制限
        box_width = 2.36  # 6cm

        for group in category.groups:
            # メンバー数に応じて高さを計算
            member_count = len(group.members)
            member_height = member_count * self.config.box_member_height_per_line_in
            calculated_height = self.config.box_title_height_in + member_height
            
            # 高さの決定ロジック：
            # - 計算高さが1.15cm以内なら1.15cm
            # - 1.15cmを超える場合は計算高さ+0.3cm
            if calculated_height <= self.config.box_base_height_in:
                total_height = self.config.box_base_height_in
            else:
                total_height = calculated_height + self.config.box_extra_height_in

            # プロジェクト責任者・オーナー系は左寄せ、それ以外は右寄せ
            is_primary = self._is_primary_role(group.title)
            if is_primary:
                box_left = left + 0.05  # 左寄せ
            else:
                box_left = left + width - box_width - 0.05  # 右寄せ

            boxes.append(
                BoxLayout(
                    left_in=box_left,
                    top_in=current_top,
                    width_in=box_width,
                    height_in=total_height,
                    title=group.title,
                    members=group.members,
                    title_color=title_color,
                    member_color=member_color,
                    border_color=border_color,
                )
            )

            current_top += total_height + self.config.group_spacing_in

        return boxes

    def _is_primary_role(self, title: str) -> bool:
        """タイトルがプロジェクト責任者・オーナー系かを判定する。"""
        primary_keywords = [
            "プロジェクト責任者",
            "プロジェクトオーナー",
            "ビジネスオーナー",
        ]
        return any(keyword in title for keyword in primary_keywords)

    def _draw_category(self, slide, layout: CategoryLayout) -> None:
        """カテゴリーを描画する（背景四角形と個別の箱）。"""
        # 背景の四角形を描画
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(layout.left_in),
            Inches(layout.top_in),
            Inches(layout.width_in),
            Inches(layout.height_in),
        )

        # 背景色を設定
        fill = bg_shape.fill
        fill.solid()
        fill.fore_color.rgb = layout.background_color

        # 枠線を灰色点線に設定
        line = bg_shape.line
        line.color.rgb = layout.border_color
        line.width = Pt(1.5)
        line.dash_style = 2  # 点線

        # カテゴリー名をテキストボックスで追加（背景の一番上）
        title_box = slide.shapes.add_textbox(
            Inches(layout.left_in),
            Inches(layout.top_in + 0.05),  # 上部に少しマージン
            Inches(layout.width_in),
            Inches(0.3),  # タイトル用の高さ
        )
        text_frame = title_box.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = layout.name
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        font = paragraph.font
        font.size = Pt(14)
        font.bold = True
        font.color.rgb = COLORS["black"]

        # 各グループの箱を描画
        for box in layout.boxes:
            self._draw_box(slide, box)

    def _draw_box(self, slide, box: BoxLayout) -> None:
        """グループの箱を描画する（タイトル部分とメンバー部分）。"""
        # タイトル部分の四角形
        title_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(box.left_in),
            Inches(box.top_in),
            Inches(box.width_in),
            Inches(self.config.box_title_height_in),
        )

        # タイトル背景色
        title_fill = title_shape.fill
        title_fill.solid()
        title_fill.fore_color.rgb = box.title_color

        # タイトル枠線
        title_line = title_shape.line
        title_line.color.rgb = box.border_color
        title_line.width = Pt(1)

        # タイトルテキスト
        title_frame = title_shape.text_frame
        title_frame.clear()
        title_frame.vertical_anchor = 1  # 中央揃え
        title_para = title_frame.paragraphs[0]
        title_para.text = box.title
        title_para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        title_font = title_para.font
        title_font.size = Pt(self.config.font_size_title_pt)
        title_font.bold = True
        title_font.color.rgb = COLORS["black"]

        # メンバー部分の四角形
        member_top = box.top_in + self.config.box_title_height_in
        member_height = box.height_in - self.config.box_title_height_in

        member_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(box.left_in),
            Inches(member_top),
            Inches(box.width_in),
            Inches(member_height),
        )

        # メンバー背景色
        member_fill = member_shape.fill
        member_fill.solid()
        member_fill.fore_color.rgb = box.member_color

        # メンバー枠線
        member_line = member_shape.line
        member_line.color.rgb = box.border_color
        member_line.width = Pt(1)

        # メンバーテキスト
        member_frame = member_shape.text_frame
        member_frame.clear()
        member_frame.word_wrap = True
        member_frame.vertical_anchor = 1  # 中央揃え

        # 各メンバーを改行区切りで追加
        for i, member in enumerate(box.members):
            if i > 0:
                member_frame.add_paragraph()
            para = member_frame.paragraphs[i]
            para.text = member
            para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            font = para.font
            font.size = Pt(self.config.font_size_member_pt)
            font.color.rgb = COLORS["black"]

    def _draw_connections(self, slide, layouts: list[CategoryLayout]) -> None:
        """カテゴリー間およびカテゴリー内の接続線を描画する。"""
        # カテゴリー内の接続線（縦方向）
        for layout in layouts:
            self._draw_category_internal_connections(slide, layout)

        # カテゴリー間の接続線（横方向）
        self._draw_category_to_category_connections(slide, layouts)

    def _draw_category_internal_connections(
        self, slide, layout: CategoryLayout
    ) -> None:
        """カテゴリー内の箱同士を接続する線を描画する。"""
        if len(layout.boxes) < 2:
            return

        for i in range(len(layout.boxes) - 1):
            box1 = layout.boxes[i]
            box2 = layout.boxes[i + 1]

            # box1の下端中央からbox2の上端中央への線
            x1 = box1.left_in + box1.width_in / 2
            y1 = box1.top_in + box1.height_in
            x2 = box2.left_in + box2.width_in / 2
            y2 = box2.top_in

            connector = slide.shapes.add_connector(
                1,  # msoConnectorStraight
                Inches(x1),
                Inches(y1),
                Inches(x2),
                Inches(y2),
            )

            line = connector.line
            line.color.rgb = COLORS["gray"]
            line.width = Pt(1.5)

    def _draw_category_to_category_connections(
        self, slide, layouts: list[CategoryLayout]
    ) -> None:
        """カテゴリー間を接続する線を描画する。"""
        if len(layouts) < 2:
            return

        # 隣接するカテゴリー間を接続
        for i in range(len(layouts) - 1):
            layout1 = layouts[i]
            layout2 = layouts[i + 1]

            # カテゴリーの中央高さで接続
            if layout1.boxes and layout2.boxes:
                # 最初の箱を基準に接続
                box1 = layout1.boxes[0]
                box2 = layout2.boxes[0]

                x1 = layout1.left_in + layout1.width_in
                y1 = box1.top_in + box1.height_in / 2
                x2 = layout2.left_in
                y2 = box2.top_in + box2.height_in / 2

                connector = slide.shapes.add_connector(
                    1,  # msoConnectorStraight
                    Inches(x1),
                    Inches(y1),
                    Inches(x2),
                    Inches(y2),
                )

                line = connector.line
                line.color.rgb = COLORS["gray"]
                line.width = Pt(1.5)

    def _remove_empty_placeholders(self, slide) -> None:
        """スライド内の空のプレースホルダーを削除する。
        
        タイトルプレースホルダー（type 1）は保持し、
        それ以外の空のプレースホルダーを削除する。
        """
        shapes_to_remove: list[int] = []
        
        for idx, shape in enumerate(slide.shapes):
            # プレースホルダーでない図形はスキップ
            if not shape.is_placeholder:
                continue
            
            phf = shape.placeholder_format
            placeholder_type = phf.type
            
            # タイトルプレースホルダー（type 1）は保持
            if placeholder_type == 1:  # PP_PLACEHOLDER.TITLE
                logger.debug(f"タイトルプレースホルダーを保持: idx={idx}")
                continue
            
            # text_frameを持たない図形はスキップ
            if not hasattr(shape, 'text_frame'):
                logger.debug(f"text_frameを持たないプレースホルダーをスキップ: idx={idx}, type={placeholder_type}")
                continue
            
            # テキストが空かどうかをチェック
            text_frame = shape.text_frame
            is_empty = True
            for paragraph in text_frame.paragraphs:
                if paragraph.text.strip():
                    is_empty = False
                    break
            
            # 空のプレースホルダーを削除対象にマーク
            if is_empty:
                shapes_to_remove.append(idx)
                logger.debug(
                    f"空のプレースホルダーを削除対象にマーク: "
                    f"idx={idx}, type={placeholder_type}, name={getattr(shape, 'name', 'unknown')}"
                )
        
        # インデックスの大きい順に削除（削除によるインデックスのずれを防ぐため）
        for idx in reversed(shapes_to_remove):
            try:
                shape = slide.shapes[idx]
                # 図形をスライドから削除
                sp = shape.element
                sp.getparent().remove(sp)
                logger.info(
                    f"空のプレースホルダーを削除しました: "
                    f"idx={idx}, name={getattr(shape, 'name', 'unknown')}"
                )
            except Exception as exc:  # noqa: BLE001
                logger.warning(
                    f"プレースホルダーの削除に失敗しました: idx={idx}, error={exc}"
                )
