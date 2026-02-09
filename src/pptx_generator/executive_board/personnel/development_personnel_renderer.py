"""開発要員計画スライドのレンダリング機能。"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt, Cm

from personnel.models import (
    DevelopmentPersonnelPlan,
    PersonnelMessage,
    PersonnelMonthSummary,
    PersonnelPhaseMonthSummary,
    PersonnelScheduleTask,
)
from common.settings import BrandingConfig

logger = logging.getLogger(__name__)


def _get_fiscal_year(year: int, month: int) -> int:
    """暦年と月から年度を計算する（4月始まり）。"""
    if 1 <= month <= 3:
        return year - 1
    return year


def _get_task_color(task_name: str) -> str:
    """タスク名からタスクの背景色を判定する。"""
    task_name_lower = task_name.lower()
    
    # 白色：拠点要件精緻化、拠点展開、計画策定、契約締結手続
    white_keywords = ['拠点要件精緻化', '拠点展開', '計画策定', '契約締結手続']
    if any(keyword in task_name for keyword in white_keywords):
        return "#FFFFFF"
    
    # 明るく薄い緑色：移行関連
    light_green_keywords = ['移行', 'データ移行']
    if any(keyword in task_name for keyword in light_green_keywords):
        return "#C6E0B4"
    
    # 明るい緑色：要件定義、設計、UAT
    green_keywords = ['要件定義', '設計', 'uat']
    if any(keyword in task_name_lower for keyword in green_keywords):
        return "#92D050"
    
    # ピンク色：開発、製造、構築、単体テスト
    pink_keywords = ['開発', '製造', '構築', '単体テスト']
    if any(keyword in task_name for keyword in pink_keywords):
        return "#FFC7CE"
    
    # 薄い黄色：結合テスト、性能テスト、障害テスト
    yellow_keywords = ['結合テスト', '性能テスト', '障害テスト']
    if any(keyword in task_name for keyword in yellow_keywords):
        return "#FFF2CC"
    
    # 薄い青色：システムテスト、運用テスト
    cyan_keywords = ['システムテスト', '運用テスト', '受入テスト']
    if any(keyword in task_name for keyword in cyan_keywords):
        return "#DEEBF7"
    
    # デフォルト：白色
    return "#FFFFFF"


@dataclass(slots=True)
class PersonnelRenderConfig:
    """開発要員計画レンダリング設定。"""

    slide_width_in: float = 13.33
    slide_height_in: float = 7.5
    table_left_in: float = 0.3
    table_top_in: float = 1.7
    table_width_in: float = 12.7
    table_height_in: float = 5.0
    phase_col_width_in: float = 1.0  # フェーズ名列の幅
    type_col_width_in: float = 0.5   # 種別列（社員/PN/SI/小計）の幅
    total_col_width_in: float = 0.6
    # ヘッダー行の高さを0.8cm以内に制限
    header_row_height_cm: float = 0.8
    # マイルストーン行の高さ（cm）- 三角形とラベル用
    milestone_row_height_cm: float = 0.35
    # マイルストーンと矢羽の間の余白（cm）
    milestone_arrow_gap_cm: float = 0.2
    # 矢羽の高さ（トラック高さ）を0.55cm以上に設定（矢羽実体は0.5cm以上）
    arrow_height_cm: float = 0.55
    # 矢羽の最小高さ（cm）
    min_arrow_height_cm: float = 0.5
    # データ行の高さ（cm）
    data_row_height_cm: float = 0.45


class DevelopmentPersonnelRenderer:
    """開発要員計画スライドのレンダラー。"""

    # 色設定
    HEADER_BG_COLOR = RGBColor(68, 114, 196)  # 青色（ヘッダー用）
    LIGHT_BLUE_BG_COLOR = RGBColor(221, 235, 247)  # 薄い水色（マイルストーン/スケジュールラベル、フェーズ名用）
    CYAN_BG_COLOR = RGBColor(155, 194, 230)  # 水色（小計行用）
    BLUE_BG_COLOR = RGBColor(68, 114, 196)  # 青色（全体、合計行用）
    WHITE_BG_COLOR = RGBColor(255, 255, 255)  # 白色（社員/PN/SI行、マイルストーン行の空白セル用）
    HIGHLIGHT_BORDER_COLOR = RGBColor(0, 32, 96)  # 紺色

    def __init__(
        self,
        *,
        template_path: Path | None = None,
        branding: BrandingConfig | None = None,
        config: PersonnelRenderConfig | None = None,
    ) -> None:
        self.template_path = template_path
        self.branding = branding or BrandingConfig.default()
        self.config = config or PersonnelRenderConfig()

    def render(
        self,
        plan: DevelopmentPersonnelPlan,
        output_path: Path,
        layout_name: str = "System_layout",
    ) -> None:
        """開発要員計画スライドをPPTXファイルとして保存する。"""
        presentation = self._load_template()
        
        # 既存スライドを削除
        self._clear_existing_slides(presentation)
        
        layout = self._find_layout(presentation, layout_name)
        slide = presentation.slides.add_slide(layout)

        # 背景を白に設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Titleを設定
        self._set_title(slide, plan.title)

        # Date_deptを設定
        self._set_date_dept(slide, plan.department)

        # Message_lineを設定
        self._set_message_line(slide, plan.messages)

        # Schedule_box（表）を描画
        self._draw_schedule_table(slide, plan)

        # 保存
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        logger.info("開発要員計画スライドを保存しました: %s", output_path)

    def _load_template(self) -> Presentation:
        """テンプレートをロードする。"""
        if self.template_path and self.template_path.exists():
            logger.debug("テンプレートを使用: %s", self.template_path)
            return Presentation(self.template_path)
        logger.debug("新規プレゼンテーションを作成")
        return Presentation()

    def _find_layout(self, presentation: Presentation, layout_name: str):
        """レイアウトを検索する。"""
        for layout in presentation.slide_layouts:
            if layout.name == layout_name:
                return layout
        logger.warning("レイアウト '%s' が見つからないため最初のレイアウトを使用", layout_name)
        if len(presentation.slide_layouts) == 0:
            msg = "テンプレートに利用可能なレイアウトが存在しません"
            raise RuntimeError(msg)
        return presentation.slide_layouts[0]

    def _clear_existing_slides(self, presentation: Presentation) -> None:
        """既存スライドを全て削除する。"""
        slide_count = len(presentation.slides)
        for i in range(slide_count - 1, -1, -1):
            rId = presentation.slides._sldIdLst[i].rId
            presentation.part.drop_rel(rId)
            del presentation.slides._sldIdLst[i]

    def _find_placeholder_by_name(self, slide, name: str):
        """プレースホルダーを名前で検索する。"""
        # デバッグ: すべてのプレースホルダーと図形をログに出力
        logger.debug("=== 検索対象: '%s' ===", name)
        logger.debug("slide.placeholdersの数: %d", len(list(slide.placeholders)))
        for i, ph in enumerate(slide.placeholders):
            ph_name = getattr(ph, 'name', 'N/A')
            ph_type = None
            try:
                ph_type = ph.placeholder_format.type
            except Exception:
                pass
            logger.debug("  placeholder[%d]: name='%s', type=%s", i, ph_name, ph_type)
        
        logger.debug("slide.shapesの数: %d", len(slide.shapes))
        for i, shape in enumerate(slide.shapes):
            shape_name = getattr(shape, 'name', 'N/A')
            has_text = hasattr(shape, 'text_frame')
            logger.debug("  shape[%d]: name='%s', has_text=%s", i, shape_name, has_text)
        
        # まずslide.placeholdersから検索
        for placeholder in slide.placeholders:
            if placeholder.name == name:
                logger.debug("プレースホルダー発見(placeholders): '%s'", name)
                return placeholder
        
        # 次にslide.shapesから検索
        for shape in slide.shapes:
            if shape.name == name and hasattr(shape, "text_frame"):
                logger.debug("プレースホルダー発見(shapes): '%s'", name)
                return shape
        
        # 部分一致で検索
        for placeholder in slide.placeholders:
            if name in placeholder.name or placeholder.name in name:
                logger.debug("プレースホルダー発見(部分一致): '%s' -> '%s'", name, placeholder.name)
                return placeholder
        
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and (name in shape.name or shape.name in name):
                logger.debug("シェイプ発見(部分一致): '%s' -> '%s'", name, shape.name)
                return shape
        
        logger.debug("プレースホルダー '%s' が見つかりませんでした", name)
        return None

    def _set_title(self, slide, title: str) -> None:
        """タイトルを設定する。"""
        # プレースホルダーを検索
        placeholder = self._find_placeholder_by_name(slide, "Title")
        
        if placeholder is not None and hasattr(placeholder, "text_frame"):
            placeholder.text_frame.clear()
            para = placeholder.text_frame.paragraphs[0]
            para.text = title
            para.font.size = Pt(28)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0, 0, 0)
            return
        
        # TITLEタイプのプレースホルダーを検索
        from pptx.enum.shapes import PP_PLACEHOLDER
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                placeholder.text_frame.clear()
                para = placeholder.text_frame.paragraphs[0]
                para.text = title
                para.font.size = Pt(28)
                para.font.bold = True
                para.font.color.rgb = RGBColor(0, 0, 0)
                return
        
        # なければ新規作成
        title_box = slide.shapes.add_textbox(
            Inches(0.3),
            Inches(0.2),
            Inches(10),
            Inches(0.5),
        )
        title_frame = title_box.text_frame
        para = title_frame.paragraphs[0]
        para.text = title
        para.font.size = Pt(28)
        para.font.bold = True
        para.font.color.rgb = RGBColor(0, 0, 0)
        title_box.name = "Title"

    def _find_body_placeholder_by_position(self, slide, position: str):
        """BODYタイプのプレースホルダーを位置で検索する。
        
        Args:
            slide: スライド
            position: "right" (右上) or "left" (左側のメッセージライン用)
        """
        from pptx.enum.shapes import PP_PLACEHOLDER
        
        body_placeholders = []
        for ph in slide.placeholders:
            try:
                if ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    body_placeholders.append(ph)
            except Exception:
                continue
        
        if not body_placeholders:
            return None
        
        if len(body_placeholders) == 1:
            return body_placeholders[0]
        
        # 位置でソート
        if position == "right":
            # 右上のプレースホルダー (leftが大きい方)
            body_placeholders.sort(key=lambda p: p.left, reverse=True)
        else:
            # 左側のプレースホルダー (topが大きい方、下の方)
            body_placeholders.sort(key=lambda p: p.top, reverse=True)
        
        return body_placeholders[0]

    def _set_date_dept(self, slide, department: str) -> None:
        """日付と部門名を設定する。"""
        today = datetime.now().strftime("%Y年%m月%d日")
        text = f"{department} {today}"
        
        # プレースホルダーを検索（まず名前で、次に位置で）
        placeholder = self._find_placeholder_by_name(slide, "Date_dept")
        
        if placeholder is None:
            # 右上のBODYプレースホルダーを探す
            placeholder = self._find_body_placeholder_by_position(slide, "right")
        
        if placeholder is not None and hasattr(placeholder, "text_frame"):
            placeholder.text_frame.clear()
            para = placeholder.text_frame.paragraphs[0]
            para.text = text
            para.font.size = Pt(10)
            para.font.color.rgb = RGBColor(128, 128, 128)
            para.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
            logger.debug("Date_deptを設定: プレースホルダー '%s'", placeholder.name)
            return
        
        # なければ新規作成
        date_box = slide.shapes.add_textbox(
            Inches(10.5),
            Inches(0.3),
            Inches(2.5),
            Inches(0.3),
        )
        date_frame = date_box.text_frame
        para = date_frame.paragraphs[0]
        para.text = text
        para.font.size = Pt(10)
        para.font.color.rgb = RGBColor(128, 128, 128)
        para.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
        date_box.name = "Date_dept"

    def _set_message_line(self, slide, messages: list[PersonnelMessage]) -> None:
        """メッセージラインを設定する。"""
        # 番号をマルイチ文字に変換
        circled_numbers = ["①", "②", "③", "④", "⑤", "⑥", "⑦", "⑧", "⑨", "⑩"]
        
        message_texts = []
        for msg in messages:
            idx = msg.number - 1
            if idx < len(circled_numbers):
                prefix = circled_numbers[idx]
            else:
                prefix = f"({msg.number})"
            message_texts.append(f"{prefix}{msg.text}")
        
        text = "\n".join(message_texts)
        
        # プレースホルダーを検索（まず名前で、次に位置で）
        placeholder = self._find_placeholder_by_name(slide, "Message_line")
        
        if placeholder is None:
            # 左下のBODYプレースホルダーを探す
            placeholder = self._find_body_placeholder_by_position(slide, "left")
        
        if placeholder is not None and hasattr(placeholder, "text_frame"):
            placeholder.text_frame.clear()
            para = placeholder.text_frame.paragraphs[0]
            para.text = text
            para.font.size = Pt(11)
            para.font.color.rgb = RGBColor(0, 0, 0)
            logger.debug("Message_lineを設定: プレースホルダー '%s'", placeholder.name)
            return
        
        # なければ新規作成
        message_box = slide.shapes.add_textbox(
            Inches(0.3),
            Inches(0.8),
            Inches(12.7),
            Inches(0.8),
        )
        message_frame = message_box.text_frame
        message_frame.word_wrap = True
        para = message_frame.paragraphs[0]
        para.text = text
        para.font.size = Pt(11)
        para.font.color.rgb = RGBColor(0, 0, 0)
        message_box.name = "Message_line"

    def _draw_schedule_table(self, slide, plan: DevelopmentPersonnelPlan) -> None:
        """スケジュール表を描画する。"""
        # 表示単位に応じて処理を分岐
        if plan.display_unit == "month":
            self._draw_schedule_table_monthly(slide, plan)
        else:
            self._draw_schedule_table_quarterly(slide, plan)

    def _draw_schedule_table_monthly(self, slide, plan: DevelopmentPersonnelPlan) -> None:
        """月単位のスケジュール表を描画する（12ヶ月以内の場合）。"""
        target_months = plan.target_months
        if not target_months:
            logger.warning("対象月データがありません")
            return
        
        # 列数を計算: フェーズ名列 + 種別列 + 各月 + 開発期間累計列
        month_cols = len(target_months)
        cols = 2 + month_cols + 1  # フェーズ名列 + 種別列 + 月列 + 累計列
        
        # 行数を計算: ヘッダー2行 + マイルストーン/スケジュール行 + (フェーズ数 × 4行) + 合計4行
        phase_count = len(plan.phase_month_summaries)
        rows = 2 + 1 + (phase_count * 4) + 4  # ヘッダー2 + スケジュール1 + フェーズ×4 + 合計4
        
        # Schedule_boxプレースホルダーの位置を取得
        table_left = Inches(self.config.table_left_in)
        table_top = Inches(self.config.table_top_in)
        table_width = Inches(self.config.table_width_in)
        table_height = Inches(self.config.table_height_in)
        
        schedule_box = self._find_placeholder_by_name(slide, "Schedule_box")
        if schedule_box is not None:
            table_left = schedule_box.left
            table_top = schedule_box.top
            table_width = schedule_box.width
            table_height = schedule_box.height
            # Schedule_boxを削除
            sp = schedule_box.element
            sp.getparent().remove(sp)
            logger.debug("Schedule_boxプレースホルダーを発見: left=%s, top=%s, width=%s, height=%s",
                        table_left, table_top, table_width, table_height)
        else:
            # OBJECTタイプのプレースホルダーを検索
            from pptx.enum.shapes import PP_PLACEHOLDER
            for placeholder in slide.placeholders:
                try:
                    if placeholder.placeholder_format.type == PP_PLACEHOLDER.OBJECT:
                        table_left = placeholder.left
                        table_top = placeholder.top
                        table_width = placeholder.width
                        table_height = placeholder.height
                        sp = placeholder.element
                        sp.getparent().remove(sp)
                        logger.debug("OBJECTプレースホルダーを発見: left=%s, top=%s", table_left, table_top)
                        break
                except Exception:
                    continue
        
        # スケジュール行の高さを事前計算（トラック数に基づく）
        all_tasks = []
        for phase in plan.phases:
            all_tasks.extend(phase.tasks)
        
        # トラック数を計算
        tasks_with_tracks = self._assign_tracks_to_tasks_monthly(all_tasks, target_months)
        if tasks_with_tracks:
            tracks_needed = max(track for _, track in tasks_with_tracks) + 1
        else:
            tracks_needed = 1
        
        # スケジュール行の高さ = マイルストーン用スペース + 間隔 + トラック数 × 矢羽の高さ + マージン
        schedule_row_height = Cm(
            self.config.milestone_row_height_cm +
            self.config.milestone_arrow_gap_cm +
            self.config.arrow_height_cm * tracks_needed +
            0.1
        )
        
        # 表を作成
        table_shape = slide.shapes.add_table(
            rows,
            cols,
            table_left,
            table_top,
            table_width,
            table_height,
        )
        table = table_shape.table
        
        # 列幅を設定
        phase_col_width = Inches(self.config.phase_col_width_in)
        type_col_width = Inches(self.config.type_col_width_in)
        total_col_width = Inches(self.config.total_col_width_in)
        month_col_width = int((table_width - phase_col_width - type_col_width - total_col_width) / month_cols)
        month_col_width_in = (table_width / 914400 - self.config.phase_col_width_in - self.config.type_col_width_in - self.config.total_col_width_in) / month_cols
        
        table.columns[0].width = phase_col_width
        table.columns[1].width = type_col_width
        for i in range(2, cols - 1):
            table.columns[i].width = month_col_width
        table.columns[cols - 1].width = total_col_width
        
        # 行の高さを設定
        header_row_height = Cm(self.config.header_row_height_cm)
        data_row_height = Cm(self.config.data_row_height_cm)
        # ヘッダー行1, 2の高さを0.8cm以内に設定
        table.rows[0].height = header_row_height
        table.rows[1].height = header_row_height
        # スケジュール行（3行目）の高さをトラック数に応じて設定
        table.rows[2].height = schedule_row_height
        # データ行の高さを設定
        for row_idx in range(3, rows):
            table.rows[row_idx].height = data_row_height
        
        # 年度ごとに月をグループ化
        fiscal_year_months: dict[int, list[tuple[int, int]]] = {}
        for year, month in target_months:
            fy = year if month >= 4 else year - 1
            if fy not in fiscal_year_months:
                fiscal_year_months[fy] = []
            fiscal_year_months[fy].append((year, month))
        
        # ヘッダー行1: 年度（フェーズ名列と種別列をマージ）
        self._style_header_cell(table.cell(0, 0), "", 8)  # 左上セルは空欄
        table.cell(0, 0).merge(table.cell(1, 1))  # 2列×2行をマージ
        
        col = 2
        for fy in sorted(fiscal_year_months.keys()):
            months_in_fy = fiscal_year_months[fy]
            start_col = col
            col += len(months_in_fy)
            # 年度セルをマージ
            table.cell(0, start_col).merge(table.cell(0, col - 1))
            self._style_header_cell(table.cell(0, start_col), f"{fy}年度", 8)
        
        # 累計列ヘッダー
        self._style_header_cell(table.cell(0, cols - 1), "開発期間", 7)
        table.cell(0, cols - 1).merge(table.cell(1, cols - 1))
        self._set_cell_text(table.cell(0, cols - 1), "開発期間\n累計", 7, bold=True, color=RGBColor(255, 255, 255))
        
        # ヘッダー行2: 月
        col = 2
        for year, month in target_months:
            self._style_header_cell(table.cell(1, col), f"{month}月", 7)
            col += 1
        
        # マイルストーン/スケジュール行（3行目）
        schedule_row = 2
        self._style_milestone_label_cell(table.cell(schedule_row, 0), "マイルストーン\n/スケジュール", 6)
        table.cell(schedule_row, 0).merge(table.cell(schedule_row, 1))  # フェーズ名列と種別列をマージ
        
        # マイルストーン/スケジュール行のセルをスタイル設定（空白セル）- 白色背景
        col = 2
        for _ in target_months:
            self._style_milestone_data_cell(table.cell(schedule_row, col), "", 6)
            col += 1
        self._style_milestone_data_cell(table.cell(schedule_row, cols - 1), "", 7)
        
        # 矢羽をスケジュール行に描画
        if all_tasks:
            self._draw_schedule_arrows_monthly(
                slide,
                table_shape,
                table,
                schedule_row,
                all_tasks,
                target_months,
                month_col_width_in,
                tasks_with_tracks,
            )
        
        # マイルストーンを三角形で描画
        self._draw_milestones_as_triangles_monthly(
            slide,
            table_shape,
            table,
            schedule_row,
            plan.milestones,
            target_months,
            month_col_width_in,
        )
        
        # 種別ラベル（社員/PN/SI/小計）
        type_labels = ["社員", "PN", "SI", "小計"]
        
        # フェーズ行（4行目以降、各フェーズ4行）
        for phase_idx, phase in enumerate(plan.phase_month_summaries):
            base_row = 3 + phase_idx * 4
            
            # フェーズ名（4行をマージ）- 水色背景
            self._style_phase_cell(table.cell(base_row, 0), phase.phase_name, 7)
            table.cell(base_row, 0).merge(table.cell(base_row + 3, 0))
            
            for type_idx, type_label in enumerate(type_labels):
                row_idx = base_row + type_idx
                
                # 種別ラベルとデータセルのスタイルを種別に応じて変更
                if type_label == "小計":
                    # 小計行: 青色背景、黒文字
                    self._style_subtotal_type_cell(table.cell(row_idx, 1), type_label, 7)
                else:
                    # 社員/PN/SI行: 白色背景、黒文字
                    self._style_employee_type_cell(table.cell(row_idx, 1), type_label, 7)
                
                # 月別工数
                col = 2
                for year, month in target_months:
                    # 該当月のデータを取得
                    month_data = next(
                        (md for md in phase.months if md.year == year and md.month == month),
                        None,
                    )
                    if month_data:
                        if type_label == "社員":
                            value = month_data.employee
                        elif type_label == "PN":
                            value = month_data.pn
                        elif type_label == "SI":
                            value = month_data.si
                        else:  # 小計
                            value = month_data.total
                        text = f"{value:.1f}" if value > 0 else "-"
                    else:
                        text = "-"
                    
                    # データセルのスタイルを種別に応じて変更
                    if type_label == "小計":
                        self._style_subtotal_data_cell(table.cell(row_idx, col), text, 7)
                    else:
                        self._style_employee_data_cell(table.cell(row_idx, col), text, 7)
                    col += 1
                
                # 累計
                if type_label == "社員":
                    total = phase.total_employee
                elif type_label == "PN":
                    total = phase.total_pn
                elif type_label == "SI":
                    total = phase.total_si
                else:  # 小計
                    total = phase.total
                
                # 累計セルのスタイルを種別に応じて変更
                if type_label == "小計":
                    self._style_subtotal_data_cell(table.cell(row_idx, cols - 1), f"{total:.1f}", 7)
                else:
                    self._style_employee_data_cell(table.cell(row_idx, cols - 1), f"{total:.1f}", 7)
        
        # 合計行（最後の4行）
        total_base_row = rows - 4
        total_type_labels = ["社員", "PN", "SI", "合計"]  # 最終行は「合計」
        
        # 全体ラベル（4行をマージ）- 青色背景、太字なし
        self._style_overall_label_cell(table.cell(total_base_row, 0), "全体", 8)
        table.cell(total_base_row, 0).merge(table.cell(total_base_row + 3, 0))
        
        for type_idx, type_label in enumerate(total_type_labels):
            row_idx = total_base_row + type_idx
            
            # 種別ラベルとデータセルのスタイルを種別に応じて変更
            if type_label == "合計":
                # 合計行: 青色背景、黒文字、太字なし
                self._style_grand_total_cell(table.cell(row_idx, 1), type_label, 7)
            else:
                # 社員/PN/SI行: 白色背景、黒文字
                self._style_employee_type_cell(table.cell(row_idx, 1), type_label, 7)
            
            col = 2
            for year, month in target_months:
                # 該当月の合計を取得
                month_total = next(
                    (mt for mt in plan.total_month_summary if mt.year == year and mt.month == month),
                    None,
                )
                if month_total:
                    if type_label == "社員":
                        value = month_total.employee
                    elif type_label == "PN":
                        value = month_total.pn
                    elif type_label == "SI":
                        value = month_total.si
                    else:  # 合計
                        value = month_total.total
                    text = f"{value:.1f}" if value > 0 else "-"
                else:
                    text = "-"
                
                # データセルのスタイルを種別に応じて変更
                if type_label == "合計":
                    self._style_grand_total_cell(table.cell(row_idx, col), text, 7)
                else:
                    self._style_employee_data_cell(table.cell(row_idx, col), text, 7)
                col += 1
            
            # 全体累計
            if type_label == "社員":
                grand_total = sum(mt.employee for mt in plan.total_month_summary)
            elif type_label == "PN":
                grand_total = sum(mt.pn for mt in plan.total_month_summary)
            elif type_label == "SI":
                grand_total = sum(mt.si for mt in plan.total_month_summary)
            else:  # 合計
                grand_total = sum(mt.total for mt in plan.total_month_summary)
            
            # 累計セルのスタイルを種別に応じて変更
            if type_label == "合計":
                self._style_grand_total_cell(table.cell(row_idx, cols - 1), f"{grand_total:.1f}", 8)
            else:
                self._style_employee_data_cell(table.cell(row_idx, cols - 1), f"{grand_total:.1f}", 8)
        
        # ハイライトボックスを追加（メッセージに対応）- 月単位の場合はスキップ
        # TODO: 月単位用のハイライトボックス実装

    def _draw_schedule_table_quarterly(self, slide, plan: DevelopmentPersonnelPlan) -> None:
        """四半期単位のスケジュール表を描画する（12ヶ月超の場合）。"""
        fiscal_years = plan.fiscal_years
        if not fiscal_years:
            logger.warning("年度データがありません")
            return
        
        # 列数を計算: フェーズ名列 + 種別列 + 各年度×4四半期 + 開発期間累計列
        quarter_cols = len(fiscal_years) * 4
        cols = 2 + quarter_cols + 1  # フェーズ名列 + 種別列 + 四半期列 + 累計列
        
        # 行数を計算: ヘッダー2行 + マイルストーン/スケジュール行 + (フェーズ数 × 4行) + 合計4行
        phase_count = len(plan.phase_summaries)
        rows = 2 + 1 + (phase_count * 4) + 4  # ヘッダー2 + スケジュール1 + フェーズ×4 + 合計4
        
        # Schedule_boxプレースホルダーの位置を取得
        table_left = Inches(self.config.table_left_in)
        table_top = Inches(self.config.table_top_in)
        table_width = Inches(self.config.table_width_in)
        table_height = Inches(self.config.table_height_in)
        
        schedule_box = self._find_placeholder_by_name(slide, "Schedule_box")
        if schedule_box is not None:
            table_left = schedule_box.left
            table_top = schedule_box.top
            table_width = schedule_box.width
            table_height = schedule_box.height
            # Schedule_boxを削除
            sp = schedule_box.element
            sp.getparent().remove(sp)
            logger.debug("Schedule_boxプレースホルダーを発見: left=%s, top=%s, width=%s, height=%s",
                        table_left, table_top, table_width, table_height)
        else:
            # OBJECTタイプのプレースホルダーを検索
            from pptx.enum.shapes import PP_PLACEHOLDER
            for placeholder in slide.placeholders:
                try:
                    if placeholder.placeholder_format.type == PP_PLACEHOLDER.OBJECT:
                        table_left = placeholder.left
                        table_top = placeholder.top
                        table_width = placeholder.width
                        table_height = placeholder.height
                        sp = placeholder.element
                        sp.getparent().remove(sp)
                        logger.debug("OBJECTプレースホルダーを発見: left=%s, top=%s", table_left, table_top)
                        break
                except Exception:
                    continue
        
        # スケジュール行の高さを事前計算（トラック数に基づく）
        all_tasks = []
        for phase in plan.phases:
            all_tasks.extend(phase.tasks)
        
        # トラック数を計算
        tasks_with_tracks = self._assign_tracks_to_tasks(all_tasks, fiscal_years)
        if tasks_with_tracks:
            tracks_needed = max(track for _, track in tasks_with_tracks) + 1
        else:
            tracks_needed = 1
        
        # スケジュール行の高さ = マイルストーン用スペース + 間隔 + トラック数 × 矢羽の高さ + マージン
        schedule_row_height = Cm(
            self.config.milestone_row_height_cm +
            self.config.milestone_arrow_gap_cm +
            self.config.arrow_height_cm * tracks_needed +
            0.1
        )
        
        # 表を作成
        table_shape = slide.shapes.add_table(
            rows,
            cols,
            table_left,
            table_top,
            table_width,
            table_height,
        )
        table = table_shape.table
        
        # 列幅を設定
        phase_col_width = Inches(self.config.phase_col_width_in)
        type_col_width = Inches(self.config.type_col_width_in)
        total_col_width = Inches(self.config.total_col_width_in)
        quarter_col_width = int((table_width - phase_col_width - type_col_width - total_col_width) / quarter_cols)
        quarter_col_width_in = (table_width / 914400 - self.config.phase_col_width_in - self.config.type_col_width_in - self.config.total_col_width_in) / quarter_cols
        
        table.columns[0].width = phase_col_width
        table.columns[1].width = type_col_width
        for i in range(2, cols - 1):
            table.columns[i].width = quarter_col_width
        table.columns[cols - 1].width = total_col_width
        
        # 行の高さを設定
        header_row_height = Cm(self.config.header_row_height_cm)
        data_row_height = Cm(self.config.data_row_height_cm)
        # ヘッダー行1, 2の高さを0.8cm以内に設定
        table.rows[0].height = header_row_height
        table.rows[1].height = header_row_height
        # スケジュール行（3行目）の高さをトラック数に応じて設定
        table.rows[2].height = schedule_row_height
        # データ行の高さを設定
        for row_idx in range(3, rows):
            table.rows[row_idx].height = data_row_height
        
        # ヘッダー行1: 年度（フェーズ名列と種別列をマージ）
        self._style_header_cell(table.cell(0, 0), "", 8)  # 左上セルは空欄
        table.cell(0, 0).merge(table.cell(1, 1))  # 2列×2行をマージ
        
        col = 2
        for fy in fiscal_years:
            start_col = col
            for _ in range(4):
                col += 1
            # 年度セルをマージ
            table.cell(0, start_col).merge(table.cell(0, col - 1))
            self._style_header_cell(table.cell(0, start_col), f"{fy}年度", 8)
        
        # 累計列ヘッダー
        self._style_header_cell(table.cell(0, cols - 1), "開発期間", 7)
        table.cell(0, cols - 1).merge(table.cell(1, cols - 1))
        self._set_cell_text(table.cell(0, cols - 1), "開発期間\n累計", 7, bold=True, color=RGBColor(255, 255, 255))
        
        # ヘッダー行2: 四半期
        col = 2
        for fy in fiscal_years:
            for q in range(1, 5):
                self._style_header_cell(table.cell(1, col), f"{q}Q", 7)
                col += 1
        
        # マイルストーン/スケジュール行（3行目）
        schedule_row = 2
        self._style_milestone_label_cell(table.cell(schedule_row, 0), "マイルストーン\n/スケジュール", 6)
        table.cell(schedule_row, 0).merge(table.cell(schedule_row, 1))  # フェーズ名列と種別列をマージ
        
        # マイルストーン/スケジュール行のセルをスタイル設定（空白セル）- 白色背景
        col = 2
        for _ in fiscal_years:
            for _ in range(1, 5):
                self._style_milestone_data_cell(table.cell(schedule_row, col), "", 6)
                col += 1
        self._style_milestone_data_cell(table.cell(schedule_row, cols - 1), "", 7)
        
        # 矢羽をスケジュール行に描画
        if all_tasks:
            self._draw_schedule_arrows(
                slide,
                table_shape,
                table,
                schedule_row,
                all_tasks,
                fiscal_years,
                quarter_col_width_in,
                tasks_with_tracks,
            )
        
        # マイルストーンを三角形で描画
        self._draw_milestones_as_triangles(
            slide,
            table_shape,
            table,
            schedule_row,
            plan.milestones,
            fiscal_years,
            quarter_col_width_in,
        )
        
        # 種別ラベル（社員/PN/SI/小計）
        type_labels = ["社員", "PN", "SI", "小計"]
        
        # フェーズ行（4行目以降、各フェーズ4行）
        for phase_idx, phase in enumerate(plan.phase_summaries):
            base_row = 3 + phase_idx * 4
            
            # フェーズ名（4行をマージ）- 水色背景
            self._style_phase_cell(table.cell(base_row, 0), phase.phase_name, 7)
            table.cell(base_row, 0).merge(table.cell(base_row + 3, 0))
            
            for type_idx, type_label in enumerate(type_labels):
                row_idx = base_row + type_idx
                
                # 種別ラベルとデータセルのスタイルを種別に応じて変更
                if type_label == "小計":
                    # 小計行: 青色背景、黒文字
                    self._style_subtotal_type_cell(table.cell(row_idx, 1), type_label, 7)
                else:
                    # 社員/PN/SI行: 白色背景、黒文字
                    self._style_employee_type_cell(table.cell(row_idx, 1), type_label, 7)
                
                # 四半期工数
                col = 2
                for fy in fiscal_years:
                    for q in range(1, 5):
                        # 該当四半期のデータを取得
                        quarter_data = next(
                            (qd for qd in phase.quarters if qd.fiscal_year == fy and qd.quarter == q),
                            None,
                        )
                        if quarter_data:
                            if type_label == "社員":
                                value = quarter_data.employee
                            elif type_label == "PN":
                                value = quarter_data.pn
                            elif type_label == "SI":
                                value = quarter_data.si
                            else:  # 小計
                                value = quarter_data.total
                            text = f"{value:.1f}" if value > 0 else "-"
                        else:
                            text = "-"
                        
                        # データセルのスタイルを種別に応じて変更
                        if type_label == "小計":
                            self._style_subtotal_data_cell(table.cell(row_idx, col), text, 7)
                        else:
                            self._style_employee_data_cell(table.cell(row_idx, col), text, 7)
                        col += 1
                
                # 累計
                if type_label == "社員":
                    total = phase.total_employee
                elif type_label == "PN":
                    total = phase.total_pn
                elif type_label == "SI":
                    total = phase.total_si
                else:  # 小計
                    total = phase.total
                
                # 累計セルのスタイルを種別に応じて変更
                if type_label == "小計":
                    self._style_subtotal_data_cell(table.cell(row_idx, cols - 1), f"{total:.1f}", 7)
                else:
                    self._style_employee_data_cell(table.cell(row_idx, cols - 1), f"{total:.1f}", 7)
        
        # 合計行（最後の4行）
        total_base_row = rows - 4
        total_type_labels = ["社員", "PN", "SI", "合計"]  # 最終行は「合計」
        
        # 全体ラベル（4行をマージ）- 青色背景、太字なし
        self._style_overall_label_cell(table.cell(total_base_row, 0), "全体", 8)
        table.cell(total_base_row, 0).merge(table.cell(total_base_row + 3, 0))
        
        for type_idx, type_label in enumerate(total_type_labels):
            row_idx = total_base_row + type_idx
            
            # 種別ラベルとデータセルのスタイルを種別に応じて変更
            if type_label == "合計":
                # 合計行: 青色背景、黒文字、太字なし
                self._style_grand_total_cell(table.cell(row_idx, 1), type_label, 7)
            else:
                # 社員/PN/SI行: 白色背景、黒文字
                self._style_employee_type_cell(table.cell(row_idx, 1), type_label, 7)
            
            col = 2
            for fy in fiscal_years:
                for q in range(1, 5):
                    # 該当四半期の合計を取得
                    quarter_total = next(
                        (qt for qt in plan.total_summary if qt.fiscal_year == fy and qt.quarter == q),
                        None,
                    )
                    if quarter_total:
                        if type_label == "社員":
                            value = quarter_total.employee
                        elif type_label == "PN":
                            value = quarter_total.pn
                        elif type_label == "SI":
                            value = quarter_total.si
                        else:  # 合計
                            value = quarter_total.total
                        text = f"{value:.1f}" if value > 0 else "-"
                    else:
                        text = "-"
                    
                    # データセルのスタイルを種別に応じて変更
                    if type_label == "合計":
                        self._style_grand_total_cell(table.cell(row_idx, col), text, 7)
                    else:
                        self._style_employee_data_cell(table.cell(row_idx, col), text, 7)
                    col += 1
            
            # 全体累計
            if type_label == "社員":
                grand_total = sum(qt.employee for qt in plan.total_summary)
            elif type_label == "PN":
                grand_total = sum(qt.pn for qt in plan.total_summary)
            elif type_label == "SI":
                grand_total = sum(qt.si for qt in plan.total_summary)
            else:  # 合計
                grand_total = sum(qt.total for qt in plan.total_summary)
            
            # 累計セルのスタイルを種別に応じて変更
            if type_label == "合計":
                self._style_grand_total_cell(table.cell(row_idx, cols - 1), f"{grand_total:.1f}", 8)
            else:
                self._style_employee_data_cell(table.cell(row_idx, cols - 1), f"{grand_total:.1f}", 8)
        
        # ハイライトボックスを追加（メッセージに対応）
        self._add_highlight_boxes(slide, plan.messages, table_shape, fiscal_years)

    def _style_employee_type_cell(self, cell, text: str, font_size: int) -> None:
        """社員/PN/SI種別セルのスタイルを設定する（白色背景、黒文字）。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.WHITE_BG_COLOR
        self._set_cell_text(cell, text, font_size, bold=False, color=RGBColor(0, 0, 0))

    def _style_employee_data_cell(self, cell, text: str, font_size: int) -> None:
        """社員/PN/SIデータセルのスタイルを設定する（白色背景、黒文字）。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.WHITE_BG_COLOR
        self._set_cell_text(cell, text, font_size, bold=False, color=RGBColor(0, 0, 0))

    def _style_subtotal_type_cell(self, cell, text: str, font_size: int) -> None:
        """小計種別セルのスタイルを設定する（水色背景、黒文字）。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.CYAN_BG_COLOR
        self._set_cell_text(cell, text, font_size, bold=False, color=RGBColor(0, 0, 0))

    def _style_subtotal_data_cell(self, cell, text: str, font_size: int) -> None:
        """小計データセルのスタイルを設定する（水色背景、黒文字）。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.CYAN_BG_COLOR
        self._set_cell_text(cell, text, font_size, bold=False, color=RGBColor(0, 0, 0))

    def _style_overall_label_cell(self, cell, text: str, font_size: int) -> None:
        """全体ラベルセルのスタイルを設定する（青色背景、黒文字、太字なし）。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.BLUE_BG_COLOR
        self._set_cell_text(cell, text, font_size, bold=False, color=RGBColor(0, 0, 0))

    def _style_grand_total_cell(self, cell, text: str, font_size: int) -> None:
        """合計セルのスタイルを設定する（青色背景、黒文字、太字なし）。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.BLUE_BG_COLOR
        self._set_cell_text(cell, text, font_size, bold=False, color=RGBColor(0, 0, 0))

    def _add_highlight_boxes(
        self,
        slide,
        messages: list[PersonnelMessage],
        table_shape,
        fiscal_years: list[int],
    ) -> None:
        """メッセージのハイライトボックスを追加する。"""
        circled_numbers = ["①", "②", "③", "④", "⑤", "⑥", "⑦", "⑧", "⑨", "⑩"]
        
        for msg in messages:
            if not msg.highlight_quarters:
                continue
            
            for fy, q in msg.highlight_quarters:
                # 対応する列を計算（2列目から始まる）
                if fy not in fiscal_years:
                    continue
                fy_idx = fiscal_years.index(fy)
                col_idx = 2 + fy_idx * 4 + (q - 1)
                
                # 表の位置を基準にハイライトボックスを配置
                table = table_shape.table
                table_left = table_shape.left
                table_top = table_shape.top
                
                # 列の左端位置を計算（フェーズ名列と種別列をスキップ）
                col_left = table_left + table.columns[0].width + table.columns[1].width
                for i in range(2, col_idx):
                    col_left += table.columns[i].width
                
                col_width = table.columns[col_idx].width
                box_height = Inches(0.3)
                box_top = table_top - box_height - Inches(0.05)
                
                # ハイライトボックス（枠線のみ）
                highlight = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    col_left,
                    box_top,
                    col_width,
                    box_height,
                )
                highlight.fill.background()  # 透明
                highlight.line.color.rgb = self.HIGHLIGHT_BORDER_COLOR
                highlight.line.width = Pt(2)
                
                # 番号ラベル
                number_idx = msg.number - 1
                if number_idx < len(circled_numbers):
                    label_text = circled_numbers[number_idx]
                else:
                    label_text = f"({msg.number})"
                
                label = slide.shapes.add_textbox(
                    col_left - Inches(0.15),
                    box_top - Inches(0.15),
                    Inches(0.2),
                    Inches(0.2),
                )
                label_frame = label.text_frame
                para = label_frame.paragraphs[0]
                para.text = label_text
                para.font.size = Pt(10)
                para.font.bold = True
                para.font.color.rgb = self.HIGHLIGHT_BORDER_COLOR

    def _style_header_cell(self, cell, text: str, font_size: int) -> None:
        """ヘッダーセルのスタイルを設定する。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.HEADER_BG_COLOR
        self._set_cell_text(cell, text, font_size, bold=True, color=RGBColor(255, 255, 255))

    def _style_milestone_label_cell(self, cell, text: str, font_size: int) -> None:
        """マイルストーン/スケジュールラベルセルのスタイルを設定する（薄い水色背景、黒文字）。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.LIGHT_BLUE_BG_COLOR
        self._set_cell_text(cell, text, font_size, bold=True, color=RGBColor(0, 0, 0))

    def _style_milestone_data_cell(self, cell, text: str, font_size: int) -> None:
        """マイルストーン/スケジュール行の空白セルのスタイルを設定する（白色背景）。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.WHITE_BG_COLOR
        self._set_cell_text(cell, text, font_size, bold=False, color=RGBColor(0, 0, 0))

    def _style_phase_cell(self, cell, text: str, font_size: int) -> None:
        """フェーズセルのスタイルを設定する（薄い水色背景、黒文字）。"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.LIGHT_BLUE_BG_COLOR
        self._set_cell_text(cell, text, font_size, bold=False, color=RGBColor(0, 0, 0))

    def _set_cell_text(
        self,
        cell,
        text: str,
        font_size: int,
        *,
        bold: bool = False,
        color: RGBColor | None = None,
    ) -> None:
        """セルにテキストを設定する。"""
        text_frame = cell.text_frame
        text_frame.clear()
        para = text_frame.paragraphs[0]
        para.text = text
        para.font.size = Pt(font_size)
        para.font.bold = bold
        if color:
            para.font.color.rgb = color
        para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        cell.vertical_anchor = 1  # 中央揃え


    def _assign_tracks_to_tasks_monthly(
        self,
        tasks: list[PersonnelScheduleTask],
        target_months: list[tuple[int, int]],
    ) -> list[tuple[PersonnelScheduleTask, int]]:
        """月単位表示用のタスクトラック割り当て。"""
        # 四半期版と同じロジックを使用
        return self._assign_tracks_to_tasks(tasks, [])

    def _draw_schedule_arrows_monthly(
        self,
        slide,
        table_shape,
        table,
        row_idx: int,
        tasks: list[PersonnelScheduleTask],
        target_months: list[tuple[int, int]],
        month_col_width_in: float,
        tasks_with_tracks: list[tuple[PersonnelScheduleTask, int]] | None = None,
    ) -> None:
        """月単位表示用のスケジュール矢羽を描画する。"""
        # 表の位置情報を取得
        table_left = table_shape.left
        table_top = table_shape.top
        phase_col_width = table.columns[0].width + table.columns[1].width
        
        # 行の上端位置を計算
        row_top = table_top
        for i in range(row_idx):
            row_top += table.rows[i].height
        
        row_height = table.rows[row_idx].height
        
        # タスクの期間重複を検出してトラックを割り当て（未計算の場合）
        if tasks_with_tracks is None:
            tasks_with_tracks = self._assign_tracks_to_tasks(tasks, [])
        
        # 必要なトラック数を計算
        if tasks_with_tracks:
            max_track = max(track for _, track in tasks_with_tracks)
            tracks_needed = max_track + 1
        else:
            tracks_needed = 1
        
        # トラックの高さを計算
        track_height = Cm(self.config.arrow_height_cm)
        arrow_actual_height = max(Cm(self.config.min_arrow_height_cm), track_height - Cm(0.05))
        
        for task, track in tasks_with_tracks:
            start_offset, end_offset, valid = self._calculate_task_position_monthly(
                task, target_months, month_col_width_in
            )
            if not valid:
                continue
            
            # 矢羽の位置を計算
            arrow_left = table_left + phase_col_width + start_offset
            arrow_width = end_offset - start_offset
            
            # トラック位置を考慮
            track_top = row_top + Cm(self.config.milestone_row_height_cm + self.config.milestone_arrow_gap_cm) + Cm(0.025) + track * track_height
            
            # 矢羽の色を取得
            task_color = _get_task_color(task.task_name)
            
            # 矢羽図形（五角形）を作成
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.PENTAGON,
                arrow_left + Cm(0.02),
                track_top,
                arrow_width - Cm(0.04),
                arrow_actual_height,
            )
            
            # 塗りつぶし
            fill = arrow.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(task_color.lstrip("#"))
            
            # 枠線（青色）
            line = arrow.line
            line.color.rgb = RGBColor(0, 112, 192)
            line.width = Pt(1)
            
            # タスク名をテキストとして追加
            if arrow_width > Cm(2):
                text_frame = arrow.text_frame
                text_frame.clear()
                text_frame.word_wrap = True
                paragraph = text_frame.paragraphs[0]
                paragraph.text = task.task_name
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                
                font = paragraph.font
                font.size = Pt(6)
                font.bold = False
                font.color.rgb = RGBColor(0, 0, 0)

    def _calculate_task_position_monthly(
        self,
        task: PersonnelScheduleTask,
        target_months: list[tuple[int, int]],
        month_col_width_in: float,
    ) -> tuple[int, int, bool]:
        """月単位表示用のタスク位置計算。"""
        start_date = datetime.strptime(task.start_date, "%Y-%m-%d")
        end_date = datetime.strptime(task.end_date, "%Y-%m-%d")
        
        # 開始月と終了月のインデックスを探す
        start_ym = (start_date.year, start_date.month)
        end_ym = (end_date.year, end_date.month)
        
        # 開始位置を計算
        start_col_idx = -1
        for idx, ym in enumerate(target_months):
            if ym == start_ym:
                start_col_idx = idx
                break
            elif ym > start_ym:
                start_col_idx = idx
                break
        
        if start_col_idx == -1:
            # 範囲外の場合
            if start_ym < target_months[0]:
                start_col_idx = 0
            else:
                return 0, 0, False
        
        # 終了位置を計算
        end_col_idx = len(target_months) - 1
        for idx, ym in enumerate(target_months):
            if ym == end_ym:
                end_col_idx = idx
                break
            elif ym > end_ym:
                end_col_idx = max(0, idx - 1)
                break
        
        # 月内の位置を計算（開始日は月の何日目か）
        days_in_month = 30  # 概算
        start_day_ratio = (start_date.day - 1) / days_in_month
        end_day_ratio = min(1.0, end_date.day / days_in_month)
        
        # EMU単位で計算
        col_width_emu = int(month_col_width_in * 914400)
        start_offset_emu = int(start_col_idx * col_width_emu + start_day_ratio * col_width_emu)
        end_offset_emu = int(end_col_idx * col_width_emu + end_day_ratio * col_width_emu)
        
        return start_offset_emu, end_offset_emu, True

    def _draw_milestones_as_triangles_monthly(
        self,
        slide,
        table_shape,
        table,
        row_idx: int,
        milestones,
        target_months: list[tuple[int, int]],
        month_col_width_in: float,
    ) -> None:
        """月単位表示用のマイルストーン三角形を描画する。"""
        table_left = table_shape.left
        table_top = table_shape.top
        phase_col_width = table.columns[0].width + table.columns[1].width
        
        row_top = table_top
        for i in range(row_idx):
            row_top += table.rows[i].height
        
        for milestone in milestones:
            date = datetime.strptime(milestone.date, "%Y-%m-%d")
            ym = (date.year, date.month)
            
            # 対象月のインデックスを探す
            col_idx = -1
            for idx, target_ym in enumerate(target_months):
                if target_ym == ym:
                    col_idx = idx
                    break
            
            if col_idx == -1:
                continue
            
            # 月内の位置を計算
            days_in_month = 30
            day_ratio = max(0, min(1, (date.day - 1) / days_in_month))
            
            # 列の位置を計算
            col_width_emu = int(month_col_width_in * 914400)
            x_offset_emu = int(col_idx * col_width_emu + day_ratio * col_width_emu)
            
            # 三角形の位置
            triangle_size = Pt(10)
            triangle_left = table_left + phase_col_width + x_offset_emu - triangle_size / 2
            triangle_top = row_top + Cm(0.05)
            
            # 下向き三角形を作成
            triangle = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                triangle_left,
                triangle_top,
                triangle_size,
                triangle_size,
            )
            
            triangle.rotation = 180
            
            fill = triangle.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            line = triangle.line
            line.color.rgb = RGBColor(0, 0, 0)
            line.width = Pt(1)
            
            # マイルストーン名
            textbox_left = triangle_left + triangle_size + Cm(0.05)
            textbox = slide.shapes.add_textbox(
                textbox_left,
                triangle_top,
                Inches(1.2),
                triangle_size,
            )
            
            text_frame = textbox.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.text = milestone.name
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            
            font = paragraph.font
            font.size = Pt(6)
            font.bold = True
            font.color.rgb = RGBColor(0, 0, 0)

    def _draw_schedule_arrows(
        self,
        slide,
        table_shape,
        table,
        row_idx: int,
        tasks: list[PersonnelScheduleTask],
        fiscal_years: list[int],
        quarter_col_width_in: float,
        tasks_with_tracks: list[tuple[PersonnelScheduleTask, int]] | None = None,
    ) -> None:
        """スケジュール行にタスクの矢羽を描画する。
        
        Args:
            slide: スライドオブジェクト
            table_shape: 表シェイプ
            table: 表オブジェクト
            row_idx: スケジュール行のインデックス
            tasks: タスク一覧
            fiscal_years: 年度一覧
            quarter_col_width_in: 四半期列の幅（インチ）
            tasks_with_tracks: タスクとトラック番号のペアリスト（事前計算済みの場合）
        """
        # 表の位置情報を取得
        table_left = table_shape.left
        table_top = table_shape.top
        phase_col_width = table.columns[0].width + table.columns[1].width  # フェーズ名列 + 種別列
        
        # 行の上端位置を計算（ヘッダー2行分をスキップ）
        row_top = table_top
        for i in range(row_idx):
            row_top += table.rows[i].height
        
        row_height = table.rows[row_idx].height
        
        # タスクの期間重複を検出してトラックを割り当て（未計算の場合）
        if tasks_with_tracks is None:
            tasks_with_tracks = self._assign_tracks_to_tasks(tasks, fiscal_years)
        
        # 必要なトラック数を計算
        if tasks_with_tracks:
            max_track = max(track for _, track in tasks_with_tracks)
            tracks_needed = max_track + 1
        else:
            tracks_needed = 1
        
        # トラックの高さを計算
        track_height = Cm(self.config.arrow_height_cm)
        # 矢羽の実際の高さを計算（最小0.5cm以上を保証）
        arrow_actual_height = max(Cm(self.config.min_arrow_height_cm), track_height - Cm(0.05))
        
        for task, track in tasks_with_tracks:
            start_offset, end_offset, valid = self._calculate_task_position(
                task, fiscal_years, quarter_col_width_in
            )
            if not valid:
                continue
            
            # 矢羽の位置を計算
            arrow_left = table_left + phase_col_width + start_offset
            arrow_width = end_offset - start_offset
            
            # トラック位置を考慮（マイルストーン用スペース + 間隔の下に配置）
            track_top = row_top + Cm(self.config.milestone_row_height_cm + self.config.milestone_arrow_gap_cm) + Cm(0.025) + track * track_height
            
            # 矢羽の色を取得
            task_color = _get_task_color(task.task_name)
            
            # 矢羽図形（五角形）を作成（高さは最小0.5cm以上）
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.PENTAGON,
                arrow_left + Cm(0.02),
                track_top,
                arrow_width - Cm(0.04),
                arrow_actual_height,
            )
            
            # 塗りつぶし
            fill = arrow.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(task_color.lstrip("#"))
            
            # 枠線（青色）
            line = arrow.line
            line.color.rgb = RGBColor(0, 112, 192)  # 青色 #0070C0
            line.width = Pt(1)
            
            # タスク名を常に表示（幅が狭い場合もフォントサイズを小さくして表示）
            text_frame = arrow.text_frame
            text_frame.clear()
            text_frame.word_wrap = False
            paragraph = text_frame.paragraphs[0]
            paragraph.text = task.task_name
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            
            font = paragraph.font
            font.size = Pt(6)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

    def _assign_tracks_to_tasks(
        self,
        tasks: list[PersonnelScheduleTask],
        fiscal_years: list[int],
    ) -> list[tuple[PersonnelScheduleTask, int]]:
        """タスクの期間重複を検出してトラック番号を割り当てる。
        
        Returns:
            タスクとトラック番号のペアのリスト
        """
        # タスクを開始日順にソート
        sorted_tasks = sorted(
            tasks,
            key=lambda t: datetime.strptime(t.start_date, "%Y-%m-%d")
        )
        
        # トラック割り当て（グリーディアルゴリズム）
        tracks: list[list[tuple[datetime, datetime]]] = []
        task_tracks: list[tuple[PersonnelScheduleTask, int]] = []
        
        for task in sorted_tasks:
            task_start = datetime.strptime(task.start_date, "%Y-%m-%d")
            task_end = datetime.strptime(task.end_date, "%Y-%m-%d")
            
            # 既存のトラックで配置可能か確認
            placed = False
            for track_idx, track_periods in enumerate(tracks):
                # このトラックで重複がないか確認
                overlap = False
                for period_start, period_end in track_periods:
                    if task_start <= period_end and task_end >= period_start:
                        overlap = True
                        break
                
                if not overlap:
                    # このトラックに配置可能
                    track_periods.append((task_start, task_end))
                    task_tracks.append((task, track_idx))
                    placed = True
                    break
            
            if not placed:
                # 新しいトラックを作成
                tracks.append([(task_start, task_end)])
                task_tracks.append((task, len(tracks) - 1))
        
        return task_tracks

    def _calculate_task_position(
        self,
        task: PersonnelScheduleTask,
        fiscal_years: list[int],
        quarter_col_width_in: float,
    ) -> tuple[int, int, bool]:
        """タスクの開始・終了位置を計算する。
        
        Returns:
            (start_offset_emu, end_offset_emu, valid): 開始・終了オフセット（EMU）と有効フラグ
        """
        start_date = datetime.strptime(task.start_date, "%Y-%m-%d")
        end_date = datetime.strptime(task.end_date, "%Y-%m-%d")
        
        # 年度と四半期を計算
        start_fy = _get_fiscal_year(start_date.year, start_date.month)
        end_fy = _get_fiscal_year(end_date.year, end_date.month)
        
        # 四半期を計算
        def get_quarter(month: int) -> int:
            if 4 <= month <= 6:
                return 1
            elif 7 <= month <= 9:
                return 2
            elif 10 <= month <= 12:
                return 3
            else:
                return 4
        
        start_q = get_quarter(start_date.month)
        end_q = get_quarter(end_date.month)
        
        # 年度が範囲外の場合はスキップ
        if start_fy > fiscal_years[-1] or end_fy < fiscal_years[0]:
            return 0, 0, False
        
        # 開始位置の列インデックスを計算
        if start_fy < fiscal_years[0]:
            start_col_idx = 0
            start_day_ratio = 0.0
        else:
            fy_idx = fiscal_years.index(start_fy) if start_fy in fiscal_years else 0
            start_col_idx = fy_idx * 4 + (start_q - 1)
            # 四半期内の位置を計算
            quarter_start_month = {1: 4, 2: 7, 3: 10, 4: 1}[start_q]
            quarter_start = datetime(
                start_date.year if start_q != 4 else (start_date.year if start_date.month >= 4 else start_date.year),
                quarter_start_month, 1
            )
            days_in_quarter = 90  # 概算
            days_from_quarter_start = (start_date - quarter_start).days
            start_day_ratio = max(0, min(1, days_from_quarter_start / days_in_quarter))
        
        # 終了位置の列インデックスを計算
        if end_fy > fiscal_years[-1]:
            end_col_idx = len(fiscal_years) * 4 - 1
            end_day_ratio = 1.0
        else:
            fy_idx = fiscal_years.index(end_fy) if end_fy in fiscal_years else len(fiscal_years) - 1
            end_col_idx = fy_idx * 4 + (end_q - 1)
            # 四半期内の位置を計算
            quarter_start_month = {1: 4, 2: 7, 3: 10, 4: 1}[end_q]
            quarter_start = datetime(
                end_date.year if end_q != 4 else (end_date.year if end_date.month >= 4 else end_date.year),
                quarter_start_month, 1
            )
            days_in_quarter = 90
            days_from_quarter_start = (end_date - quarter_start).days + 1
            end_day_ratio = max(0, min(1, days_from_quarter_start / days_in_quarter))
        
        # EMU単位で計算（1インチ = 914400 EMU）
        col_width_emu = int(quarter_col_width_in * 914400)
        start_offset_emu = int(start_col_idx * col_width_emu + start_day_ratio * col_width_emu)
        end_offset_emu = int(end_col_idx * col_width_emu + end_day_ratio * col_width_emu)
        
        return start_offset_emu, end_offset_emu, True

    def _draw_milestones_as_triangles(
        self,
        slide,
        table_shape,
        table,
        row_idx: int,
        milestones,
        fiscal_years: list[int],
        quarter_col_width_in: float,
    ) -> None:
        """マイルストーンを下向き三角形で描画する。"""
        # 表の位置情報を取得
        table_left = table_shape.left
        table_top = table_shape.top
        phase_col_width = table.columns[0].width + table.columns[1].width  # フェーズ名列 + 種別列
        
        # 行の上端位置を計算
        row_top = table_top
        for i in range(row_idx):
            row_top += table.rows[i].height
        
        row_height = table.rows[row_idx].height
        
        for milestone in milestones:
            date = datetime.strptime(milestone.date, "%Y-%m-%d")
            fy = _get_fiscal_year(date.year, date.month)
            
            # 年度が範囲外の場合はスキップ
            if fy not in fiscal_years:
                continue
            
            # 四半期を計算
            if 4 <= date.month <= 6:
                q = 1
            elif 7 <= date.month <= 9:
                q = 2
            elif 10 <= date.month <= 12:
                q = 3
            else:
                q = 4
            
            # 列インデックスを計算
            fy_idx = fiscal_years.index(fy)
            col_idx = fy_idx * 4 + (q - 1)
            
            # 四半期内の位置を計算
            quarter_start_month = {1: 4, 2: 7, 3: 10, 4: 1}[q]
            quarter_start_year = date.year if q != 4 else (date.year if date.month >= 4 else date.year)
            quarter_start = datetime(quarter_start_year, quarter_start_month, 1)
            days_in_quarter = 90
            days_from_quarter_start = (date - quarter_start).days
            day_ratio = max(0, min(1, days_from_quarter_start / days_in_quarter))
            
            # 列の位置を計算
            col_width_emu = int(quarter_col_width_in * 914400)
            x_offset_emu = int(col_idx * col_width_emu + day_ratio * col_width_emu)
            
            # 三角形の位置
            triangle_size = Pt(10)
            triangle_left = table_left + phase_col_width + x_offset_emu - triangle_size / 2
            triangle_top = row_top + Cm(0.05)
            
            # 下向き三角形を作成
            triangle = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                triangle_left,
                triangle_top,
                triangle_size,
                triangle_size,
            )
            
            # 三角形を180度回転（下向きに）
            triangle.rotation = 180
            
            # 塗りつぶし（白色）
            fill = triangle.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            # 枠線（黒色）
            line = triangle.line
            line.color.rgb = RGBColor(0, 0, 0)
            line.width = Pt(1)
            
            # マイルストーン名をテキストボックスで追加
            textbox_left = triangle_left + triangle_size + Cm(0.05)
            textbox = slide.shapes.add_textbox(
                textbox_left,
                triangle_top,
                Inches(1.2),
                triangle_size,
            )
            
            text_frame = textbox.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.text = milestone.name
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            
            font = paragraph.font
            font.size = Pt(6)
            font.bold = True
            font.color.rgb = RGBColor(0, 0, 0)
