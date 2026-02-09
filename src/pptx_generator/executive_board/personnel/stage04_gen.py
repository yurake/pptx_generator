#!/usr/bin/env python3
"""Stage4 hook: 開発要員計画スライドを既存PPTXへ追加する。"""
from __future__ import annotations

import json
import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from common.stage_shared import load_context, resolve_local_path  # noqa: E402


def _configure_sys_path() -> None:
    src_dir = Path(__file__).resolve().parents[3] / "src"
    if str(src_dir) not in sys.path:
        sys.path.insert(0, str(src_dir))


DEFAULT_TEMPLATE = Path(__file__).resolve().parents[1] / "input/executive_board.pptx"
DEFAULT_LAYOUT = "2_System_layout"


_configure_sys_path()

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from personnel.models import DevelopmentPersonnelPlan  # noqa: E402
from schedule.models import ScheduleGantt  # noqa: E402
from schedule.schedule_renderer import ScheduleGanttRenderer  # noqa: E402
from personnel.development_personnel_renderer import DevelopmentPersonnelRenderer  # noqa: E402
from common.rendering.title_style import apply_title_font_size, apply_title_prefix  # noqa: E402
from common.rendering.text_cleanup import apply_strip_bullet_prefix  # noqa: E402
from common.rendering.text_style import apply_text_font_size_by_name  # noqa: E402


class AppendScheduleRenderer(ScheduleGanttRenderer):  # type: ignore[misc]
    """既存のPPTXにスケジュールスライドを追加するレンダラー（既存スライドは保持）。"""

    def _clear_existing_slides(self, presentation) -> None:  # override
        return


def _load_plan(context: dict[str, str]) -> DevelopmentPersonnelPlan:
    plan_candidates = [
        os.environ.get("PPTX_PERSONNEL_PLAN_JSON"),
        context.get("personnel_plan_path"),
    ]
    for candidate in plan_candidates:
        if not candidate:
            continue
        path = Path(candidate).expanduser().resolve()
        if path.exists():
            data = json.loads(path.read_text(encoding="utf-8"))
            return DevelopmentPersonnelPlan.model_validate(data)
    raise FileNotFoundError("development_personnel_plan.json が見つかりません。prepare フックを先に実行してください。")


def _resolve_template_path(context: dict[str, str], output_path: Path) -> Path:
    if output_path.exists():
        return output_path
    env_template = os.environ.get("PPTX_TEMPLATE_PATH")
    candidates = [
        env_template,
        context.get("template_path"),
        str(DEFAULT_TEMPLATE),
    ]
    for cand in candidates:
        if not cand:
            continue
        resolved = resolve_local_path(cand, Path(__file__).resolve().parent)
        if resolved.exists():
            return resolved
    raise FileNotFoundError("Template PPTX path could not be resolved for personnel slide.")


def _apply_title_and_messages(slide, plan: DevelopmentPersonnelPlan) -> None:
    """タイトルとメッセージ、日付をプレースホルダーに反映する。"""
    for shape in slide.shapes:
        if (
            getattr(shape, "is_placeholder", False)
            and getattr(shape, "placeholder_format", None)
            and shape.placeholder_format.type == 1
            and hasattr(shape, "text_frame")
        ):
            tf = shape.text_frame
            tf.text = plan.title
            if tf.paragraphs:
                tf.paragraphs[0].font.size = Pt(28)
                tf.paragraphs[0].font.bold = True
            break

    bodies = [
        sh
        for sh in slide.shapes
        if getattr(sh, "is_placeholder", False)
        and getattr(sh, "placeholder_format", None)
        and sh.placeholder_format.type == 2
        and hasattr(sh, "text_frame")
    ]
    circled = ["①", "②", "③", "④", "⑤", "⑥", "⑦", "⑧", "⑨"]
    lines: list[str] = []
    for idx, msg in enumerate(plan.messages):
        prefix = circled[idx] if idx < len(circled) else f"{idx+1}."
        lines.append(f"{prefix}{msg.text}")
    message_text = "\n".join(lines)

    date_text = plan.department or ""
    try:
        from datetime import datetime

        dt = datetime.fromisoformat(plan.generated_at).date()
        date_text = f"{plan.department} {dt.strftime('%Y年%m月%d日')}".strip()
    except Exception:
        pass

    if bodies:
        tf = bodies[0].text_frame
        tf.clear()
        parts = message_text.split("\n") if message_text else []
        if parts:
            tf.paragraphs[0].text = parts[0]
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.bold = None
            tf.paragraphs[0].alignment = None
            for line in parts[1:]:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(11)
                p.font.bold = None
                p.alignment = None
    if len(bodies) >= 2 and date_text:
        tf = bodies[1].text_frame
        tf.text = date_text
        if tf.paragraphs:
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT


def _add_personnel_table(slide, plan: DevelopmentPersonnelPlan) -> None:
    """期待レイアウトに合わせた人月表を追加する。"""
    years = plan.fiscal_years
    quarters = [(fy, q) for fy in years for q in range(1, 5)]
    cols = 2 + len(quarters) + 1  # phase/role + quarters + total

    rows = []
    header1 = ["", ""] + [str(fy) + "年度" if i % 4 == 0 else "" for i, (fy, _) in enumerate(quarters)] + ["開発期間\n累計"]
    header2 = ["", ""] + [f"{q}Q" for _, q in quarters] + [""]
    rows.append(header1)
    rows.append(header2)
    rows.append(["マイルストーン\n/スケジュール"] + [""] * (cols - 1))

    def quarter_map(summary_list):
        return {(item.fiscal_year, item.quarter): item for item in summary_list}

    def fmt(val: float | None) -> str:
        if val is None or val == 0:
            return "-"
        if float(val).is_integer():
            return f"{int(val)}"
        return f"{val:.1f}"

    def row_from_summary(name: str, data_map, attr: str) -> list[str]:
        vals = []
        for key in quarters:
            item = data_map.get(key)
            num = getattr(item, attr, None) if item else None
            vals.append(fmt(num))
        total = sum(getattr(v, attr, 0) for v in data_map.values() if v)
        vals.append(fmt(total))
        return [name] + vals

    for phase in plan.phase_summaries:
        qmap = quarter_map(phase.quarters)
        rows.append([phase.phase_name, "社員"] + row_from_summary("", qmap, "employee")[1:])
        rows.append(["", "PN"] + row_from_summary("", qmap, "pn")[1:])
        rows.append(["", "SI"] + row_from_summary("", qmap, "si")[1:])
        subtotal = []
        for key in quarters:
            item = qmap.get(key)
            total = (item.employee if item else 0) + (item.pn if item else 0) + (item.si if item else 0)
            subtotal.append(fmt(total))
        subtotal.append(fmt(sum(float(x) for x in subtotal if x not in {"-", ""})))
        rows.append(["", "小計"] + subtotal)

    total_map = quarter_map(plan.total_summary)
    rows.append(["全体", "社員"] + row_from_summary("", total_map, "employee")[1:])
    rows.append(["", "PN"] + row_from_summary("", total_map, "pn")[1:])
    rows.append(["", "SI"] + row_from_summary("", total_map, "si")[1:])
    overall_vals = []
    for key in quarters:
        item = total_map.get(key)
        total = (item.employee if item else 0) + (item.pn if item else 0) + (item.si if item else 0)
        overall_vals.append(fmt(total))
    overall_vals.append(fmt(sum(float(x) for x in overall_vals if x not in {"-", ""})))
    rows.append(["", "合計"] + overall_vals)

    row_count = len(rows)
    table_shape = slide.shapes.add_table(
        row_count,
        cols,
        Inches(0.18),
        Inches(2.12),
        Inches(10.472),
        Inches(3.445),
    )
    table = table_shape.table

    # 列幅を調整（期待レイアウトに合わせる）
    col_widths_in = [1.0, 0.5] + [1.0465] * 8 + [0.6]
    for idx, w in enumerate(col_widths_in):
        table.columns[idx].width = Inches(w)

    # 行高さを調整（ヘッダー2行0.315in、スケジュール行0.689in、データ行0.177in）
    header_h = Inches(0.315)
    schedule_h = Inches(0.689)
    data_h = Inches(0.177)
    if len(table.rows) >= 3:
        table.rows[0].height = header_h
        table.rows[1].height = header_h
        table.rows[2].height = schedule_h
        for r in range(3, len(table.rows)):
            table.rows[r].height = data_h

    header_blue = RGBColor(68, 114, 196)
    light_blue = RGBColor(221, 235, 247)
    cyan = RGBColor(155, 194, 230)
    white = RGBColor(255, 255, 255)

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val

            # 背景色の決定
            fill_color = None
            font_color = None
            if r_idx in (0, 1):  # ヘッダー
                if val.strip() or (r_idx == 0 and c_idx == 0):
                    fill_color = header_blue
                    font_color = white
            elif "マイルストーン" in row[0]:
                if c_idx == 0:
                    fill_color = light_blue
                elif c_idx == 1:
                    fill_color = None
                else:
                    fill_color = white
            elif row[1] == "小計":
                fill_color = cyan if c_idx >= 1 else None
            elif row[1] == "合計":
                fill_color = header_blue if c_idx >= 1 else None
                font_color = None
            elif row[0] and c_idx == 0:
                fill_color = light_blue
            else:
                fill_color = white if r_idx >= 3 else fill_color

            _style_table_cell(cell, header=r_idx <= 2, fill_color=fill_color, font_color=font_color)
    return table_shape


def _style_table_cell(cell, header: bool, fill_color=None, font_color=None) -> None:
    cell.text_frame.word_wrap = True
    for para in cell.text_frame.paragraphs:
        para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        if header:
            para.font.bold = True
            para.font.size = Pt(8)
        else:
            para.font.size = Pt(8)
        if font_color:
            para.font.color.rgb = font_color
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def main(argv: list[str] | None = None) -> int:
    _ = argv  # unused

    context = load_context()

    output_dir_env = os.environ.get("PPTX_OUTPUT_DIR")
    pptx_name_env = os.environ.get("PPTX_PPTX_NAME")
    output_pptx_env = os.environ.get("PPTX_OUTPUT_PPTX_PATH")

    if output_pptx_env:
        output_path = Path(output_pptx_env).expanduser().resolve()
    else:
        output_dir = Path(output_dir_env) if output_dir_env else Path(".pptx/gen")
        output_dir = output_dir.expanduser().resolve()
        output_dir.mkdir(parents=True, exist_ok=True)
        pptx_name = pptx_name_env or "executive_board_generated.pptx"
        output_path = output_dir / pptx_name

    plan = _load_plan(context)
    template_path = _resolve_template_path(context, output_path)
    layout_name = os.environ.get("PPTX_PERSONNEL_LAYOUT") or context.get("personnel_layout") or DEFAULT_LAYOUT

    # 1) スケジュールスライドを追加（従来どおり）
    schedule_json = context.get("schedule_json_path")
    if not schedule_json or not Path(schedule_json).expanduser().exists():
        raise FileNotFoundError("schedule_json_path が context にありません。prepare ステージを確認してください。")
    schedule = ScheduleGantt.parse_file(schedule_json)

    schedule_renderer = AppendScheduleRenderer(template_path=template_path)
    schedule_renderer.render(schedule, output_path, layout_name=layout_name)
    prs = Presentation(output_path)
    if prs.slides:
        apply_title_font_size(prs.slides[-1])
        apply_title_prefix(prs.slides[-1], 4)
        apply_strip_bullet_prefix(prs.slides[-1])
        prs.save(output_path)

    # 2) 旧実装の DevelopmentPersonnelRenderer を用いて人月スライドを追加
    class AppendPersonnelRenderer(DevelopmentPersonnelRenderer):  # type: ignore[misc]
        def _clear_existing_slides(self, presentation):  # override
            return

        def _load_template(self):  # override: 既存PPTXをテンプレートとして読み込む
            if self.template_path and Path(self.template_path).exists():
                return Presentation(self.template_path)
            return super()._load_template()

    personnel_renderer = AppendPersonnelRenderer(template_path=output_path)
    personnel_renderer.render(plan, output_path, layout_name=layout_name)
    prs = Presentation(output_path)
    if prs.slides:
        apply_title_font_size(prs.slides[-1])
        apply_title_prefix(prs.slides[-1], 6)
        apply_text_font_size_by_name(prs.slides[-1], "Text Placeholder 1", 14)
        apply_strip_bullet_prefix(prs.slides[-1])
        prs.save(output_path)

    # スケジュールスライド（9列テーブル）を複数生成した場合、最初の1枚だけ残す
    prs = Presentation(output_path)
    schedule_indices = []
    for idx, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if getattr(sh, "has_table", False) and len(sh.table.columns) == 9:
                schedule_indices.append(idx)
                break
    if len(schedule_indices) > 1:
        keep = min(schedule_indices)
        for idx in sorted([i for i in schedule_indices if i != keep], reverse=True):
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]
        prs.save(output_path)

    # overview の不要テーブルをクリア
    prs = Presentation(output_path)
    if prs.slides:
        slide = prs.slides[0]
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            if getattr(shape, "name", "").strip() != "表 10":
                continue
            for row in shape.table.rows:
                for cell in row.cells:
                    cell.text = ""
        prs.save(output_path)

    print(f"[stage04_gen] personnel slide appended -> {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
