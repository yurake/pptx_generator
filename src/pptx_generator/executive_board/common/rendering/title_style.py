#!/usr/bin/env python3
from __future__ import annotations

from typing import Iterable

from pptx.util import Pt


def _is_title_shape(name: str) -> bool:
    if not name:
        return False
    lowered = name.lower()
    return "タイトル" in name or lowered.startswith("title")


def apply_title_font_size(slide, font_size: int = 22) -> None:
    for shape in slide.shapes:
        name = getattr(shape, "name", "")
        if not _is_title_shape(name):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            for run in paragraph.runs:
                run.font.size = Pt(font_size)


def apply_title_prefix(slide, page_no: int) -> None:
    fullwidth_digits = "０１２３４５６７８９"
    if 0 <= page_no <= 9:
        prefix = f"{fullwidth_digits[page_no]}．"
    else:
        prefix = f"{page_no}．"
    for shape in slide.shapes:
        name = getattr(shape, "name", "")
        if not _is_title_shape(name):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        text_frame = shape.text_frame
        if not text_frame.paragraphs:
            continue
        paragraph = text_frame.paragraphs[0]
        current = paragraph.text or ""
        if current.startswith(prefix):
            continue
        if len(current) >= 2 and current[1] == "．":
            if current[0].isdigit() or current[0] in fullwidth_digits:
                continue
        paragraph.text = f"{prefix}{current}"
