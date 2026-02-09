#!/usr/bin/env python3
from __future__ import annotations

from pptx.util import Pt


def apply_text_font_size_by_name(slide, shape_name: str, font_size: int) -> None:
    for shape in slide.shapes:
        if getattr(shape, "name", "").strip() != shape_name:
            continue
        if not getattr(shape, "has_text_frame", False):
            return
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
        return


def apply_text_style_by_name(
    slide,
    shape_name: str,
    *,
    font_size: int | None = None,
    bold: bool | None = None,
) -> None:
    for shape in slide.shapes:
        if getattr(shape, "name", "").strip() != shape_name:
            continue
        if not getattr(shape, "has_text_frame", False):
            return
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            if font_size is not None:
                paragraph.font.size = Pt(font_size)
            if bold is not None:
                paragraph.font.bold = bold
            for run in paragraph.runs:
                if font_size is not None:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold
        return
