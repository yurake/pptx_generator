"""pptx_generator CLI."""

from __future__ import annotations

import hashlib
import json
import logging
import os
import shutil
from dataclasses import dataclass, replace
from datetime import datetime, timezone
from importlib import resources as importlib_resources
from pathlib import Path
from typing import Any, Optional

import click
from dotenv import load_dotenv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pydantic import BaseModel, ValidationError

from .branding_extractor import (BrandingExtractionError,
                                 extract_branding_config)
from .draft_intel import load_return_reasons
from .generate_ready import generate_ready_to_jobspec
from .layout_validation import (LayoutValidationError, LayoutValidationOptions,
                                LayoutValidationResult, LayoutValidationSuite)
from .models import (ContentApprovalDocument, DraftDocument,
                     GenerateReadyDocument, JobSpec, JobSpecScaffold,
                     SpecValidationError, TemplateRelease,
                     TemplateReleaseDiagnostics, TemplateReleaseGoldenRun,
                     TemplateReleaseReport, TemplateSpec)
from .pipeline import (AnalyzerOptions, ContentApprovalOptions,
                       ContentApprovalStep, DraftStructuringOptions,
                       DraftStructuringStep, MappingOptions, MappingStep,
                       MonitoringIntegrationOptions, MonitoringIntegrationStep,
                       PdfExportError, PdfExportOptions, PdfExportStep,
                       PipelineContext, PipelineRunner, PipelineStep,
                       PolisherError, PolisherOptions, PolisherStep,
                       PrepareNormalizationError, PrepareNormalizationOptions,
                       PrepareNormalizationStep, RefinerOptions,
                       RenderingAuditOptions, RenderingAuditStep,
                       RenderingOptions, SimpleAnalyzerStep, SimpleRefinerStep,
                       SimpleRendererStep, SpecValidatorStep,
                       TemplateExtractor, TemplateExtractorOptions)
from .pipeline.analyzer import SlideSnapshot
from .pipeline.draft_structuring import DraftStructuringError
from .prepare import (PrepareAIOrchestrationError, PrepareAIOrchestrator,
                      PrepareDocument, PreparePolicyError,
                      PrepareSourceDocument, load_prepare_policy_set)
from .review_engine import AnalyzerReviewEngineAdapter
from .settings import BrandingConfig, RulesConfig
from .spec_loader import load_jobspec_from_path
from .template_audit import (build_release_report, build_template_release,
                             load_template_release)

DEFAULT_RULES_PATH = Path("config/rules.json")
DEFAULT_BRANDING_PATH = Path("config/branding.json")
DEFAULT_CHAPTER_TEMPLATES_DIR = Path("config/chapter_templates")
DEFAULT_RETURN_REASONS_PATH = Path("config/return_reasons.json")
DEFAULT_PREPARE_POLICY_PATH = Path("config/prepare_policies/default.json")
DEFAULT_PREPARE_OUTPUT_DIR = Path(".pptx/prepare")
DEFAULT_JOBSPEC_PATH = Path(".pptx/extract/jobspec.json")

logger = logging.getLogger(__name__)

_LOG_LEVEL_ALIASES = {
    "debug": logging.DEBUG,
    "info": logging.INFO,
    "warning": logging.WARNING,
    "warn": logging.WARNING,
    "error": logging.ERROR,
    "err": logging.ERROR,
    "fatal": logging.CRITICAL,
    "critical": logging.CRITICAL,
}


@dataclass(slots=True)
class OutlineResult:
    """アウトライン stage の実行結果。"""

    context: PipelineContext
    draft_path: Path
    approved_path: Path
    log_path: Path
    meta_path: Path
    generate_ready_path: Path
    generate_ready_meta_path: Path


_DEFAULT_DRAFT_OPTIONS = DraftStructuringOptions()
DEFAULT_DRAFT_FILENAME = _DEFAULT_DRAFT_OPTIONS.draft_filename
DEFAULT_APPROVED_FILENAME = _DEFAULT_DRAFT_OPTIONS.approved_filename
DEFAULT_DRAFT_LOG_FILENAME = _DEFAULT_DRAFT_OPTIONS.log_filename
DEFAULT_GENERATE_READY_FILENAME = _DEFAULT_DRAFT_OPTIONS.generate_ready_filename
DEFAULT_GENERATE_READY_META_FILENAME = _DEFAULT_DRAFT_OPTIONS.generate_ready_meta_filename
DEFAULT_DRAFT_META_FILENAME = "draft_meta.json"


load_dotenv()


def _parse_log_level(value: str | None) -> int | None:
    if value is None:
        return None
    candidate = value.strip()
    if not candidate:
        return None
    lowered = candidate.lower()
    if lowered in _LOG_LEVEL_ALIASES:
        return _LOG_LEVEL_ALIASES[lowered]
    try:
        numeric_level = int(candidate)
    except ValueError:
        return None
    return numeric_level


def _determine_log_level(verbose: bool, debug: bool) -> tuple[int, list[tuple[int, str]]]:
    """CLI 全体のログレベルを決定する。"""
    deferred_logs: list[tuple[int, str]] = []

    if debug:
        return logging.DEBUG, deferred_logs
    if verbose:
        return logging.INFO, deferred_logs

    env_level = os.getenv("LOG_LEVEL")
    parsed_level = _parse_log_level(env_level)
    if env_level:
        if parsed_level is not None:
            return parsed_level, deferred_logs
        deferred_logs.append(
            (
                logging.WARNING,
                f"LOG_LEVEL='{env_level}' を解釈できません。WARNING レベルにフォールバックします。",
            )
        )

    if os.getenv("OPENAI_LOG"):
        deferred_logs.append(
            (
                logging.WARNING,
                "OPENAI_LOG 環境変数は廃止されました。LOG_LEVEL を利用してください。",
            )
        )

    return logging.WARNING, deferred_logs


def _discover_template_ai_policy() -> Path | None:
    """Return the default template AI policy path if available."""

    env_policy = os.getenv("PPTX_TEMPLATE_AI_POLICY")
    if env_policy:
        candidate = Path(env_policy)
        if candidate.is_file():
            return candidate.resolve()

    cwd_candidate = Path.cwd() / "config" / "template_ai_policies.json"
    if cwd_candidate.exists():
        logger.info("Detected default template AI policy: %s", cwd_candidate)
        return cwd_candidate.resolve()

    try:
        resource = importlib_resources.files("pptx_generator").joinpath(
            "config/template_ai_policies.json"
        )
        if resource.is_file():
            try:
                text = resource.read_text(encoding="utf-8")
            except FileNotFoundError:
                text = None
            if text is not None:
                cache_dir = Path(".pptx/cache")
                cache_dir.mkdir(parents=True, exist_ok=True)
                cache_path = cache_dir / "template_ai_policies.json"
                cache_path.write_text(text, encoding="utf-8")
                logger.info(
                    "Detected default template AI policy (package resource cached to %s)",
                    cache_path,
                )
                return cache_path.resolve()
    except (ModuleNotFoundError, FileNotFoundError):
        pass

    for parent in Path(__file__).resolve().parents:
        candidate = parent / "config" / "template_ai_policies.json"
        if candidate.exists():
            logger.info("Detected default template AI policy: %s", candidate)
            return candidate.resolve()

    return None


def _configure_llm_logger() -> None:
    log_dir = Path("logs")
    log_dir.mkdir(parents=True, exist_ok=True)
    llm_logger = logging.getLogger("pptx_generator.content_ai.llm")

    class _LLMLogFormatter(logging.Formatter):
        """content_ai ログ用の安全なフォーマッタ。"""

        def __init__(self, *args, max_chars: int = 2000, **kwargs) -> None:
            super().__init__(*args, **kwargs)
            self._max_chars = max_chars

        def _sanitize(self, value: object) -> str:
            if value is None:
                return "-"
            text = str(value).replace("\n", "\\n")
            if len(text) > self._max_chars:
                return f"{text[: self._max_chars]}...(truncated)"
            return text

        def format(self, record: logging.LogRecord) -> str:  # noqa: D401
            for attr in ("slide_id", "card_id", "model", "intent", "reason", "finish_reason", "refusal"):
                if not hasattr(record, attr):
                    setattr(record, attr, "-")

            record.warnings = self._sanitize(getattr(record, "warnings", None))
            record.prompt_excerpt = self._sanitize(getattr(record, "prompt", None))
            record.raw_response_excerpt = self._sanitize(getattr(record, "raw_response", None))

            return super().format(record)

    formatter = _LLMLogFormatter(
        fmt=(
            "%(asctime)s %(levelname)s %(name)s "
            "slide_id=%(slide_id)s card_id=%(card_id)s model=%(model)s intent=%(intent)s "
            "reason=%(reason)s finish=%(finish_reason)s refusal=%(refusal)s warnings=%(warnings)s "
                "message=%(message)s prompt=%(prompt_excerpt)s raw_response=%(raw_response_excerpt)s"
        ),
    )
    class _LLMLogFilter(logging.Filter):
        def filter(self, record: logging.LogRecord) -> bool:
            return bool(getattr(record, "raw_response", None) or getattr(record, "prompt", None))

    if not any(isinstance(f, _LLMLogFilter) for f in llm_logger.filters):
        llm_logger.addFilter(_LLMLogFilter())

    existing_handler = next(
        (
            handler
            for handler in llm_logger.handlers
            if isinstance(handler, logging.FileHandler)
            and getattr(handler, "baseFilename", None) == str(log_dir / "out.log")
        ),
        None,
    )
    if existing_handler:
        existing_handler.setFormatter(formatter)
    else:
        handler = logging.FileHandler(log_dir / "out.log", encoding="utf-8")
        handler.setFormatter(formatter)
        llm_logger.addHandler(handler)

    stream_handler_exists = any(
        isinstance(handler, logging.StreamHandler) and not isinstance(handler, logging.FileHandler)
        for handler in llm_logger.handlers
    )
    if not stream_handler_exists:
        stream_handler = logging.StreamHandler()
        stream_handler.setFormatter(formatter)
        llm_logger.addHandler(stream_handler)
    llm_logger.setLevel(logging.INFO)
    llm_logger.propagate = False


def _log_current_llm_provider(context: str) -> None:
    provider_env = os.getenv("PPTX_LLM_PROVIDER")
    provider = provider_env.strip().lower() if provider_env else "mock"
    source = "env" if provider_env else "default"
    logging.getLogger("pptx_generator.cli.llm").info(
        "LLM provider (%s): %s (source=%s)", context, provider, source
    )


def _configure_file_logging() -> None:
    log_dir = Path("logs")
    log_dir.mkdir(parents=True, exist_ok=True)
    file_path = log_dir / "out.log"
    root_logger = logging.getLogger()
    if not any(
        isinstance(handler, logging.FileHandler)
        and getattr(handler, "baseFilename", None) == str(file_path)
        for handler in root_logger.handlers
    ):
        handler = logging.FileHandler(file_path, encoding="utf-8")
        formatter = logging.Formatter(
            "%(asctime)s %(levelname)s %(name)s %(message)s")
        handler.setFormatter(formatter)
        root_logger.addHandler(handler)


def _resolve_config_path(value: str, *, base_dir: Path | None = None) -> Path:
    """設定ファイルで指定されたパスを解決する。"""
    candidate = Path(value)
    if candidate.is_absolute():
        resolved = candidate
    else:
        if candidate.parts and candidate.parts[0] == "config":
            resolved = Path.cwd() / candidate
        elif base_dir is not None:
            resolved = base_dir / candidate
        else:
            resolved = Path.cwd() / candidate
    resolved = resolved.resolve()
    if not resolved.exists():
        msg = f"設定ファイルで指定されたパスが見つかりません: {resolved}"
        raise FileNotFoundError(msg)
    return resolved


@click.group(
    help="JSON 仕様から PPTX を生成する CLI",
    context_settings={"help_option_names": ["-h", "--help"]},
)
@click.option("-v", "--verbose", is_flag=True, help="INFO レベルの冗長ログを出力する")
@click.option("--debug", is_flag=True, help="DEBUG レベルで詳細ログを出力する")
def app(verbose: bool, debug: bool) -> None:
    """CLI ルートエントリ。"""
    level, deferred_logs = _determine_log_level(verbose, debug)
    logging.basicConfig(
        level=level,
        format="%(asctime)s %(levelname)s %(name)s - %(message)s",
        force=True,
    )
    logging.getLogger("openai").setLevel(level)
    cli_logger = logging.getLogger("pptx_generator.cli")
    for message_level, message in deferred_logs:
        cli_logger.log(message_level, message)
    _configure_llm_logger()
    _configure_file_logging()


def _prepare_branding(
    template: Optional[Path], branding: Optional[Path]
) -> tuple[BrandingConfig, dict[str, object]]:
    def _load_default_branding() -> tuple[BrandingConfig, dict[str, object]]:
        try:
            config = BrandingConfig.load(DEFAULT_BRANDING_PATH)
            return config, {"type": "default", "path": str(DEFAULT_BRANDING_PATH)}
        except FileNotFoundError:
            click.echo(
                f"デフォルトのブランド設定が見つかりません: {DEFAULT_BRANDING_PATH}. 内蔵設定を使用します。",
                err=True,
            )
            return BrandingConfig.default(), {"type": "builtin"}

    branding_payload: dict[str, object] | None = None
    if branding is not None:
        branding_config = BrandingConfig.load(branding)
        branding_source = {"type": "file", "path": str(branding)}
    elif template is not None:
        try:
            extraction = extract_branding_config(template)
        except BrandingExtractionError as exc:
            click.echo(f"ブランド設定の抽出に失敗しました: {exc}", err=True)
            branding_config, branding_source = _load_default_branding()
            branding_source["error"] = str(exc)
        else:
            branding_config = extraction.to_branding_config()
            branding_payload = extraction.to_branding_payload()
            branding_source = {"type": "template", "template": str(template)}
    else:
        branding_config, branding_source = _load_default_branding()

    artifact = {"source": branding_source}
    if branding_payload is not None:
        artifact["config"] = branding_payload
    return branding_config, artifact


def _resolve_template_path(
    *,
    spec: JobSpec,
    spec_source: Path,
) -> Path:
    """ジョブスペックとオプションからテンプレートパスを決定する。"""

    template_path_value: str | None = None
    meta = getattr(spec, "meta", None)
    if meta is not None:
        template_path_value = getattr(meta, "template_path", None)
        if template_path_value is None and isinstance(meta, BaseModel):
            extra = getattr(meta, "model_extra", None)
            if isinstance(extra, dict):
                template_path_value = extra.get("template_path")
        if template_path_value is None and isinstance(meta, dict):
            template_path_value = meta.get("template_path")

    if not template_path_value:
        try:
            raw_spec = json.loads(spec_source.read_text(encoding="utf-8"))
            template_path_value = raw_spec.get("meta", {}).get("template_path")
        except Exception:  # noqa: BLE001
            template_path_value = None

    if not template_path_value:
        raise ValueError(
            "jobspec.meta.template_path にテンプレートパスを設定してください。"
        )

    candidate_raw = Path(template_path_value)
    if candidate_raw.is_absolute():
        resolved = candidate_raw
    else:
        spec_relative = (spec_source.parent / candidate_raw).resolve()
        cwd_relative = (Path.cwd() / candidate_raw).resolve()
        if spec_relative.exists():
            resolved = spec_relative
        elif cwd_relative.exists():
            resolved = cwd_relative
        else:
            raise ValueError(
                "jobspec.meta.template_path にテンプレートパスを設定してください。"
                f"（確認したパス: {spec_relative}, {cwd_relative}）"
            )
    if not resolved.exists():
        raise ValueError(f"テンプレートファイルが見つかりません: {resolved}")
    return resolved


def _resolve_layouts_path(
    *,
    spec: JobSpec,
    spec_source: Path,
) -> Path | None:
    """ジョブスペックとオプションから layouts.jsonl のパスを決定する。"""

    layouts_path_value: str | None = None
    meta = getattr(spec, "meta", None)
    if meta is not None:
        layouts_path_value = getattr(meta, "layouts_path", None)
        if layouts_path_value is None and isinstance(meta, BaseModel):
            extra = getattr(meta, "model_extra", None)
            if isinstance(extra, dict):
                layouts_path_value = extra.get("layouts_path")
        if layouts_path_value is None and isinstance(meta, dict):
            layouts_path_value = meta.get("layouts_path")

    if layouts_path_value is None:
        try:
            raw_spec = json.loads(spec_source.read_text(encoding="utf-8"))
            layouts_path_value = raw_spec.get("meta", {}).get("layouts_path")
        except Exception:  # noqa: BLE001
            layouts_path_value = None

    if not layouts_path_value:
        return None

    candidate_raw = Path(layouts_path_value)
    candidates: list[Path]
    if candidate_raw.is_absolute():
        candidates = [candidate_raw]
    else:
        spec_relative = (spec_source.parent / candidate_raw).resolve()
        cwd_relative = (Path.cwd() / candidate_raw).resolve()
        candidates = [spec_relative, cwd_relative]

    for candidate in candidates:
        if candidate.exists():
            return candidate

    message = (
        "jobspec.meta.layouts_path に有効なパスを設定してください。"
        f"（確認したパス: {', '.join(str(path) for path in candidates)}）"
    )
    raise ValueError(message)


@dataclass(slots=True)
class TemplateExtractionResult:
    template_spec: TemplateSpec
    jobspec_scaffold: JobSpecScaffold
    template_spec_path: Path
    branding_path: Path
    jobspec_path: Path
    validation_result: LayoutValidationResult | None
    output_dir: Path
    slide_snapshot_path: Path | None


@dataclass(slots=True)
class TemplateReleaseExecutionResult:
    release: TemplateRelease
    report: TemplateReleaseReport
    release_path: Path
    report_path: Path
    golden_runs_path: Path | None
    baseline_release: Path | None


def _run_template_extraction(
    *,
    template_path: Path,
    output_dir: Path,
    layout: str | None,
    anchor: str | None,
    output_format: str,
    template_ai_policy: Path | None,
    template_ai_policy_id: str | None,
    disable_template_ai: bool,
    layout_mode: str,
    skip_validation: bool = False,
    emit_slide_snapshot: bool = False,
) -> TemplateExtractionResult:
    fmt = output_format.lower()
    extractor_options = TemplateExtractorOptions(
        template_path=template_path,
        output_path=None,
        layout_filter=layout,
        anchor_filter=anchor,
        format=fmt,
        layout_mode=layout_mode,
    )

    output_dir.mkdir(parents=True, exist_ok=True)

    ai_policy_path = template_ai_policy
    if ai_policy_path is None and not disable_template_ai:
        env_policy = os.getenv("PPTX_TEMPLATE_AI_POLICY")
        if env_policy:
            ai_policy_path = Path(env_policy)
        else:
            ai_policy_path = _discover_template_ai_policy()
    ai_policy_id = template_ai_policy_id or os.getenv("PPTX_TEMPLATE_AI_POLICY_ID")
    effective_disable = disable_template_ai or ai_policy_path is None
    if effective_disable:
        ai_policy_path = None
        if not disable_template_ai:
            logger.info("Template AI validation disabled: no policy file available")

    extractor = TemplateExtractor(extractor_options)
    template_spec = extractor.extract()
    jobspec_scaffold = extractor.build_jobspec_scaffold(template_spec)
    branding_result = extract_branding_config(template_path)

    if fmt == "yaml":
        import yaml

        spec_path = output_dir / "template_spec.yaml"
        spec_content = yaml.dump(
            template_spec.model_dump(mode="json", exclude_none=True),
            allow_unicode=True,
            default_flow_style=False,
            indent=2,
        )
    else:
        spec_path = output_dir / "template_spec.json"
        spec_content = json.dumps(
            template_spec.model_dump(mode="json", exclude_none=True),
            indent=2,
            ensure_ascii=False,
        )

    spec_path.write_text(spec_content, encoding="utf-8")
    logger.info("Saved template spec to %s", spec_path.resolve())

    branding_path = output_dir / "branding.json"
    branding_payload = branding_result.to_branding_payload()
    branding_text = json.dumps(branding_payload, ensure_ascii=False, indent=2)
    branding_path.write_text(branding_text, encoding="utf-8")
    logger.info("Saved branding payload to %s", branding_path.resolve())

    validation_result: LayoutValidationResult | None = None
    if not skip_validation:
        logger.info("Starting layout validation for %s", template_path)
        validation_options = LayoutValidationOptions(
            template_path=template_path,
            output_dir=output_dir,
            template_ai_policy_path=ai_policy_path,
            template_ai_policy_id=ai_policy_id,
            disable_template_ai=effective_disable,
        )
        validation_suite = LayoutValidationSuite(validation_options)
        validation_result = validation_suite.run()
        logger.info(
            "Layout validation finished: warnings=%d errors=%d",
            validation_result.warnings_count,
            validation_result.errors_count,
        )

        try:
            layouts_relative = str(
                validation_result.layouts_path.relative_to(output_dir)
            )
        except ValueError:
            layouts_relative = str(validation_result.layouts_path)
    else:
        validation_result = None
        layouts_relative = None

    template_spec_relative: str | None = None
    try:
        template_spec_relative = str(spec_path.relative_to(output_dir))
    except ValueError:
        template_spec_relative = str(spec_path)

    meta_update: dict[str, str | None] = {
        "template_spec_path": template_spec_relative,
    }
    if layouts_relative is not None:
        meta_update["layouts_path"] = layouts_relative

    jobspec_scaffold.meta = jobspec_scaffold.meta.model_copy(update=meta_update)

    jobspec_path = output_dir / "jobspec.json"
    extractor.save_jobspec_scaffold(jobspec_scaffold, jobspec_path)
    logger.info("Saved jobspec scaffold to %s", jobspec_path.resolve())

    slide_snapshot_path: Path | None = None
    if emit_slide_snapshot:
        slide_snapshot_path = _generate_slide_snapshot(
            template_path=template_path,
            output_dir=output_dir,
        )

    return TemplateExtractionResult(
        template_spec=template_spec,
        jobspec_scaffold=jobspec_scaffold,
        template_spec_path=spec_path,
        branding_path=branding_path,
        jobspec_path=jobspec_path,
        validation_result=validation_result,
        output_dir=output_dir,
        slide_snapshot_path=slide_snapshot_path,
    )


def _generate_slide_snapshot(*, template_path: Path, output_dir: Path) -> Path | None:
    logger.info("Generating slide snapshot for %s", template_path)
    try:
        presentation = Presentation(template_path)
    except Exception as exc:  # noqa: BLE001
        logger.warning("スライドスナップショットの生成に失敗しました: %s", exc)
        return None

    slides_payload: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides):
        snapshot = SlideSnapshot.from_slide(slide, index)
        slides_payload.append(_serialize_slide_snapshot(slide, snapshot))

    if not slides_payload:
        logger.info("スライドが存在しないためスナップショットを出力しません: %s", template_path)
        return None

    payload = {
        "template_path": str(template_path),
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "slides": slides_payload,
    }

    output_dir.mkdir(parents=True, exist_ok=True)
    path = output_dir / "slide_snapshot.json"
    path.write_text(json.dumps(payload, ensure_ascii=False,
                    indent=2), encoding="utf-8")
    logger.info("Saved slide snapshot to %s", path.resolve())
    return path


def _serialize_slide_snapshot(slide, snapshot: SlideSnapshot) -> dict[str, Any]:
    layout_name = getattr(getattr(slide, "slide_layout", None), "name", None)
    slide_identifier = getattr(slide, "slide_id", None)

    shapes_payload: list[dict[str, Any]] = []
    for shape in snapshot.shapes:
        paragraphs = [
            {
                "index": paragraph.paragraph_index,
                "text": paragraph.text,
                "level": paragraph.level,
                "font_size_pt": paragraph.font_size_pt,
                "color_hex": paragraph.color_hex,
            }
            for paragraph in shape.paragraphs
        ]
        shapes_payload.append(
            {
                "shape_id": shape.shape_id,
                "name": shape.name or "",
                "shape_type": _shape_type_name(shape.shape_type),
                "left_in": shape.left_in,
                "top_in": shape.top_in,
                "width_in": shape.width_in,
                "height_in": shape.height_in,
                "is_placeholder": shape.is_placeholder,
                "placeholder_type": _placeholder_type_name(shape.placeholder_type),
                "paragraphs": paragraphs,
            }
        )

    return {
        "index": snapshot.index,
        "slide_id": slide_identifier,
        "layout": layout_name,
        "shapes": shapes_payload,
    }


def _shape_type_name(shape_type: int | None) -> str:
    if shape_type is None:
        return "unknown"
    try:
        return MSO_SHAPE_TYPE(shape_type).name
    except ValueError:
        return str(shape_type)


def _placeholder_type_name(placeholder_type: int | None) -> str | None:
    if placeholder_type is None:
        return None
    try:
        return PP_PLACEHOLDER(placeholder_type).name
    except ValueError:
        return str(placeholder_type)


def _echo_template_extraction_result(result: TemplateExtractionResult) -> None:
    template_spec = result.template_spec
    jobspec_scaffold = result.jobspec_scaffold
    validation_result = result.validation_result

    click.echo(f"テンプレート抽出が完了しました: {result.template_spec_path}")
    click.echo(f"ブランド設定を出力しました: {result.branding_path}")
    click.echo(f"ジョブスペック雛形を出力しました: {result.jobspec_path}")
    click.echo(f"抽出されたレイアウト数: {len(template_spec.layouts)}")
    click.echo(f"Layout Mode: {template_spec.layout_mode}")
    if template_spec.blueprint:
        click.echo(f"Blueprint Slides: {len(template_spec.blueprint.slides)}")

    total_anchors = sum(len(layout.anchors)
                        for layout in template_spec.layouts)
    click.echo(f"抽出された図形・アンカー数: {total_anchors}")
    click.echo(f"ジョブスペックのスライド数: {len(jobspec_scaffold.slides)}")

    if validation_result is not None:
        click.echo(f"Layouts: {validation_result.layouts_path}")
        click.echo(f"Diagnostics: {validation_result.diagnostics_path}")
        if validation_result.diff_report_path is not None:
            click.echo(f"Diff: {validation_result.diff_report_path}")
        click.echo(
            "検出結果: warnings=%d, errors=%d"
            % (validation_result.warnings_count, validation_result.errors_count)
        )
    else:
        click.echo("検証をスキップしました (--force)")

    if result.slide_snapshot_path is not None:
        click.echo(f"スライドスナップショットを出力しました: {result.slide_snapshot_path}")

    if template_spec.warnings:
        click.echo(f"警告: {len(template_spec.warnings)} 件")
        for warning in template_spec.warnings:
            click.echo(f"  - {warning}", err=True)

    if template_spec.errors:
        click.echo(f"エラー: {len(template_spec.errors)} 件")
        for error in template_spec.errors:
            click.echo(f"  - {error}", err=True)


def _run_template_release(
    *,
    template_path: Path,
    brand: str,
    version: str,
    template_id: str | None,
    output_dir: Path,
    generated_by: str | None,
    reviewed_by: str | None,
    baseline_release: Path | None,
    golden_specs: tuple[Path, ...],
    layout_mode: str,
) -> TemplateReleaseExecutionResult:
    resolved_template_id = _resolve_template_id(template_id, brand, version)

    extractor = TemplateExtractor(
        TemplateExtractorOptions(template_path=template_path, layout_mode=layout_mode))
    spec = extractor.extract()

    output_dir.mkdir(parents=True, exist_ok=True)

    baseline = load_template_release(
        baseline_release) if baseline_release else None

    resolved_golden_specs, auto_golden_warnings = _resolve_golden_specs(
        user_specs=list(golden_specs),
        baseline=baseline,
        baseline_release=baseline_release,
    )

    golden_runs: list[TemplateReleaseGoldenRun] = []
    golden_warnings: list[str] = []
    golden_errors: list[str] = []
    if resolved_golden_specs:
        golden_runs, golden_warnings, golden_errors = _run_golden_specs(
            template_path=template_path,
            golden_specs=resolved_golden_specs,
            output_dir=output_dir,
        )

    combined_warnings = golden_warnings + auto_golden_warnings

    release = build_template_release(
        template_path=template_path,
        spec=spec,
        template_id=resolved_template_id,
        brand=brand,
        version=version,
        generated_by=generated_by,
        reviewed_by=reviewed_by,
        golden_runs=golden_runs,
        extra_warnings=combined_warnings,
        extra_errors=golden_errors,
    )
    release_path = output_dir / "template_release.json"
    release_path.write_text(
        json.dumps(release.model_dump(), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    logger.info("Saved template release to %s", release_path.resolve())

    golden_runs_path: Path | None = None
    if golden_runs:
        golden_runs_path = output_dir / "golden_runs.json"
        golden_runs_path.write_text(
            json.dumps(
                [run.model_dump() for run in golden_runs],
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )
        logger.info("Saved golden run log to %s", golden_runs_path.resolve())

    baseline_model = baseline
    report = build_release_report(current=release, baseline=baseline_model)
    report_path = output_dir / "release_report.json"
    report_path.write_text(
        json.dumps(report.model_dump(), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    logger.info("Saved release report to %s", report_path.resolve())

    return TemplateReleaseExecutionResult(
        release=release,
        report=report,
        release_path=release_path,
        report_path=report_path,
        golden_runs_path=golden_runs_path,
        baseline_release=baseline_release,
    )


def _echo_template_release_result(result: TemplateReleaseExecutionResult) -> None:
    click.echo(f"テンプレートリリースメタを出力しました: {result.release_path}")
    click.echo(f"差分レポートを出力しました: {result.report_path}")
    if result.golden_runs_path is not None:
        click.echo(f"ゴールデンサンプル結果を出力しました: {result.golden_runs_path}")
    if result.baseline_release is not None:
        click.echo(f"比較対象: {result.baseline_release}")
    _print_diagnostics(result.release.diagnostics)


def _build_reference_text(document: ContentApprovalDocument) -> tuple[str | None, bool]:
    lines: list[str] = []
    for slide in document.slides:
        if slide.elements.title:
            lines.append(str(slide.elements.title))
        lines.extend(str(item) for item in slide.elements.body)
        if slide.elements.note:
            lines.append(str(slide.elements.note))

    reference_text = "\n".join(line.strip() for line in lines if line.strip())
    if not reference_text:
        return None, False

    return reference_text, False


def _build_analyzer_options(
    rules_config: RulesConfig,
    branding_config: BrandingConfig,
    emit_structure_snapshot: bool,
) -> AnalyzerOptions:
    analyzer_rules = rules_config.analyzer
    analyzer_defaults = AnalyzerOptions()
    body_font_size = branding_config.body_font.size_pt
    body_font_color = branding_config.body_font.color_hex
    primary_color = branding_config.primary_color
    background_color = branding_config.background_color

    max_bullet_level = (
        rules_config.max_bullet_level
        if rules_config.max_bullet_level is not None
        else analyzer_defaults.max_bullet_level
    )

    return AnalyzerOptions(
        min_font_size=analyzer_rules.min_font_size
        if analyzer_rules.min_font_size is not None
        else body_font_size,
        default_font_size=analyzer_rules.default_font_size
        if analyzer_rules.default_font_size is not None
        else body_font_size,
        default_font_color=analyzer_rules.default_font_color or body_font_color,
        preferred_text_color=analyzer_rules.preferred_text_color or primary_color,
        background_color=analyzer_rules.background_color or background_color,
        min_contrast_ratio=analyzer_rules.min_contrast_ratio
        if analyzer_rules.min_contrast_ratio is not None
        else analyzer_defaults.min_contrast_ratio,
        large_text_min_contrast=analyzer_rules.large_text_min_contrast
        if analyzer_rules.large_text_min_contrast is not None
        else analyzer_defaults.large_text_min_contrast,
        large_text_threshold_pt=analyzer_rules.large_text_threshold_pt
        if analyzer_rules.large_text_threshold_pt is not None
        else body_font_size,
        margin_in=analyzer_rules.margin_in
        if analyzer_rules.margin_in is not None
        else analyzer_defaults.margin_in,
        slide_width_in=analyzer_rules.slide_width_in
        if analyzer_rules.slide_width_in is not None
        else analyzer_defaults.slide_width_in,
        slide_height_in=analyzer_rules.slide_height_in
        if analyzer_rules.slide_height_in is not None
        else analyzer_defaults.slide_height_in,
        max_bullet_level=max_bullet_level,
        snapshot_output_filename="analysis_snapshot.json"
        if emit_structure_snapshot
        else None,
    )


def _build_refiner_options(
    rules_config: RulesConfig, branding_config: BrandingConfig
) -> RefinerOptions:
    analyzer_rules = rules_config.analyzer
    refiner_rules = rules_config.refiner
    body_font_size = branding_config.body_font.size_pt
    body_font_color = branding_config.body_font.color_hex
    primary_color = branding_config.primary_color

    refiner_defaults = RefinerOptions()
    max_bullet_level = (
        rules_config.max_bullet_level
        if rules_config.max_bullet_level is not None
        else refiner_defaults.max_bullet_level
    )

    return RefinerOptions(
        max_bullet_level=max_bullet_level,
        enable_bullet_reindent=refiner_rules.enable_bullet_reindent,
        enable_font_raise=refiner_rules.enable_font_raise,
        min_font_size=refiner_rules.min_font_size
        if refiner_rules.min_font_size is not None
        else body_font_size,
        enable_color_adjust=refiner_rules.enable_color_adjust,
        preferred_text_color=refiner_rules.preferred_text_color
        or analyzer_rules.preferred_text_color
        or primary_color,
        fallback_font_color=refiner_rules.fallback_font_color or body_font_color,
        default_font_name=branding_config.body_font.name,
    )


def _build_polisher_options(
    rules_config: RulesConfig,
    *,
    polisher_toggle: bool | None,
    polisher_path: Optional[Path],
    polisher_rules: Optional[Path],
    polisher_timeout: Optional[int],
    polisher_args: tuple[str, ...],
    polisher_cwd: Optional[Path],
    rules_path: Path,
) -> PolisherOptions:
    config = rules_config.polisher
    enabled = polisher_toggle if polisher_toggle is not None else config.enabled

    executable: Path | None = polisher_path
    if executable is None and config.executable:
        executable = _resolve_config_path(
            config.executable, base_dir=rules_path.parent)

    rules_file: Path | None = polisher_rules
    if rules_file is None and config.rules_path:
        rules_file = _resolve_config_path(
            config.rules_path, base_dir=rules_path.parent)

    timeout_sec = polisher_timeout or config.timeout_sec
    arguments = tuple(config.arguments) + tuple(polisher_args)

    return PolisherOptions(
        enabled=enabled,
        executable=executable,
        rules_path=rules_file,
        timeout_sec=timeout_sec,
        arguments=arguments,
        working_dir=polisher_cwd,
    )


def _dump_json(path: Path, payload: object) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    text = json.dumps(payload, ensure_ascii=False, indent=2)
    path.write_text(text, encoding="utf-8")
    logger.info("Saved JSON to %s", path.resolve())


def _build_prepare_story_outline(document: PrepareDocument) -> dict[str, Any]:
    chapter_cards: dict[str, list[str]] = {}

    def resolve_bucket(card: PrepareCard) -> str:
        if isinstance(card.meta, dict):
            source_chapter = card.meta.get("source_chapter")
            if isinstance(source_chapter, dict):
                source_id = source_chapter.get("id")
                source_title = source_chapter.get("title")
                if isinstance(source_id, str) and source_id.strip():
                    return source_id.strip()
                if isinstance(source_title, str) and source_title.strip():
                    return source_title.strip()
        blueprint = card.blueprint_meta()
        if blueprint and blueprint.get("slide_id"):
            return str(blueprint.get("slide_id"))
        return card.role.story_phase

    for card in document.cards:
        bucket = resolve_bucket(card)
        chapter_cards.setdefault(bucket, []).append(card.card_id)

    chapters_payload: list[dict[str, Any]] = []
    for chapter in document.story_context.chapters:
        cards = chapter_cards.pop(chapter.id, [])
        if not cards:
            cards = chapter_cards.pop(chapter.title, [])
        chapters_payload.append(
            {
                "id": chapter.id,
                "title": chapter.title,
                "cards": cards,
            }
        )

    for title, cards in chapter_cards.items():
        chapters_payload.append({"id": title, "title": title, "cards": cards})

    return {
        "prepare_id": document.prepare_id,
        "chapters": chapters_payload,
        "narrative_theme": None,
        "summary": None,
    }


def _load_jobspec(path: Path) -> JobSpec:
    logger.info("Loading JobSpec from %s", path.resolve())
    return load_jobspec_from_path(path)


def _run_content_approval_pipeline(
    *,
    spec: JobSpec,
    output_dir: Path,
    content_approved: Path | None,
    content_review_log: Path | None,
    require_document: bool,
) -> PipelineContext:
    output_dir.mkdir(parents=True, exist_ok=True)
    context = PipelineContext(spec=spec, workdir=output_dir)

    step = ContentApprovalStep(
        ContentApprovalOptions(
            approved_path=content_approved,
            review_log_path=content_review_log,
            require_document=require_document,
            require_all_approved=True,
        )
    )
    PipelineRunner([step]).execute(context)
    return context


def _write_content_outputs(
    *,
    context: PipelineContext,
    output_dir: Path,
    spec_filename: str,
    content_filename: str,
    review_filename: str,
    meta_filename: str,
) -> tuple[Path, Path | None, Path | None, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)

    spec_path = output_dir / spec_filename
    _dump_json(spec_path, context.spec.model_dump(mode="json"))

    content_document = context.artifacts.get("content_approved")
    content_path: Path | None = None
    if isinstance(content_document, ContentApprovalDocument):
        content_path = output_dir / content_filename
        _dump_json(content_path, content_document.model_dump(mode="json"))

    review_logs = context.artifacts.get("content_review_log")
    review_path: Path | None = None
    if review_logs:
        review_payload: list[dict[str, object]] = []
        for entry in review_logs:
            if hasattr(entry, "model_dump"):
                review_payload.append(entry.model_dump(
                    mode="json"))  # type: ignore[call-arg]
            else:
                review_payload.append(entry)
        review_path = output_dir / review_filename
        _dump_json(review_path, review_payload)

    content_meta = context.artifacts.get("content_approved_meta")
    review_meta = context.artifacts.get("content_review_log_meta")
    meta_payload: dict[str, object] = {
        "spec": {
            "slides": len(context.spec.slides),
            "output_path": str(spec_path),
        }
    }
    if isinstance(content_meta, dict):
        meta_payload["content_approved"] = {
            **content_meta,
            "output_path": str(content_path) if content_path else None,
        }
    if isinstance(review_meta, dict):
        meta_payload["content_review_log"] = {
            **review_meta,
            "output_path": str(review_path) if review_path else None,
        }

    meta_path = output_dir / meta_filename
    _dump_json(meta_path, meta_payload)

    return spec_path, content_path, review_path, meta_path


def _run_draft_pipeline(
    *,
    spec: JobSpec,
    output_dir: Path,
    prepare_cards: Path | None,
    require_prepare: bool,
    draft_options: DraftStructuringOptions,
) -> PipelineContext:
    output_dir.mkdir(parents=True, exist_ok=True)

    if prepare_cards is not None:
        resolved_cards = prepare_cards if prepare_cards.is_absolute() else (
            Path.cwd() / prepare_cards
        )
        resolved_cards = resolved_cards.resolve()
        try:
            display_path = resolved_cards.relative_to(Path.cwd())
        except ValueError:
            display_path = resolved_cards
        logger.info("prepare_card.json を読み込みます: %s", display_path)

    steps: list[PipelineStep] = []
    steps.append(
        PrepareNormalizationStep(
            PrepareNormalizationOptions(
                cards_path=prepare_cards,
                require_document=require_prepare,
            )
        )
    )
    steps.append(DraftStructuringStep(draft_options))

    context = PipelineContext(spec=spec, workdir=output_dir)
    PipelineRunner(steps).execute(context)
    return context


def _write_draft_meta(
    *,
    context: PipelineContext,
    output_dir: Path,
    meta_filename: str,
    draft_filename: str,
    approved_filename: str,
    log_filename: str,
) -> Path:
    draft_document = context.artifacts.get("draft_document")
    sections = 0
    slides = 0
    approved_sections: list[str] = []
    section_status: dict[str, str] = {}
    appendix_limit: int | None = None
    structure_pattern: str | None = None
    target_length: int | None = None
    template_id: str | None = None
    template_match_score: float | None = None
    template_mismatch: list[dict[str, object]] = []
    analyzer_summary: dict[str, int] = {}
    return_reason_stats: dict[str, int] = {}

    if isinstance(draft_document, DraftDocument):
        sections = len(draft_document.sections)
        for section in draft_document.sections:
            section_status[section.name] = section.status
            if section.status == "approved":
                approved_sections.append(section.name)
            slides += len(section.slides)
        appendix_limit = draft_document.meta.appendix_limit
        structure_pattern = draft_document.meta.structure_pattern
        target_length = draft_document.meta.target_length
        template_id = draft_document.meta.template_id
        template_match_score = draft_document.meta.template_match_score
        template_mismatch = [item.model_dump(
            mode="json") for item in draft_document.meta.template_mismatch]
        analyzer_summary = draft_document.meta.analyzer_summary
        return_reason_stats = draft_document.meta.return_reason_stats

    paths = {
        "draft_draft": str((output_dir / draft_filename).resolve()),
        "draft_approved": str((output_dir / approved_filename).resolve()),
        "draft_review_log": str((output_dir / log_filename).resolve()),
    }

    approved_path = context.artifacts.get("draft_document_path")
    if isinstance(approved_path, str):
        paths["draft_approved"] = str(Path(approved_path).resolve())
    log_path = context.artifacts.get("draft_review_log_path")
    if isinstance(log_path, str):
        paths["draft_review_log"] = str(Path(log_path).resolve())
    ready_path = context.artifacts.get("generate_ready_path")
    if isinstance(ready_path, str):
        paths["generate_ready"] = str(Path(ready_path).resolve())
    ready_meta_path = context.artifacts.get("generate_ready_meta_path")
    if isinstance(ready_meta_path, str):
        paths["generate_ready_meta"] = str(Path(ready_meta_path).resolve())

    meta_payload = {
        "spec_id": context.artifacts.get("draft_spec_id"),
        "sections": sections,
        "slides": slides,
        "approved_sections": approved_sections,
        "section_status": section_status,
        "appendix_limit": appendix_limit,
        "structure_pattern": structure_pattern,
        "target_length": target_length,
        "paths": paths,
        "template": {
            "template_id": template_id,
            "match_score": template_match_score,
            "mismatch": template_mismatch,
        },
        "analyzer_summary": analyzer_summary,
        "return_reason_stats": return_reason_stats,
    }

    meta_path = output_dir / meta_filename
    _dump_json(meta_path, meta_payload)
    return meta_path


def _execute_outline(
    *,
    spec: JobSpec,
    layouts: Path | None,
    output_dir: Path,
    spec_source_path: Path,
    target_length: int | None,
    structure_pattern: str | None,
    appendix_limit: int,
    chapter_templates_dir: Path | None,
    chapter_template: str | None,
    analysis_summary_path: Path | None,
    prepare_cards: Path | None,
    require_prepare: bool,
) -> OutlineResult:
    draft_options = DraftStructuringOptions(
        layouts_path=layouts,
        output_dir=output_dir,
        spec_source_path=spec_source_path,
        target_length=target_length,
        structure_pattern=structure_pattern,
        appendix_limit=appendix_limit,
        chapter_templates_dir=chapter_templates_dir,
        chapter_template_id=chapter_template,
        analysis_summary_path=analysis_summary_path,
    )

    context = _run_draft_pipeline(
        spec=spec,
        output_dir=output_dir,
        prepare_cards=prepare_cards,
        require_prepare=require_prepare,
        draft_options=draft_options,
    )

    meta_path = _write_draft_meta(
        context=context,
        output_dir=output_dir,
        meta_filename=DEFAULT_DRAFT_META_FILENAME,
        draft_filename=DEFAULT_DRAFT_FILENAME,
        approved_filename=DEFAULT_APPROVED_FILENAME,
        log_filename=DEFAULT_DRAFT_LOG_FILENAME,
    )

    ready_artifact = context.artifacts.get("generate_ready_path")
    ready_meta_artifact = context.artifacts.get("generate_ready_meta_path")
    ready_path = (
        Path(ready_artifact)
        if isinstance(ready_artifact, str)
        else (output_dir / DEFAULT_GENERATE_READY_FILENAME)
    )
    ready_meta_path = (
        Path(ready_meta_artifact)
        if isinstance(ready_meta_artifact, str)
        else (output_dir / DEFAULT_GENERATE_READY_META_FILENAME)
    )

    return OutlineResult(
        context=context,
        draft_path=output_dir / DEFAULT_DRAFT_FILENAME,
        approved_path=output_dir / DEFAULT_APPROVED_FILENAME,
        log_path=output_dir / DEFAULT_DRAFT_LOG_FILENAME,
        meta_path=meta_path,
        generate_ready_path=ready_path,
        generate_ready_meta_path=ready_meta_path,
    )


def _print_outline_result(result: OutlineResult, *, show_layout_reasons: bool) -> None:
    click.echo(f"Outline Draft: {result.draft_path}")
    click.echo(f"Outline Approved: {result.approved_path}")
    click.echo(f"Outline Review Log: {result.log_path}")
    click.echo(f"Outline Meta: {result.meta_path}")
    click.echo(f"Outline Generate Ready: {result.generate_ready_path}")
    click.echo(f"Outline Ready Meta: {result.generate_ready_meta_path}")

    if not show_layout_reasons:
        return

    draft_document = result.context.artifacts.get("draft_document")
    if not isinstance(draft_document, DraftDocument):
        return

    click.echo("layout_hint 候補スコア内訳:")
    for section in draft_document.sections:
        for slide in section.slides:
            detail = slide.layout_score_detail
            if not detail:
                continue
            click.echo(
                f"- {slide.ref_id} -> {slide.layout_hint} "
                f"(uses_tag={detail.uses_tag:.2f}, "
                f"capacity={detail.content_capacity:.2f}, "
                f"diversity={detail.diversity:.2f}, "
                f"analyzer={detail.analyzer_support:.2f})"
            )



def _run_mapping_pipeline(
    *,
    spec: JobSpec,
    output_dir: Path,
    spec_source_path: Path,
    rules_config: RulesConfig,
    refiner_options: RefinerOptions,
    branding_artifact: dict[str, object],
    prepare_cards: Path | None,
    require_prepare: bool,
    layouts: Optional[Path],
    draft_output: Path,
    template: Optional[Path],
    draft_context: PipelineContext | None = None,
    draft_options: DraftStructuringOptions | None = None,
) -> PipelineContext:
    if template is None:
        msg = "jobspec.meta.template_path を設定し、テンプレートパスを埋め込んでください。"
        raise ValueError(msg)

    output_dir.mkdir(parents=True, exist_ok=True)
    draft_output.mkdir(parents=True, exist_ok=True)

    mapping_options = MappingOptions(
        layouts_path=layouts,
        output_dir=output_dir,
        template_path=template,
    )

    if draft_context is None:
        draft_context = _run_draft_pipeline(
            spec=spec,
            output_dir=draft_output,
            prepare_cards=prepare_cards,
            require_prepare=require_prepare,
            draft_options=draft_options
            or DraftStructuringOptions(
                layouts_path=layouts,
                output_dir=draft_output,
                spec_source_path=spec_source_path,
            ),
        )
    else:
        if draft_context.workdir != draft_output:
            logger.debug(
                "draft_context.workdir と draft_output が一致しません: %s != %s",
                draft_context.workdir,
                draft_output,
            )

    draft_generate_ready = draft_context.artifacts.get("generate_ready")
    if isinstance(draft_generate_ready, GenerateReadyDocument) and draft_generate_ready.meta.layout_mode == "static":
        return _pass_through_static_generate_ready(
            spec=spec,
            draft_context=draft_context,
            output_dir=output_dir,
            mapping_options=mapping_options,
        )

    context = PipelineContext(
        spec=spec, workdir=output_dir, artifacts=dict(draft_context.artifacts))
    context.add_artifact("branding", branding_artifact)

    spec_validator = SpecValidatorStep(
        max_title_length=rules_config.max_title_length,
        max_bullet_length=rules_config.max_bullet_length,
        max_bullet_level=rules_config.max_bullet_level,
        forbidden_words=rules_config.forbidden_words,
    )
    refiner = SimpleRefinerStep(refiner_options)
    mapping = MappingStep(mapping_options)

    PipelineRunner([spec_validator, refiner, mapping]).execute(context)

    meta_source = context.artifacts.get("generate_ready_meta_path")
    if isinstance(meta_source, str):
        source_path = Path(meta_source)
        destination = output_dir / DEFAULT_GENERATE_READY_META_FILENAME
        try:
            if source_path.exists():
                if destination.resolve() != source_path.resolve():
                    destination.parent.mkdir(parents=True, exist_ok=True)
                    shutil.copy2(source_path, destination)
                context.add_artifact(
                    "generate_ready_meta_path", str(destination))
        except OSError as exc:  # noqa: PERF203 - unexpected copy failure should bubble up later
            raise RuntimeError(
                f"generate_ready_meta.json のコピーに失敗しました: {exc}"
            ) from exc

    return context


def _pass_through_static_generate_ready(
    *,
    spec: JobSpec,
    draft_context: PipelineContext,
    output_dir: Path,
    mapping_options: MappingOptions,
) -> PipelineContext:
    output_dir.mkdir(parents=True, exist_ok=True)

    artifacts = dict(draft_context.artifacts)
    context = PipelineContext(
        spec=spec, workdir=output_dir, artifacts=artifacts)

    generate_ready = artifacts.get("generate_ready")
    if not isinstance(generate_ready, GenerateReadyDocument):
        raise RuntimeError("static モードの成果物に generate_ready が存在しません")

    template_path = None
    if mapping_options.template_path is not None:
        template_path = str(mapping_options.template_path)
    elif isinstance(spec.meta, JobMeta) and spec.meta.template_path:
        template_path = spec.meta.template_path
    else:
        raw_meta = getattr(spec, "meta", None)
        if isinstance(raw_meta, dict):
            template_path = raw_meta.get("template_path")
    if template_path:
        updated_job_meta = None
        if generate_ready.meta.job_meta:
            job_meta = generate_ready.meta.job_meta
            if not job_meta.template_path:
                updated_job_meta = job_meta.model_copy(
                    update={"template_path": template_path})
            else:
                updated_job_meta = job_meta
        meta_updates = {"template_path": template_path}
        if updated_job_meta is not None:
            meta_updates["job_meta"] = updated_job_meta
        generate_ready = generate_ready.model_copy(
            update={"meta": generate_ready.meta.model_copy(
                update=meta_updates)}
        )
        artifacts["generate_ready"] = generate_ready
        logger.debug(
            "static pass-through: template_path set to %s", template_path)

    ready_dest = output_dir / mapping_options.generate_ready_filename
    ready_dest.parent.mkdir(parents=True, exist_ok=True)
    ready_dest.write_text(
        json.dumps(
            generate_ready.model_dump(mode="json", exclude_none=True),
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    context.add_artifact("generate_ready", generate_ready)
    context.add_artifact("generate_ready_path", str(ready_dest))

    meta_dest = output_dir / DEFAULT_GENERATE_READY_META_FILENAME
    meta_src_value = artifacts.get("generate_ready_meta_path")
    meta_dest.parent.mkdir(parents=True, exist_ok=True)
    meta_payload: dict[str, object] | None = None
    if isinstance(meta_src_value, str):
        meta_src = Path(meta_src_value)
        if meta_src.exists():
            try:
                meta_payload = json.loads(meta_src.read_text(encoding="utf-8"))
            except json.JSONDecodeError:
                meta_payload = None
    if meta_payload is None:
        meta_payload = {
            "mode": generate_ready.meta.layout_mode,
            "slot_summary": generate_ready.meta.slot_summary or {},
            "generate_ready_generated_at": generate_ready.meta.generated_at,
            "template_path": generate_ready.meta.template_path,
            "blueprint_path": generate_ready.meta.blueprint_path,
            "blueprint_hash": generate_ready.meta.blueprint_hash,
        }
    else:
        if template_path:
            meta_payload["template_path"] = template_path
        if "mode" not in meta_payload and generate_ready.meta.layout_mode:
            meta_payload["mode"] = generate_ready.meta.layout_mode
        if generate_ready.meta.slot_summary:
            meta_payload.setdefault(
                "slot_summary", generate_ready.meta.slot_summary)
        meta_payload.setdefault(
            "generate_ready_generated_at", generate_ready.meta.generated_at)
    meta_dest.write_text(json.dumps(
        meta_payload, ensure_ascii=False, indent=2), encoding="utf-8")
    context.add_artifact("generate_ready_meta_path", str(meta_dest))

    mapping_log_dest = output_dir / mapping_options.mapping_log_filename
    draft_mapping_log_path = artifacts.get("draft_mapping_log_path")
    mapping_src = Path(draft_mapping_log_path) if isinstance(
        draft_mapping_log_path, str) else None
    if mapping_src and mapping_src.exists() and mapping_src.resolve() != mapping_log_dest.resolve():
        shutil.copy2(mapping_src, mapping_log_dest)
    elif mapping_src and mapping_src.exists():
        # same path, nothing to copy
        pass
    else:
        mapping_log_dest.write_text(
            json.dumps(
                artifacts.get("draft_mapping_log", {}),
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )
    context.add_artifact("mapping_log_path", str(mapping_log_dest))

    mapping_meta = {
        "mode": "static",
        "slot_summary": generate_ready.meta.slot_summary or {},
        "generate_ready_generated_at": generate_ready.meta.generated_at,
        "template_path": generate_ready.meta.template_path,
        "blueprint_path": generate_ready.meta.blueprint_path,
        "blueprint_hash": generate_ready.meta.blueprint_hash,
    }
    context.add_artifact("mapping_meta", mapping_meta)

    return context


def _run_render_pipeline(
    *,
    generate_ready: GenerateReadyDocument,
    generate_ready_path: Optional[Path],
    output_dir: Path,
    template: Optional[Path],
    pptx_name: str,
    branding_config: BrandingConfig,
    branding_artifact: dict[str, object],
    analyzer_options: AnalyzerOptions,
    pdf_options: PdfExportOptions,
    polisher_options: PolisherOptions | None = None,
    base_artifacts: dict[str, object] | None = None,
) -> PipelineContext:
    output_dir.mkdir(parents=True, exist_ok=True)

    render_spec = generate_ready_to_jobspec(generate_ready)
    artifacts = dict(base_artifacts or {})

    context = PipelineContext(
        spec=render_spec,
        workdir=output_dir,
        artifacts=artifacts,
    )
    context.add_artifact("branding", branding_artifact)
    context.add_artifact("generate_ready", generate_ready)
    if generate_ready_path is not None:
        context.add_artifact("generate_ready_path", str(generate_ready_path))

    renderer = SimpleRendererStep(
        RenderingOptions(
            template_path=template,
            output_filename=pptx_name,
            branding=branding_config,
        )
    )
    baseline_analyzer_options = replace(
        analyzer_options,
        output_filename="analysis_pre_polisher.json",
        snapshot_output_filename=None,
    )
    baseline_analyzer = SimpleAnalyzerStep(
        baseline_analyzer_options,
        artifact_key="analysis_pre_polisher_path",
        register_default_artifact=False,
        allow_missing_artifact=True,
    )
    analyzer = SimpleAnalyzerStep(analyzer_options)

    polisher_step = PolisherStep(polisher_options or PolisherOptions())
    audit_step = RenderingAuditStep(RenderingAuditOptions())

    monitoring_step = MonitoringIntegrationStep(MonitoringIntegrationOptions())

    steps: list[PipelineStep] = [
        renderer,
        baseline_analyzer,
        polisher_step,
        audit_step,
    ]
    if pdf_options.enabled:
        steps.append(PdfExportStep(pdf_options))
    steps.extend([analyzer, monitoring_step])

    PipelineRunner(steps).execute(context)
    return context


def _echo_mapping_outputs(context: PipelineContext) -> None:
    draft_path = context.artifacts.get("draft_document_path")
    if draft_path is not None:
        click.echo(f"Draft: {draft_path}")
    draft_log_path = context.artifacts.get("draft_review_log_path")
    if draft_log_path is not None:
        click.echo(f"Draft Log: {draft_log_path}")
    generate_ready_path = context.artifacts.get("generate_ready_path")
    if generate_ready_path is not None:
        click.echo(f"Generate Ready: {generate_ready_path}")
    generate_ready_meta_path = context.artifacts.get(
        "generate_ready_meta_path")
    if generate_ready_meta_path is not None:
        click.echo(f"Generate Ready Meta: {generate_ready_meta_path}")
    mapping_log_path = context.artifacts.get("mapping_log_path")
    if mapping_log_path is not None:
        click.echo(f"Mapping Log: {mapping_log_path}")
    fallback_report_path = context.artifacts.get(
        "mapping_fallback_report_path")
    if fallback_report_path is not None:
        click.echo(f"Fallback Report: {fallback_report_path}")


def _echo_render_outputs(context: PipelineContext, audit_path: Path | None) -> None:
    pptx_path = context.artifacts.get("pptx_path")
    if pptx_path is not None:
        click.echo(f"PPTX: {pptx_path}")
    else:
        click.echo("PPTX: --pdf-mode=only のため保存しませんでした")
    analysis_path = context.artifacts.get("analysis_path")
    click.echo(f"Analysis: {analysis_path}")
    baseline_analysis_path = context.artifacts.get(
        "analysis_pre_polisher_path")
    if baseline_analysis_path is not None:
        click.echo(f"Analysis (Pre-Polisher): {baseline_analysis_path}")
    rendering_log_path = context.artifacts.get("rendering_log_path")
    if rendering_log_path is not None:
        click.echo(f"Rendering Log: {rendering_log_path}")
    rendering_summary = context.artifacts.get("rendering_summary")
    if isinstance(rendering_summary, dict):
        click.echo(
            "Rendering Warnings: %s" % rendering_summary.get(
                "warnings_total", 0)
        )
    review_engine_path = context.artifacts.get("review_engine_analysis_path")
    if review_engine_path is not None:
        click.echo(f"ReviewEngine Analysis: {review_engine_path}")
    snapshot_path = context.artifacts.get("analyzer_snapshot_path")
    if snapshot_path is not None:
        click.echo(f"Analyzer Snapshot: {snapshot_path}")
    pdf_path = context.artifacts.get("pdf_path")
    if pdf_path is not None:
        click.echo(f"PDF: {pdf_path}")
    polisher_meta = context.artifacts.get("polisher_metadata")
    if isinstance(polisher_meta, dict):
        status = polisher_meta.get("status", "unknown")
        click.echo(f"Polisher: {status}")
        summary = polisher_meta.get("summary")
        if isinstance(summary, dict) and summary:
            click.echo(
                "Polisher Summary: "
                + json.dumps(summary, ensure_ascii=False, sort_keys=True)
            )
    draft_path = context.artifacts.get("draft_document_path")
    if draft_path is not None:
        click.echo(f"Draft: {draft_path}")
    draft_log_path = context.artifacts.get("draft_review_log_path")
    if draft_log_path is not None:
        click.echo(f"Draft Log: {draft_log_path}")
    generate_ready_path = context.artifacts.get("generate_ready_path")
    if generate_ready_path is not None:
        click.echo(f"Generate Ready: {generate_ready_path}")
    mapping_log_path = context.artifacts.get("mapping_log_path")
    if mapping_log_path is not None:
        click.echo(f"Mapping Log: {mapping_log_path}")
    fallback_report_path = context.artifacts.get(
        "mapping_fallback_report_path")
    if fallback_report_path is not None:
        click.echo(f"Fallback Report: {fallback_report_path}")
    monitoring_report_path = context.artifacts.get("monitoring_report_path")
    if monitoring_report_path is not None:
        click.echo(f"Monitoring Report: {monitoring_report_path}")
    if audit_path is not None:
        click.echo(f"Audit: {audit_path}")


@app.command("gen")
@click.argument(
    "generate_ready_path",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
)
@click.option(
    "--output",
    "-o",
    "output_dir",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=Path(".pptx/gen"),
    show_default=True,
    help="生成物を保存するディレクトリ",
)
@click.option(
    "--pptx-name",
    default="proposal.pptx",
    show_default=True,
    help="出力 PPTX のファイル名",
)
@click.option(
    "--rules",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=DEFAULT_RULES_PATH,
    show_default=True,
    help="検証ルール設定ファイル",
)
@click.option(
    "--branding",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    show_default=str(DEFAULT_BRANDING_PATH),
    help="ブランド設定ファイル",
)
@click.option(
    "--export-pdf",
    is_flag=True,
    help="LibreOffice を利用して PDF を追加出力する",
)
@click.option(
    "--pdf-mode",
    type=click.Choice(["both", "only"], case_sensitive=False),
    default="both",
    show_default=True,
    help="PDF 出力時の挙動。only では PPTX を保存しない",
)
@click.option(
    "--pdf-output",
    type=str,
    default="proposal.pdf",
    show_default=True,
    help="出力 PDF ファイル名",
)
@click.option(
    "--libreoffice-path",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    help="LibreOffice (soffice) 実行ファイルのパス",
)
@click.option(
    "--pdf-timeout",
    type=int,
    default=120,
    show_default=True,
    help="LibreOffice 変換のタイムアウト秒",
)
@click.option(
    "--pdf-retries",
    type=int,
    default=2,
    show_default=True,
    help="LibreOffice 変換の最大リトライ回数",
)
@click.option(
    "--polisher/--no-polisher",
    "polisher_toggle",
    default=None,
    help="Open XML Polisher を実行するかを明示的に指定する（設定ファイル値を上書き）",
)
@click.option(
    "--polisher-path",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    help="Open XML Polisher (.exe / .dll) もしくはラッパースクリプトのパス",
)
@click.option(
    "--polisher-rules",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    help="Polisher へ渡すルール設定ファイル",
)
@click.option(
    "--polisher-timeout",
    type=int,
    default=None,
    help="Polisher 実行のタイムアウト秒（設定ファイル値を上書き）",
)
@click.option(
    "--polisher-arg",
    "polisher_args",
    multiple=True,
    help="Polisher へ渡す追加引数（{pptx}, {rules} をプレースホルダーとして利用可能）",
)
@click.option(
    "--polisher-cwd",
    type=click.Path(exists=True, file_okay=False,
                    dir_okay=True, path_type=Path),
    default=None,
    help="Polisher 実行時のカレントディレクトリ",
)
@click.option(
    "--emit-structure-snapshot",
    is_flag=True,
    help="Analyzer の構造スナップショット (analysis_snapshot.json) を出力する",
)
def gen(  # noqa: PLR0913
    generate_ready_path: Path,
    output_dir: Path,
    pptx_name: str,
    rules: Path,
    branding: Optional[Path],
    export_pdf: bool,
    pdf_mode: str,
    pdf_output: str,
    libreoffice_path: Optional[Path],
    pdf_timeout: int,
    pdf_retries: int,
    polisher_toggle: bool | None,
    polisher_path: Optional[Path],
    polisher_rules: Optional[Path],
    polisher_timeout: Optional[int],
    polisher_args: tuple[str, ...],
    polisher_cwd: Optional[Path],
    emit_structure_snapshot: bool,
) -> None:
    """generate_ready.json から PPTX / PDF / 監査ログを生成する。"""

    if not export_pdf and pdf_mode != "both":
        click.echo("--pdf-mode は --export-pdf と併用してください", err=True)
        raise click.exceptions.Exit(code=2)

    try:
        generate_ready = GenerateReadyDocument.parse_file(generate_ready_path)
    except Exception as exc:  # noqa: BLE001
        click.echo(f"generate_ready.json の読み込みに失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc

    template_path_str = generate_ready.meta.template_path
    if not template_path_str:
        click.echo(
            "generate_ready.json に template_path が含まれていません。stage 4 を最新仕様で再実行するか、テンプレート情報を埋め込んでください。",
            err=True,
        )
        raise click.exceptions.Exit(code=2)

    template_path = Path(template_path_str)
    if not template_path.is_absolute():
        candidate = (generate_ready_path.parent / template_path).resolve()
        template_path = candidate if candidate.exists() else template_path
    if not template_path.exists():
        click.echo(f"テンプレートファイルが見つかりません: {template_path}", err=True)
        raise click.exceptions.Exit(code=4)

    rules_config = RulesConfig.load(rules)
    branding_config, branding_artifact = _prepare_branding(
        template_path, branding)
    analyzer_options = _build_analyzer_options(
        rules_config, branding_config, emit_structure_snapshot
    )
    pdf_options = PdfExportOptions(
        enabled=export_pdf,
        mode=pdf_mode,
        output_filename=pdf_output,
        soffice_path=libreoffice_path,
        timeout_sec=pdf_timeout,
        max_retries=pdf_retries,
    )
    polisher_options = _build_polisher_options(
        rules_config,
        polisher_toggle=polisher_toggle,
        polisher_path=polisher_path,
        polisher_rules=polisher_rules,
        polisher_timeout=polisher_timeout,
        polisher_args=polisher_args,
        polisher_cwd=polisher_cwd,
        rules_path=rules,
    )

    mapping_meta: dict[str, object] = {
        "generate_ready_path": str(generate_ready_path),
        "generate_ready_generated_at": generate_ready.meta.generated_at,
        "template_version": generate_ready.meta.template_version,
        "template_path": str(template_path),
    }

    base_artifacts: dict[str, object] = {
        "generate_ready": generate_ready,
        "generate_ready_path": str(generate_ready_path),
        "mapping_meta": mapping_meta,
    }

    mapping_log_path = generate_ready_path.with_name("mapping_log.json")
    if mapping_log_path.exists():
        base_artifacts["mapping_log_path"] = str(mapping_log_path)
        try:
            mapping_log = json.loads(
                mapping_log_path.read_text(encoding="utf-8"))
        except Exception as exc:  # noqa: BLE001
            logger.warning("mapping_log.json の読み込みに失敗しました: %s", exc)
        else:
            meta_payload = mapping_log.get("meta")
            if isinstance(meta_payload, dict):
                mapping_meta.update(meta_payload)
    fallback_path = generate_ready_path.with_name("fallback_report.json")
    if fallback_path.exists():
        base_artifacts["mapping_fallback_report_path"] = str(fallback_path)

    try:
        render_context = _run_render_pipeline(
            generate_ready=generate_ready,
            generate_ready_path=generate_ready_path,
            output_dir=output_dir,
            template=template_path,
            pptx_name=pptx_name,
            branding_config=branding_config,
            branding_artifact=branding_artifact,
            analyzer_options=analyzer_options,
            pdf_options=pdf_options,
            polisher_options=polisher_options,
            base_artifacts=base_artifacts,
        )
    except PdfExportError as exc:
        click.echo(f"PDF 出力に失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=5) from exc
    except PolisherError as exc:
        click.echo(f"Polisher の実行に失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=6) from exc
    except FileNotFoundError as exc:
        click.echo(f"ファイルが見つかりません: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except Exception as exc:  # noqa: BLE001
        logging.exception("パイプライン実行中にエラーが発生しました")
        raise click.exceptions.Exit(code=1) from exc

    analysis_path = render_context.artifacts.get("analysis_path")
    _emit_review_engine_analysis(render_context, analysis_path)
    audit_path = _write_audit_log(render_context)
    _echo_render_outputs(render_context, audit_path)


@app.command("prepare")
@click.argument(
    "prepare_path",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
)
@click.option(
    "--output",
    "-o",
    "output_dir",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=DEFAULT_PREPARE_OUTPUT_DIR,
    show_default=True,
    help="コンテンツ準備成果物を保存するディレクトリ",
)
@click.option(
    "--mode",
    type=click.Choice(["dynamic", "static"], case_sensitive=False),
    required=True,
    help="カード生成モード。static は Blueprint を利用する",
)
@click.option(
    "--jobspec",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    help="静的モードで参照する jobspec.json (未指定時は .pptx/extract/jobspec.json を探索)",
)
@click.option(
    "-p",
    "--page-limit",
    "--card-limit",
    type=click.IntRange(1, None),
    default=None,
    help="生成するカード枚数の上限",
)
def prepare(
    prepare_path: Path,
    output_dir: Path,
    jobspec: Path | None,
    mode: str,
    page_limit: int | None,
) -> None:
    """stage 2 コンテンツ準備: PrepareCard 成果物を生成する。"""

    try:
        source = PrepareSourceDocument.parse_file(prepare_path)
    except FileNotFoundError as exc:
        click.echo(f"プレペア入力ファイルが見つかりません: {exc}", err=True)
        raise click.exceptions.Exit(code=2) from exc
    except (json.JSONDecodeError, ValidationError) as exc:
        click.echo(f"プレペア入力の解析に失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=2) from exc

    policy_path = DEFAULT_PREPARE_POLICY_PATH
    try:
        policy_set = load_prepare_policy_set(policy_path)
    except PreparePolicyError as exc:
        click.echo(f"プレペアポリシーの読み込みに失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc

    blueprint_spec: TemplateSpec | None = None
    blueprint_ref: dict[str, str] | None = None
    template_spec_path: Path | None = None
    jobspec_path: Path | None = jobspec
    normalized_mode = mode.lower()
    if normalized_mode not in {"dynamic", "static"}:
        raise click.exceptions.BadParameter(
            "--mode には dynamic か static を指定してください")

    if normalized_mode == "static":
        if page_limit is not None:
            click.echo("static モードでは --page-limit を利用できません", err=True)
            raise click.exceptions.Exit(code=2)
        resolved_jobspec = jobspec_path or DEFAULT_JOBSPEC_PATH
        if not resolved_jobspec.exists():
            click.echo(
                "static モードでは --jobspec で jobspec.json のパスを指定するか、.pptx/extract/jobspec.json を用意してください",
                err=True,
            )
            raise click.exceptions.Exit(code=2)
        try:
            spec_for_static = load_jobspec_from_path(resolved_jobspec)
        except (FileNotFoundError, ValidationError, SpecValidationError) as exc:
            click.echo(f"jobspec.json の読み込みに失敗しました: {exc}", err=True)
            raise click.exceptions.Exit(code=2) from exc

        template_spec_ref = getattr(
            spec_for_static.meta, "template_spec_path", None)
        if not template_spec_ref:
            click.echo(
                "jobspec.meta.template_spec_path が設定されていません。テンプレ抽出を再実行してください",
                err=True,
            )
            raise click.exceptions.Exit(code=2)

        template_spec_path = Path(template_spec_ref)
        if not template_spec_path.is_absolute():
            template_spec_path = (
                resolved_jobspec.parent / template_spec_path).resolve()

        if not template_spec_path.exists():
            click.echo(
                f"template_spec.json が見つかりません: {template_spec_path}",
                err=True,
            )
            raise click.exceptions.Exit(code=2)

        try:
            template_spec_text = template_spec_path.read_text(encoding="utf-8")
        except FileNotFoundError as exc:
            click.echo(f"template_spec の読み込みに失敗しました: {exc}", err=True)
            raise click.exceptions.Exit(code=2) from exc
        try:
            if template_spec_path.suffix.lower() in {".yaml", ".yml"}:
                import yaml

                payload = yaml.safe_load(template_spec_text)
                blueprint_spec = TemplateSpec.model_validate(payload)
            else:
                blueprint_spec = TemplateSpec.model_validate_json(
                    template_spec_text)
        except ValueError as exc:
            click.echo(f"template_spec の検証に失敗しました: {exc}", err=True)
            raise click.exceptions.Exit(code=2) from exc
        if blueprint_spec.layout_mode != "static":
            click.echo("template_spec の layout_mode が static ではありません", err=True)
            raise click.exceptions.Exit(code=2)
        if blueprint_spec.blueprint is None:
            click.echo("template_spec に blueprint が含まれていません", err=True)
            raise click.exceptions.Exit(code=2)

        blueprint_hash = hashlib.sha256(
            json.dumps(
                blueprint_spec.blueprint.model_dump(mode="json"),
                ensure_ascii=False,
                sort_keys=True,
            ).encode("utf-8")
        ).hexdigest()
        blueprint_ref = {
            "path": str(template_spec_path),
            "hash": f"sha256:{blueprint_hash}",
        }

    orchestrator = PrepareAIOrchestrator(policy_set)
    try:
        document, meta, ai_logs = orchestrator.generate_document(
            source,
            policy_id=None,
            page_limit=page_limit,
            mode=normalized_mode,  # type: ignore[arg-type]
            blueprint=blueprint_spec.blueprint if blueprint_spec else None,
            blueprint_ref=blueprint_ref,
        )
    except PrepareAIOrchestrationError as exc:
        click.echo(f"プレペアカードの生成に失敗しました: {exc}", err=True)
        exit_code = 6 if normalized_mode == "static" else 4
        raise click.exceptions.Exit(code=exit_code) from exc

    output_dir.mkdir(parents=True, exist_ok=True)
    cards_path = output_dir / "prepare_card.json"
    log_path = output_dir / "prepare_log.json"
    ai_log_path = output_dir / "prepare_ai_log.json"
    meta_path = output_dir / "ai_generation_meta.json"
    story_outline_path = output_dir / "prepare_story_outline.json"
    audit_path = output_dir / "audit_log.json"

    def _relativize(path: Path) -> str:
        try:
            return str(path.relative_to(output_dir))
        except ValueError:
            return str(path)

    document.meta = dict(document.meta or {})
    document.meta.update(
        {
            "prepare_card_path": _relativize(cards_path),
            "prepare_log_path": _relativize(log_path),
            "prepare_ai_log_path": _relativize(ai_log_path),
            "ai_generation_meta_path": _relativize(meta_path),
            "prepare_story_outline_path": _relativize(story_outline_path),
            "prepare_audit_log_path": _relativize(audit_path),
        }
    )

    _dump_json(cards_path, document.model_dump(mode="json", exclude_none=True))
    _dump_json(log_path, [])
    _dump_json(
        ai_log_path,
        [record.model_dump(mode="json", exclude_none=True)
         for record in ai_logs],
    )
    _dump_json(meta_path, meta.model_dump(mode="json", exclude_none=True))
    _dump_json(story_outline_path, _build_prepare_story_outline(document))

    audit_payload = {
        "prepare_normalization": {
            "generated_at": meta.generated_at.isoformat(),
            "policy_id": meta.policy_id,
            "input_hash": meta.input_hash,
            "mode": meta.mode,
            "outputs": {
                "prepare_card": str(cards_path.resolve()),
                "prepare_log": str(log_path.resolve()),
                "prepare_ai_log": str(ai_log_path.resolve()),
                "ai_generation_meta": str(meta_path.resolve()),
                "prepare_story_outline": str(story_outline_path.resolve()),
            },
            "statistics": meta.statistics,
        }
    }
    if template_spec_path is not None:
        audit_payload["prepare_normalization"]["outputs"]["template_spec"] = str(
            template_spec_path)
    if blueprint_ref:
        audit_payload["prepare_normalization"]["blueprint"] = blueprint_ref
    if meta.slot_coverage:
        audit_payload["prepare_normalization"]["slot_summary"] = meta.slot_coverage
    _dump_json(audit_path, audit_payload)

    click.echo(f"Prepare Card: {cards_path}")
    click.echo(f"Prepare Log: {log_path}")
    click.echo(f"Prepare AI Log: {ai_log_path}")
    click.echo(f"AI Generation Meta: {meta_path}")
    click.echo(f"Prepare Story Outline: {story_outline_path}")
    click.echo(f"Audit Log: {audit_path}")


@app.command("outline")
@click.argument(
    "spec_path",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
)
@click.option(
    "--output",
    "-o",
    "output_dir",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=Path(".pptx/draft"),
    show_default=True,
    help="ドラフト成果物を保存するディレクトリ",
)
@click.option(
    "--target-length",
    type=int,
    default=None,
    help="目標スライド枚数",
)
@click.option(
    "--structure-pattern",
    type=str,
    default=None,
    help="章構成パターン名",
)
@click.option(
    "--appendix-limit",
    type=int,
    default=5,
    show_default=True,
    help="付録枚数の上限",
)
@click.option(
    "--chapter-templates-dir",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=DEFAULT_CHAPTER_TEMPLATES_DIR,
    show_default=True,
    help="章テンプレート辞書ディレクトリ",
)
@click.option(
    "--chapter-template",
    type=str,
    default=None,
    help="適用する章テンプレート ID",
)
@click.option(
    "--import-analysis",
    "analysis_summary_path",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    help="analysis_summary.json のパス",
)
@click.option(
    "--return-reasons-path",
    type=click.Path(dir_okay=False, writable=False, path_type=Path),
    default=DEFAULT_RETURN_REASONS_PATH,
    show_default=True,
    help="差戻し理由テンプレート辞書のパス",
)
@click.option(
    "--return-reasons",
    is_flag=True,
    default=False,
    help="差戻し理由テンプレート一覧を表示して終了する",
)
@click.option(
    "--show-layout-reasons",
    is_flag=True,
    default=False,
    help="layout_hint 候補のスコア内訳を表示する",
)
@click.option(
    "--prepare-cards",
    type=click.Path(exists=False, dir_okay=False,
                    readable=True, path_type=Path),
    default=DEFAULT_PREPARE_OUTPUT_DIR / "prepare_card.json",
    show_default=True,
    help="stage 2 の prepare_card.json",
)
def outline(
    spec_path: Path,
    output_dir: Path,
    target_length: int | None,
    structure_pattern: str | None,
    appendix_limit: int,
    chapter_templates_dir: Path,
    chapter_template: str | None,
    analysis_summary_path: Path | None,
    return_reasons_path: Path,
    return_reasons: bool,
    show_layout_reasons: bool,
    prepare_cards: Path,
) -> None:
    """stage 4 ドラフト構成（アウトライン）を生成する。"""

    if return_reasons:
        reasons = load_return_reasons(return_reasons_path)
        if not reasons:
            click.echo(f"差戻し理由テンプレートが見つかりません: {return_reasons_path}")
        else:
            click.echo("差戻し理由テンプレート一覧:")
            for reason in reasons:
                label = f"{reason.code} ({reason.severity})"
                if reason.label and reason.label != reason.code:
                    label += f" - {reason.label}"
                click.echo(f"- {label}")
        return

    try:
        spec = _load_jobspec(spec_path)
    except SpecValidationError as exc:
        _echo_errors("スキーマ検証に失敗しました", exc.errors)
        raise click.exceptions.Exit(code=2) from exc

    templates_dir = chapter_templates_dir if chapter_templates_dir.exists() else None

    try:
        resolved_layouts = _resolve_layouts_path(
            spec=spec,
            spec_source=spec_path,
        )
    except ValueError as exc:
        click.echo(str(exc), err=True)
        raise click.exceptions.Exit(code=2) from exc

    try:
        result = _execute_outline(
            spec=spec,
            layouts=resolved_layouts,
            output_dir=output_dir,
            spec_source_path=spec_path,
            target_length=target_length,
            structure_pattern=structure_pattern,
            appendix_limit=appendix_limit,
            chapter_templates_dir=templates_dir,
            chapter_template=chapter_template,
            analysis_summary_path=analysis_summary_path,
            prepare_cards=prepare_cards,
            require_prepare=True,
        )
    except PrepareNormalizationError as exc:
        click.echo(f"プレペア成果物の読み込みに失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except DraftStructuringError as exc:
        click.echo(f"ドラフト構成の生成に失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except FileNotFoundError as exc:
        click.echo(f"ファイルが見つかりません: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except Exception as exc:  # noqa: BLE001
        logging.exception("outline 実行中にエラーが発生しました")
        raise click.exceptions.Exit(code=1) from exc
    _print_outline_result(result, show_layout_reasons=show_layout_reasons)


@app.command("compose")
@click.argument(
    "spec_path",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
)
@click.option(
    "--draft-output",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=Path(".pptx/draft"),
    show_default=True,
    help="ドラフト成果物を保存するディレクトリ",
)
@click.option(
    "--target-length",
    type=int,
    default=None,
    help="目標スライド枚数",
)
@click.option(
    "--structure-pattern",
    type=str,
    default=None,
    help="章構成パターン名",
)
@click.option(
    "--appendix-limit",
    type=int,
    default=5,
    show_default=True,
    help="付録枚数の上限",
)
@click.option(
    "--chapter-templates-dir",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=DEFAULT_CHAPTER_TEMPLATES_DIR,
    show_default=True,
    help="章テンプレート辞書ディレクトリ",
)
@click.option(
    "--chapter-template",
    type=str,
    default=None,
    help="適用する章テンプレート ID",
)
@click.option(
    "--import-analysis",
    "analysis_summary_path",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    help="analysis_summary.json のパス",
)
@click.option(
    "--show-layout-reasons",
    is_flag=True,
    default=False,
    help="layout_hint 候補のスコア内訳を表示する",
)
@click.option(
    "--output",
    "-o",
    "output_dir",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=Path(".pptx/compose"),
    show_default=True,
    help="generate_ready.json 等の出力ディレクトリ",
)
@click.option(
    "--rules",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=DEFAULT_RULES_PATH,
    show_default=True,
    help="検証ルール設定ファイル",
)
@click.option(
    "--branding",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    show_default=str(DEFAULT_BRANDING_PATH),
    help="ブランド設定ファイル（任意）",
)
@click.option(
    "--prepare-cards",
    "prepare_cards",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=DEFAULT_PREPARE_OUTPUT_DIR / "prepare_card.json",
    show_default=True,
    help="stage 2 の prepare_card.json",
)
def compose(  # noqa: PLR0913
    spec_path: Path,
    draft_output: Path,
    target_length: int | None,
    structure_pattern: str | None,
    appendix_limit: int,
    chapter_templates_dir: Path,
    chapter_template: str | None,
    analysis_summary_path: Path | None,
    show_layout_reasons: bool,
    output_dir: Path,
    rules: Path,
    branding: Optional[Path],
    prepare_cards: Path,
) -> None:
    """stage 4+5 を連続実行しドラフトとマッピング成果物を生成する。"""

    try:
        spec = _load_jobspec(spec_path)
    except SpecValidationError as exc:
        _echo_errors("スキーマ検証に失敗しました", exc.errors)
        raise click.exceptions.Exit(code=2) from exc

    try:
        resolved_template = _resolve_template_path(
            spec=spec,
            spec_source=spec_path,
        )
    except ValueError as exc:
        click.echo(str(exc), err=True)
        raise click.exceptions.Exit(code=2) from exc
    try:
        resolved_layouts = _resolve_layouts_path(
            spec=spec,
            spec_source=spec_path,
        )
    except ValueError as exc:
        click.echo(str(exc), err=True)
        raise click.exceptions.Exit(code=2) from exc

    templates_dir = chapter_templates_dir if chapter_templates_dir.exists() else None

    try:
        outline_result = _execute_outline(
            spec=spec,
            layouts=resolved_layouts,
            output_dir=draft_output,
            spec_source_path=spec_path,
            target_length=target_length,
            structure_pattern=structure_pattern,
            appendix_limit=appendix_limit,
            chapter_templates_dir=templates_dir,
            chapter_template=chapter_template,
            analysis_summary_path=analysis_summary_path,
            prepare_cards=prepare_cards,
            require_prepare=True,
        )
    except PrepareNormalizationError as exc:
        click.echo(f"プレペア成果物の読み込みに失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except DraftStructuringError as exc:
        click.echo(f"ドラフト構成の生成に失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except FileNotFoundError as exc:
        click.echo(f"ファイルが見つかりません: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except Exception as exc:  # noqa: BLE001
        logging.exception("compose 実行中にアウトライン stage でエラーが発生しました")
        raise click.exceptions.Exit(code=1) from exc

    _print_outline_result(outline_result, show_layout_reasons=show_layout_reasons)

    rules_config = RulesConfig.load(rules)
    branding_config, branding_artifact = _prepare_branding(
        resolved_template, branding
    )
    refiner_options = _build_refiner_options(rules_config, branding_config)

    try:
        mapping_context = _run_mapping_pipeline(
            spec=spec,
            output_dir=output_dir,
            spec_source_path=spec_path,
            rules_config=rules_config,
            refiner_options=refiner_options,
            branding_artifact=branding_artifact,
            prepare_cards=prepare_cards,
            require_prepare=True,
            layouts=resolved_layouts,
            draft_output=draft_output,
            template=resolved_template,
            draft_context=outline_result.context,
            draft_options=DraftStructuringOptions(
                layouts_path=resolved_layouts,
                output_dir=draft_output,
                spec_source_path=spec_path,
                target_length=target_length,
                structure_pattern=structure_pattern,
                appendix_limit=appendix_limit,
                chapter_templates_dir=chapter_templates_dir,
                chapter_template_id=chapter_template,
                analysis_summary_path=analysis_summary_path,
            ),
        )
    except ValueError as exc:
        click.echo(str(exc), err=True)
        raise click.exceptions.Exit(code=2) from exc
    except SpecValidationError as exc:
        _echo_errors("業務ルール検証に失敗しました", exc.errors)
        raise click.exceptions.Exit(code=3) from exc
    except PrepareNormalizationError as exc:
        click.echo(f"プレペア成果物の読み込みに失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except Exception as exc:  # noqa: BLE001
        logging.exception("compose 実行中にマッピング stage でエラーが発生しました")
        raise click.exceptions.Exit(code=1) from exc

    _echo_mapping_outputs(mapping_context)


@app.command("mapping")
@click.argument(
    "spec_path",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
)
@click.option(
    "--output",
    "-o",
    "output_dir",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=Path(".pptx/gen"),
    show_default=True,
    help="generate_ready.json 等の出力ディレクトリ",
)
@click.option(
    "--rules",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=DEFAULT_RULES_PATH,
    show_default=True,
    help="検証ルール設定ファイル",
)
@click.option(
    "--draft-output",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=Path(".pptx/draft"),
    show_default=True,
    help="draft_draft.json / draft_approved.json の出力先",
)
@click.option(
    "--branding",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    show_default=str(DEFAULT_BRANDING_PATH),
    help="ブランド設定ファイル（任意）",
)
@click.option(
    "--prepare-cards",
    "prepare_cards",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=DEFAULT_PREPARE_OUTPUT_DIR / "prepare_card.json",
    show_default=True,
    help="stage 2 の prepare_card.json",
)
def mapping(  # noqa: PLR0913
    spec_path: Path,
    output_dir: Path,
    rules: Path,
    draft_output: Path,
    branding: Optional[Path],
    prepare_cards: Path,
) -> None:
    """stage 5 マッピングを実行し generate_ready.json を生成する。"""
    try:
        spec = _load_jobspec(spec_path)
    except SpecValidationError as exc:
        _echo_errors("スキーマ検証に失敗しました", exc.errors)
        raise click.exceptions.Exit(code=2) from exc

    try:
        resolved_template = _resolve_template_path(
            spec=spec,
            spec_source=spec_path,
        )
    except ValueError as exc:
        click.echo(str(exc), err=True)
        raise click.exceptions.Exit(code=2) from exc

    try:
        resolved_layouts = _resolve_layouts_path(
            spec=spec,
            spec_source=spec_path,
        )
    except ValueError as exc:
        click.echo(str(exc), err=True)
        raise click.exceptions.Exit(code=2) from exc

    rules_config = RulesConfig.load(rules)
    branding_config, branding_artifact = _prepare_branding(
        resolved_template, branding
    )
    refiner_options = _build_refiner_options(rules_config, branding_config)

    try:
        context = _run_mapping_pipeline(
            spec=spec,
            output_dir=output_dir,
            spec_source_path=spec_path,
            rules_config=rules_config,
            refiner_options=refiner_options,
            branding_artifact=branding_artifact,
            prepare_cards=prepare_cards,
            require_prepare=True,
            layouts=resolved_layouts,
            draft_output=draft_output,
            template=resolved_template,
        )
    except ValueError as exc:
        click.echo(str(exc), err=True)
        raise click.exceptions.Exit(code=2) from exc
    except SpecValidationError as exc:
        _echo_errors("業務ルール検証に失敗しました", exc.errors)
        raise click.exceptions.Exit(code=3) from exc
    except PrepareNormalizationError as exc:
        click.echo(f"プレペア成果物の読み込みに失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except Exception as exc:  # noqa: BLE001
        logging.exception("マッピング実行中にエラーが発生しました")
        raise click.exceptions.Exit(code=1) from exc

    _echo_mapping_outputs(context)


@app.command("template")
@click.argument(
    "template_path",
    type=click.Path(dir_okay=False, readable=True, path_type=Path),
)
@click.option(
    "--output",
    "-o",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=Path(".pptx/extract"),
    show_default=True,
    help="抽出・検証結果を保存するディレクトリ",
)
@click.option(
    "--format",
    type=click.Choice(["json", "yaml"], case_sensitive=False),
    default="json",
    show_default=True,
    help="テンプレート仕様の出力形式",
)
@click.option(
    "--layout",
    type=str,
    default=None,
    help="抽出対象レイアウト名のフィルタ（前方一致）",
)
@click.option(
    "--anchor",
    type=str,
    default=None,
    help="抽出対象アンカー名のフィルタ（前方一致）",
)
@click.option(
    "--layout-mode",
    type=click.Choice(["dynamic", "static"], case_sensitive=False),
    default="dynamic",
    show_default=True,
    help="テンプレートの想定運用モード。static を指定すると Blueprint を出力する",
)
@click.option(
    "--with-release",
    is_flag=True,
    help="抽出・検証後にテンプレートリリースメタも生成する",
)
@click.option(
    "--brand",
    type=str,
    default=None,
    help="--with-release 時のブランド名",
)
@click.option(
    "--version",
    type=str,
    default=None,
    help="--with-release 時のテンプレートバージョン",
)
@click.option(
    "--template-id",
    type=str,
    default=None,
    help="--with-release 時のテンプレート識別子。未指定時は <brand>_<version> を使用",
)
@click.option(
    "--release-output",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=Path(".pptx/release"),
    show_default=True,
    help="テンプレートリリース成果物の出力ディレクトリ",
)
@click.option(
    "--generated-by",
    type=str,
    default=None,
    help="テンプレートリリースメタの生成者",
)
@click.option(
    "--reviewed-by",
    type=str,
    default=None,
    help="テンプレートリリースメタのレビュー担当者",
)
@click.option(
    "--baseline-release",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    help="比較対象となる過去の template_release.json",
)
@click.option(
    "--golden-spec",
    "golden_specs",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    multiple=True,
    help="テンプレ互換性検証に使用する spec ファイル（複数指定可）",
)
@click.option(
    "--template-ai-policy",
    type=click.Path(dir_okay=False, readable=True, path_type=Path),
    default=None,
    help="テンプレート usage_tags 推定に使用する AI ポリシー JSON",
)
@click.option(
    "--template-ai-policy-id",
    type=str,
    default=None,
    help="テンプレート AI ポリシーセット内の利用対象 ID",
)
@click.option(
    "--disable-template-ai",
    is_flag=True,
    default=False,
    help="生成AIによる usage_tags 推定を無効化する",
)
@click.option(
    "--slide",
    is_flag=True,
    default=False,
    help="実スライドの図形・段落情報を slide_snapshot.json として出力する",
)
@click.option(
    "--force",
    "-f",
    is_flag=True,
    default=False,
    help="レイアウト検証をスキップして強制的にテンプレ stage を継続する（緊急時のみ使用）",
)
def template(  # noqa: PLR0913
    template_path: Path,
    output: Path,
    format: str,
    layout: str | None,
    anchor: str | None,
    layout_mode: str,
    with_release: bool,
    brand: str | None,
    version: str | None,
    template_id: str | None,
    release_output: Path,
    generated_by: str | None,
    reviewed_by: str | None,
    baseline_release: Path | None,
    golden_specs: tuple[Path, ...],
    template_ai_policy: Path | None,
    template_ai_policy_id: str | None,
    disable_template_ai: bool,
    slide: bool,
    force: bool,
) -> None:
    """テンプレ stage（抽出・検証・必要に応じてリリース）を実行する。"""
    _log_current_llm_provider("template")
    try:
        extraction_result = _run_template_extraction(
            template_path=template_path,
            output_dir=output,
            layout=layout,
            anchor=anchor,
            output_format=format,
            template_ai_policy=template_ai_policy,
            template_ai_policy_id=template_ai_policy_id,
            disable_template_ai=disable_template_ai,
            layout_mode=layout_mode,
            skip_validation=force,
            emit_slide_snapshot=slide,
        )
    except FileNotFoundError as exc:
        click.echo(f"ファイルが見つかりません: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except LayoutValidationError as exc:
        click.echo(f"レイアウト検証に失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=6) from exc
    except Exception as exc:  # noqa: BLE001
        logging.exception("テンプレート抽出中にエラーが発生しました")
        click.echo(f"テンプレート抽出に失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=1) from exc

    _echo_template_extraction_result(extraction_result)

    validation_result = extraction_result.validation_result
    if validation_result is not None and validation_result.errors_count > 0:
        click.echo(
            "レイアウト検証でエラーが検出されました。Diagnostics を確認してください。",
            err=True,
        )
        raise click.exceptions.Exit(code=6)
    if validation_result is None and not force:
        click.echo(
            "レイアウト検証を実施できませんでした。--force を使用しない場合は出力を確認してください。",
            err=True,
        )
        raise click.exceptions.Exit(code=6)

    if extraction_result.template_spec.errors:
        click.echo(
            "テンプレート仕様にエラーが含まれています。出力ファイルを確認してください。",
            err=True,
        )
        raise click.exceptions.Exit(code=6)

    if not with_release:
        click.echo("テンプレ stage（抽出＋検証）が完了しました。")
        return

    if brand is None or version is None:
        raise click.UsageError(
            "--with-release を使用する場合は --brand と --version を指定してください。")

    try:
        release_result = _run_template_release(
            template_path=template_path,
            brand=brand,
            version=version,
            template_id=template_id,
            output_dir=release_output,
            generated_by=generated_by,
            reviewed_by=reviewed_by,
            baseline_release=baseline_release,
            golden_specs=golden_specs,
            layout_mode=layout_mode,
        )
    except FileNotFoundError as exc:
        click.echo(f"ファイルが見つかりません: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except click.exceptions.Exit:
        raise
    except Exception as exc:  # noqa: BLE001
        logging.exception("テンプレートリリース生成中にエラーが発生しました")
        click.echo(f"テンプレートリリースの生成に失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=1) from exc

    _echo_template_release_result(release_result)
    if release_result.release.diagnostics.errors:
        raise click.exceptions.Exit(code=6)

    click.echo("テンプレ stage（抽出＋検証＋リリース）が完了しました。")


@app.command("tpl-extract")
@click.option(
    "--template",
    "-t",
    "template_path",
    type=click.Path(dir_okay=False, readable=True, path_type=Path),
    required=True,
    help="抽出対象の PPTX テンプレートファイル",
)
@click.option(
    "--layout",
    type=str,
    default=None,
    help="抽出対象レイアウト名のフィルタ（前方一致）",
)
@click.option(
    "--anchor",
    type=str,
    default=None,
    help="抽出対象アンカー名のフィルタ（前方一致）",
)
@click.option(
    "--format",
    type=click.Choice(["json", "yaml"], case_sensitive=False),
    default="json",
    show_default=True,
    help="出力形式",
)
@click.option(
    "--layout-mode",
    type=click.Choice(["dynamic", "static"], case_sensitive=False),
    default="dynamic",
    show_default=True,
    help="テンプレートの想定運用モード。static を指定すると Blueprint を出力する",
)
@click.option(
    "--output",
    "-o",
    "output_dir",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=Path(".pptx/extract"),
    show_default=True,
    help="テンプレート仕様とブランド設定を保存するディレクトリ",
)
@click.option(
    "--template-ai-policy",
    type=click.Path(dir_okay=False, readable=True, path_type=Path),
    default=None,
    help="テンプレート usage_tags 推定に使用する AI ポリシー JSON",
)
@click.option(
    "--template-ai-policy-id",
    type=str,
    default=None,
    help="テンプレート AI ポリシーセット内の利用対象 ID",
)
@click.option(
    "--disable-template-ai",
    is_flag=True,
    default=False,
    help="生成AIによる usage_tags 推定を無効化する",
)
def tpl_extract(
    template_path: Path,
    output_dir: Path,
    layout: Optional[str],
    anchor: Optional[str],
    format: str,
    layout_mode: str,
    template_ai_policy: Path | None,
    template_ai_policy_id: str | None,
    disable_template_ai: bool,
) -> None:
    """テンプレートファイルから図形・プレースホルダー情報を抽出してJSON仕様の雛形を生成する。"""
    try:
        extraction_result = _run_template_extraction(
            template_path=template_path,
            output_dir=output_dir,
            layout=layout,
            anchor=anchor,
            output_format=format,
            template_ai_policy=template_ai_policy,
            template_ai_policy_id=template_ai_policy_id,
            disable_template_ai=disable_template_ai,
            layout_mode=layout_mode,
        )
    except FileNotFoundError as exc:
        click.echo(f"ファイルが見つかりません: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except LayoutValidationError as exc:
        click.echo(f"レイアウト検証に失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=6) from exc
    except Exception as exc:  # noqa: BLE001
        if isinstance(exc, click.exceptions.Exit):
            raise
        logging.exception("テンプレート抽出中にエラーが発生しました")
        click.echo(f"テンプレート抽出に失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=1) from exc
    else:
        _echo_template_extraction_result(extraction_result)
        validation_result = extraction_result.validation_result
        if validation_result is not None and validation_result.errors_count > 0:
            click.echo(
                "レイアウト検証でエラーが検出されました。Diagnostics を確認してください。",
                err=True,
            )
            raise click.exceptions.Exit(code=6)
        if extraction_result.template_spec.errors:
            click.echo(
                "テンプレート仕様にエラーが含まれています。出力ファイルを確認してください。",
                err=True,
            )
            raise click.exceptions.Exit(code=6)


@app.command("layout-validate")
@click.option(
    "--template",
    "-t",
    "template_path",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    required=True,
    help="検証対象の PPTX テンプレートファイル",
)
@click.option(
    "--output",
    "-o",
    "output_dir",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=Path(".pptx/validation"),
    show_default=True,
    help="検証成果物の出力ディレクトリ",
)
@click.option(
    "--template-id",
    type=str,
    default=None,
    help="layouts.jsonl に記録するテンプレート ID",
)
@click.option(
    "--baseline",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    help="比較対象となる過去の layouts.jsonl",
)
@click.option(
    "--analyzer-snapshot",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    help="Analyzer が出力した構造スナップショット JSON",
)
def layout_validate(
    template_path: Path,
    output_dir: Path,
    template_id: Optional[str],
    baseline: Optional[Path],
    analyzer_snapshot: Optional[Path],
) -> None:
    """テンプレート構造の検証スイートを実行する。"""

    options = LayoutValidationOptions(
        template_path=template_path,
        output_dir=output_dir,
        template_id=template_id,
        baseline_path=baseline,
        analyzer_snapshot_path=analyzer_snapshot,
    )
    suite = LayoutValidationSuite(options)

    try:
        result = suite.run()
    except LayoutValidationError as exc:
        click.echo(f"レイアウト検証に失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=6) from exc

    click.echo(f"Layouts: {result.layouts_path}")
    click.echo(f"Diagnostics: {result.diagnostics_path}")
    if result.diff_report_path is not None:
        click.echo(f"Diff: {result.diff_report_path}")
    click.echo(
        "検出結果: warnings=%d, errors=%d" % (
            result.warnings_count, result.errors_count)
    )


@app.command("tpl-release")
@click.option(
    "--template",
    "-t",
    "template_path",
    type=click.Path(dir_okay=False, readable=True, path_type=Path),
    required=True,
    help="リリース対象の PPTX テンプレートファイル",
)
@click.option("--brand", type=str, required=True, help="ブランド名")
@click.option("--version", type=str, required=True, help="テンプレートバージョン")
@click.option(
    "--template-id",
    type=str,
    default=None,
    help="テンプレート識別子（未指定時は brand_version を使用）",
)
@click.option(
    "--output",
    "-o",
    "output_dir",
    type=click.Path(file_okay=False, dir_okay=True, path_type=Path),
    default=Path(".pptx/release"),
    show_default=True,
    help="リリース成果物を保存するディレクトリ",
)
@click.option(
    "--generated-by",
    type=str,
    default=None,
    help="リリースメタの生成者",
)
@click.option(
    "--reviewed-by",
    type=str,
    default=None,
    help="レビュー担当者",
)
@click.option(
    "--baseline-release",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    default=None,
    help="比較対象となる過去の template_release.json",
)
@click.option(
    "--golden-spec",
    "golden_specs",
    type=click.Path(exists=True, dir_okay=False,
                    readable=True, path_type=Path),
    multiple=True,
    help="テンプレ互換性検証に使用する spec ファイル（複数指定可）",
)
@click.option(
    "--layout-mode",
    type=click.Choice(["dynamic", "static"], case_sensitive=False),
    default="dynamic",
    show_default=True,
    help="テンプレートの想定運用モード。static を指定すると Blueprint を出力する",
)
def tpl_release(
    template_path: Path,
    brand: str,
    version: str,
    template_id: Optional[str],
    output_dir: Path,
    generated_by: Optional[str],
    reviewed_by: Optional[str],
    baseline_release: Optional[Path],
    golden_specs: tuple[Path, ...],
    layout_mode: str,
) -> None:
    """テンプレート受け渡しメタと差分レポートを生成する。"""

    try:
        result = _run_template_release(
            template_path=template_path,
            brand=brand,
            version=version,
            template_id=template_id,
            output_dir=output_dir,
            generated_by=generated_by,
            reviewed_by=reviewed_by,
            baseline_release=baseline_release,
            golden_specs=golden_specs,
            layout_mode=layout_mode,
        )
    except click.exceptions.Exit:
        raise
    except FileNotFoundError as exc:
        click.echo(f"ファイルが見つかりません: {exc}", err=True)
        raise click.exceptions.Exit(code=4) from exc
    except Exception as exc:  # noqa: BLE001
        logging.exception("テンプレートリリース生成中にエラーが発生しました")
        click.echo(f"テンプレートリリースの生成に失敗しました: {exc}", err=True)
        raise click.exceptions.Exit(code=1) from exc
    else:
        _echo_template_release_result(result)
        if result.release.diagnostics.errors:
            raise click.exceptions.Exit(code=6)


def _run_golden_specs(
    *, template_path: Path, golden_specs: list[Path], output_dir: Path
) -> tuple[list[TemplateReleaseGoldenRun], list[str], list[str]]:
    results: list[TemplateReleaseGoldenRun] = []
    warnings: list[str] = []
    errors: list[str] = []

    if not golden_specs:
        return results, warnings, errors

    rules_config = RulesConfig.load(DEFAULT_RULES_PATH)
    branding_config = _load_branding_for_template(template_path, warnings)

    golden_root = output_dir / "golden_runs"

    for spec_path in golden_specs:
        run_dir = golden_root / spec_path.stem
        result = TemplateReleaseGoldenRun(
            spec_path=str(spec_path),
            status="passed",
            output_dir=str(run_dir),
        )

        try:
            spec = _load_jobspec(spec_path)
        except SpecValidationError as exc:
            detail = json.dumps(exc.errors, ensure_ascii=False)
            message = (
                f"golden spec {spec_path} のスキーマ検証に失敗しました"
            )
            result.status = "failed"
            result.errors.extend([message, detail])
            errors.append(message)
            results.append(result)
            continue
        except Exception as exc:  # noqa: BLE001
            message = f"golden spec {spec_path} の読み込みに失敗しました: {exc}"
            result.status = "failed"
            result.errors.append(message)
            errors.append(message)
            results.append(result)
            continue

        run_dir.mkdir(parents=True, exist_ok=True)
        context = PipelineContext(spec=spec, workdir=run_dir)

        renderer = SimpleRendererStep(
            RenderingOptions(
                template_path=template_path,
                output_filename=f"{spec_path.stem}.pptx",
                branding=branding_config,
            )
        )
        refiner_bullet_level = (
            rules_config.max_bullet_level
            if rules_config.max_bullet_level is not None
            else RefinerOptions().max_bullet_level
        )
        refiner = SimpleRefinerStep(
            RefinerOptions(
                max_bullet_level=refiner_bullet_level,
            )
        )
        analyzer = SimpleAnalyzerStep(
            AnalyzerOptions(
                min_font_size=branding_config.body_font.size_pt,
                default_font_size=branding_config.body_font.size_pt,
                default_font_color=branding_config.body_font.color_hex,
                preferred_text_color=branding_config.primary_color,
                background_color=branding_config.background_color,
                max_bullet_level=(
                    rules_config.max_bullet_level
                    if rules_config.max_bullet_level is not None
                    else AnalyzerOptions().max_bullet_level
                ),
                large_text_threshold_pt=branding_config.body_font.size_pt,
            )
        )

        steps = [
            SpecValidatorStep(
                max_title_length=rules_config.max_title_length,
                max_bullet_length=rules_config.max_bullet_length,
                max_bullet_level=rules_config.max_bullet_level,
                forbidden_words=rules_config.forbidden_words,
            ),
            refiner,
            renderer,
            analyzer,
        ]
        runner = PipelineRunner(steps)

        try:
            runner.execute(context)
        except SpecValidationError as exc:
            detail = json.dumps(exc.errors, ensure_ascii=False)
            message = (
                f"golden spec {spec_path} の業務ルール検証に失敗しました"
            )
            result.status = "failed"
            result.errors.extend([message, detail])
            errors.append(message)
        except Exception as exc:  # noqa: BLE001
            logging.exception(
                "ゴールデンサンプル実行中にエラーが発生しました: %s", spec_path
            )
            message = f"golden spec {spec_path} の実行に失敗しました: {exc}"
            result.status = "failed"
            result.errors.append(message)
            errors.append(message)
        else:
            pptx_path = context.artifacts.get("pptx_path")
            if pptx_path is not None:
                result.pptx_path = str(pptx_path)
            analysis_path = context.artifacts.get("analysis_path")
            if analysis_path is not None:
                result.analysis_path = str(analysis_path)
            pdf_path = context.artifacts.get("pdf_path")
            if pdf_path is not None:
                result.pdf_path = str(pdf_path)

            analyzer_warnings = context.artifacts.get("analyzer_warnings")
            if isinstance(analyzer_warnings, list):
                new_warnings = [str(item) for item in analyzer_warnings]
                result.warnings.extend(new_warnings)
                for warning in new_warnings:
                    warnings.append(f"golden spec {spec_path}: {warning}")

        results.append(result)

    return results, warnings, errors


def _resolve_golden_specs(
    *,
    user_specs: list[Path],
    baseline: TemplateRelease | None,
    baseline_release: Path | None,
) -> tuple[list[Path], list[str]]:
    resolved: list[Path] = []
    warnings: list[str] = []
    seen: set[Path] = set()

    def _add_spec(path: Path) -> None:
        try:
            normalized = path.resolve()
        except OSError:
            normalized = path
        if normalized in seen:
            return
        resolved.append(path)
        seen.add(normalized)

    for spec in user_specs:
        _add_spec(spec)

    if baseline is None:
        return resolved, warnings

    if user_specs:
        return resolved, warnings

    base_dir = baseline_release.parent if baseline_release is not None else Path.cwd()
    for run in baseline.golden_runs:
        candidate = _resolve_golden_spec_path(run.spec_path, base_dir)
        if candidate is None:
            warnings.append(
                f"baseline のゴールデンスペックを解決できませんでした: {run.spec_path}"
            )
            continue
        _add_spec(candidate)

    return resolved, warnings


def _resolve_golden_spec_path(spec_path: str, base_dir: Path) -> Path | None:
    candidate = Path(spec_path)
    if candidate.is_absolute() and candidate.exists():
        return candidate

    cwd_candidate = Path.cwd() / candidate
    if cwd_candidate.exists():
        return cwd_candidate

    fallback = base_dir / candidate
    if fallback.exists():
        return fallback

    return None


def _load_branding_for_template(
    template_path: Path, warnings: list[str]
) -> BrandingConfig:
    try:
        extraction = extract_branding_config(template_path)
    except BrandingExtractionError as exc:
        warnings.append(
            f"テンプレートからブランド設定を抽出できなかったためデフォルト設定を使用します: {exc}"
        )
        try:
            return BrandingConfig.load(DEFAULT_BRANDING_PATH)
        except FileNotFoundError:
            warnings.append(
                f"デフォルトのブランド設定が見つからないため内蔵設定を使用します: {DEFAULT_BRANDING_PATH}"
            )
            return BrandingConfig.default()
    else:
        return extraction.to_branding_config()


def _resolve_template_id(
    template_id: Optional[str], brand: str, version: str
) -> str:
    if template_id and template_id.strip():
        return template_id.strip()
    base = f"{brand}_{version}"
    return base.replace(" ", "_")


def _print_diagnostics(diagnostics: TemplateReleaseDiagnostics) -> None:
    if diagnostics.warnings:
        click.echo(f"警告: {len(diagnostics.warnings)} 件", err=True)
        for warning in diagnostics.warnings:
            click.echo(f"  - {warning}", err=True)
    if diagnostics.errors:
        click.echo(f"エラー: {len(diagnostics.errors)} 件", err=True)
        for error in diagnostics.errors:
            click.echo(f"  - {error}", err=True)


def _echo_errors(message: str, errors: list[dict[str, object]] | None) -> None:
    click.echo(message, err=True)
    if not errors:
        return
    formatted = json.dumps(errors, ensure_ascii=False, indent=2)
    click.echo(formatted, err=True)


def _emit_review_engine_analysis(
    context: PipelineContext, analysis_path: object | None
) -> Path | None:
    if analysis_path is None:
        return None

    path = Path(str(analysis_path))
    if not path.exists():
        logger.warning(
            "Review Engine 連携ファイル生成のため analysis.json が見つかりません: %s", path
        )
        return None

    adapter = AnalyzerReviewEngineAdapter()
    try:
        logger.info("Loading analysis payload from %s", path.resolve())
        analysis_payload = json.loads(path.read_text(encoding="utf-8"))
    except Exception as exc:  # noqa: BLE001
        logger.warning(
            "analysis.json の読み込みに失敗したため Review Engine 連携ファイルを生成しません: %s",
            exc,
        )
        return None

    try:
        payload = adapter.build_payload(analysis_payload, context.spec)
    except Exception as exc:  # noqa: BLE001
        logger.warning(
            "Review Engine 連携ペイロードの生成に失敗しました: %s",
            exc,
        )
        return None

    output_path = path.with_name("review_engine_analyzer.json")
    output_path.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    logger.info("Saved review engine payload to %s", output_path.resolve())
    context.add_artifact("review_engine_analysis_path", output_path)
    return output_path


def _write_audit_log(context: PipelineContext) -> Path:
    outputs_dir = context.workdir
    outputs_dir.mkdir(parents=True, exist_ok=True)
    artifacts_payload = {
        "pptx": _artifact_str(context.artifacts.get("pptx_path")),
        "analysis": _artifact_str(context.artifacts.get("analysis_path")),
        "analysis_pre_polisher": _artifact_str(
            context.artifacts.get("analysis_pre_polisher_path")
        ),
        "review_engine_analysis": _artifact_str(
            context.artifacts.get("review_engine_analysis_path")
        ),
        "pdf": _artifact_str(context.artifacts.get("pdf_path")),
        "generate_ready": _artifact_str(
            context.artifacts.get("generate_ready_path")
        ),
        "rendering_log": _artifact_str(
            context.artifacts.get("rendering_log_path")
        ),
        "mapping_log": _artifact_str(context.artifacts.get("mapping_log_path")),
        "mapping_fallback_report": _artifact_str(
            context.artifacts.get("mapping_fallback_report_path")
        ),
        "monitoring_report": _artifact_str(
            context.artifacts.get("monitoring_report_path")
        ),
    }

    pdf_meta = context.artifacts.get("pdf_export_metadata")
    if isinstance(pdf_meta, dict):
        pdf_payload = {
            "enabled": True,
            "status": pdf_meta.get("status", "success"),
            "attempts": pdf_meta.get("attempts", 0),
            "elapsed_ms": int(pdf_meta.get("elapsed_sec", 0.0) * 1000),
            "converter": pdf_meta.get("converter"),
        }
    else:
        pdf_payload = None

    polisher_meta = context.artifacts.get("polisher_metadata")
    if isinstance(polisher_meta, dict):
        polisher_payload = {
            "enabled": bool(polisher_meta.get("enabled")),
            "status": polisher_meta.get("status"),
            "elapsed_ms": int(polisher_meta.get("elapsed_sec", 0.0) * 1000)
            if polisher_meta.get("elapsed_sec") is not None
            else None,
            "rules_path": polisher_meta.get("rules_path"),
            "summary": polisher_meta.get("summary"),
        }
    else:
        polisher_payload = None

    hashes: dict[str, str] = {}
    for label, key in (
        ("generate_ready", "generate_ready_path"),
        ("pptx", "pptx_path"),
        ("analysis", "analysis_path"),
        ("analysis_pre_polisher", "analysis_pre_polisher_path"),
        ("pdf", "pdf_path"),
        ("rendering_log", "rendering_log_path"),
        ("monitoring_report", "monitoring_report_path"),
        ("mapping_log", "mapping_log_path"),
        ("mapping_fallback_report", "mapping_fallback_report_path"),
    ):
        digest = _sha256_of(context.artifacts.get(key))
        if digest:
            hashes[label] = digest

    audit_payload = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "spec_meta": context.spec.meta.model_dump(),
        "slides": len(context.spec.slides),
        "artifacts": artifacts_payload,
        "rendering": context.artifacts.get("rendering_summary"),
        "pdf_export": pdf_payload,
        "refiner_adjustments": context.artifacts.get("refiner_adjustments"),
        "branding": context.artifacts.get("branding"),
        "polisher": polisher_payload,
    }
    monitoring_summary = context.artifacts.get("monitoring_summary")
    if monitoring_summary is not None:
        audit_payload["monitoring"] = monitoring_summary
    if hashes:
        audit_payload["hashes"] = hashes
    content_meta = context.artifacts.get("content_approved_meta")
    if content_meta is not None:
        audit_payload["content_approval"] = content_meta
    review_meta = context.artifacts.get("content_review_log_meta")
    if review_meta is not None:
        audit_payload["content_review_log"] = review_meta
    mapping_meta = context.artifacts.get("mapping_meta")
    if mapping_meta is not None:
        audit_payload["mapping"] = mapping_meta
    audit_path = outputs_dir / "audit_log.json"
    audit_path.write_text(json.dumps(
        audit_payload, ensure_ascii=False, indent=2), encoding="utf-8")
    logger.info("Saved audit log to %s", audit_path.resolve())
    context.add_artifact("audit_path", audit_path)
    return audit_path


def _artifact_str(value: object | None) -> str | None:
    if value is None:
        return None
    return str(value)


def _sha256_of(value: object | None) -> str | None:
    if value is None:
        return None
    path = Path(str(value))
    if not path.exists():
        return None
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(8192), b""):
            if not chunk:
                break
            digest.update(chunk)
    return f"sha256:{digest.hexdigest()}"


if __name__ == "__main__":
    app()
