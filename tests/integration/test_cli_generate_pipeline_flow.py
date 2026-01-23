"""CLI の統合テスト (generate_ready フロー確認用)."""

from __future__ import annotations

import logging
import json
import os
import shutil
import subprocess
import sys
from pathlib import Path

import click
import pytest
from click.testing import CliRunner
from pptx import Presentation

from pptx_generator import cli
from pptx_generator.template import BrandingExtractionError
from pptx_generator.cli import DEFAULT_GENERATE_READY_META_FILENAME, app
from pptx_generator.layout_validation import LayoutValidationSuite
from pptx_generator.models import (JobAuth, JobMeta, JobSpec, Slide, TemplateStyle,
                                   PipelineFallbackError)
from pptx_generator.pipeline import pdf_exporter
from pptx_generator.pipeline.base import PipelineContext
from pptx_generator.cli_handlers.mapping import TemplateStylePayload
from pptx_generator.cli_handlers.outline import OutlineResult

SAMPLE_TEMPLATE = Path("samples/templates/templates.pptx")
PREPARE_SOURCE = Path("samples/input/pitch.md")


def _libreoffice_available() -> bool:
    env_path = os.environ.get("LIBREOFFICE_PATH")
    if env_path and Path(env_path).exists():
        return True
    return shutil.which("soffice") is not None


def _collect_paragraph_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                texts.append(text)
    return texts


def _create_template_with_slide(path: Path) -> None:
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Snapshot Title"
    body_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if body_placeholder is not None and getattr(body_placeholder, "text_frame", None):
        body_placeholder.text = "Snapshot Body"
    presentation.save(path)


def _prepare_inputs(runner: CliRunner, temp_dir: Path) -> dict[str, Path]:
    prepare_dir = temp_dir / "prepare"
    result = runner.invoke(
        app,
        [
            "prepare",
            str(PREPARE_SOURCE),
            "--mode",
            "dynamic",
            "--output",
            str(prepare_dir),
        ],
        catch_exceptions=False,
    )
    assert result.exit_code == 0, result.output
    return {
        "dir": prepare_dir,
        "cards": prepare_dir / "prepare_card.json",
        "log": prepare_dir / "prepare_log.json",
        "meta": prepare_dir / "ai_generation_meta.json",
    }


def _prepare_args(paths: dict[str, Path]) -> list[str]:
    return [
        "--prepare-cards",
        str(paths["cards"]),
    ]


def _create_matching_jobspec(root: Path, prepare_paths: dict[str, Path], *, filename: str = "matching_jobspec.json") -> Path:
    base_spec = JobSpec.parse_file(Path("samples/json/sample_jobspec.json"))
    cards_payload = json.loads(
        prepare_paths["cards"].read_text(encoding="utf-8"))
    cards = cards_payload.get("cards", [])

    slides: list[Slide] = []
    for index, card in enumerate(cards, start=1):
        card_id = card.get("card_id") or f"card-{index:03d}"
        content = card.get("content") or {}
        title = (content.get("title") or content.get(
            "headline") or card_id)[:25]
        notes = None
        note_entries = content.get("notes") or []
        if note_entries:
            note = note_entries[0]
            if isinstance(note, dict):
                notes = note.get("text")
            elif isinstance(note, str):
                notes = note
        slides.append(
            Slide(
                id=card_id,
                layout="Content" if index > 1 else "Title",
                title=title,
                notes=notes,
            )
        )

    jobspec = JobSpec(
        meta=base_spec.meta.model_copy(deep=True),
        auth=base_spec.auth.model_copy(deep=True),
        slides=slides,
    )

    spec_path = root / filename
    spec_path.parent.mkdir(parents=True, exist_ok=True)
    payload = jobspec.model_dump(mode="json")
    payload.setdefault("meta", {})
    payload["meta"]["template_path"] = "templates/templates.pptx"
    payload["meta"]["layouts_path"] = "layouts.jsonl"
    spec_path.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    template_dst = spec_path.parent / "templates" / "templates.pptx"
    template_dst.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy(Path("samples/templates/templates.pptx"), template_dst)

    layouts_path = spec_path.parent / "layouts.jsonl"
    layouts_path.write_text(
        json.dumps(
            {
                "layout_id": "Title",
                "usage_tags": ["intro"],
                "text_hint": {"max_lines": 5},
                "media_hint": {"allow_table": False},
            }
        )
        + "\n",
        encoding="utf-8",
    )
    layouts_path.write_text(
        layouts_path.read_text(encoding="utf-8")
        + json.dumps(
            {
                "layout_id": "Content",
                "usage_tags": ["body"],
                "text_hint": {"max_lines": 8},
                "media_hint": {"allow_table": True},
            }
        )
        + "\n",
        encoding="utf-8",
    )
    return spec_path


def test_cli_template_basic(tmp_path: Path) -> None:
    runner = CliRunner()
    extract_dir = tmp_path / "extract"

    result = runner.invoke(
        app,
        [
            "template",
            str(SAMPLE_TEMPLATE),
            "--output",
            str(extract_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0, result.output
    assert (extract_dir / "template_spec.json").exists()
    assert (extract_dir / "branding.json").exists()
    assert (extract_dir / "jobspec.json").exists()
    assert (extract_dir / "layouts.jsonl").exists()
    assert (extract_dir / "diagnostics.json").exists()
    assert "テンプレ stage（抽出＋検証）が完了しました。" in result.output


def test_compose_logs_outline_stage_error(monkeypatch: pytest.MonkeyPatch, tmp_path: Path, caplog: pytest.LogCaptureFixture) -> None:
    spec_path = tmp_path / "jobspec.json"
    spec_path.write_text("{}", encoding="utf-8")
    (tmp_path / "pipeline_rules.json").write_text("{}", encoding="utf-8")

    monkeypatch.setattr(
        "pptx_generator.cli_handlers.compose.load_jobspec",
        lambda path: object(),
    )
    monkeypatch.setattr(
        "pptx_generator.cli_handlers.compose.resolve_template_path",
        lambda spec, spec_source: tmp_path / "template.pptx",
    )
    monkeypatch.setattr(
        "pptx_generator.cli_handlers.compose.resolve_layouts_path",
        lambda spec, spec_source: tmp_path / "layouts.jsonl",
    )

    def boom(**_: object) -> None:
        raise RuntimeError("outline failure")

    monkeypatch.setattr("pptx_generator.cli_handlers.compose.execute_outline", boom)

    with caplog.at_level(logging.ERROR):
        with pytest.raises(click.exceptions.Exit) as exc:
            cli.compose.callback(
                spec_path=spec_path,
                target_length=None,
                structure_pattern=None,
                appendix_limit=5,
                analysis_summary_path=None,
                show_layout_reasons=False,
                output_dir=tmp_path / "compose",
                rules=tmp_path / "pipeline_rules.json",
                prepare_cards=tmp_path / "prepare_card.json",
            )

    assert exc.value.exit_code == 1
    assert any(
        "アウトライン stage" in record.getMessage() for record in caplog.records
    )


def test_compose_logs_mapping_stage_error(monkeypatch: pytest.MonkeyPatch, tmp_path: Path, caplog: pytest.LogCaptureFixture) -> None:
    spec_path = tmp_path / "jobspec.json"
    spec_path.write_text("{}", encoding="utf-8")
    rules_path = tmp_path / "pipeline_rules.json"
    rules_path.write_text("{}", encoding="utf-8")

    monkeypatch.setattr(
        "pptx_generator.cli_handlers.compose.load_jobspec",
        lambda path: object(),
    )
    monkeypatch.setattr(
        "pptx_generator.cli_handlers.compose.resolve_template_path",
        lambda spec, spec_source: tmp_path / "template.pptx",
    )
    monkeypatch.setattr(
        "pptx_generator.cli_handlers.compose.resolve_layouts_path",
        lambda spec, spec_source: tmp_path / "layouts.jsonl",
    )

    outline_result = OutlineResult(
        context=PipelineContext(spec=object(), workdir=tmp_path),
        draft_path=tmp_path / "draft.json",
        approved_path=tmp_path / "approved.json",
        log_path=tmp_path / "log.json",
        meta_path=tmp_path / "meta.json",
        generate_ready_path=tmp_path / "ready.json",
        generate_ready_meta_path=tmp_path / "ready_meta.json",
    )
    monkeypatch.setattr(
        "pptx_generator.cli_handlers.compose.execute_outline",
        lambda **_: outline_result,
    )
    monkeypatch.setattr(
        "pptx_generator.cli_handlers.compose.print_outline_result",
        lambda *args, **kwargs: None,
    )

    dummy_rules = object()

    def fake_rules_load(cls: object, path: Path) -> object:
        return dummy_rules

    monkeypatch.setattr(
        cli.RulesConfig,
        "load",
        classmethod(fake_rules_load),
    )
    monkeypatch.setattr(
        "pptx_generator.cli_handlers.compose.prepare_template_style",
        lambda template: TemplateStylePayload(style=TemplateStyle.default(), artifact={}),
    )
    monkeypatch.setattr(
        "pptx_generator.cli_handlers.compose.build_refiner_options",
        lambda rules_config, template_style: object(),
    )
    monkeypatch.setattr(
        "pptx_generator.cli_handlers.compose.run_mapping_pipeline",
        lambda **_: (_ for _ in ()).throw(RuntimeError("mapping failure")),
    )

    with caplog.at_level(logging.ERROR):
        with pytest.raises(click.exceptions.Exit) as exc:
            cli.compose.callback(
                spec_path=spec_path,
                target_length=None,
                structure_pattern=None,
                appendix_limit=5,
                analysis_summary_path=None,
                show_layout_reasons=False,
                output_dir=tmp_path / "compose",
                rules=rules_path,
                prepare_cards=tmp_path / "prepare_card.json",
            )

    assert exc.value.exit_code == 1
    assert any(
        "マッピング stage" in record.getMessage() for record in caplog.records
    )


def test_cli_template_with_release(tmp_path: Path) -> None:
    runner = CliRunner()
    extract_dir = tmp_path / "extract"
    release_dir = tmp_path / "release"

    result = runner.invoke(
        app,
        [
            "template",
            str(SAMPLE_TEMPLATE),
            "--output",
            str(extract_dir),
            "--with-release",
            "--brand",
            "Sample",
            "--version",
            "1.0.0",
            "--release-output",
            str(release_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0, result.output
    assert (extract_dir / "template_spec.json").exists()
    assert (extract_dir / "branding.json").exists()
    assert (release_dir / "template_release.json").exists()
    assert (release_dir / "release_report.json").exists()
    assert "テンプレ stage（抽出＋検証＋リリース）が完了しました。" in result.output


def test_cli_template_force_skips_validation(tmp_path: Path, monkeypatch) -> None:
    runner = CliRunner()
    extract_dir = tmp_path / "extract"
    template_path = tmp_path / "template_with_slide.pptx"
    _create_template_with_slide(template_path)

    def _fail_run(self):  # noqa: D401
        raise AssertionError(
            "validation should be skipped when --force is specified")

    monkeypatch.setattr(LayoutValidationSuite, "run", _fail_run)

    result = runner.invoke(
        app,
        [
            "template",
            str(template_path),
            "--output",
            str(extract_dir),
            "--slide",
            "--force",
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0, result.output
    assert (extract_dir / "template_spec.json").exists()
    assert "検証をスキップしました" in result.output
    snapshot_path = extract_dir / "slide_snapshot.json"
    assert snapshot_path.exists(), "--slide 指定時にスナップショットが出力されていません"


def test_cli_template_emits_slide_snapshot(tmp_path: Path) -> None:
    runner = CliRunner()
    extract_dir = tmp_path / "extract"
    template_path = tmp_path / "template_with_slide.pptx"
    _create_template_with_slide(template_path)

    result = runner.invoke(
        app,
        [
            "template",
            str(template_path),
            "--output",
            str(extract_dir),
            "--slide",
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0, result.output
    snapshot_path = extract_dir / "slide_snapshot.json"
    assert snapshot_path.exists(), result.output

    payload = json.loads(snapshot_path.read_text(encoding="utf-8"))
    assert payload["slides"], "スライドスナップショットが空です"
    first_slide = payload["slides"][0]
    assert "shapes" in first_slide and first_slide["shapes"], "図形情報が存在しません"
    first_shape = first_slide["shapes"][0]
    paragraph_texts = [
        paragraph
        for paragraph in first_shape.get("paragraphs", [])
        if paragraph.get("text")
    ]
    assert paragraph_texts, "段落テキストが取得できていません"
    sample_paragraph = paragraph_texts[0]
    assert "font_name" in sample_paragraph, "フォント属性が出力されていません"
    assert "alignment" in sample_paragraph, "段落属性が出力されていません"
    assert "placeholder_type" in first_shape, "プレースホルダー種別がシリアライズされていません"


def _prepare_generate_ready(
    runner: CliRunner,
    spec_path: Path,
    mapping_dir: Path,
    *,
    prepare_paths: dict[str, Path],
) -> Path:
    draft_dir = mapping_dir / "draft"
    args = [
        "mapping",
        str(spec_path),
        "--output",
        str(mapping_dir),
        *_prepare_args(prepare_paths),
    ]

    result = runner.invoke(app, args, catch_exceptions=False)
    assert result.exit_code == 0, result.output

    assert draft_dir.exists()
    ready_path = mapping_dir / "generate_ready.json"
    assert ready_path.exists()
    meta_path = mapping_dir / "generate_ready_meta.json"
    assert meta_path.exists()
    payload = json.loads(ready_path.read_text(encoding="utf-8"))
    meta = payload.get("meta", {})
    template_path = meta.get("template_path")
    assert template_path is not None
    template_path_obj = Path(template_path)
    if template_path_obj.is_absolute():
        assert template_path_obj.exists()
    else:
        resolved = (ready_path.parent / template_path_obj).resolve()
        assert resolved.exists()
    return ready_path


def test_cli_gen_generates_outputs(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    output_dir = tmp_path / "gen"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    prepare_paths = _prepare_inputs(runner, tmp_path)
    generate_ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        prepare_paths=prepare_paths,
    )

    result = runner.invoke(
        app,
        [
            "gen",
            str(generate_ready_path),
            "--output",
            str(output_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0, result.output
    assert "Polisher: disabled" in result.output

    spec = JobSpec.parse_file(spec_path)
    pptx_path = output_dir / "proposal.pptx"
    analysis_path = output_dir / "analysis.json"
    baseline_analysis_path = output_dir / "analysis_pre_polisher.json"
    audit_path = output_dir / "audit_log.json"
    rendering_log_path = output_dir / "rendering_log.json"

    assert pptx_path.exists()
    assert analysis_path.exists()
    assert baseline_analysis_path.exists()
    assert audit_path.exists()
    assert rendering_log_path.exists()

    audit_payload = json.loads(audit_path.read_text(encoding="utf-8"))
    hashes = audit_payload["hashes"]
    assert hashes.get("generate_ready", "").startswith("sha256:")
    mapping_info = audit_payload.get("mapping")
    assert mapping_info is not None
    assert mapping_info.get("generate_ready_path") == str(generate_ready_path)

    cards_payload = json.loads(
        prepare_paths["cards"].read_text(encoding="utf-8"))
    cards = cards_payload["cards"]

    presentation = Presentation(pptx_path)
    assert len(presentation.slides) == len(cards) == len(spec.slides)
    for card, slide in zip(cards, presentation.slides, strict=False):
        content = card.get("content") or {}
        expected_title = (content.get("title") or content.get(
            "headline") or card.get("card_id") or "").strip()
        if not expected_title:
            continue
        title_shape = slide.shapes.title
        if title_shape is None:
            continue
        assert title_shape.text == expected_title


def test_cli_prepare_generates_outputs(tmp_path: Path) -> None:
    runner = CliRunner()
    prepare_dir = tmp_path / "prepare"

    result = runner.invoke(
        app,
        [
            "prepare",
            str(PREPARE_SOURCE),
            "--mode",
            "dynamic",
            "--output",
            str(prepare_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0, result.output

    cards_path = prepare_dir / "prepare_card.json"
    meta_path = prepare_dir / "ai_generation_meta.json"
    log_path = prepare_dir / "prepare_ai_log.json"
    audit_path = prepare_dir / "audit_log.json"

    for path in (cards_path, meta_path, log_path, audit_path):
        assert path.exists()

    cards_payload = json.loads(cards_path.read_text(encoding="utf-8"))
    assert len(cards_payload["cards"]) >= 1
    audit_payload = json.loads(audit_path.read_text(encoding="utf-8"))
    assert audit_payload["prepare_normalization"]["statistics"]["cards_total"] == len(
        cards_payload["cards"])


def test_cli_mapping_then_gen(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    output_dir = tmp_path / "render"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        prepare_paths=prepare_paths,
    )

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads(
        (output_dir / "audit_log.json").read_text(encoding="utf-8"))
    artifacts = audit_payload.get("artifacts", {})
    assert artifacts.get("generate_ready") == str(ready_path)


def test_cli_mapping_requires_template(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)
    payload = json.loads(spec_path.read_text(encoding="utf-8"))
    payload.get("meta", {}).pop("template_path", None)
    spec_path.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    result = runner.invoke(
        app,
        [
            "mapping",
            str(spec_path),
            "--output",
            str(mapping_dir),
            *_prepare_args(prepare_paths),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 2
    assert (
        "jobspec.meta.template_path にテンプレートパスを設定してください。"
        in result.output
    )


def test_cli_compose_missing_layouts_path(tmp_path: Path) -> None:
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)
    payload = json.loads(spec_path.read_text(encoding="utf-8"))
    payload["meta"]["layouts_path"] = "missing/layouts.jsonl"
    spec_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

    result = runner.invoke(
        app,
        [
            "compose",
            str(spec_path),
            "--output",
            str(tmp_path / "compose"),
            *_prepare_args(prepare_paths),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 2
    assert "layouts_path" in result.output


def test_cli_compose_schema_validation_failure(tmp_path: Path) -> None:
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = tmp_path / "invalid_jobspec.json"
    spec_path.write_text(json.dumps({"meta": {}}, ensure_ascii=False, indent=2), encoding="utf-8")

    result = runner.invoke(
        app,
        [
            "compose",
            str(spec_path),
            "--output",
            str(tmp_path / "compose"),
            *_prepare_args(prepare_paths),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 2
    assert "スキーマ検証に失敗しました" in result.output


def test_cli_compose_fails_when_layouts_mismatch(tmp_path: Path) -> None:
    output_dir = tmp_path / "compose-layout-mismatch"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    payload = json.loads(spec_path.read_text(encoding="utf-8"))
    payload["slides"][0]["layout"] = "UnknownLayout"
    spec_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

    result = runner.invoke(
        app,
        [
            "compose",
            str(spec_path),
            "--output",
            str(output_dir),
            *_prepare_args(prepare_paths),
        ],
        env={"PPTX_STRICT_LAYOUTS": "1"},
        catch_exceptions=False,
    )

    assert result.exit_code == 4
    assert "レイアウト" in result.output


def test_cli_mapping_fails_when_layouts_mismatch(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping-layout-mismatch"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    payload = json.loads(spec_path.read_text(encoding="utf-8"))
    payload["slides"][0]["layout"] = "UnknownLayout"
    spec_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

    result = runner.invoke(
        app,
        [
            "mapping",
            str(spec_path),
            "--output",
            str(mapping_dir),
            *_prepare_args(prepare_paths),
        ],
        env={"PPTX_STRICT_LAYOUTS": "1"},
        catch_exceptions=False,
    )

    assert result.exit_code == 4
    assert "レイアウト" in result.output


def test_cli_compose_generates_stage45_outputs(tmp_path: Path) -> None:
    output_dir = tmp_path / "compose-gen"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    result = runner.invoke(
        app,
        [
            "compose",
            str(spec_path),
            "--output",
            str(output_dir),
            *_prepare_args(prepare_paths),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0, result.output
    draft_dir = output_dir / "draft"
    assert draft_dir.exists()
    assert (output_dir / "generate_ready.json").exists()
    assert (output_dir / "mapping_log.json").exists()
    assert (output_dir / "generate_ready_meta.json").exists()


def test_cli_gen_missing_template_path(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        prepare_paths=prepare_paths,
    )

    payload = json.loads(ready_path.read_text(encoding="utf-8"))
    payload["meta"].pop("template_path", None)
    stripped_ready = mapping_dir / "generate_ready_no_template.json"
    stripped_ready.write_text(json.dumps(
        payload, ensure_ascii=False), encoding="utf-8")

    result = runner.invoke(
        app,
        ["gen", str(stripped_ready), "--output", str(tmp_path / "out")],
        catch_exceptions=False,
    )

    assert result.exit_code == 2
    assert "template_path" in result.output


def test_cli_mapping_invalid_prepare_fails(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    invalid_cards = tmp_path / "prepare_card.json"
    invalid_cards.write_text("{}", encoding="utf-8")

    result = runner.invoke(
        app,
        [
            "mapping",
            str(spec_path),
            "--output",
            str(mapping_dir),
            "--prepare-cards",
            str(invalid_cards),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 4
    assert "プレペア成果物の読み込みに失敗しました" in result.output


def test_cli_compose_invalid_prepare_fails(tmp_path: Path) -> None:
    output_dir = tmp_path / "compose-invalid-prepare"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    invalid_cards = tmp_path / "prepare_card.json"
    invalid_cards.write_text(json.dumps({"cards": "invalid"}, ensure_ascii=False), encoding="utf-8")

    result = runner.invoke(
        app,
        [
            "compose",
            str(spec_path),
            "--output",
            str(output_dir),
            "--prepare-cards",
            str(invalid_cards),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 4
    assert "プレペア成果物の読み込みに失敗しました" in result.output


def test_cli_compose_propagates_mapping_fallback(tmp_path: Path, monkeypatch) -> None:
    output_dir = tmp_path / "compose-fallback"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    monkeypatch.setattr(
        "pptx_generator.cli_handlers.compose.run_mapping_pipeline",
        lambda **kwargs: (_ for _ in ()).throw(PipelineFallbackError("forced fallback")),
    )

    result = runner.invoke(
        app,
        [
            "compose",
            str(spec_path),
            "--output",
            str(output_dir),
            *_prepare_args(prepare_paths),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 4
    assert "forced fallback" in result.output


def test_cli_gen_exports_pdf(tmp_path: Path, monkeypatch) -> None:
    mapping_dir = tmp_path / "mapping"
    output_dir = tmp_path / "gen-pdf"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        prepare_paths=prepare_paths,
    )

    def fake_which(cmd: str) -> str | None:
        if cmd == "soffice":
            return sys.executable
        return shutil.which(cmd)

    def fake_run(*args, **kwargs):  # noqa: ANN401
        (Path(output_dir) / "proposal.pdf").write_bytes(b"%PDF-1.4 fake")
        return subprocess.CompletedProcess(args, returncode=0)

    monkeypatch.setattr(pdf_exporter.shutil, "which", fake_which)
    monkeypatch.setattr(pdf_exporter.subprocess, "run", fake_run)

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
            "--export-pdf",
            "--pdf-output",
            "custom.pdf",
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads(
        (output_dir / "audit_log.json").read_text(encoding="utf-8"))
    pdf_meta = audit_payload.get("pdf_export")
    assert pdf_meta is not None
    assert pdf_meta.get("status") == "success"


def test_cli_gen_pdf_only(tmp_path: Path, monkeypatch) -> None:
    if not _libreoffice_available():
        pytest.skip("LibreOffice が利用できないためスキップします")

    mapping_dir = tmp_path / "mapping"
    output_dir = tmp_path / "gen-pdf-only"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        prepare_paths=prepare_paths,
    )

    def fake_run(*args, **kwargs):  # noqa: ANN401
        (Path(output_dir) / "proposal.pdf").write_bytes(b"%PDF-1.4 fake")
        return subprocess.CompletedProcess(args, returncode=0)

    monkeypatch.setattr(pdf_exporter.subprocess, "run", fake_run)

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
            "--export-pdf",
            "--pdf-mode",
            "only",
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    assert not (output_dir / "proposal.pptx").exists()


def test_cli_gen_pdf_skip_env(tmp_path: Path, monkeypatch) -> None:
    mapping_dir = tmp_path / "mapping"
    output_dir = tmp_path / "gen-pdf-skip"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        prepare_paths=prepare_paths,
    )

    def fail_run(*args, **kwargs):  # noqa: ANN401
        raise subprocess.CalledProcessError(cmd=args, returncode=1)

    monkeypatch.setattr(pdf_exporter.subprocess, "run", fail_run)

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
            "--export-pdf",
        ],
        env={"PPTXGEN_SKIP_PDF_CONVERT": "1"},
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads(
        (output_dir / "audit_log.json").read_text(encoding="utf-8"))
    pdf_meta = audit_payload.get("pdf_export")
    assert pdf_meta is not None
    assert pdf_meta.get("status") == "skipped"


def test_cli_gen_pdf_retries_after_timeout(tmp_path: Path, monkeypatch) -> None:
    mapping_dir = tmp_path / "mapping"
    output_dir = tmp_path / "gen-pdf-timeout"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        prepare_paths=prepare_paths,
    )

    calls = {"count": 0}

    def fake_run(*args, **kwargs):  # noqa: ANN401
        calls["count"] += 1
        if calls["count"] == 1:
            raise subprocess.TimeoutExpired(cmd=args, timeout=1)
        (Path(output_dir) / "proposal.pdf").write_bytes(b"%PDF-1.4 fake")
        return subprocess.CompletedProcess(args, returncode=0)

    monkeypatch.setattr(pdf_exporter.subprocess, "run", fake_run)
    monkeypatch.setattr(
        pdf_exporter.shutil,
        "which",
        lambda cmd: sys.executable if cmd == "soffice" else shutil.which(cmd),
    )

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
            "--export-pdf",
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads(
        (output_dir / "audit_log.json").read_text(encoding="utf-8"))
    pdf_meta = audit_payload.get("pdf_export")
    assert pdf_meta is not None
    assert pdf_meta.get("status") == "success"
    assert pdf_meta.get("attempts") == 2


def test_cli_gen_pdf_skips_when_converter_unavailable(tmp_path: Path, monkeypatch) -> None:
    mapping_dir = tmp_path / "mapping"
    output_dir = tmp_path / "gen-pdf-fallback"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        prepare_paths=prepare_paths,
    )

    monkeypatch.setattr(
        pdf_exporter.LibreOfficeConverter,
        "convert",
        lambda self, pptx_path, output_dir: (_ for _ in ()).throw(pdf_exporter.PdfExportError("libreoffice missing")),
    )

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
            "--export-pdf",
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    assert (output_dir / "proposal.pptx").exists()
    audit_payload = json.loads(
        (output_dir / "audit_log.json").read_text(encoding="utf-8"))
    pdf_meta = audit_payload.get("pdf_export")
    assert pdf_meta is not None
    assert pdf_meta.get("status") == "skipped"
    assert "libreoffice" in pdf_meta.get("converter", "")


def test_pdf_export_step_skips_on_error(tmp_path: Path, monkeypatch) -> None:
    pptx_path = tmp_path / "input.pptx"
    pptx_path.write_bytes(b"pptx")
    spec = JobSpec(meta=JobMeta(schema_version="1.0", title="deck"), auth=JobAuth(created_by="tester"), slides=[])
    context = PipelineContext(spec=spec, workdir=tmp_path)
    context.add_artifact("pptx_path", pptx_path)

    monkeypatch.setattr(
        pdf_exporter.LibreOfficeConverter,
        "convert",
        lambda self, pptx_path, output_dir: (_ for _ in ()).throw(pdf_exporter.PdfExportError("fail")),
    )

    step = pdf_exporter.PdfExportStep(pdf_exporter.PdfExportOptions(enabled=True))
    step.run(context)

    pdf_meta = context.artifacts.get("pdf_export_metadata")
    assert pdf_meta is not None
    assert pdf_meta.get("status") == "skipped"


def test_cli_template_reports_validation_errors(tmp_path: Path, monkeypatch) -> None:
    template_path = tmp_path / "broken_template.pptx"
    _create_template_with_slide(template_path)
    output_dir = tmp_path / "extract"
    runner = CliRunner()

    class DummyValidationResult:
        def __init__(self, base_dir: Path) -> None:
            self.layouts_path = base_dir / "layouts.jsonl"
            self.layouts_path.write_text("{}", encoding="utf-8")
            self.diagnostics_path = base_dir / "diagnostics.json"
            self.diagnostics_path.write_text("{}", encoding="utf-8")
            self.diff_report_path = None
            self.record_count = 0
            self.warnings_count = 0
            self.errors_count = 1

    def fake_run(self):  # noqa: ANN001
        return DummyValidationResult(output_dir)

    monkeypatch.setattr(LayoutValidationSuite, "run", fake_run)

    result = runner.invoke(
        app,
        [
            "template",
            str(template_path),
            "--output",
            str(output_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 6
    assert "レイアウト検証でエラーが検出されました" in result.output


def test_cli_gen_with_polisher_stub(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    output_dir = tmp_path / "gen-polisher"
    rules_path = tmp_path / "polisher-rules.json"
    rules_path.write_text(json.dumps(
        {"min_font_size_pt": 18.0}), encoding="utf-8")

    script_path = tmp_path / "polisher_stub.py"
    script_path.write_text(
        "\n".join(
            [
                "import argparse",
                "import json",
                "from pathlib import Path",
                "",
                "parser = argparse.ArgumentParser()",
                "parser.add_argument('--input', required=True)",
                "parser.add_argument('--rules', required=True)",
                "args = parser.parse_args()",
                "path = Path(args.input)",
                "path.read_bytes()",
                "Path(args.rules).touch(exist_ok=True)",
                "print(json.dumps({'stub': 'ok'}))",
            ]
        ),
        encoding="utf-8",
    )

    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)
    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        prepare_paths=prepare_paths,
    )

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
            "--polisher",
            "--polisher-path",
            sys.executable,
            "--polisher-arg",
            str(script_path),
            "--polisher-arg",
            "--input",
            "--polisher-arg",
            "{pptx}",
            "--polisher-arg",
            "--rules",
            "--polisher-arg",
            "{rules}",
            "--polisher-rules",
            str(rules_path),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads(
        (output_dir / "audit_log.json").read_text(encoding="utf-8"))
    polisher_meta = audit_payload.get("polisher")
    assert polisher_meta is not None
    assert polisher_meta.get("status") == "success"


def test_cli_gen_template_branding_fallback(tmp_path, monkeypatch) -> None:
    mapping_dir = tmp_path / "mapping"
    output_dir = tmp_path / "gen-template-fallback"
    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        prepare_paths=prepare_paths,
    )

    monkeypatch.setattr("pptx_generator.template.template_style.extract_branding_config", lambda _: (
        _ for _ in ()).throw(BrandingExtractionError("boom")))

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads(
        (output_dir / "audit_log.json").read_text(encoding="utf-8"))
    style_info = audit_payload.get("template_style")
    assert style_info is not None
    source_info = style_info.get("source", {})
    assert source_info.get("type") == "default"
    assert "error" in source_info


def test_cli_gen_template_uses_template_branding(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    output_dir = tmp_path / "gen-template-branding"
    template_path = tmp_path / "template.pptx"
    shutil.copyfile(SAMPLE_TEMPLATE, template_path)

    runner = CliRunner()
    prepare_paths = _prepare_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, prepare_paths)
    spec_payload = json.loads(spec_path.read_text(encoding="utf-8"))
    spec_payload.setdefault("meta", {})
    spec_payload["meta"]["template_path"] = str(template_path)
    spec_path.write_text(
        json.dumps(spec_payload, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        prepare_paths=prepare_paths,
    )

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads(
        (output_dir / "audit_log.json").read_text(encoding="utf-8"))
    style_info = audit_payload.get("template_style")
    assert style_info is not None
    source_info = style_info.get("source", {})
    assert source_info.get("type") == "template"
    assert Path(source_info.get("template", "")).resolve() == template_path.resolve()


def test_cli_gen_default_output_directory(tmp_path) -> None:
    runner = CliRunner()
    repo_root = Path.cwd()
    with runner.isolated_filesystem(temp_dir=tmp_path) as fs_root:
        shutil.copytree(repo_root / "samples", Path(fs_root) / "samples")

        prepare_paths = _prepare_inputs(runner, Path(fs_root))
        spec_path = _create_matching_jobspec(
            Path(fs_root) / "samples/json",
            prepare_paths,
            filename="jobspec_matching_cards.json",
        )
        mapping_result = runner.invoke(
            app,
            [
                "mapping",
                str(spec_path.relative_to(Path(fs_root))),
                "--output",
                "samples/gen-ready",
                *_prepare_args(prepare_paths),
            ],
            catch_exceptions=False,
        )
        assert mapping_result.exit_code == 0, mapping_result.output

        ready_path = Path("samples/gen-ready/generate_ready.json")
        result = runner.invoke(
            app,
            [
                "gen",
                str(ready_path),
            ],
            catch_exceptions=False,
        )
        assert result.exit_code == 0


def test_static_mode_pipeline(tmp_path: Path) -> None:
    runner = CliRunner()
    template_path = Path("samples/templates/templates.pptx")

    blueprint_payload = {
        "template_path": str(template_path),
        "extracted_at": "2025-11-09T00:00:00Z",
        "layouts": [],
        "warnings": [],
        "errors": [],
        "layout_mode": "static",
        "blueprint": {
            "slides": [
                {
                    "slide_id": "title-01",
                    "layout": "Title",
                    "required": True,
                    "intent_tags": ["opening"],
                    "slots": [
                        {
                            "slot_id": "title-01.slot01",
                            "anchor": "Title",
                            "content_type": "text",
                            "required": True,
                            "intent_tags": ["headline"],
                        }
                    ],
                },
                {
                    "slide_id": "section_covor_left-01",
                    "layout": "Section Covor Left",
                    "required": True,
                    "intent_tags": ["body"],
                    "slots": [
                        {
                            "slot_id": "section_covor_left-01.slot01",
                            "anchor": "Section Title",
                            "content_type": "text",
                            "required": True,
                            "intent_tags": ["headline"],
                        }
                    ],
                },
            ]
        },
    }
    blueprint_path = tmp_path / "static_template_spec.json"
    blueprint_path.write_text(json.dumps(
        blueprint_payload, ensure_ascii=False, indent=2), encoding="utf-8")

    prepare_payload = {
        "meta": {"prepare_id": "static-prepare", "title": "Static Prepare"},
        "chapters": [
            {
                "id": "intro",
                "title": "Intro",
                "message": "Welcome",
                "details": ["Welcome"],
                "supporting_points": [],
                "intent_tags": ["opening"],
            },
            {
                "id": "overview",
                "title": "Overview",
                "message": "Overview message",
                "details": ["Overview detail"],
                "supporting_points": [],
                "intent_tags": ["body"],
            },
        ],
    }
    prepare_path = tmp_path / "prepare_static.json"
    prepare_path.write_text(json.dumps(
        prepare_payload, ensure_ascii=False, indent=2), encoding="utf-8")

    jobspec_payload = {
        "meta": {
            "schema_version": "1.1",
            "title": "Static Deck",
            "template_path": str(template_path),
            "layouts_path": str(Path("samples/extract/layouts.jsonl")),
            "locale": "ja-JP",
            "template_spec_path": str(blueprint_path),
        },
        "auth": {"created_by": "tester"},
        "slides": [
            {
                "id": "title-01",
                "layout": "Title",
                "title": "Intro",
                "auto_draw_anchors": ["Num"],
                "auto_draw_boxes": {
                    "Num": {
                        "left_in": 9.4,
                        "top_in": 6.9,
                        "width_in": 1.0,
                        "height_in": 0.4,
                    }
                },
            },
            {"id": "section_covor_left-01",
                "layout": "Section Covor Left", "title": "Overview"},
        ],
    }
    jobspec_path = tmp_path / "jobspec_static.json"
    jobspec_path.write_text(json.dumps(
        jobspec_payload, ensure_ascii=False, indent=2), encoding="utf-8")

    prepare_dir = tmp_path / "prepare_static"
    result = runner.invoke(
        app,
        [
            "prepare",
            str(prepare_path),
            "--mode",
            "static",
            "--jobspec",
            str(jobspec_path),
            "--output",
            str(prepare_dir),
        ],
        catch_exceptions=False,
    )
    assert result.exit_code == 0, result.output

    meta_payload = json.loads(
        (prepare_dir / "ai_generation_meta.json").read_text(encoding="utf-8"))
    assert meta_payload["mode"] == "static"
    assert meta_payload["slot_coverage"]["required_total"] == 2

    mapping_dir = tmp_path / "mapping_static"
    result = runner.invoke(
        app,
        [
            "mapping",
            str(jobspec_path),
            "--prepare-cards",
            str(prepare_dir / "prepare_card.json"),
            "--output",
            str(mapping_dir),
        ],
        catch_exceptions=False,
    )
    assert result.exit_code == 0, result.output
    draft_dir = mapping_dir / "draft"
    assert draft_dir.exists()

    ready_payload = json.loads(
        (mapping_dir / "generate_ready.json").read_text(encoding="utf-8"))
    assert ready_payload["meta"]["layout_mode"] == "static"
    first_slot = ready_payload["slides"][0]["meta"]["blueprint_slots"][0]
    assert first_slot["fulfilled"] is True
    first_auto_draw = ready_payload["slides"][0]["meta"]["auto_draw"][0]
    assert first_auto_draw["anchor"] == "Num"
    assert first_auto_draw["left_in"] == pytest.approx(9.4, rel=1e-3)

    mapping_log_payload = json.loads(
        (mapping_dir / "mapping_log.json").read_text(encoding="utf-8"))
    meta_payload = mapping_log_payload["meta"]
    assert meta_payload["mode"] == "static"
    assert meta_payload["static_slot_checks"]["unused_slots"] == []
    assert meta_payload["slot_summary"]["required_total"] == 2

    generate_ready_meta_payload = json.loads(
        (mapping_dir / DEFAULT_GENERATE_READY_META_FILENAME).read_text(encoding="utf-8"))
    assert generate_ready_meta_payload["mode"] == "static"
    assert Path(generate_ready_meta_payload["blueprint_path"]).resolve(
    ) == blueprint_path.resolve()
