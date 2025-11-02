"""SimpleAnalyzerStep の診断結果を検証するテスト。"""

from __future__ import annotations

import base64
import json

import pytest

from pptx import Presentation
from pptx.util import Inches

from pptx_generator.models import (
    FontSpec,
    JobAuth,
    JobMeta,
    JobSpec,
    Slide,
    SlideBullet,
    SlideBulletGroup,
    SlideImage,
)
from pptx_generator.pipeline import (
    AnalyzerOptions,
    MappingOptions,
    MappingStep,
    PipelineContext,
    RenderingOptions,
    SimpleAnalyzerStep,
    SimpleRendererStep,
)
from pptx_generator.generate_ready import generate_ready_to_jobspec


def _group(*bullets: SlideBullet, anchor: str | None = None) -> SlideBulletGroup:
    return SlideBulletGroup(anchor=anchor, items=list(bullets))


def _write_dummy_png(path) -> None:
    # 1px x 1px の透明 PNG
    payload = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNgYAAAAAMAASsJTYQAAAAASUVORK5CYII="
    )
    path.write_bytes(payload)


def _render_spec(spec: JobSpec, workdir, template_path=None) -> PipelineContext:
    context = PipelineContext(spec=spec, workdir=workdir)
    options = RenderingOptions(output_filename="test.pptx")
    if template_path is not None:
        options.template_path = template_path
    renderer = SimpleRendererStep(options)
    renderer.run(context)
    return context


def test_simple_analyzer_detects_quality_issues(tmp_path) -> None:
    image_path = tmp_path / "image.png"
    _write_dummy_png(image_path)

    spec = JobSpec(
        meta=JobMeta(
            schema_version="1.1",
            title="テスト案件",
            client="Zeta",
            author="営業部",
            created_at="2025-10-05",
            theme="corporate",
        ),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="slide-1",
                layout="Title and Content",
                title="テストスライド",
                bullets=[
                    _group(
                        SlideBullet(
                            id="bullet-1",
                            text="本文",
                            level=4,
                            font=FontSpec(
                                name="Yu Gothic",
                                size_pt=12.0,
                                color_hex="#FFFFFF",
                            ),
                        )
                    )
                ],
                images=[
                    SlideImage(
                        id="img-1",
                        source=str(image_path),
                        left_in=0.1,
                        top_in=0.2,
                        width_in=9.5,
                        height_in=7.0,
                    )
                ],
            )
        ],
    )

    context = _render_spec(spec, tmp_path)
    analyzer = SimpleAnalyzerStep(
        AnalyzerOptions(
            min_font_size=16.0,
            default_font_size=16.0,
            max_bullet_level=3,
            default_font_color="#CCCCCC",
            preferred_text_color="#005BAC",
            background_color="#FFFFFF",
            margin_in=0.5,
        )
    )

    analyzer.run(context)

    analysis_path = context.require_artifact("analysis_path")
    payload = json.loads(analysis_path.read_text(encoding="utf-8"))

    issue_types = {issue["type"] for issue in payload["issues"]}
    assert {
        "margin",
        "font_min",
        "contrast_low",
        "bullet_depth",
        "layout_consistency",
        "grid_misaligned",
    } <= issue_types

    fix_types = {fix["type"] for fix in payload["fixes"]}
    assert {"move", "font_raise", "color_adjust", "bullet_cap", "bullet_reindent"} <= fix_types

    margin_issue = next(issue for issue in payload["issues"] if issue["type"] == "margin")
    assert margin_issue["fix"]["payload"]

    layout_issue = next(issue for issue in payload["issues"] if issue["type"] == "layout_consistency")
    assert layout_issue["fix"]["payload"]["level"] == 0
    contrast_issue = next(issue for issue in payload["issues"] if issue["type"] == "contrast_low")
    assert contrast_issue["metrics"]["required_ratio"] == pytest.approx(4.5)
    assert contrast_issue["metrics"]["font_size_pt"] == pytest.approx(12.0)


def test_analyzer_updates_mapping_log(tmp_path) -> None:
    spec = JobSpec(
        meta=JobMeta(
            schema_version="1.1",
            title="Analyzer 連携テスト",
            client="Zeta",
            author="営業部",
            created_at="2025-10-21",
            theme="corporate",
        ),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="slide-1",
                layout="Title and Content",
                title="概要",
                bullets=[
                    _group(
                        SlideBullet(
                            id="bullet-1",
                            text="最初のポイント",
                            level=0,
                        )
                    )
                ],
            )
        ],
    )

    mapping_context = PipelineContext(spec=spec, workdir=tmp_path)
    mapping_step = MappingStep(MappingOptions(output_dir=tmp_path))
    mapping_step.run(mapping_context)

    generate_ready = mapping_context.artifacts["generate_ready"]
    render_spec = generate_ready_to_jobspec(generate_ready)
    render_context = PipelineContext(
        spec=render_spec,
        workdir=tmp_path,
        artifacts=dict(mapping_context.artifacts),
    )

    renderer = SimpleRendererStep(RenderingOptions(output_filename="rm031.pptx"))
    renderer.run(render_context)

    analyzer = SimpleAnalyzerStep(
        AnalyzerOptions(
            min_font_size=40.0,
            default_font_size=18.0,
            max_bullet_level=3,
            default_font_color="#777777",
            preferred_text_color="#005BAC",
            background_color="#FFFFFF",
        )
    )
    analyzer.run(render_context)

    mapping_log_path = tmp_path / "mapping_log.json"
    payload = json.loads(mapping_log_path.read_text(encoding="utf-8"))

    slide_summary = payload["slides"][0]["analyzer"]
    assert slide_summary["issue_count"] >= 1
    assert "font_min" in slide_summary["issue_counts_by_type"]
    issue_types = {issue["issue_type"] for issue in slide_summary["issues"]}
    assert "font_min" in issue_types
    target_slide_ids = {issue["target"].get("slide_id") for issue in slide_summary["issues"]}
    assert spec.slides[0].id in target_slide_ids

    meta = payload["meta"]
    assert meta["analyzer_issue_count"] == slide_summary["issue_count"]
    assert meta["analyzer_issue_counts_by_type"]["font_min"] == slide_summary["issue_counts_by_type"]["font_min"]


def test_margin_check_respects_actual_slide_size(tmp_path) -> None:
    template_path = tmp_path / "widescreen.pptx"
    template = Presentation()
    template.slide_width = Inches(13.333333)
    template.slide_height = Inches(7.5)
    template.save(template_path)

    image_path = tmp_path / "wide.png"
    _write_dummy_png(image_path)

    spec = JobSpec(
        meta=JobMeta(
            schema_version="1.1",
            title="ワイドスライド",
            client="Zeta",
            author="営業部",
            created_at="2025-10-08",
            theme="corporate",
        ),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="wide-1",
                layout="Title and Content",
                images=[
                    SlideImage(
                        id="wide-img",
                        source=str(image_path),
                        left_in=1.0,
                        top_in=1.0,
                        width_in=11.0,
                        height_in=5.5,
                    )
                ],
            )
        ],
    )

    context = _render_spec(spec, tmp_path, template_path=template_path)
    analyzer = SimpleAnalyzerStep(
        AnalyzerOptions(
            margin_in=0.5,
            grid_size_in=0.125,
            grid_tolerance_in=0.02,
        )
    )

    analyzer.run(context)

    analysis_path = context.require_artifact("analysis_path")
    payload = json.loads(analysis_path.read_text(encoding="utf-8"))

    margin_issues = [issue for issue in payload["issues"] if issue["type"] == "margin"]
    assert not margin_issues


def test_simple_analyzer_allows_large_text_with_lower_contrast(tmp_path) -> None:
    spec = JobSpec(
        meta=JobMeta(
            schema_version="1.1",
            title="コントラスト調整テスト",
            client="Zeta",
            author="営業部",
            created_at="2025-10-07",
            theme="corporate",
        ),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="slide-large",
                layout="Title and Content",
                bullets=[
                    _group(
                        SlideBullet(
                            id="bullet-large",
                            text="セカンダリカラーの本文",
                            level=0,
                            font=FontSpec(
                                name="Yu Gothic",
                                size_pt=24.0,
                                color_hex="#0097A7",
                            ),
                        )
                    )
                ],
            )
        ],
    )

    context = _render_spec(spec, tmp_path)
    analyzer = SimpleAnalyzerStep(
        AnalyzerOptions(
            min_font_size=16.0,
            default_font_size=16.0,
            max_bullet_level=3,
            default_font_color="#333333",
            preferred_text_color="#005BAC",
            background_color="#FFFFFF",
            large_text_threshold_pt=18.0,
            large_text_min_contrast=3.0,
        )
    )

    analyzer.run(context)

    analysis_path = context.require_artifact("analysis_path")
    payload = json.loads(analysis_path.read_text(encoding="utf-8"))

    issue_types = {issue["type"] for issue in payload["issues"]}
    assert "contrast_low" not in issue_types


def test_analyzer_outputs_structure_snapshot(tmp_path) -> None:
    spec = JobSpec(
        meta=JobMeta(
            schema_version="1.1",
            title="スナップショット検証",
            client="Zeta",
            author="営業部",
            created_at="2025-10-17",
            theme="corporate",
        ),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="slide-structure",
                layout="Title and Content",
                bullets=[
                    _group(
                        SlideBullet(
                            id="bullet-structure",
                            text="アンカー検証",
                            level=1,
                        ),
                    )
                ],
            )
        ],
    )

    context = _render_spec(spec, tmp_path)
    analyzer = SimpleAnalyzerStep(
        AnalyzerOptions(snapshot_output_filename="analysis_snapshot.json")
    )

    analyzer.run(context)

    snapshot_path = context.require_artifact("analyzer_snapshot_path")
    payload = json.loads(snapshot_path.read_text(encoding="utf-8"))

    assert payload["schema_version"] == "1.0.0"
    assert payload["slides"], "スナップショットにスライドが含まれていません"
    slide_entry = payload["slides"][0]
    assert slide_entry["slide_id"] == "slide-structure"
    assert slide_entry["layout"] == "Title and Content"
    assert "placeholders" in slide_entry
