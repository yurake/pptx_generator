"""CLI の統合テスト (generate_ready フロー確認用)."""

from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from collections import Counter

import pytest
from click.testing import CliRunner
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_generator.cli import app
from pptx_generator.branding_extractor import BrandingExtractionError
from pptx_generator.models import JobSpec, Slide, TemplateRelease, TemplateReleaseGoldenRun, TemplateReleaseReport, TemplateSpec
from pptx_generator.pipeline import pdf_exporter
from pptx_generator.layout_validation import LayoutValidationResult, LayoutValidationSuite

SAMPLE_TEMPLATE = Path("samples/templates/templates.pptx")
BRIEF_SOURCE = Path("samples/contents/sample_import_content_summary.txt")


def _libreoffice_available() -> bool:
    env_path = os.environ.get("LIBREOFFICE_PATH")
    if env_path and Path(env_path).exists():
        return True
    return shutil.which("soffice") is not None


def _collect_paragraph_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                texts.append(text)
    return texts


def _prepare_brief_inputs(runner: CliRunner, temp_dir: Path) -> dict[str, Path]:
    brief_dir = temp_dir / "prepare"
    result = runner.invoke(
        app,
        [
            "prepare",
            str(BRIEF_SOURCE),
            "--output",
            str(brief_dir),
        ],
        catch_exceptions=False,
    )
    assert result.exit_code == 0, result.output
    return {
        "dir": brief_dir,
        "cards": brief_dir / "prepare_card.json",
        "log": brief_dir / "brief_log.json",
        "meta": brief_dir / "ai_generation_meta.json",
    }


def _brief_args(paths: dict[str, Path]) -> list[str]:
    return [
        "--brief-cards",
        str(paths["cards"]),
        "--brief-log",
        str(paths["log"]),
        "--brief-meta",
        str(paths["meta"]),
    ]


def _create_matching_jobspec(root: Path, brief_paths: dict[str, Path], *, filename: str = "matching_jobspec.json") -> Path:
    base_spec = JobSpec.parse_file(Path("samples/json/sample_jobspec.json"))
    cards_payload = json.loads(brief_paths["cards"].read_text(encoding="utf-8"))
    cards = cards_payload.get("cards", [])

    slides: list[Slide] = []
    for index, card in enumerate(cards, start=1):
        card_id = card.get("card_id") or f"card-{index:03d}"
        title = card.get("chapter") or card.get("message") or card_id
        slides.append(
            Slide(
                id=card_id,
                layout="Content" if index > 1 else "Title",
                title=title[:120],
                notes=card.get("message"),
            )
        )

    jobspec = JobSpec(
        meta=base_spec.meta.model_copy(deep=True),
        auth=base_spec.auth.model_copy(deep=True),
        slides=slides,
    )

    spec_path = root / filename
    spec_path.parent.mkdir(parents=True, exist_ok=True)
    payload = jobspec.model_dump(mode="json")
    spec_path.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return spec_path


def test_cli_template_basic(tmp_path: Path) -> None:
    runner = CliRunner()
    extract_dir = tmp_path / "extract"

    result = runner.invoke(
        app,
        [
            "template",
            str(SAMPLE_TEMPLATE),
            "--output",
            str(extract_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0, result.output
    assert (extract_dir / "template_spec.json").exists()
    assert (extract_dir / "branding.json").exists()
    assert (extract_dir / "jobspec.json").exists()
    assert (extract_dir / "layouts.jsonl").exists()
    assert (extract_dir / "diagnostics.json").exists()
    assert "テンプレ工程（抽出＋検証）が完了しました。" in result.output


def test_cli_template_with_release(tmp_path: Path) -> None:
    runner = CliRunner()
    extract_dir = tmp_path / "extract"
    release_dir = tmp_path / "release"

    result = runner.invoke(
        app,
        [
            "template",
            str(SAMPLE_TEMPLATE),
            "--output",
            str(extract_dir),
            "--with-release",
            "--brand",
            "Sample",
            "--version",
            "1.0.0",
            "--release-output",
            str(release_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0, result.output
    assert (extract_dir / "template_spec.json").exists()
    assert (extract_dir / "branding.json").exists()
    assert (release_dir / "template_release.json").exists()
    assert (release_dir / "release_report.json").exists()
    assert "テンプレ工程（抽出＋検証＋リリース）が完了しました。" in result.output


def _prepare_generate_ready(
    runner: CliRunner,
    spec_path: Path,
    mapping_dir: Path,
    *,
    draft_dir: Path,
    brief_paths: dict[str, Path],
    extra_args: list[str] | None = None,
) -> Path:
    args = [
        "mapping",
        str(spec_path),
        "--output",
        str(mapping_dir),
        "--template",
        str(SAMPLE_TEMPLATE),
        "--draft-output",
        str(draft_dir),
        *_brief_args(brief_paths),
    ]
    if extra_args:
        args.extend(extra_args)

    result = runner.invoke(app, args, catch_exceptions=False)
    assert result.exit_code == 0, result.output

    ready_path = mapping_dir / "generate_ready.json"
    assert ready_path.exists()
    meta_path = mapping_dir / "generate_ready_meta.json"
    assert meta_path.exists()
    payload = json.loads(ready_path.read_text(encoding="utf-8"))
    meta = payload.get("meta", {})
    template_path = meta.get("template_path")
    assert template_path is not None
    assert Path(template_path).is_absolute()
    return ready_path


def test_cli_gen_generates_outputs(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    draft_dir = tmp_path / "draft"
    output_dir = tmp_path / "gen"
    runner = CliRunner()
    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, brief_paths)

    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    generate_ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        draft_dir=draft_dir,
        brief_paths=brief_paths,
    )

    result = runner.invoke(
        app,
        [
            "gen",
            str(generate_ready_path),
            "--output",
            str(output_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0, result.output
    assert "Polisher: disabled" in result.output

    spec = JobSpec.parse_file(spec_path)
    pptx_path = output_dir / "proposal.pptx"
    analysis_path = output_dir / "analysis.json"
    baseline_analysis_path = output_dir / "analysis_pre_polisher.json"
    audit_path = output_dir / "audit_log.json"
    rendering_log_path = output_dir / "rendering_log.json"

    assert pptx_path.exists()
    assert analysis_path.exists()
    assert baseline_analysis_path.exists()
    assert audit_path.exists()
    assert rendering_log_path.exists()

    audit_payload = json.loads(audit_path.read_text(encoding="utf-8"))
    hashes = audit_payload["hashes"]
    assert hashes.get("generate_ready", "").startswith("sha256:")
    mapping_info = audit_payload.get("mapping")
    assert mapping_info is not None
    assert mapping_info.get("generate_ready_path") == str(generate_ready_path)

    cards_payload = json.loads(brief_paths["cards"].read_text(encoding="utf-8"))
    cards = cards_payload["cards"]

    presentation = Presentation(pptx_path)
    assert len(presentation.slides) == len(cards) == len(spec.slides)
    for card, slide in zip(cards, presentation.slides, strict=False):
        expected_title = (card.get("message") or card.get("chapter") or card.get("card_id") or "").strip()
        if not expected_title:
            continue
        title_shape = slide.shapes.title
        if title_shape is None:
            continue
        assert title_shape.text == expected_title


def test_cli_prepare_generates_outputs(tmp_path: Path) -> None:
    runner = CliRunner()
    brief_dir = tmp_path / "prepare"

    result = runner.invoke(
        app,
        [
            "prepare",
            str(BRIEF_SOURCE),
            "--output",
            str(brief_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0, result.output

    cards_path = brief_dir / "prepare_card.json"
    meta_path = brief_dir / "ai_generation_meta.json"
    log_path = brief_dir / "brief_ai_log.json"
    audit_path = brief_dir / "audit_log.json"

    for path in (cards_path, meta_path, log_path, audit_path):
        assert path.exists()

    cards_payload = json.loads(cards_path.read_text(encoding="utf-8"))
    assert len(cards_payload["cards"]) >= 1
    audit_payload = json.loads(audit_path.read_text(encoding="utf-8"))
    assert audit_payload["brief_normalization"]["statistics"]["cards_total"] == len(cards_payload["cards"])


def test_cli_mapping_then_gen(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    draft_dir = tmp_path / "draft"
    output_dir = tmp_path / "render"
    runner = CliRunner()
    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, brief_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        draft_dir=draft_dir,
        brief_paths=brief_paths,
    )

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads((output_dir / "audit_log.json").read_text(encoding="utf-8"))
    artifacts = audit_payload.get("artifacts", {})
    assert artifacts.get("generate_ready") == str(ready_path)


def test_cli_mapping_requires_template(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    draft_dir = tmp_path / "draft"
    runner = CliRunner()
    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, brief_paths)

    result = runner.invoke(
        app,
        [
            "mapping",
            str(spec_path),
            "--output",
            str(mapping_dir),
            "--draft-output",
            str(draft_dir),
            *_brief_args(brief_paths),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 2
    assert (
        "テンプレートファイルを --template で指定するか、jobspec.meta.template_path にテンプレートパスを設定してください。"
        in result.output
    )


def test_cli_compose_generates_stage45_outputs(tmp_path: Path) -> None:
    draft_dir = tmp_path / "compose-draft"
    output_dir = tmp_path / "compose-gen"
    runner = CliRunner()
    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, brief_paths)

    result = runner.invoke(
        app,
        [
            "compose",
            str(spec_path),
            "--draft-output",
            str(draft_dir),
            "--output",
            str(output_dir),
            "--template",
            str(SAMPLE_TEMPLATE),
            *_brief_args(brief_paths),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0, result.output
    assert (output_dir / "generate_ready.json").exists()
    assert (output_dir / "mapping_log.json").exists()
    assert (output_dir / "generate_ready_meta.json").exists()


def test_cli_gen_missing_template_path(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    draft_dir = tmp_path / "draft"
    runner = CliRunner()
    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, brief_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        draft_dir=draft_dir,
        brief_paths=brief_paths,
    )

    payload = json.loads(ready_path.read_text(encoding="utf-8"))
    payload["meta"].pop("template_path", None)
    stripped_ready = mapping_dir / "generate_ready_no_template.json"
    stripped_ready.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")

    result = runner.invoke(
        app,
        ["gen", str(stripped_ready), "--output", str(tmp_path / "out")],
        catch_exceptions=False,
    )

    assert result.exit_code == 2
    assert "template_path" in result.output


def test_cli_mapping_invalid_brief_fails(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    draft_dir = tmp_path / "draft"
    runner = CliRunner()
    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, brief_paths)

    invalid_cards = tmp_path / "prepare_card.json"
    invalid_cards.write_text("{}", encoding="utf-8")

    result = runner.invoke(
        app,
        [
            "mapping",
            str(spec_path),
            "--output",
            str(mapping_dir),
            "--template",
            str(SAMPLE_TEMPLATE),
            "--draft-output",
            str(draft_dir),
            "--brief-cards",
            str(invalid_cards),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 4
    assert "ブリーフ成果物の読み込みに失敗しました" in result.output


def test_cli_gen_exports_pdf(tmp_path: Path, monkeypatch) -> None:
    mapping_dir = tmp_path / "mapping"
    draft_dir = tmp_path / "draft"
    output_dir = tmp_path / "gen-pdf"
    runner = CliRunner()
    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, brief_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        draft_dir=draft_dir,
        brief_paths=brief_paths,
    )

    def fake_which(cmd: str) -> str | None:
        if cmd == "soffice":
            return sys.executable
        return shutil.which(cmd)

    def fake_run(*args, **kwargs):  # noqa: ANN401
        (Path(output_dir) / "proposal.pdf").write_bytes(b"%PDF-1.4 fake")
        return subprocess.CompletedProcess(args, returncode=0)

    monkeypatch.setattr(pdf_exporter.shutil, "which", fake_which)
    monkeypatch.setattr(pdf_exporter.subprocess, "run", fake_run)

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
            "--export-pdf",
            "--pdf-output",
            "custom.pdf",
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads((output_dir / "audit_log.json").read_text(encoding="utf-8"))
    pdf_meta = audit_payload.get("pdf_export")
    assert pdf_meta is not None
    assert pdf_meta.get("status") == "success"


def test_cli_gen_pdf_only(tmp_path: Path, monkeypatch) -> None:
    if not _libreoffice_available():
        pytest.skip("LibreOffice が利用できないためスキップします")

    mapping_dir = tmp_path / "mapping"
    draft_dir = tmp_path / "draft"
    output_dir = tmp_path / "gen-pdf-only"
    runner = CliRunner()
    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, brief_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        draft_dir=draft_dir,
        brief_paths=brief_paths,
    )

    def fake_run(*args, **kwargs):  # noqa: ANN401
        (Path(output_dir) / "proposal.pdf").write_bytes(b"%PDF-1.4 fake")
        return subprocess.CompletedProcess(args, returncode=0)

    monkeypatch.setattr(pdf_exporter.subprocess, "run", fake_run)

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
            "--export-pdf",
            "--pdf-mode",
            "only",
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    assert not (output_dir / "proposal.pptx").exists()


def test_cli_gen_pdf_skip_env(tmp_path: Path, monkeypatch) -> None:
    mapping_dir = tmp_path / "mapping"
    draft_dir = tmp_path / "draft"
    output_dir = tmp_path / "gen-pdf-skip"
    runner = CliRunner()
    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, brief_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        draft_dir=draft_dir,
        brief_paths=brief_paths,
    )

    def fail_run(*args, **kwargs):  # noqa: ANN401
        raise subprocess.CalledProcessError(cmd=args, returncode=1)

    monkeypatch.setattr(pdf_exporter.subprocess, "run", fail_run)

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
            "--export-pdf",
        ],
        env={"PPTXGEN_SKIP_PDF_CONVERT": "1"},
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads((output_dir / "audit_log.json").read_text(encoding="utf-8"))
    pdf_meta = audit_payload.get("pdf_export")
    assert pdf_meta is not None
    assert pdf_meta.get("status") == "skipped"


def test_cli_gen_with_polisher_stub(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    draft_dir = tmp_path / "draft"
    output_dir = tmp_path / "gen-polisher"
    rules_path = tmp_path / "polisher-rules.json"
    rules_path.write_text(json.dumps({"min_font_size_pt": 18.0}), encoding="utf-8")

    script_path = tmp_path / "polisher_stub.py"
    script_path.write_text(
        "\n".join(
            [
                "import argparse",
                "import json",
                "from pathlib import Path",
                "",
                "parser = argparse.ArgumentParser()",
                "parser.add_argument('--input', required=True)",
                "parser.add_argument('--rules', required=True)",
                "args = parser.parse_args()",
                "path = Path(args.input)",
                "path.read_bytes()",
                "Path(args.rules).touch(exist_ok=True)",
                "print(json.dumps({'stub': 'ok'}))",
            ]
        ),
        encoding="utf-8",
    )

    runner = CliRunner()
    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, brief_paths)
    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        draft_dir=draft_dir,
        brief_paths=brief_paths,
    )

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
            "--polisher",
            "--polisher-path",
            sys.executable,
            "--polisher-arg",
            str(script_path),
            "--polisher-arg",
            "--input",
            "--polisher-arg",
            "{pptx}",
            "--polisher-arg",
            "--rules",
            "--polisher-arg",
            "{rules}",
            "--polisher-rules",
            str(rules_path),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads((output_dir / "audit_log.json").read_text(encoding="utf-8"))
    polisher_meta = audit_payload.get("polisher")
    assert polisher_meta is not None
    assert polisher_meta.get("status") == "success"


def test_cli_gen_template_with_explicit_branding(tmp_path: Path) -> None:
    mapping_dir = tmp_path / "mapping"
    draft_dir = tmp_path / "draft"
    output_dir = tmp_path / "gen-template-branding"
    template_path = tmp_path / "template.pptx"
    branding_path = tmp_path / "custom-branding.json"
    shutil.copyfile(SAMPLE_TEMPLATE, template_path)

    branding_payload = {
        "fonts": {
            "heading": {"name": "Test Font", "size_pt": 30, "color_hex": "#123456"},
            "body": {"name": "Test Font", "size_pt": 16, "color_hex": "#654321"},
        },
        "colors": {
            "primary": "#111111",
            "secondary": "#222222",
            "accent": "#333333",
            "background": "#FFFFFF",
        },
    }
    branding_path.write_text(json.dumps(branding_payload), encoding="utf-8")

    runner = CliRunner()
    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, brief_paths)
    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        draft_dir=draft_dir,
        brief_paths=brief_paths,
        extra_args=["--template", str(template_path)],
    )

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
            "--branding",
            str(branding_path),
            *_brief_args(brief_paths),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads((output_dir / "audit_log.json").read_text(encoding="utf-8"))
    branding_info = audit_payload.get("branding")
    assert branding_info is not None
    assert branding_info.get("source", {}).get("type") == "file"


def test_cli_gen_template_branding_fallback(tmp_path, monkeypatch) -> None:
    mapping_dir = tmp_path / "mapping"
    draft_dir = tmp_path / "draft"
    output_dir = tmp_path / "gen-template-fallback"
    runner = CliRunner()
    brief_paths = _prepare_brief_inputs(runner, tmp_path)
    spec_path = _create_matching_jobspec(tmp_path, brief_paths)

    ready_path = _prepare_generate_ready(
        runner,
        spec_path,
        mapping_dir,
        draft_dir=draft_dir,
        brief_paths=brief_paths,
    )

    monkeypatch.setattr("pptx_generator.cli.extract_branding_config", lambda _: (_ for _ in ()).throw(BrandingExtractionError("boom")))

    result = runner.invoke(
        app,
        [
            "gen",
            str(ready_path),
            "--output",
            str(output_dir),
        ],
        catch_exceptions=False,
    )

    assert result.exit_code == 0
    audit_payload = json.loads((output_dir / "audit_log.json").read_text(encoding="utf-8"))
    branding_info = audit_payload.get("branding")
    assert branding_info is not None
    source_info = branding_info.get("source", {})
    assert source_info.get("type") == "default"
    assert "error" in source_info


def test_cli_gen_default_output_directory(tmp_path) -> None:
    runner = CliRunner()
    repo_root = Path.cwd()
    with runner.isolated_filesystem(temp_dir=tmp_path) as fs_root:
        shutil.copytree(repo_root / "samples", Path(fs_root) / "samples")
        shutil.copytree(repo_root / "config", Path(fs_root) / "config")

        brief_paths = _prepare_brief_inputs(runner, Path(fs_root))
        spec_path = _create_matching_jobspec(
            Path(fs_root) / "samples/json",
            brief_paths,
            filename="jobspec_matching_cards.json",
        )
        mapping_result = runner.invoke(
            app,
            [
                "mapping",
                str(spec_path.relative_to(Path(fs_root))),
                "--output",
                "samples/gen-ready",
                "--template",
                "samples/templates/templates.pptx",
                *_brief_args(brief_paths),
            ],
            catch_exceptions=False,
        )
        assert mapping_result.exit_code == 0, mapping_result.output

        ready_path = Path("samples/gen-ready/generate_ready.json")
        result = runner.invoke(
            app,
            [
                "gen",
                str(ready_path),
            ],
            catch_exceptions=False,
        )
        assert result.exit_code == 0
