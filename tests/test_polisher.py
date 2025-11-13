"""PolisherStep の動作検証テスト。"""

from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest
from pptx import Presentation

from pptx_generator.models import JobAuth, JobMeta, JobSpec
from pptx_generator.pipeline import (PipelineContext, PolisherError,
                                     PolisherOptions, PolisherStep)


def _build_context(tmp_path: Path) -> tuple[PipelineContext, Path]:
    spec = JobSpec(
        meta=JobMeta(
            schema_version="1.1",
            title="Polisher テスト",
            client="Test",
            author="営業部",
            created_at="2025-10-18",
            theme="corporate",
        ),
        auth=JobAuth(created_by="tester"),
        slides=[],
    )
    context = PipelineContext(spec=spec, workdir=tmp_path)
    pptx_path = tmp_path / "output.pptx"
    pptx_path.write_bytes(b"pptx-initial")
    context.add_artifact("pptx_path", pptx_path)
    return context, pptx_path


def test_polisher_disabled_skips_execution(tmp_path: Path) -> None:
    context, _ = _build_context(tmp_path)
    step = PolisherStep(PolisherOptions(enabled=False))

    step.run(context)

    metadata = context.require_artifact("polisher_metadata")
    assert metadata["status"] == "disabled"
    assert metadata["enabled"] is False


def test_polisher_missing_executable_raises(tmp_path: Path) -> None:
    context, _ = _build_context(tmp_path)
    step = PolisherStep(PolisherOptions(enabled=True))

    with pytest.raises(PolisherError):
        step.run(context)


def test_polisher_executes_stub_command(tmp_path: Path) -> None:
    context, pptx_path = _build_context(tmp_path)
    rules_path = tmp_path / "polisher-rules.json"
    rules_path.write_text("{}", encoding="utf-8")

    script_path = tmp_path / "polisher_stub.py"
    script_path.write_text(
        "\n".join(
            [
                "import argparse",
                "import json",
                "from pathlib import Path",
                "",
                "parser = argparse.ArgumentParser()",
                "parser.add_argument('--input', required=True)",
                "parser.add_argument('--rules', required=True)",
                "args = parser.parse_args()",
                "path = Path(args.input)",
                "data = path.read_bytes()",
                "path.write_bytes(data)",
                "Path(args.rules).touch(exist_ok=True)",
                "print(json.dumps({'slides': 0, 'adjusted_font_size': 0, 'adjusted_color': 0}))",
            ]
        ),
        encoding="utf-8",
    )

    step = PolisherStep(
        PolisherOptions(
            enabled=True,
            executable=Path(sys.executable),
            rules_path=rules_path,
            timeout_sec=30,
            arguments=(str(script_path), "--input",
                       "{pptx}", "--rules", "{rules}"),
        )
    )

    step.run(context)

    metadata = context.require_artifact("polisher_metadata")
    assert metadata["status"] == "success"
    assert metadata["enabled"] is True
    assert metadata["returncode"] == 0
    assert metadata["command"][0] == str(Path(sys.executable))
    summary = metadata.get("summary")
    assert isinstance(summary, dict)
    assert summary.get("slides") == 0
    assert pptx_path.read_bytes() == b"pptx-initial"


def test_ai_footer_disabled_when_no_config(tmp_path: Path) -> None:
    """AI フッタ設定がない場合、フッタ付与がスキップされることを確認。"""
    context, pptx_path = _build_context(tmp_path)

    # 最小限のPPTXファイルを作成
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(pptx_path))

    step = PolisherStep(PolisherOptions(enabled=False))
    result = step._add_ai_footer_if_needed(pptx_path, context)

    assert result["enabled"] is False
    assert result["slides_modified"] == 0
    assert result["error"] is None


def test_ai_footer_enabled_with_config(tmp_path: Path) -> None:
    """AI フッタ設定がある場合、フッタが付与されることを確認。"""
    context, pptx_path = _build_context(tmp_path)

    # 最小限のPPTXファイルを作成
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(pptx_path))

    # ルール設定ファイルを作成
    rules_path = tmp_path / "polisher-rules.json"
    rules_path.write_text(
        json.dumps({
            "ai_footer": {
                "enabled": True,
                "text": "※テストフッタ",
                "font_size_pt": 8.0,
                "color": "#666666",
                "position": "bottom_right",
                "margin_in": 0.25
            }
        }),
        encoding="utf-8"
    )

    # AI生成メタデータを作成
    meta_path = tmp_path / "ai_generation_meta.json"
    meta_path.write_text(
        json.dumps({
            "cards": [
                {"card_id": "slide1"},
                {"card_id": "slide2"}
            ]
        }),
        encoding="utf-8"
    )

    step = PolisherStep(PolisherOptions(enabled=False, rules_path=rules_path))
    result = step._add_ai_footer_if_needed(pptx_path, context)

    assert result["enabled"] is True
    assert result["slides_modified"] == 1
    assert result["error"] is None

    # PPTXファイルを開いてフッタが追加されていることを確認
    prs_modified = Presentation(str(pptx_path))
    assert len(prs_modified.slides) == 1
    slide = prs_modified.slides[0]
    # テキストボックスが追加されていることを確認
    textboxes = [shape for shape in slide.shapes if hasattr(
        shape, "text_frame")]
    assert len(textboxes) > 0


def test_ai_footer_skipped_when_disabled(tmp_path: Path) -> None:
    """AI フッタが無効化されている場合、処理がスキップされることを確認。"""
    context, pptx_path = _build_context(tmp_path)

    # 最小限のPPTXファイルを作成
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(pptx_path))

    # ルール設定ファイルを作成（enabled=false）
    rules_path = tmp_path / "polisher-rules.json"
    rules_path.write_text(
        json.dumps({
            "ai_footer": {
                "enabled": False
            }
        }),
        encoding="utf-8"
    )

    step = PolisherStep(PolisherOptions(enabled=False, rules_path=rules_path))
    result = step._add_ai_footer_if_needed(pptx_path, context)

    assert result["enabled"] is False
    assert result["slides_modified"] == 0


def test_ai_footer_skipped_when_no_metadata(tmp_path: Path) -> None:
    """AI生成メタデータが存在しない場合、処理がスキップされることを確認。"""
    context, pptx_path = _build_context(tmp_path)

    # 最小限のPPTXファイルを作成
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(pptx_path))

    # ルール設定ファイルを作成
    rules_path = tmp_path / "polisher-rules.json"
    rules_path.write_text(
        json.dumps({
            "ai_footer": {
                "enabled": True,
                "text": "※テストフッタ"
            }
        }),
        encoding="utf-8"
    )

    # AI生成メタデータは作成しない

    step = PolisherStep(PolisherOptions(enabled=False, rules_path=rules_path))
    result = step._add_ai_footer_if_needed(pptx_path, context)

    assert result["enabled"] is True
    assert result["slides_modified"] == 0
    assert result["error"] is None
