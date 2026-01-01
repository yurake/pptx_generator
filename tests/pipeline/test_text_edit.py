from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

import json

from pptx_generator.pipeline.text_edit import (
    apply_shape_text_edits,
    generate_edits_template,
    overwrite_text_frame_preserving_style,
)
from pptx_generator.pipeline.analyzer.snapshot import table_cell_shape_id


def test_overwrite_text_frame_preserves_run_style() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))

    text_frame = textbox.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.text = "before"
    paragraph.level = 1
    paragraph.line_spacing = Pt(30)
    paragraph.space_after = Pt(6)
    paragraph.alignment = None
    run = paragraph.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x12, 0x34, 0x56)
    text_frame.word_wrap = False

    overwrite_text_frame_preserving_style(text_frame, "after line1\nafter line2")

    assert [p.text for p in text_frame.paragraphs] == ["after line1", "after line2"]
    for para in text_frame.paragraphs:
        assert para.level == 1
        assert para.line_spacing == Pt(30)
        assert para.space_after == Pt(6)
        assert para.runs[0].font.name == "Calibri"
        assert para.runs[0].font.size == Pt(18)
        assert para.runs[0].font.bold is True
        assert para.runs[0].font.color.rgb == RGBColor(0x12, 0x34, 0x56)
    assert text_frame.word_wrap is False


def test_apply_shape_text_edits_updates_textbox(tmp_path) -> None:
    pptx_path = tmp_path / "input.pptx"
    output_path = tmp_path / "output.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.text = "before"
    paragraph.runs[0].font.name = "Calibri"
    presentation.save(pptx_path)

    applied, missing = apply_shape_text_edits(
        pptx_path,
        [{"shape_id": textbox.shape_id, "contents": "after"}],
        output_path=output_path,
    )

    assert applied == 1
    assert missing == []
    reloaded = Presentation(output_path)
    reloaded_shape = next(s for s in reloaded.slides[0].shapes if s.shape_id == textbox.shape_id)
    assert reloaded_shape.text == "after"


def test_apply_shape_text_edits_updates_table_cell(tmp_path) -> None:
    pptx_path = tmp_path / "input_table.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    table = table_shape.table
    table.cell(1, 0).text = "cell"
    presentation.save(pptx_path)

    target_id = table_cell_shape_id(int(table_shape.shape_id), 1, 0)
    applied, missing = apply_shape_text_edits(
        pptx_path,
        [{"shape_id": target_id, "contents": "updated"}],
    )

    assert applied == 1
    assert missing == []
    reloaded = Presentation(pptx_path)
    assert reloaded.slides[0].shapes[0].table.cell(1, 0).text == "updated"


def test_apply_shape_text_edits_with_slide_index(tmp_path) -> None:
    pptx_path = tmp_path / "input_slide_index.pptx"
    presentation = Presentation()
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    shape1.text = "s1"
    shape1.name = "slide1_box"
    shape2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    shape2.text = "s2"
    shape2.name = "slide2_box"
    presentation.save(pptx_path)

    applied, missing = apply_shape_text_edits(
            pptx_path,
            [
                {"slide_index": 1, "shape_id": shape1.shape_id, "contents": "wrong slide", "name": shape1.name},
                {"slide_index": 1, "shape_id": shape2.shape_id, "contents": "target slide", "name": shape2.name},
            ],
        )

    assert applied == 1
    assert missing == ["1:" + str(shape1.shape_id)]
    reloaded = Presentation(pptx_path)
    # slide_index を指定した場合、指定スライド内の shape_id にマッチした場合のみ適用される
    assert reloaded.slides[0].shapes[0].text == "s1"
    assert reloaded.slides[1].shapes[0].text == "target slide"


def test_generate_edits_template_outputs_json(tmp_path) -> None:
    pptx_path = tmp_path / "input_template.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "hello"
    presentation.save(pptx_path)

    output = generate_edits_template(pptx_path)

    payload = json.loads(output.read_text(encoding="utf-8"))
    assert "edits" in payload
    assert any(edit["contents"] == "hello" for edit in payload["edits"])


def test_pptx_edit_cli_runs_with_mock_llm(tmp_path) -> None:
    from click.testing import CliRunner
    from pptx_generator.cli_commands.edit import create_edit_command

    pptx_path = tmp_path / "input_cli.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "keep"
    presentation.save(pptx_path)

    runner = CliRunner()
    cmd = create_edit_command()
    output_path = tmp_path / "output_cli.pptx"
    result = runner.invoke(cmd, [str(pptx_path), "--output", str(output_path)])

    assert result.exit_code == 0
    assert output_path.exists()
    reloaded = Presentation(output_path)
    reloaded_text = reloaded.slides[0].shapes[0].text
    assert reloaded_text == "keep"


def test_pptx_edit_cli_fails_with_legacy_option(tmp_path) -> None:
    from click.testing import CliRunner
    from pptx_generator.cli_commands.edit import create_edit_command

    pptx_path = tmp_path / "input_cli.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "keep"
    presentation.save(pptx_path)

    runner = CliRunner()
    cmd = create_edit_command()
    result = runner.invoke(cmd, ["--pptx-path", str(pptx_path)])

    assert result.exit_code != 0
    assert "No such option" in result.output or "no such option" in result.output


def test_pptx_edit_cli_llm_failure(tmp_path, monkeypatch) -> None:
    from click.testing import CliRunner
    from pptx_generator.cli_commands import edit as edit_module
    from pptx_generator.cli_commands.edit import create_edit_command

    pptx_path = tmp_path / "input_cli_fail.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "keep"
    presentation.save(pptx_path)

    class DummyClient:
        def rewrite(self, request):
            raise RuntimeError("llm fail")

    monkeypatch.setattr(edit_module, "create_edit_ai_client", lambda: DummyClient())

    runner = CliRunner()
    cmd = create_edit_command()
    output_path = pptx_path.with_name(f"{pptx_path.stem}_edited{pptx_path.suffix}")
    # 出力先を固定できるようにルートを指定
    out_root = tmp_path / "out"
    monkeypatch.setenv("PPTX_OUTPUT_ROOT", str(out_root))

    result = runner.invoke(cmd, [str(pptx_path)])

    assert result.exit_code != 0
    # エラー時は出力PPTXが作成されない
    assert not any(out_root.rglob("*.pptx"))


def test_pptx_edit_cli_default_output_dir(tmp_path, monkeypatch) -> None:
    from click.testing import CliRunner
    from pptx_generator.cli_commands.edit import create_edit_command
    from pathlib import Path

    pptx_path = tmp_path / "input_cli_default.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "keep"
    presentation.save(pptx_path)

    monkeypatch.delenv("PPTX_OUTPUT_ROOT", raising=False)
    runner = CliRunner()
    cmd = create_edit_command()
    result = runner.invoke(cmd, [str(pptx_path)])

    assert result.exit_code == 0
    # 出力は .pptx/<tx>/edit/<job_id>/input_cli_default.pptx 配下
    out_dir = Path(".pptx")
    outputs = list(out_dir.glob("*/edit/*/input_cli_default.pptx"))
    assert outputs, "expected output file under .pptx/<tx>/edit/<job_id>/"


def test_pptx_edit_cli_same_name_multiple_runs(tmp_path, monkeypatch) -> None:
    from click.testing import CliRunner
    from pptx_generator.cli_commands.edit import create_edit_command

    pptx_path = tmp_path / "input_cli_multi.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "keep"
    presentation.save(pptx_path)

    out_root = tmp_path / "out"
    monkeypatch.setenv("PPTX_OUTPUT_ROOT", str(out_root))
    runner = CliRunner()
    cmd = create_edit_command()

    res1 = runner.invoke(cmd, [str(pptx_path)])
    res2 = runner.invoke(cmd, [str(pptx_path)])

    assert res1.exit_code == 0
    assert res2.exit_code == 0

    outputs = list(out_root.glob("*/edit/*/input_cli_multi.pptx"))
    assert len(outputs) == 2, "should create separate files under different job_id"
