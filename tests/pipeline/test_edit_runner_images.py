from pptx import Presentation
from pptx.util import Inches

from pptx_generator.pipeline import edit_runner
from pptx_generator.pipeline.edit_runner import generate_edits_via_llm
from pptx_generator.pipeline.slide_image_exporter import (
    SlideImageAsset,
    SlideImageExportResult,
)


def test_generate_edits_via_llm_includes_image(tmp_path):
    pptx_path = tmp_path / "edit_source.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "keep"
    presentation.save(pptx_path)

    image_path = tmp_path / "slide.png"
    image_path.write_bytes(b"fake")
    output_dir = tmp_path / "images"

    class DummyExporter:
        def export(self, _pptx_path, _output_dir):
            return SlideImageExportResult(
                images_by_slide={
                    0: [
                        SlideImageAsset(
                            slide_index=0,
                            format="png",
                            path=image_path,
                            media_type="image/png",
                        )
                    ]
                },
                status="success",
            )

    called = {}

    class DummyClient:
        def rewrite(self, request):
            called["images"] = request.images
            called["prompt"] = request.prompt
            return type("Resp", (), {"edits": [{"shape_id": 1, "contents": "y"}], "model": "m"})

    def snapshot_fn(_):
        return [{"shape_id": 1, "slide_index": 0, "text": "x"}]

    edits, models = generate_edits_via_llm(
        pptx_path,
        snapshot_fn=snapshot_fn,
        client_factory=lambda: DummyClient(),
        image_exporter=DummyExporter(),
        image_output_dir=output_dir,
    )

    assert edits
    assert models == {"m"}
    assert called["images"]
    assert "screenshot" in called["prompt"]
    assert (output_dir / "edit_slide_images.json").exists()


def test_resolve_edit_image_exporter(monkeypatch, tmp_path):
    monkeypatch.setenv("PPTX_EDIT_IMAGE_INPUT", "1")
    monkeypatch.setenv("PPTX_EDIT_IMAGE_FORMATS", "png")
    exporter = edit_runner._resolve_edit_image_exporter()
    assert exporter is not None
    assert exporter.options.formats == ("png",)

    output = edit_runner._resolve_edit_image_output_dir(tmp_path / "out.pptx")
    assert output.exists()
