"""SimpleAnalyzerStep の診断結果を検証するテスト。"""

from __future__ import annotations

import base64
import json
from typing import Any

import pytest

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Inches, Pt

from pptx_generator.prepare import (
    PrepareBodyBlock,
    PrepareCard,
    PrepareCardContent,
    PrepareCardRole,
    PrepareDocument,
    PrepareStoryContext,
)
from pptx_generator.models import (
    DraftDocument,
    DraftSection,
    DraftSlideCard,
    FontSpec,
    JobAuth,
    JobMeta,
    JobSpec,
    Slide,
    SlideBullet,
    SlideBulletGroup,
    SlideImage,
    SlideTextbox,
)
from pptx_generator.pipeline import (
    AnalyzerOptions,
    MappingOptions,
    MappingStep,
    PipelineContext,
    RenderingOptions,
    SimpleAnalyzerStep,
    SimpleRendererStep,
)
from pptx_generator.pipeline.generate_ready import generate_ready_to_jobspec
from pptx_generator.pipeline.analyzer import (
    BulletParagraphResolver,
    ShapeSnapshot,
    SlideSnapshot,
    ParagraphSnapshot,
    _contrast_ratio,
    _extract_font_info,
    _extract_paragraph_style,
    _extract_text_frame_padding,
    _length_to_inches,
    _length_to_pt,
    _normalize_hex,
    _color_to_hex,
    _hex_to_rgb,
    _enum_name,
)


def _group(*bullets: SlideBullet, anchor: str | None = None) -> SlideBulletGroup:
    return SlideBulletGroup(anchor=anchor, items=list(bullets))


def _write_dummy_png(path) -> None:
    # 1px x 1px の透明 PNG
    payload = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNgYAAAAAMAASsJTYQAAAAASUVORK5CYII="
    )
    path.write_bytes(payload)


def _build_snapshot_with_textbox() -> SlideSnapshot:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])

    body_shape = slide.shapes.placeholders[1]
    body_shape.name = "Body"
    text_frame = body_shape.text_frame
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_top = Inches(0.1)
    paragraph = text_frame.paragraphs[0]
    paragraph.text = "本文"
    paragraph.level = 0
    paragraph.font.size = Pt(24)
    paragraph.font.name = "Calibri"
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(0x12, 0x34, 0x56)
    paragraph.line_spacing = 36
    paragraph.space_before = Pt(6)
    paragraph.space_after = Pt(12)

    textbox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(4), Inches(1))
    textbox.name = "anchored-shape"
    textbox.text_frame.text = "アンカー"
    additional = textbox.text_frame.add_paragraph()
    additional.text = "サブ"

    return SlideSnapshot.from_slide(slide, 0)


def _render_spec(spec: JobSpec, workdir, template_path=None) -> PipelineContext:
    context = PipelineContext(spec=spec, workdir=workdir)
    options = RenderingOptions(output_filename="test.pptx")
    if template_path is not None:
        options.template_path = template_path
    renderer = SimpleRendererStep(options)
    renderer.run(context)
    return context


def _attach_minimal_draft_document(context: PipelineContext, spec: JobSpec) -> None:
    cards = [
        DraftSlideCard(ref_id=slide.id, order=index, layout_hint=slide.layout)
        for index, slide in enumerate(spec.slides, start=1)
    ]
    section = DraftSection(name="auto", order=1, slides=cards)
    context.add_artifact("draft_document", DraftDocument(sections=[section]))


def test_simple_analyzer_detects_quality_issues(tmp_path) -> None:
    image_path = tmp_path / "image.png"
    _write_dummy_png(image_path)

    spec = JobSpec(
        meta=JobMeta(
            schema_version="1.1",
            title="テスト案件",
            client="Zeta",
            author="営業部",
            created_at="2025-10-05",
            theme="corporate",
        ),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="slide-1",
                layout="Title and Content",
                title="テストスライド",
                bullets=[
                    _group(
                        SlideBullet(
                            id="bullet-1",
                            text="本文",
                            level=4,
                            font=FontSpec(
                                name="Meiryo UI",
                                size_pt=12.0,
                                color_hex="#FFFFFF",
                            ),
                        )
                    )
                ],
                images=[
                    SlideImage(
                        id="img-1",
                        source=str(image_path),
                        left_in=0.1,
                        top_in=0.2,
                        width_in=9.5,
                        height_in=7.0,
                    )
                ],
            )
        ],
    )

    context = _render_spec(spec, tmp_path)
    analyzer = SimpleAnalyzerStep(
        AnalyzerOptions(
            min_font_size=16.0,
            default_font_size=16.0,
            max_bullet_level=3,
            default_font_color="#CCCCCC",
            preferred_text_color="#005BAC",
            background_color="#FFFFFF",
            margin_in=0.5,
        )
    )

    analyzer.run(context)

    analysis_path = context.require_artifact("analysis_path")
    payload = json.loads(analysis_path.read_text(encoding="utf-8"))

    issue_types = {issue["type"] for issue in payload["issues"]}
    assert {
        "margin",
        "font_min",
        "contrast_low",
        "bullet_depth",
        "layout_consistency",
        "grid_misaligned",
    } <= issue_types

    fix_types = {fix["type"] for fix in payload["fixes"]}
    assert {"move", "font_raise", "color_adjust", "bullet_cap", "bullet_reindent"} <= fix_types

    margin_issue = next(issue for issue in payload["issues"] if issue["type"] == "margin")
    assert margin_issue["fix"]["payload"]

    layout_issue = next(issue for issue in payload["issues"] if issue["type"] == "layout_consistency")
    assert layout_issue["fix"]["payload"]["level"] == 0
    contrast_issue = next(issue for issue in payload["issues"] if issue["type"] == "contrast_low")
    assert contrast_issue["metrics"]["required_ratio"] == pytest.approx(4.5)
    assert contrast_issue["metrics"]["font_size_pt"] == pytest.approx(12.0)


def test_analyzer_updates_mapping_log(tmp_path) -> None:
    spec = JobSpec(
        meta=JobMeta(
            schema_version="1.1",
            title="Analyzer 連携テスト",
            client="Zeta",
            author="営業部",
            created_at="2025-10-21",
            theme="corporate",
        ),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="slide-1",
                layout="Title and Content",
                title="概要",
                bullets=[
                    _group(
                        SlideBullet(
                            id="bullet-1",
                            text="最初のポイント",
                            level=0,
                        )
                    )
                ],
            )
        ],
    )

    mapping_context = PipelineContext(spec=spec, workdir=tmp_path)
    _attach_minimal_draft_document(mapping_context, spec)
    prepare_doc = PrepareDocument(
        prepare_id="prepare-test",
        cards=[
            PrepareCard(
                card_id="slide-1",
                order=1,
                role=PrepareCardRole(story_phase="introduction", intent_tags=["intro"]),
                content=PrepareCardContent(
                    headline="概要",
                    body=[PrepareBodyBlock(type="paragraph", text="最初のポイント")],
                ),
            )
        ],
        story_context=PrepareStoryContext(chapters=[]),
    )
    mapping_context.add_artifact("prepare_document", prepare_doc)
    mapping_step = MappingStep(MappingOptions(output_dir=tmp_path))
    mapping_step.run(mapping_context)

    generate_ready = mapping_context.artifacts["generate_ready"]
    render_spec = generate_ready_to_jobspec(generate_ready)
    render_context = PipelineContext(
        spec=render_spec,
        workdir=tmp_path,
        artifacts=dict(mapping_context.artifacts),
    )

    renderer = SimpleRendererStep(RenderingOptions(output_filename="rm031.pptx"))
    renderer.run(render_context)

    analyzer = SimpleAnalyzerStep(
        AnalyzerOptions(
            min_font_size=40.0,
            default_font_size=18.0,
            max_bullet_level=3,
            default_font_color="#777777",
            preferred_text_color="#005BAC",
            background_color="#FFFFFF",
        )
    )
    analyzer.run(render_context)

    mapping_log_path = tmp_path / "mapping_log.json"
    payload = json.loads(mapping_log_path.read_text(encoding="utf-8"))

    slide_summary = payload["slides"][0]["analyzer"]
    assert slide_summary["issue_count"] >= 1
    assert "font_min" in slide_summary["issue_counts_by_type"]
    issue_types = {issue["issue_type"] for issue in slide_summary["issues"]}
    assert "font_min" in issue_types
    target_slide_ids = {issue["target"].get("slide_id") for issue in slide_summary["issues"]}
    assert spec.slides[0].id in target_slide_ids

    meta = payload["meta"]
    assert meta["analyzer_issue_count"] == slide_summary["issue_count"]
    assert meta["analyzer_issue_counts_by_type"]["font_min"] == slide_summary["issue_counts_by_type"]["font_min"]


def test_margin_check_respects_actual_slide_size(tmp_path) -> None:
    template_path = tmp_path / "widescreen.pptx"
    template = Presentation()
    template.slide_width = Inches(13.333333)
    template.slide_height = Inches(7.5)
    template.save(template_path)

    image_path = tmp_path / "wide.png"
    _write_dummy_png(image_path)

    spec = JobSpec(
        meta=JobMeta(
            schema_version="1.1",
            title="ワイドスライド",
            client="Zeta",
            author="営業部",
            created_at="2025-10-08",
            theme="corporate",
        ),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="wide-1",
                layout="Title and Content",
                images=[
                    SlideImage(
                        id="wide-img",
                        source=str(image_path),
                        left_in=1.0,
                        top_in=1.0,
                        width_in=11.0,
                        height_in=5.5,
                    )
                ],
            )
        ],
    )

    context = _render_spec(spec, tmp_path, template_path=template_path)
    analyzer = SimpleAnalyzerStep(
        AnalyzerOptions(
            margin_in=0.5,
            grid_size_in=0.125,
            grid_tolerance_in=0.02,
        )
    )

    analyzer.run(context)

    analysis_path = context.require_artifact("analysis_path")
    payload = json.loads(analysis_path.read_text(encoding="utf-8"))

    margin_issues = [issue for issue in payload["issues"] if issue["type"] == "margin"]
    assert not margin_issues


def test_simple_analyzer_allows_large_text_with_lower_contrast(tmp_path) -> None:
    spec = JobSpec(
        meta=JobMeta(
            schema_version="1.1",
            title="コントラスト調整テスト",
            client="Zeta",
            author="営業部",
            created_at="2025-10-07",
            theme="corporate",
        ),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="slide-large",
                layout="Title and Content",
                bullets=[
                    _group(
                        SlideBullet(
                            id="bullet-large",
                            text="セカンダリカラーの本文",
                            level=0,
                            font=FontSpec(
                                name="Meiryo UI",
                                size_pt=24.0,
                                color_hex="#0097A7",
                            ),
                        )
                    )
                ],
            )
        ],
    )

    context = _render_spec(spec, tmp_path)
    analyzer = SimpleAnalyzerStep(
        AnalyzerOptions(
            min_font_size=16.0,
            default_font_size=16.0,
            max_bullet_level=3,
            default_font_color="#333333",
            preferred_text_color="#005BAC",
            background_color="#FFFFFF",
            large_text_threshold_pt=18.0,
            large_text_min_contrast=3.0,
        )
    )

    analyzer.run(context)

    analysis_path = context.require_artifact("analysis_path")
    payload = json.loads(analysis_path.read_text(encoding="utf-8"))

    issue_types = {issue["type"] for issue in payload["issues"]}
    assert "contrast_low" not in issue_types


def test_analyzer_outputs_structure_snapshot(tmp_path) -> None:
    spec = JobSpec(
        meta=JobMeta(
            schema_version="1.1",
            title="スナップショット検証",
            client="Zeta",
            author="営業部",
            created_at="2025-10-17",
            theme="corporate",
        ),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="slide-structure",
                layout="Title and Content",
                bullets=[
                    _group(
                        SlideBullet(
                            id="bullet-structure",
                            text="アンカー検証",
                            level=1,
                        ),
                    )
                ],
            )
        ],
    )

    context = _render_spec(spec, tmp_path)
    analyzer = SimpleAnalyzerStep(
        AnalyzerOptions(snapshot_output_filename="analysis_snapshot.json")
    )

    analyzer.run(context)

    snapshot_path = context.require_artifact("analyzer_snapshot_path")
    payload = json.loads(snapshot_path.read_text(encoding="utf-8"))

    assert payload["schema_version"] == "1.0.0"
    assert payload["slides"], "スナップショットにスライドが含まれていません"
    slide_entry = payload["slides"][0]
    assert slide_entry["slide_id"] == "slide-structure"
    assert slide_entry["layout"] == "Title and Content"
    assert "placeholders" in slide_entry
    placeholder_paragraphs = [
        entry
        for entry in slide_entry["placeholders"]
        if entry.get("paragraphs")
    ]
    assert placeholder_paragraphs, "プレースホルダーに段落情報が含まれていません"
    paragraph_texts = [
        paragraph["text"]
        for entry in placeholder_paragraphs
        for paragraph in entry["paragraphs"]
        if paragraph["text"]
    ]
    assert paragraph_texts, "段落テキストが空です"
    assert "アンカー検証" in paragraph_texts


def test_slide_snapshot_from_slide_extracts_metadata() -> None:
    snapshot = _build_snapshot_with_textbox()
    assert snapshot.index == 0
    assert snapshot.body_placeholder_id is not None

    body_snapshot = snapshot.shape_by_id(snapshot.body_placeholder_id)
    assert body_snapshot is not None
    assert body_snapshot.is_placeholder
    assert body_snapshot.placeholder_type is not None
    assert body_snapshot.text_frame_padding is not None
    assert pytest.approx(body_snapshot.text_frame_padding["left_in"], rel=1e-3) == 0.2
    assert snapshot.find_shape_by_name("anchored-shape") is not None

    resolver = BulletParagraphResolver(snapshot)
    fallback = resolver.resolve(None)
    assert fallback is not None
    assert fallback.text == "本文"

    anchored_first = resolver.resolve("anchored-shape")
    anchored_second = resolver.resolve("anchored-shape")
    assert anchored_first is not None and anchored_first.text == "アンカー"
    assert anchored_second is not None and anchored_second.text == "サブ"


def test_analyzer_locates_shapes_via_snapshot_helpers(tmp_path) -> None:
    picture_shape = ShapeSnapshot(
        shape_id=101,
        name="anchor-picture",
        shape_type=int(MSO_SHAPE_TYPE.PICTURE),
        left_in=1.0,
        top_in=1.5,
        width_in=2.5,
        height_in=1.0,
    )
    textbox_shape = ShapeSnapshot(
        shape_id=102,
        name="anchor-textbox",
        shape_type=int(MSO_SHAPE_TYPE.TEXT_BOX),
        left_in=0.5,
        top_in=0.75,
        width_in=3.5,
        height_in=1.2,
    )
    placeholder_shape = ShapeSnapshot(
        shape_id=103,
        name="textbox-2",
        shape_type=int(MSO_SHAPE_TYPE.PLACEHOLDER),
        left_in=0.4,
        top_in=1.0,
        width_in=3.0,
        height_in=1.1,
        is_placeholder=True,
        placeholder_type=int(PP_PLACEHOLDER.BODY),
    )
    snapshot = SlideSnapshot(
        index=0,
        shapes=[picture_shape, textbox_shape, placeholder_shape],
        body_placeholder_id=placeholder_shape.shape_id,
    )

    analyzer = SimpleAnalyzerStep()
    image_spec = SlideImage(
        id="image-1",
        source=str((tmp_path / "dummy.png").as_posix()),
        anchor="anchor-picture",
    )
    assert analyzer._locate_image_shape(snapshot, image_spec) is picture_shape

    fallback_image_spec = SlideImage(
        id="image-2",
        source=str((tmp_path / "dummy2.png").as_posix()),
    )
    assert analyzer._locate_image_shape(snapshot, fallback_image_spec) is picture_shape

    textbox_spec = SlideTextbox(id="textbox-1", text="内容", anchor="anchor-textbox")
    assert analyzer._locate_textbox_shape(snapshot, textbox_spec) is textbox_shape

    fallback_textbox = SlideTextbox(id="textbox-2", text="詳細")
    assert analyzer._locate_textbox_shape(snapshot, fallback_textbox) is placeholder_shape


def test_slide_snapshot_handles_placeholder_errors() -> None:
    class BrokenPlaceholderFormat:
        @property
        def type(self):
            raise ValueError("boom")

    class DummyShape:
        shape_id = 1
        left = 0
        top = 0
        width = 0
        height = 0
        name = None
        shape_type = int(MSO_SHAPE_TYPE.AUTO_SHAPE)
        is_placeholder = True
        placeholder_format = BrokenPlaceholderFormat()
        z_order_position = "invalid"
        rotation = "oops"
        has_text_frame = False

    class DummySlide:
        shapes = [DummyShape()]

    snapshot = SlideSnapshot.from_slide(DummySlide(), 0)
    assert snapshot.shape_by_id(999) is None
    assert snapshot.find_shape_by_name("") is None

    resolver = BulletParagraphResolver(snapshot)
    assert resolver.resolve("missing-anchor") is None
    assert resolver.resolve(None) is None


def test_analyzer_internal_helpers_edge_cases(tmp_path) -> None:
    analyzer = SimpleAnalyzerStep()
    spec = Slide(id="edge", layout="Title and Content")
    image_spec = SlideImage(id="img-missing", source="dummy.png", anchor="nothing")
    text_spec = SlideTextbox(id="txt-missing", text="value", anchor="lost")
    dummy_snapshot = SlideSnapshot(index=0, shapes=[], body_placeholder_id=None)

    with pytest.raises(ValueError):
        analyzer._save_snapshot([], tmp_path)

    issues, fixes = analyzer._analyze_images(
        Slide(id=spec.id, layout=spec.layout, images=[image_spec]),
        dummy_snapshot,
        slide_width_in=10.0,
        slide_height_in=7.5,
    )
    assert not issues and not fixes

    issues, fixes = analyzer._analyze_textboxes(
        Slide(id=spec.id, layout=spec.layout, textboxes=[text_spec]),
        dummy_snapshot,
    )
    assert not issues and not fixes

    assert analyzer._shape_type_name(None) == "unknown"
    assert analyzer._shape_type_name(9999) == "9999"
    assert analyzer._placeholder_type_name(None) is None
    assert analyzer._placeholder_type_name(9999) == "9999"

    slide = Slide(id="s", layout="L")
    bullet = SlideBullet(id="b", text="t", level=0)
    target = {"slide_id": "s", "element_id": "b", "element_type": "bullet"}

    size_result = analyzer._check_font_size(slide, bullet, None, target)
    assert size_result is None

    bad_color_paragraph = ParagraphSnapshot(
        shape_id=0,
        shape_name=None,
        shape_type=int(MSO_SHAPE_TYPE.TEXT_BOX),
        paragraph_index=0,
        text="bad",
        level=0,
        font_size_pt=12.0,
        color_hex="#GGGGGG",
    )
    assert (
        analyzer._check_contrast(slide, bullet, bad_color_paragraph, target) is None
    )


def test_paragraph_and_length_helpers() -> None:
    class DummyLength:
        def __init__(self, *, pt: float | None = None, inches: float | None = None, value: float | None = None):
            self.pt = pt
            self.inches = inches
            self._value = value

        def __float__(self) -> float:
            if self._value is None:
                raise TypeError
            return float(self._value)

    class DummyParagraphFormat:
        line_spacing = DummyLength(pt=18.0)
        space_before = DummyLength(pt=2.0)
        space_after = DummyLength(pt=4.0)
        left_indent = DummyLength(inches=0.25)
        right_indent = DummyLength(inches=0.5)
        first_line_indent = DummyLength(inches=0.1)

    class DummyFont:
        def __init__(self, *, size=None, color=None, name=None, bold=None, italic=None):
            self.size = size
            self.color = color
            self.name = name
            self.bold = bold
            self.italic = italic

    class DummyColor:
        def __init__(self, rgb=None):
            self.rgb = rgb

    class DummyParagraph:
        def __init__(self):
            self.font = DummyFont()
            self.paragraph_format = DummyParagraphFormat()
            self.alignment = type("Align", (), {"name": "CENTER"})()
            self.runs = [
                type(
                    "Run",
                    (),
                    {
                        "font": DummyFont(
                            size=DummyLength(pt=20.0),
                            color=DummyColor((0x12, 0x34, 0x56)),
                            name="Fallback",
                            bold=True,
                            italic=False,
                        )
                    },
                )()
            ]

    paragraph = DummyParagraph()
    style = _extract_paragraph_style(paragraph)
    assert style["alignment"] == "center"
    assert style["left_indent_in"] == 0.25

    empty_padding = type("Frame", (), {"margin_left": None, "margin_right": None, "margin_top": None, "margin_bottom": None})()
    assert _extract_text_frame_padding(None) is None
    assert _extract_text_frame_padding(empty_padding) is None

    class FallbackLength:
        def __float__(self):
            return 9144.0

    frame = type(
        "Frame",
        (),
        {
            "margin_left": DummyLength(inches=0.1),
            "margin_right": FallbackLength(),
            "margin_top": DummyLength(inches=0.2),
            "margin_bottom": DummyLength(inches=0.3),
        },
    )()
    padding = _extract_text_frame_padding(frame)
    assert padding["left_in"] == pytest.approx(0.1)
    assert padding["right_in"] == pytest.approx(0.01)

    enum_value = type("EnumValue", (), {"name": "VALUE"})()
    assert _enum_name(enum_value) == "value"
    assert _enum_name("  Mixed ") == "mixed"

    class StrangeStr:
        def __str__(self):
            raise RuntimeError("boom")

    assert _enum_name(StrangeStr()) is None

    assert _length_to_inches(DummyLength(inches=1.5)) == 1.5
    assert _length_to_inches(914400) == pytest.approx(1.0)
    assert _length_to_inches(object()) is None

    assert _length_to_pt(DummyLength(pt=12.0)) == 12.0
    assert _length_to_pt(14.0) == 14.0
    assert _length_to_pt(object()) is None

    assert _normalize_hex("abc123") == "#abc123"
    assert _color_to_hex(DummyColor()) is None
    assert _hex_to_rgb("#336699") == pytest.approx((0x33 / 255, 0x66 / 255, 0x99 / 255))
    with pytest.raises(ValueError):
        _hex_to_rgb("short")

    size, color, name, bold, italic = _extract_font_info(paragraph)
    assert size == 20.0
    assert color == "#123456"
    assert name == "Fallback"
    assert bold is True
    assert italic is False

    assert _contrast_ratio("#000000", "#FFFFFF") == pytest.approx(21.0)


def test_analyzer_run_handles_missing_artifacts(tmp_path) -> None:
    spec = JobSpec(
        meta=JobMeta(schema_version="1.1", title="missing", client="Z", author="X", created_at="2025-01-01", theme="corp"),
        auth=JobAuth(created_by="tester"),
        slides=[],
    )
    context = PipelineContext(spec=spec, workdir=tmp_path)

    analyzer = SimpleAnalyzerStep(allow_missing_artifact=True)
    analyzer.run(context)

    analyzer_strict = SimpleAnalyzerStep()
    with pytest.raises(RuntimeError):
        analyzer_strict.run(context)


def test_locate_textbox_shape_fallback_to_any_textbox() -> None:
    floating_shape = ShapeSnapshot(
        shape_id=50,
        name="other",
        shape_type=int(MSO_SHAPE_TYPE.TEXT_BOX),
        left_in=0.5,
        top_in=0.5,
        width_in=1.0,
        height_in=1.0,
    )
    snapshot = SlideSnapshot(index=0, shapes=[floating_shape])
    analyzer = SimpleAnalyzerStep()
    textbox_spec = SlideTextbox(id="no-match", text="value")
    assert analyzer._locate_textbox_shape(snapshot, textbox_spec) is floating_shape


def test_analyze_textboxes_reports_grid_misalignment(tmp_path) -> None:
    snapshot = SlideSnapshot(
        index=0,
        shapes=[
            ShapeSnapshot(
                shape_id=10,
                name="textbox-anchor",
                shape_type=int(MSO_SHAPE_TYPE.TEXT_BOX),
                left_in=0.18,
                top_in=0.19,
                width_in=1.0,
                height_in=1.0,
            )
        ],
    )
    analyzer = SimpleAnalyzerStep()
    spec = Slide(
        id="slide-grid",
        layout="Title and Content",
        textboxes=[
            SlideTextbox(id="textbox-anchor", text="Grid", anchor="textbox-anchor")
        ],
    )
    issues, fixes = analyzer._analyze_textboxes(spec, snapshot)
    assert issues and fixes
    assert issues[0]["type"] == "grid_misaligned"


def test_find_shape_by_name_enforces_type() -> None:
    shape = ShapeSnapshot(
        shape_id=99,
        name="target",
        shape_type=int(MSO_SHAPE_TYPE.TEXT_BOX),
        left_in=0,
        top_in=0,
        width_in=1,
        height_in=1,
    )
    snapshot = SlideSnapshot(index=0, shapes=[shape])
    assert snapshot.find_shape_by_name("target", shape_type=int(MSO_SHAPE_TYPE.PICTURE)) is None


def test_evaluate_bullet_logs_missing_paragraph() -> None:
    analyzer = SimpleAnalyzerStep()
    slide = Slide(id="slide", layout="Layout")
    bullet = SlideBullet(id="b", text="text", level=0)
    empty_snapshot = SlideSnapshot(index=0, shapes=[])
    resolver = BulletParagraphResolver(empty_snapshot)
    issues, fixes, updated_level, actual_level = analyzer._evaluate_bullet(
        slide,
        bullet,
        anchor="missing",
        resolver=resolver,
        applied_level=None,
        previous_level=None,
    )
    assert updated_level == actual_level
    assert isinstance(issues, list)
    assert isinstance(fixes, list)


def test_check_margins_detects_violation() -> None:
    analyzer = SimpleAnalyzerStep()
    slide = Slide(id="slide", layout="Layout")
    image = SlideImage(id="img", source="dummy.png")
    shape = ShapeSnapshot(
        shape_id=1,
        name="img",
        shape_type=int(MSO_SHAPE_TYPE.PICTURE),
        left_in=0.0,
        top_in=0.0,
        width_in=9.8,
        height_in=7.4,
    )
    issue, fix = analyzer._check_margins(
        slide,
        image,
        shape,
        slide_width_in=10.0,
        slide_height_in=7.5,
    )
    assert issue["type"] == "margin"
    assert issue["metrics"]["violations"]


def test_sync_mapping_log_edge_cases(tmp_path) -> None:
    spec = JobSpec(
        meta=JobMeta(schema_version="1.1", title="sync", client="c", author="a", created_at="2025-01-01", theme="corp"),
        auth=JobAuth(created_by="tester"),
        slides=[],
    )
    context = PipelineContext(spec=spec, workdir=tmp_path)
    analyzer = SimpleAnalyzerStep()

    class BadStr:
        def __str__(self):
            raise ValueError("bad")

    context.add_artifact("mapping_log_path", BadStr())
    analyzer._sync_mapping_log(context, {})

    missing_path = tmp_path / "missing.json"
    context.add_artifact("mapping_log_path", missing_path)
    analyzer._sync_mapping_log(context, {})

    invalid_path = tmp_path / "invalid.json"
    invalid_path.write_text("{", encoding="utf-8")
    context.add_artifact("mapping_log_path", invalid_path)
    analyzer._sync_mapping_log(context, {})


def test_extend_results_handles_optional_fix() -> None:
    analyzer = SimpleAnalyzerStep()
    issues: list[dict[str, Any]] = []
    fixes: list[dict[str, Any]] = []

    analyzer._extend_results(issues, fixes, None)
    assert not issues
    assert not fixes

    issue_only = {"id": "issue-1"}
    analyzer._extend_results(issues, fixes, (issue_only, None))
    assert issues == [issue_only]
    assert not fixes

    issue_with_fix = {"id": "issue-2"}
    fix_payload = {"id": "fix-2"}
    analyzer._extend_results(issues, fixes, (issue_with_fix, fix_payload))
    assert issues[-1] is issue_with_fix
    assert fixes[-1] is fix_payload
