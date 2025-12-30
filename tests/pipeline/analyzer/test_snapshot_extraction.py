from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_generator.pipeline.analyzer.snapshot import SlideSnapshot
from pptx_generator.pipeline.analyzer.utils import EMU_PER_INCH


class _DummyTextFrame:
    def __init__(self, text: str):
        self.paragraphs = [_DummyParagraph(text)]
        self.word_wrap = None
        self.vertical_anchor = None
        self.auto_size = None
        self.margin_left = None
        self.margin_right = None
        self.margin_top = None
        self.margin_bottom = None

    def clear(self) -> None:  # pragma: no cover - unused in this test
        self.paragraphs = []


class _DummyFont:
    def __init__(self):
        self.size = None
        self.color = None
        self.name = None
        self.bold = None
        self.italic = None


class _DummyRun:
    def __init__(self, text: str):
        self.text = text
        self.font = _DummyFont()


class _DummyParagraph:
    def __init__(self, text: str):
        self.text = text
        self.level = 0
        self.font = _DummyFont()
        self.runs = [_DummyRun(text)]
        self.paragraph_format = _DummyParagraphFormat()
        self.alignment = None


class _DummyParagraphFormat:
    def __init__(self):
        self.line_spacing = None
        self.space_before = None
        self.space_after = None
        self.left_indent = None
        self.right_indent = None
        self.first_line_indent = None


class _DummyCell:
    def __init__(self, row_idx: int, col_idx: int, text: str, left: int, top: int):
        self.row_idx = row_idx
        self.col_idx = col_idx
        self.left = left
        self.top = top
        self.width = EMU_PER_INCH
        self.height = EMU_PER_INCH
        self.text_frame = _DummyTextFrame(text)


class _DummyRow:
    def __init__(self, row_idx: int, texts: list[str], origin_left: int, origin_top: int):
        self.cells = [
            _DummyCell(row_idx, col_idx, text, origin_left + col_idx * EMU_PER_INCH, origin_top)
            for col_idx, text in enumerate(texts)
        ]


class _DummyTable:
    def __init__(self, rows: list[list[str]], left: int, top: int):
        self.rows = [_DummyRow(idx, row, left, top + idx * EMU_PER_INCH) for idx, row in enumerate(rows)]


class _DummyShape:
    def __init__(
        self,
        *,
        shape_id: int,
        name: str,
        shape_type: int,
        left: int,
        top: int,
        width: int = EMU_PER_INCH,
        height: int = EMU_PER_INCH,
        text: str | None = None,
        children: list["_DummyShape"] | None = None,
        table: _DummyTable | None = None,
    ):
        self.shape_id = shape_id
        self.name = name
        self.shape_type = shape_type
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.shapes = children or []
        self.table = table
        self.has_table = table is not None
        self.is_placeholder = False
        self.placeholder_format = None
        self.z_order_position = None
        self.rotation = None
        if text is not None:
            self.has_text_frame = True
            self.text_frame = _DummyTextFrame(text)
        else:
            self.has_text_frame = False
            self.text_frame = None


class _DummySlide:
    def __init__(self, shapes):
        self.shapes = shapes


def test_slide_snapshot_collects_group_children_with_offsets() -> None:
    group_left = 2 * EMU_PER_INCH
    group_top = 3 * EMU_PER_INCH
    child = _DummyShape(
        shape_id=2,
        name="child",
        shape_type=int(MSO_SHAPE_TYPE.TEXT_BOX),
        left=1 * EMU_PER_INCH,
        top=0,
        text="inside group",
    )
    group = _DummyShape(
        shape_id=1,
        name="group",
        shape_type=int(MSO_SHAPE_TYPE.GROUP),
        left=group_left,
        top=group_top,
        children=[child],
    )
    slide = _DummySlide([group])

    snapshot = SlideSnapshot.from_slide(slide, 0)

    assert len(snapshot.shapes) == 2
    child_snapshot = next(s for s in snapshot.shapes if s.name == "child")
    assert child_snapshot.parent_shape_id == group.shape_id
    assert child_snapshot.left_in == 3.0  # 2in offset + 1in child
    assert child_snapshot.top_in == 3.0
    assert child_snapshot.paragraphs[0].text == "inside group"


def test_slide_snapshot_includes_table_cells_with_coordinates() -> None:
    table = _DummyTable(rows=[["r0c0", "r0c1"], ["r1c0", "r1c1"]], left=EMU_PER_INCH, top=EMU_PER_INCH)
    table_shape = _DummyShape(
        shape_id=10,
        name="table-shape",
        shape_type=int(MSO_SHAPE_TYPE.TABLE),
        left=EMU_PER_INCH,
        top=EMU_PER_INCH,
        table=table,
    )
    slide = _DummySlide([table_shape])

    snapshot = SlideSnapshot.from_slide(slide, 0)

    cell_snapshots = [shape for shape in snapshot.shapes if shape.table_cell is not None]
    assert len(cell_snapshots) == 4
    sample_cell = next(shape for shape in cell_snapshots if shape.table_cell == {"row": 1, "col": 0})
    assert sample_cell.parent_shape_id == table_shape.shape_id
    assert sample_cell.left_in == 1.0  # absolute座標
    assert sample_cell.top_in == 2.0  # table top (1in) + row offset (1in)
    assert sample_cell.paragraphs[0].text == "r1c0"
