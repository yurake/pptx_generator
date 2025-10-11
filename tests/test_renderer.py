"""SimpleRendererStep の拡張動作テスト。"""

from __future__ import annotations

import base64
import logging
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Inches

from pptx_generator.models import (
    ChartOptions,
    ChartSeries,
    JobAuth,
    JobMeta,
    JobSpec,
    Slide,
    SlideBullet,
    SlideBulletGroup,
    SlideChart,
    SlideImage,
    SlideTable,
    TableStyle,
)
from pptx_generator.pipeline.base import PipelineContext
from pptx_generator.pipeline.renderer import RenderingOptions, SimpleRendererStep
from pptx_generator.settings import BrandingConfig, BrandingFont
from pydantic import ValidationError


def _write_dummy_png(path: Path) -> None:
    payload = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR4nGNgYAAAAAMAASsJTYQAAAAASUVORK5CYII="
    )
    path.write_bytes(payload)


def _load_placeholder_boxes(
    template_path: Path, layout_name: str
) -> dict[str, tuple[int, int, int, int, int]]:
    presentation = Presentation(template_path)
    layout = next(
        layout for layout in presentation.slide_layouts if layout.name == layout_name
    )
    placeholders: dict[str, tuple[int, int, int, int, int]] = {}
    for shape in layout.shapes:
        if getattr(shape, "is_placeholder", False):
            placeholders[shape.name] = (
                int(shape.left),
                int(shape.top),
                int(shape.width),
                int(shape.height),
                shape.placeholder_format.idx,
            )
    return placeholders


def _emu_box_from_inches(
    box: tuple[float, float, float, float],
) -> tuple[int, int, int, int]:
    left_in, top_in, width_in, height_in = box
    return (
        int(Inches(left_in)),
        int(Inches(top_in)),
        int(Inches(width_in)),
        int(Inches(height_in)),
    )


def test_renderer_renders_rich_content(tmp_path: Path) -> None:
    image_path = tmp_path / "image.png"
    _write_dummy_png(image_path)

    spec = JobSpec(
        meta=JobMeta(schema_version="1.0", title="提案"),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="slide-1",
                layout="Title and Content",
                title="サンプルタイトル",
                bullets=[
                    SlideBulletGroup(
                        items=[
                            SlideBullet(id="b1", text="箇条書き", level=0, font=None),
                        ]
                    ),
                ],
                images=[
                    SlideImage(
                        id="img",
                        source=str(image_path),
                        left_in=6.0,
                        top_in=1.5,
                        width_in=3.0,
                        height_in=3.0,
                    )
                ],
                tables=[
                    SlideTable(
                        id="tbl",
                        columns=["指標", "値"],
                        rows=[["A", 10], ["B", 20]],
                        style=TableStyle(header_fill="#334455", zebra=True),
                    )
                ],
                charts=[
                    SlideChart(
                        id="chart",
                        type="column",
                        categories=["Before", "After"],
                        series=[
                            ChartSeries(
                                name="効果", values=[10, 5], color_hex="#123456"
                            ),
                        ],
                        options=ChartOptions(data_labels=True, y_axis_format="0%"),
                    )
                ],
            )
        ],
    )

    branding = BrandingConfig(
        heading_font=BrandingFont(
            name="HeadingBrand", size_pt=30.0, color_hex="#101010"
        ),
        body_font=BrandingFont(name="BodyBrand", size_pt=20.0, color_hex="#202020"),
        primary_color="#445566",
        secondary_color="#DDEEFF",
        accent_color="#CC5500",
        background_color="#FFFFFF",
    )

    context = PipelineContext(spec=spec, workdir=tmp_path)
    renderer = SimpleRendererStep(
        RenderingOptions(output_filename="out.pptx", branding=branding)
    )
    renderer.run(context)

    pptx_path = context.artifacts["pptx_path"]
    presentation = Presentation(pptx_path)
    slide = presentation.slides[0]

    pictures = [
        shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert pictures, "画像が描画されていること"

    tables = [shape for shape in slide.shapes if getattr(shape, "has_table", False)]
    assert tables, "テーブルが描画されていること"
    table = tables[0].table

    header_cell = table.rows[0].cells[0]
    assert header_cell.fill.fore_color.rgb == RGBColor.from_string("334455")
    data_cell = table.rows[2].cells[0]
    assert data_cell.fill.fore_color.rgb == RGBColor.from_string("DDEEFF")

    text_paragraph = None
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text == "箇条書き":
                    text_paragraph = paragraph
                    break
        if text_paragraph:
            break
    assert text_paragraph is not None
    assert text_paragraph.font.name == "BodyBrand"

    charts = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
    assert charts, "チャートが描画されていること"
    chart = charts[0].chart
    plot = chart.plots[0]
    assert plot.has_data_labels is True
    assert plot.data_labels.show_value is True
    if hasattr(chart, "value_axis"):
        assert chart.value_axis.tick_labels.number_format == "0%"
    series = chart.series[0]
    assert series.format.fill.fore_color.rgb == RGBColor.from_string("123456")


def test_renderer_falls_back_when_anchor_not_specified(tmp_path: Path) -> None:
    spec = JobSpec(
        meta=JobMeta(schema_version="1.0", title="Fallback Test"),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="fallback-table",
                layout="Title and Content",
                tables=[
                    SlideTable(
                        id="table",
                        columns=["A"],
                        rows=[["value"]],
                    )
                ],
            )
        ],
    )

    context = PipelineContext(spec=spec, workdir=tmp_path)
    renderer = SimpleRendererStep(RenderingOptions(output_filename="fallback.pptx"))
    renderer.run(context)

    presentation = Presentation(context.require_artifact("pptx_path"))
    slide = presentation.slides[0]
    table_shape = next(
        shape for shape in slide.shapes if getattr(shape, "has_table", False)
    )

    expected_box = _emu_box_from_inches((1.0, 1.5, 8.5, 3.0))
    assert (
        table_shape.left,
        table_shape.top,
        table_shape.width,
        table_shape.height,
    ) == expected_box


def test_renderer_falls_back_when_anchor_unknown(tmp_path: Path) -> None:
    spec = JobSpec(
        meta=JobMeta(schema_version="1.0", title="Fallback Unknown"),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="missing-anchor",
                layout="Title and Content",
                tables=[
                    SlideTable(
                        id="table",
                        anchor="Unknown Anchor",
                        columns=["A"],
                        rows=[["value"]],
                    )
                ],
            )
        ],
    )

    context = PipelineContext(spec=spec, workdir=tmp_path)
    renderer = SimpleRendererStep(
        RenderingOptions(output_filename="fallback-unknown.pptx")
    )
    renderer.run(context)

    presentation = Presentation(context.require_artifact("pptx_path"))
    slide = presentation.slides[0]
    table_shape = next(
        shape for shape in slide.shapes if getattr(shape, "has_table", False)
    )

    expected_box = _emu_box_from_inches((1.0, 1.5, 8.5, 3.0))
    assert (
        table_shape.left,
        table_shape.top,
        table_shape.width,
        table_shape.height,
    ) == expected_box


def test_renderer_placeholder_centering_tolerance(tmp_path: Path) -> None:
    (
        template_path,
        two_content_layout_name,
        _,
        left_box,
        right_box,
        _picture_box,
    ) = _build_template_with_named_placeholders(tmp_path)

    left_in = left_box[0] / Inches(1)
    top_in = left_box[1] / Inches(1)
    width_in = left_box[2] / Inches(1)
    height_in = left_box[3] / Inches(1)

    override_width = width_in * 0.6
    override_height = height_in * 0.75
    centered_left = left_in + (width_in - override_width) / 2
    centered_top = top_in + (height_in - override_height) / 2

    image_path = tmp_path / "tolerance.png"
    _write_dummy_png(image_path)

    spec = JobSpec(
        meta=JobMeta(schema_version="1.0", title="Tolerance Test"),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="centered-image",
                layout=two_content_layout_name,
                title="許容誤差テスト",
                images=[
                    SlideImage(
                        id="image",
                        anchor="Left Content Placeholder",
                        source=str(image_path),
                        left_in=centered_left,
                        top_in=centered_top,
                        width_in=override_width,
                        height_in=override_height,
                    )
                ],
                tables=[
                    SlideTable(
                        id="table",
                        anchor="Right Content Placeholder",
                        columns=["A"],
                        rows=[["1"]],
                    )
                ],
            )
        ],
    )

    context = PipelineContext(spec=spec, workdir=tmp_path)
    renderer = SimpleRendererStep(
        RenderingOptions(
            template_path=template_path,
            output_filename="tolerance.pptx",
            branding=BrandingConfig.default(),
        )
    )
    renderer.run(context)

    presentation = Presentation(context.require_artifact("pptx_path"))
    slide = presentation.slides[0]

    table_shape = next(
        shape for shape in slide.shapes if getattr(shape, "has_table", False)
    )
    assert (
        table_shape.left,
        table_shape.top,
        table_shape.width,
        table_shape.height,
    ) == right_box

    picture_shape = next(
        shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    placeholder_center_x = left_box[0] + left_box[2] // 2
    placeholder_center_y = left_box[1] + left_box[3] // 2
    picture_center_x = picture_shape.left + picture_shape.width // 2
    picture_center_y = picture_shape.top + picture_shape.height // 2

    assert abs(picture_center_x - placeholder_center_x) <= 1000
    assert abs(picture_center_y - placeholder_center_y) <= 1000


def _build_template_with_named_placeholders(tmp_path: Path):
    presentation = Presentation()

    two_content_layout = presentation.slide_layouts[3]
    left_placeholder_box: tuple[int, int, int, int] | None = None
    right_placeholder_box: tuple[int, int, int, int] | None = None

    supported_types = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
    for shape in two_content_layout.shapes:
        if (
            getattr(shape, "is_placeholder", False)
            and shape.placeholder_format.type in supported_types
        ):
            if left_placeholder_box is None:
                shape.name = "Left Content Placeholder"
                left_placeholder_box = (
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height,
                )
            else:
                shape.name = "Right Content Placeholder"
                right_placeholder_box = (
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height,
                )

    picture_layout = presentation.slide_layouts[8]
    picture_placeholder_box: tuple[int, int, int, int] | None = None
    for shape in picture_layout.shapes:
        if (
            getattr(shape, "is_placeholder", False)
            and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE
        ):
            shape.name = "Picture Content Placeholder"
            picture_placeholder_box = (
                shape.left,
                shape.top,
                shape.width,
                shape.height,
            )
            break

    template_path = tmp_path / "template_placeholders.pptx"
    presentation.save(template_path)

    assert left_placeholder_box and right_placeholder_box and picture_placeholder_box

    return (
        template_path,
        two_content_layout.name,
        picture_layout.name,
        left_placeholder_box,
        right_placeholder_box,
        picture_placeholder_box,
    )


def test_renderer_uses_layout_placeholder_names_for_anchors(tmp_path: Path) -> None:
    (
        template_path,
        two_content_layout_name,
        picture_layout_name,
        left_box,
        right_box,
        picture_box,
    ) = _build_template_with_named_placeholders(tmp_path)

    image_path = tmp_path / "placeholder_image.png"
    _write_dummy_png(image_path)

    spec = JobSpec(
        meta=JobMeta(schema_version="1.0", title="Placeholder Anchor Test"),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="table-slide",
                layout=two_content_layout_name,
                title="表スライド",
                tables=[
                    SlideTable(
                        id="table",
                        anchor="Left Content Placeholder",
                        columns=["ヘッダ1", "ヘッダ2"],
                        rows=[["A", "B"]],
                    )
                ],
            ),
            Slide(
                id="chart-slide",
                layout=two_content_layout_name,
                title="チャートスライド",
                charts=[
                    SlideChart(
                        id="chart",
                        anchor="Right Content Placeholder",
                        type="column",
                        categories=["X", "Y"],
                        series=[ChartSeries(name="Series", values=[1, 2])],
                    )
                ],
            ),
            Slide(
                id="image-slide",
                layout=picture_layout_name,
                title="画像スライド",
                images=[
                    SlideImage(
                        id="image",
                        anchor="Picture Content Placeholder",
                        source=str(image_path),
                    )
                ],
            ),
        ],
    )

    context = PipelineContext(spec=spec, workdir=tmp_path)
    renderer = SimpleRendererStep(
        RenderingOptions(
            template_path=template_path,
            output_filename="placeholders.pptx",
            branding=BrandingConfig.default(),
        )
    )
    renderer.run(context)

    pptx_path = context.require_artifact("pptx_path")
    presentation = Presentation(pptx_path)

    table_slide = presentation.slides[0]
    table_shape = next(
        shape for shape in table_slide.shapes if getattr(shape, "has_table", False)
    )
    assert (
        table_shape.left,
        table_shape.top,
        table_shape.width,
        table_shape.height,
    ) == left_box

    chart_slide = presentation.slides[1]
    chart_shape = next(
        shape for shape in chart_slide.shapes if getattr(shape, "has_chart", False)
    )
    assert (
        chart_shape.left,
        chart_shape.top,
        chart_shape.width,
        chart_shape.height,
    ) == right_box

    image_slide = presentation.slides[2]
    picture_shape = next(
        shape
        for shape in image_slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    expected_left, expected_top, expected_width, expected_height = picture_box
    assert picture_shape.top == expected_top
    assert picture_shape.left >= expected_left
    assert picture_shape.width <= expected_width
    assert picture_shape.height <= expected_height
    expected_center_x = expected_left + expected_width // 2
    actual_center_x = picture_shape.left + picture_shape.width // 2
    assert abs(actual_center_x - expected_center_x) <= 2000

    def placeholder_names(slide):
        return {
            shape.name
            for shape in slide.shapes
            if getattr(shape, "is_placeholder", False)
        }

    assert "Left Content Placeholder" not in placeholder_names(table_slide)
    assert "Right Content Placeholder" not in placeholder_names(chart_slide)
    assert "Picture Content Placeholder" not in placeholder_names(image_slide)


def test_renderer_fallback_when_placeholder_removed(tmp_path: Path, caplog) -> None:
    (
        template_path,
        two_content_layout_name,
        _,
        _left_box,
        _right_box,
        _picture_box,
    ) = _build_template_with_named_placeholders(tmp_path)

    presentation = Presentation(template_path)
    layout = next(
        layout
        for layout in presentation.slide_layouts
        if layout.name == two_content_layout_name
    )
    for shape in list(layout.shapes):
        if (
            getattr(shape, "is_placeholder", False)
            and shape.name == "Left Content Placeholder"
        ):
            shape.element.getparent().remove(shape.element)
            break
    removed_template_path = tmp_path / "template_placeholder_removed.pptx"
    presentation.save(removed_template_path)

    spec = JobSpec(
        meta=JobMeta(schema_version="1.0", title="Removed Placeholder"),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="table-slide",
                layout=two_content_layout_name,
                tables=[
                    SlideTable(
                        id="table",
                        anchor="Left Content Placeholder",
                        columns=["A"],
                        rows=[["1"]],
                    )
                ],
            )
        ],
    )

    context = PipelineContext(spec=spec, workdir=tmp_path)
    renderer = SimpleRendererStep(
        RenderingOptions(
            template_path=removed_template_path, output_filename="removed.pptx"
        )
    )

    caplog.clear()
    with caplog.at_level(logging.ERROR, logger="pptx_generator.pipeline.renderer"):
        renderer.run(context)

    presentation = Presentation(context.require_artifact("pptx_path"))
    slide = presentation.slides[0]
    table_shape = next(
        shape for shape in slide.shapes if getattr(shape, "has_table", False)
    )

    expected_box = _emu_box_from_inches((1.0, 1.5, 8.5, 3.0))
    assert (
        table_shape.left,
        table_shape.top,
        table_shape.width,
        table_shape.height,
    ) == expected_box
    assert not any(record.levelno >= logging.ERROR for record in caplog.records)


def test_renderer_fallback_when_placeholder_renamed(tmp_path: Path, caplog) -> None:
    (
        template_path,
        two_content_layout_name,
        _,
        _left_box,
        _right_box,
        _picture_box,
    ) = _build_template_with_named_placeholders(tmp_path)

    presentation = Presentation(template_path)
    layout = next(
        layout
        for layout in presentation.slide_layouts
        if layout.name == two_content_layout_name
    )
    for shape in layout.shapes:
        if (
            getattr(shape, "is_placeholder", False)
            and shape.name == "Left Content Placeholder"
        ):
            shape.name = "Renamed Placeholder"
            break
    renamed_template_path = tmp_path / "template_placeholder_renamed.pptx"
    presentation.save(renamed_template_path)

    spec = JobSpec(
        meta=JobMeta(schema_version="1.0", title="Renamed Placeholder"),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="table-slide",
                layout=two_content_layout_name,
                tables=[
                    SlideTable(
                        id="table",
                        anchor="Left Content Placeholder",
                        columns=["A"],
                        rows=[["1"]],
                    )
                ],
            )
        ],
    )

    context = PipelineContext(spec=spec, workdir=tmp_path)
    renderer = SimpleRendererStep(
        RenderingOptions(
            template_path=renamed_template_path, output_filename="renamed.pptx"
        )
    )

    caplog.clear()
    with caplog.at_level(logging.ERROR, logger="pptx_generator.pipeline.renderer"):
        renderer.run(context)

    presentation = Presentation(context.require_artifact("pptx_path"))
    slide = presentation.slides[0]
    table_shape = next(
        shape for shape in slide.shapes if getattr(shape, "has_table", False)
    )

    expected_box = _emu_box_from_inches((1.0, 1.5, 8.5, 3.0))
    assert (
        table_shape.left,
        table_shape.top,
        table_shape.width,
        table_shape.height,
    ) == expected_box
    assert not any(record.levelno >= logging.ERROR for record in caplog.records)


def test_renderer_handles_object_placeholders(tmp_path: Path) -> None:
    template_path = Path("samples/templates/templates.pptx")
    placeholders = _load_placeholder_boxes(template_path, "Two Column Detail")
    assert (
        "Body Left" in placeholders
        and "Body Right" in placeholders
        and "Logo" in placeholders
    )

    image_path = Path("samples/assets/logo.png")
    spec = JobSpec(
        meta=JobMeta(schema_version="1.0", title="Object Placeholder Test"),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="object-placeholder",
                layout="Two Column Detail",
                title="配置テスト",
                tables=[
                    SlideTable(
                        id="table",
                        anchor="Body Left",
                        columns=["A", "B"],
                        rows=[["x", "y"]],
                    )
                ],
                charts=[
                    SlideChart(
                        id="chart",
                        anchor="Body Right",
                        type="column",
                        categories=["c1"],
                        series=[ChartSeries(name="s", values=[1])],
                    )
                ],
                images=[
                    SlideImage(
                        id="image",
                        anchor="Logo",
                        source=str(image_path),
                        sizing="stretch",
                    )
                ],
            )
        ],
    )

    context = PipelineContext(spec=spec, workdir=tmp_path)
    renderer = SimpleRendererStep(
        RenderingOptions(
            template_path=template_path,
            output_filename="object-placeholder.pptx",
            branding=BrandingConfig.default(),
        )
    )
    renderer.run(context)

    presentation = Presentation(context.require_artifact("pptx_path"))
    slide = presentation.slides[0]

    table_shape = next(
        shape for shape in slide.shapes if getattr(shape, "has_table", False)
    )
    assert (
        int(table_shape.left),
        int(table_shape.top),
        int(table_shape.width),
        int(table_shape.height),
    ) == placeholders["Body Left"][:4]

    chart_shape = next(
        shape for shape in slide.shapes if getattr(shape, "has_chart", False)
    )
    assert (
        int(chart_shape.left),
        int(chart_shape.top),
        int(chart_shape.width),
        int(chart_shape.height),
    ) == placeholders["Body Right"][:4]

    picture_shape = next(
        shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    assert (
        int(picture_shape.left),
        int(picture_shape.top),
        int(picture_shape.width),
        int(picture_shape.height),
    ) == placeholders["Logo"][:4]

    remaining_placeholder_names = {
        shape.name for shape in slide.shapes if getattr(shape, "is_placeholder", False)
    }
    assert "Body Left" not in remaining_placeholder_names
    assert "Body Right" not in remaining_placeholder_names
    assert "Logo" not in remaining_placeholder_names


def test_renderer_removes_bullet_placeholder_when_anchor_specified(
    tmp_path: Path,
) -> None:
    """グループでアンカー指定時にプレースホルダーが削除されることを確認するテスト。"""
    (
        template_path,
        two_content_layout_name,
        _,
        left_box,
        _right_box,
        _picture_box,
    ) = _build_template_with_named_placeholders(tmp_path)

    spec = JobSpec(
        meta=JobMeta(schema_version="1.0", title="Bullet Placeholder Removal Test"),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="bullet-slide",
                layout=two_content_layout_name,
                title="箇条書きプレースホルダー削除テスト",
                bullets=[
                    SlideBulletGroup(
                        anchor="Left Content Placeholder",
                        items=[
                            SlideBullet(
                                id="bullet1",
                                text="アンカー指定の箇条書き1",
                                level=0,
                            ),
                            SlideBullet(
                                id="bullet2",
                                text="アンカー指定の箇条書き2",
                                level=1,
                            ),
                        ],
                    )
                ],
            )
        ],
    )

    context = PipelineContext(spec=spec, workdir=tmp_path)
    renderer = SimpleRendererStep(
        RenderingOptions(
            template_path=template_path,
            output_filename="bullet-placeholder-removal.pptx",
            branding=BrandingConfig.default(),
        )
    )
    renderer.run(context)

    presentation = Presentation(context.require_artifact("pptx_path"))
    slide = presentation.slides[0]

    # 箇条書きテキストが正しく配置されていることを確認
    bullet_texts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text in [
                    "アンカー指定の箇条書き1",
                    "アンカー指定の箇条書き2",
                ]:
                    bullet_texts.append(paragraph.text)
                    # テキストフレームが左側プレースホルダーの位置に配置されていることを確認
                    text_shape = shape
                    assert abs(text_shape.left - left_box[0]) < 100  # 許容誤差
                    assert abs(text_shape.top - left_box[1]) < 100

    assert "アンカー指定の箇条書き1" in bullet_texts
    assert "アンカー指定の箇条書き2" in bullet_texts

    # プレースホルダーが削除されていることを確認
    remaining_placeholder_names = {
        shape.name for shape in slide.shapes if getattr(shape, "is_placeholder", False)
    }
    assert "Left Content Placeholder" not in remaining_placeholder_names


def test_renderer_bullet_fallback_when_no_anchor(tmp_path: Path) -> None:
    """アンカー未指定グループが本文プレースホルダーへ配置されることを確認するテスト。"""
    spec = JobSpec(
        meta=JobMeta(schema_version="1.0", title="Bullet Fallback Test"),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="bullet-fallback",
                layout="Title and Content",
                title="箇条書きフォールバックテスト",
                bullets=[
                    SlideBulletGroup(
                        items=[
                            SlideBullet(id="b1", text="アンカー未指定の箇条書き", level=0),
                        ]
                    )
                ],
            )
        ],
    )

    context = PipelineContext(spec=spec, workdir=tmp_path)
    renderer = SimpleRendererStep(
        RenderingOptions(output_filename="bullet-fallback.pptx")
    )
    renderer.run(context)

    presentation = Presentation(context.require_artifact("pptx_path"))
    slide = presentation.slides[0]

    # 箇条書きテキストが配置されていることを確認
    bullet_found = False
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text == "アンカー未指定の箇条書き":
                    bullet_found = True
                    break
        if bullet_found:
            break

def test_job_spec_rejects_legacy_bullets_schema() -> None:
    legacy_spec = {
        "meta": {"schema_version": "1.1", "title": "Legacy"},
        "auth": {"created_by": "tester"},
        "slides": [
            {
                "id": "legacy-slide",
                "layout": "Title and Content",
                # 旧形式: bullets は SlideBullet のリスト
                "bullets": [
                    {"id": "b1", "text": "旧形式", "level": 0},
                ],
            }
        ],
    }

    with pytest.raises(ValidationError):
        JobSpec.model_validate(legacy_spec)


def test_slide_bullet_rejects_anchor_field() -> None:
    with pytest.raises(ValidationError):
        SlideBullet.model_validate(
            {"id": "b1", "text": "anchor", "level": 0, "anchor": "Body"}
        )


def test_renderer_renders_multiple_bullet_groups(tmp_path: Path) -> None:
    """グループ形式の複数アンカーが正しく描画されることを確認するテスト。"""
    (
        template_path,
        two_content_layout_name,
        _picture_layout_name,
        left_box,
        right_box,
        _picture_box,
    ) = _build_template_with_named_placeholders(tmp_path)

    spec = JobSpec(
        meta=JobMeta(schema_version="1.0", title="Grouped Bullets Test"),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="grouped-bullets",
                layout=two_content_layout_name,
                bullets=[
                    SlideBulletGroup(
                        anchor="Left Content Placeholder",
                        items=[
                            SlideBullet(id="l1", text="左側の箇条書き1", level=0),
                            SlideBullet(id="l2", text="左側の箇条書き2", level=1),
                        ],
                    ),
                    SlideBulletGroup(
                        anchor="Right Content Placeholder",
                        items=[
                            SlideBullet(id="r1", text="右側の箇条書き1", level=0),
                        ],
                    ),
                ],
            )
        ],
    )

    context = PipelineContext(spec=spec, workdir=tmp_path)
    renderer = SimpleRendererStep(
        RenderingOptions(
            template_path=template_path,
            output_filename="grouped-bullets.pptx",
            branding=BrandingConfig.default(),
        )
    )
    renderer.run(context)

    presentation = Presentation(context.require_artifact("pptx_path"))
    slide = presentation.slides[0]

    left_text_shape = None
    right_text_shape = None
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            texts = {paragraph.text for paragraph in shape.text_frame.paragraphs}
            if "左側の箇条書き1" in texts:
                left_text_shape = shape
            if "右側の箇条書き1" in texts:
                right_text_shape = shape

    assert left_text_shape is not None
    assert right_text_shape is not None

    assert abs(left_text_shape.left - left_box[0]) < 100
    assert abs(left_text_shape.top - left_box[1]) < 100
    assert abs(right_text_shape.left - right_box[0]) < 100
    assert abs(right_text_shape.top - right_box[1]) < 100

    remaining_placeholder_names = {
        shape.name for shape in slide.shapes if getattr(shape, "is_placeholder", False)
    }
    assert "Left Content Placeholder" not in remaining_placeholder_names
    assert "Right Content Placeholder" not in remaining_placeholder_names


def test_renderer_raises_error_for_duplicate_group_anchor(tmp_path: Path) -> None:
    (
        template_path,
        two_content_layout_name,
        _picture_layout_name,
        _left_box,
        _right_box,
        _picture_box,
    ) = _build_template_with_named_placeholders(tmp_path)

    spec = JobSpec(
        meta=JobMeta(schema_version="1.1", title="Duplicate Anchor"),
        auth=JobAuth(created_by="tester"),
        slides=[
            Slide(
                id="duplicate-anchors",
                layout=two_content_layout_name,
                bullets=[
                    SlideBulletGroup(
                        anchor="Left Content Placeholder",
                        items=[
                            SlideBullet(id="b1", text="左", level=0),
                        ],
                    ),
                    SlideBulletGroup(
                        anchor="Left Content Placeholder",
                        items=[
                            SlideBullet(id="b2", text="右", level=0),
                        ],
                    ),
                ],
            )
        ],
    )

    context = PipelineContext(spec=spec, workdir=tmp_path)
    renderer = SimpleRendererStep(
        RenderingOptions(
            template_path=template_path,
            output_filename="duplicate-anchors.pptx",
            branding=BrandingConfig.default(),
        )
    )

    with pytest.raises(ValueError, match="箇条書きのアンカー 'Left Content Placeholder'"):
        renderer.run(context)
